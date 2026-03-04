"""
Tests pour utils/pptx_generator.py
Cible: 7% -> couverture maximale (708 stmts)

python-pptx est installe (requirements.txt), donc on peut importer le module
normalement puis patcher les objets au niveau du module.
"""
import json
import os
import sys
from pathlib import Path
from unittest.mock import MagicMock, patch, call, PropertyMock

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import at module level so the module is cached BEFORE any fixture patches pptx
from utils.pptx_generator import (  # noqa: E402
    ProposalPPTXGenerator,
    build_proposal_pptx,
    _hex_to_rgb,
    _load_branding,
    BRANDING,
    COLORS,
    FONT_TITLE,
    FONT_BODY,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_placeholder(idx, has_text_frame=True):
    """Build a mock placeholder with the given idx."""
    ph = MagicMock()
    ph.placeholder_format.idx = idx
    ph.has_text_frame = has_text_frame
    tf = MagicMock()
    para = MagicMock()
    para.runs = []
    tf.paragraphs = [para]
    ph.text_frame = tf
    return ph


def _make_slide_mock(placeholder_idxs=(0, 1, 2)):
    """Return a mock slide with placeholders at the given indices."""
    slide = MagicMock()
    phs = [_make_placeholder(i) for i in placeholder_idxs]
    slide.placeholders = phs
    slide.shapes = MagicMock()
    return slide


def _make_presentation_mock(num_existing_slides=0, num_layouts=20):
    """Return a fully-mocked Presentation object."""
    prs = MagicMock()

    # slide list behaviour
    prs.slides._sldIdLst = []

    # layouts
    layouts = [MagicMock() for _ in range(num_layouts)]
    prs.slide_layouts = layouts

    # by default adding a slide returns a slide with placeholders 0, 1, 2
    default_slide = _make_slide_mock()
    prs.slides.add_slide.return_value = default_slide

    return prs


# ---------------------------------------------------------------------------
# _load_branding
# ---------------------------------------------------------------------------

class TestLoadBranding:
    def test_returns_dict_with_colors_and_fonts(self):
        """_load_branding() retourne au moins colors et fonts."""
        result = _load_branding()
        assert "colors" in result
        assert "fonts" in result

    def test_returns_dict_when_branding_file_missing(self, tmp_path):
        """_load_branding() retourne les defaults si branding.json manquant."""
        with patch("utils.pptx_generator.BRANDING_FILE", tmp_path / "nonexistent.json"):
            result = _load_branding()
        assert "colors" in result
        assert "anthracite" in result["colors"]

    def test_returns_dict_when_branding_file_invalid(self, tmp_path):
        """_load_branding() retourne les defaults si JSON invalide."""
        bad_file = tmp_path / "branding.json"
        bad_file.write_text("NOT JSON {{{", encoding="utf-8")
        with patch("utils.pptx_generator.BRANDING_FILE", bad_file):
            result = _load_branding()
        assert "colors" in result

    def test_reads_custom_branding_file(self, tmp_path):
        """_load_branding() lit un branding.json valide."""
        custom_branding = {
            "colors": {"primary": "AABBCC"},
            "fonts": {"title": "CustomFont", "body": "AnotherFont"},
        }
        branding_file = tmp_path / "branding.json"
        branding_file.write_text(json.dumps(custom_branding), encoding="utf-8")
        with patch("utils.pptx_generator.BRANDING_FILE", branding_file):
            result = _load_branding()
        assert result["fonts"]["title"] == "CustomFont"


# ---------------------------------------------------------------------------
# _hex_to_rgb
# ---------------------------------------------------------------------------

class TestHexToRgb:
    def test_converts_hex_string(self):
        """_hex_to_rgb() convertit une chaine hex en RGBColor."""
        from pptx.dml.color import RGBColor
        result = _hex_to_rgb("FF6B58")
        assert isinstance(result, RGBColor)
        assert result[0] == 0xFF
        assert result[1] == 0x6B
        assert result[2] == 0x58

    def test_strips_hash_prefix(self):
        """_hex_to_rgb() gere le prefixe #."""
        from pptx.dml.color import RGBColor
        result = _hex_to_rgb("#FFFFFF")
        assert isinstance(result, RGBColor)

    def test_converts_black(self):
        """_hex_to_rgb('000000') retourne RGBColor(0, 0, 0)."""
        from pptx.dml.color import RGBColor
        result = _hex_to_rgb("000000")
        assert result[0] == 0
        assert result[1] == 0
        assert result[2] == 0

    def test_converts_white(self):
        """_hex_to_rgb('FFFFFF') retourne RGBColor(255, 255, 255)."""
        from pptx.dml.color import RGBColor
        result = _hex_to_rgb("FFFFFF")
        assert result[0] == 255
        assert result[1] == 255
        assert result[2] == 255


# ---------------------------------------------------------------------------
# Module-level constants
# ---------------------------------------------------------------------------

class TestModuleConstants:
    def test_branding_has_colors(self):
        assert "colors" in BRANDING

    def test_colors_dict_has_corail(self):
        assert "corail" in COLORS

    def test_colors_dict_has_blanc(self):
        assert "blanc" in COLORS

    def test_font_title_is_string(self):
        assert isinstance(FONT_TITLE, str)
        assert len(FONT_TITLE) > 0

    def test_font_body_is_string(self):
        assert isinstance(FONT_BODY, str)
        assert len(FONT_BODY) > 0

    def test_slide_width_is_set(self):
        from pptx.util import Inches
        assert SLIDE_WIDTH == Inches(13.333)

    def test_slide_height_is_set(self):
        from pptx.util import Inches
        assert SLIDE_HEIGHT == Inches(7.5)


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.__init__
# ---------------------------------------------------------------------------

class TestProposalPPTXGeneratorInit:
    def test_init_calls_presentation(self, tmp_path):
        """__init__ charge le template via Presentation()."""
        fake_template = str(tmp_path / "template.pptx")
        mock_prs = _make_presentation_mock()

        with patch("utils.pptx_generator.Presentation", return_value=mock_prs) as MockPrs:
            gen = ProposalPPTXGenerator(fake_template)
            MockPrs.assert_called_once_with(fake_template)

    def test_init_clears_existing_slides(self, tmp_path):
        """__init__ supprime les slides existantes du template."""
        fake_template = str(tmp_path / "template.pptx")
        mock_prs = _make_presentation_mock()

        # Simulate 2 existing slides
        mock_id1 = MagicMock()
        mock_id1.rId = "rId1"
        mock_id2 = MagicMock()
        mock_id2.rId = "rId2"
        sld_list = [mock_id1, mock_id2]
        mock_prs.slides._sldIdLst = sld_list

        # __len__ tracks the actual sld_list so the while loop terminates naturally
        mock_prs.slides.__len__ = MagicMock(side_effect=lambda: len(sld_list))

        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            gen = ProposalPPTXGenerator(fake_template)

        assert gen.template_path == fake_template
        assert gen.prs is mock_prs

    def test_init_stores_template_path(self, tmp_path):
        """__init__ stocke template_path."""
        fake_template = str(tmp_path / "template.pptx")
        mock_prs = _make_presentation_mock()

        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            gen = ProposalPPTXGenerator(fake_template)

        assert gen.template_path == fake_template


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator._set_text_style
# ---------------------------------------------------------------------------

class TestSetTextStyle:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx"))

    def test_sets_font_name(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.util import Pt
        gen._set_text_style(run, font_name="Arial", size=Pt(12))
        assert run.font.name == "Arial"

    def test_sets_font_size(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.util import Pt
        gen._set_text_style(run, size=Pt(18))
        assert run.font.size == Pt(18)

    def test_sets_color_when_provided(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.dml.color import RGBColor
        color = RGBColor(255, 0, 0)
        from pptx.util import Pt
        gen._set_text_style(run, color=color, size=Pt(12))
        assert run.font.color.rgb == color

    def test_skips_color_when_none(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.util import Pt
        gen._set_text_style(run, color=None, size=Pt(12))
        # color.rgb should NOT be set
        run.font.color.rgb.assert_not_called()

    def test_sets_bold(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.util import Pt
        gen._set_text_style(run, bold=True, size=Pt(12))
        assert run.font.bold is True

    def test_sets_italic(self, tmp_path):
        gen = self._make_gen(tmp_path)
        run = MagicMock()
        from pptx.util import Pt
        gen._set_text_style(run, italic=True, size=Pt(12))
        assert run.font.italic is True


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator._add_paragraph
# ---------------------------------------------------------------------------

class TestAddParagraph:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx"))

    def test_adds_paragraph_to_text_frame(self, tmp_path):
        gen = self._make_gen(tmp_path)
        tf = MagicMock()
        new_para = MagicMock()
        run = MagicMock()
        new_para.add_run.return_value = run
        tf.add_paragraph.return_value = new_para

        from pptx.util import Pt
        gen._add_paragraph(tf, "Hello World", size=Pt(14))

        tf.add_paragraph.assert_called_once()
        assert run.text == "Hello World"

    def test_sets_text_on_run(self, tmp_path):
        gen = self._make_gen(tmp_path)
        tf = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tf.add_paragraph.return_value = para

        from pptx.util import Pt
        result = gen._add_paragraph(tf, "Test Text", size=Pt(12))

        assert run.text == "Test Text"

    def test_returns_paragraph(self, tmp_path):
        gen = self._make_gen(tmp_path)
        tf = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tf.add_paragraph.return_value = para

        from pptx.util import Pt
        result = gen._add_paragraph(tf, "x", size=Pt(12))

        assert result is para


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_cover_slide
# ---------------------------------------------------------------------------

class TestAddCoverSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_add_cover_slide_adds_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock(placeholder_idxs=[0, 1])
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cover_slide("ClientCorp", "My Project", "2026-01-01", "Jane Doe")

        mock_prs.slides.add_slide.assert_called_once()

    def test_add_cover_slide_sets_title_text(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cover_slide("ClientCorp", "Project Title", "2026-01-01", "Consultant Name")

        # The first paragraph's run should have been set (PROPOSITION COMMERCIALE)
        tf_title = ph_title.text_frame
        assert tf_title.clear.called

    def test_add_cover_slide_uses_section_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cover_slide("C", "P", "D", "N")

        # LAYOUT_SECTION = 15
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[15])

    def test_add_cover_slide_sets_subtitle_text(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cover_slide("Corp", "Proj", "2026", "Jane")

        # Sub placeholder must have been cleared
        assert ph_sub.text_frame.clear.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_section_slide
# ---------------------------------------------------------------------------

class TestAddSectionSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_section_slide_adds_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock(placeholder_idxs=[0, 1])
        mock_prs.slides.add_slide.return_value = slide

        gen.add_section_slide("Notre Approche", 1)
        mock_prs.slides.add_slide.assert_called()

    def test_section_number_below_10_uses_leading_zero(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        # Capture runs added to paragraphs
        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph_title.text_frame.paragraphs = [para0]

        para1 = MagicMock()
        run1 = MagicMock()
        para1.add_run.return_value = run1
        ph_title.text_frame.add_paragraph.return_value = para1

        gen.add_section_slide("Section One", 3)
        assert run0.text == "03"

    def test_section_number_10_plus_no_leading_zero(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph_title.text_frame.paragraphs = [para0]

        para1 = MagicMock()
        run1 = MagicMock()
        para1.add_run.return_value = run1
        ph_title.text_frame.add_paragraph.return_value = para1

        gen.add_section_slide("Big Section", 12)
        assert run0.text == "12"

    def test_section_slide_clears_subtitle(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph_title.text_frame.paragraphs = [para0]
        ph_title.text_frame.add_paragraph.return_value = MagicMock()

        gen.add_section_slide("S", 1)
        assert ph_sub.text == ""


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_content_slide
# ---------------------------------------------------------------------------

class TestAddContentSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_content_slide_calls_add_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock(placeholder_idxs=[0, 1, 2])
        mock_prs.slides.add_slide.return_value = slide

        gen.add_content_slide("Titre", ["Point 1", "Point 2"])
        mock_prs.slides.add_slide.assert_called()

    def test_content_slide_uses_content_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock(placeholder_idxs=[0])
        mock_prs.slides.add_slide.return_value = slide

        gen.add_content_slide("T", [])
        # LAYOUT_CONTENT = 0
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[0])

    def test_content_slide_sets_title(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        slide = MagicMock()
        slide.placeholders = [ph_title]
        mock_prs.slides.add_slide.return_value = slide

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        ph_title.text_frame.paragraphs = [para]

        gen.add_content_slide("My Title", [])
        assert run.text == "My Title"

    def test_content_slide_sets_subtitle_when_provided(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        para_t = MagicMock()
        run_t = MagicMock()
        para_t.add_run.return_value = run_t
        ph_title.text_frame.paragraphs = [para_t]

        para_s = MagicMock()
        run_s = MagicMock()
        para_s.add_run.return_value = run_s
        ph_sub.text_frame.paragraphs = [para_s]

        gen.add_content_slide("Title", [], subtitle="My Subtitle")
        assert run_s.text == "My Subtitle"

    def test_content_slide_empty_subtitle_clears_placeholder(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        ph_sub = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph_title, ph_sub]
        mock_prs.slides.add_slide.return_value = slide

        para_t = MagicMock()
        run_t = MagicMock()
        para_t.add_run.return_value = run_t
        ph_title.text_frame.paragraphs = [para_t]

        gen.add_content_slide("T", [], subtitle="")
        assert ph_sub.text == ""

    def test_content_slide_adds_bullet_points(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_body = _make_placeholder(2)
        slide = MagicMock()
        slide.placeholders = [ph_body]
        mock_prs.slides.add_slide.return_value = slide

        para0 = MagicMock()
        run_bullet = MagicMock()
        run_text = MagicMock()
        para0.add_run.side_effect = [run_bullet, run_text]
        ph_body.text_frame.paragraphs = [para0]

        para_extra = MagicMock()
        run_b2 = MagicMock()
        run_t2 = MagicMock()
        para_extra.add_run.side_effect = [run_b2, run_t2]
        ph_body.text_frame.add_paragraph.return_value = para_extra

        gen.add_content_slide("T", ["Point A", "Point B"])
        assert run_text.text == "Point A"
        assert run_t2.text == "Point B"


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_two_column_slide
# ---------------------------------------------------------------------------

class TestAddTwoColumnSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_two_column_slide_uses_correct_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = _make_slide_mock(placeholder_idxs=[0, 1, 2, 3])
        mock_prs.slides.add_slide.return_value = slide

        gen.add_two_column_slide("T", "Left Title", ["L1"], "Right Title", ["R1"])
        # LAYOUT_TWO_COL = 4
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[4])

    def test_two_column_slide_sets_title(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        slide = MagicMock()
        slide.placeholders = [ph_title]
        mock_prs.slides.add_slide.return_value = slide

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        ph_title.text_frame.paragraphs = [para]

        gen.add_two_column_slide("Main Title", "L", [], "R", [])
        assert run.text == "Main Title"

    def test_two_column_slide_populates_left_column(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_left = _make_placeholder(2)
        slide = MagicMock()
        slide.placeholders = [ph_left]
        mock_prs.slides.add_slide.return_value = slide

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph_left.text_frame.paragraphs = [para0]

        extra_para = MagicMock()
        extra_run = MagicMock()
        extra_para.add_run.return_value = extra_run
        ph_left.text_frame.add_paragraph.return_value = extra_para

        gen.add_two_column_slide("T", "Left Header", ["Item L"], "R", [])
        assert run0.text == "Left Header"

    def test_two_column_slide_populates_right_column(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_right = _make_placeholder(3)
        slide = MagicMock()
        slide.placeholders = [ph_right]
        mock_prs.slides.add_slide.return_value = slide

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph_right.text_frame.paragraphs = [para0]

        extra_para = MagicMock()
        extra_run = MagicMock()
        extra_para.add_run.return_value = extra_run
        ph_right.text_frame.add_paragraph.return_value = extra_para

        gen.add_two_column_slide("T", "L", [], "Right Header", ["Item R"])
        assert run0.text == "Right Header"


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_table_slide
# ---------------------------------------------------------------------------

class TestAddTableSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_table_slide_adds_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        mock_para = MagicMock()
        mock_run = MagicMock()
        mock_para.add_run.return_value = mock_run
        slide.placeholders[0].text_frame.paragraphs = [mock_para]

        mock_table = MagicMock()
        header_cell = MagicMock()
        header_cell.text_frame.paragraphs = [MagicMock()]
        for p in header_cell.text_frame.paragraphs:
            p.runs = [MagicMock()]
        mock_table.cell.return_value = header_cell

        table_shape = MagicMock()
        table_shape.table = mock_table
        slide.shapes.add_table.return_value = table_shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_table_slide("Budget", ["Phase", "Cost"], [["Phase 1", "50k"]])
        mock_prs.slides.add_slide.assert_called()

    def test_table_slide_uses_title_only_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []

        mock_table = MagicMock()
        mock_table.cell.return_value = MagicMock()
        mock_table.cell.return_value.text_frame.paragraphs = []
        table_shape = MagicMock()
        table_shape.table = mock_table
        slide.shapes.add_table.return_value = table_shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_table_slide("T", ["H1"], [])
        # LAYOUT_TITLE_ONLY = 17
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[17])

    def test_table_slide_alternates_row_colors(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        mock_para = MagicMock()
        mock_run = MagicMock()
        mock_para.add_run.return_value = mock_run
        slide.placeholders[0].text_frame.paragraphs = [mock_para]

        mock_table = MagicMock()
        cells = {}
        for row_i in range(3):
            for col_i in range(2):
                cell = MagicMock()
                cell.text_frame.paragraphs = [MagicMock()]
                for p in cell.text_frame.paragraphs:
                    p.runs = [MagicMock()]
                cells[(row_i, col_i)] = cell

        mock_table.cell.side_effect = lambda r, c: cells[(r, c)]
        table_shape = MagicMock()
        table_shape.table = mock_table
        slide.shapes.add_table.return_value = table_shape
        mock_prs.slides.add_slide.return_value = slide

        headers = ["H1", "H2"]
        rows = [["A", "B"], ["C", "D"]]
        gen.add_table_slide("T", headers, rows)

        # Even rows (i=0) get rose_poudre fill; odd rows don't set fill solid
        # Check that fill was called on even row cell
        even_cell = cells[(1, 0)]  # row index 1 in table = data row 0 (even)
        assert even_cell.fill.solid.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_cv_slide
# ---------------------------------------------------------------------------

class TestAddCvSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_cv_slide_adds_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        bar_tf = MagicMock()
        bar_para = MagicMock()
        bar_run = MagicMock()
        bar_para.add_run.return_value = bar_run
        bar_tf.paragraphs = [bar_para]
        bar_shape = MagicMock()
        bar_shape.text_frame = bar_tf

        exp_tf = MagicMock()
        exp_para = MagicMock()
        exp_run = MagicMock()
        exp_para.add_run.return_value = exp_run
        exp_tf.paragraphs = [exp_para]
        exp_box = MagicMock()
        exp_box.text_frame = exp_tf

        skills_tf = MagicMock()
        skills_para = MagicMock()
        skills_run = MagicMock()
        skills_para.add_run.return_value = skills_run
        skills_tf.paragraphs = [skills_para]
        skills_box = MagicMock()
        skills_box.text_frame = skills_tf

        slide.shapes.add_shape.return_value = bar_shape
        slide.shapes.add_textbox.side_effect = [exp_box, skills_box]
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cv_slide("Alice", "Data Scientist", ["Led migration"], ["Python"])
        mock_prs.slides.add_slide.assert_called()

    def test_cv_slide_sets_title(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph_title = _make_placeholder(0)
        slide = MagicMock()
        slide.placeholders = [ph_title]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        ph_title.text_frame.paragraphs = [para]

        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cv_slide("Bob Smith", "Consultant", [], [])
        assert "Bob Smith" in run.text

    def test_cv_slide_uses_title_only_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_cv_slide("N", "T", [], [])
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[17])


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_image_slide
# ---------------------------------------------------------------------------

class TestAddImageSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_image_slide_adds_slide(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        tb = MagicMock()
        tb_para = MagicMock()
        tb_run = MagicMock()
        tb_para.add_run.return_value = tb_run
        tb.text_frame.paragraphs = [tb_para]
        slide.shapes.add_textbox.return_value = tb

        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("Image Title", "/nonexistent/img.png", caption="A caption")
        mock_prs.slides.add_slide.assert_called()

    def test_image_slide_uses_content_image_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("T", "/no/img.png")
        # LAYOUT_CONTENT_IMAGE = 1
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[1])

    def test_image_slide_adds_picture_if_file_exists(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        # Create a real (fake) image file
        img_file = tmp_path / "test.png"
        img_file.write_bytes(b"\x89PNG fake content")

        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("T", str(img_file))
        slide.shapes.add_picture.assert_called()

    def test_image_slide_skips_picture_if_file_missing(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("T", "/does/not/exist.png")
        slide.shapes.add_picture.assert_not_called()

    def test_image_slide_adds_caption(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []

        tb = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tb.text_frame.paragraphs = [para]
        extra_para = MagicMock()
        extra_run = MagicMock()
        extra_para.add_run.return_value = extra_run
        tb.text_frame.add_paragraph.return_value = extra_para
        slide.shapes.add_textbox.return_value = tb

        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("T", "/no/img.png", caption="My caption")
        assert run.text == "My caption"

    def test_image_slide_adds_bullets(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []

        tb = MagicMock()
        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        tb.text_frame.paragraphs = [para0]

        para1 = MagicMock()
        run1 = MagicMock()
        para1.add_run.return_value = run1
        tb.text_frame.add_paragraph.return_value = para1
        slide.shapes.add_textbox.return_value = tb

        mock_prs.slides.add_slide.return_value = slide

        gen.add_image_slide("T", "/no/img.png", bullets=["Bullet A"])
        # add_paragraph is called for each bullet
        assert tb.text_frame.add_paragraph.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_diagram_slide (flow)
# ---------------------------------------------------------------------------

class TestAddDiagramSlideFlow:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_flow_diagram_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Flow Diagram", "flow", ["Step A", "Step B", "Step C"])
        assert slide.shapes.add_shape.called

    def test_pyramid_diagram_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Pyramid", "pyramid", ["Top", "Middle", "Bottom"])
        assert slide.shapes.add_shape.called

    def test_cycle_diagram_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Cycle", "cycle", ["Step 1", "Step 2", "Step 3"])
        assert slide.shapes.add_shape.called

    def test_timeline_diagram_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Timeline", "timeline", ["2023", "2024", "2025"])
        assert slide.shapes.add_shape.called

    def test_matrix_diagram_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Matrix", "matrix", ["Q1", "Q2", "Q3", "Q4"])
        assert slide.shapes.add_shape.called

    def test_default_diagram_type_adds_shapes(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("Unknown", "unknown_type", ["A", "B"])
        assert slide.shapes.add_shape.called

    def test_diagram_with_description_adds_textbox(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]

        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        desc_tb = MagicMock()
        desc_para = MagicMock()
        desc_run = MagicMock()
        desc_para.add_run.return_value = desc_run
        desc_tb.text_frame.paragraphs = [desc_para]
        slide.shapes.add_textbox.return_value = desc_tb

        mock_prs.slides.add_slide.return_value = slide

        gen.add_diagram_slide("T", "flow", ["A"], description="Some description")
        slide.shapes.add_textbox.assert_called()


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_stat_slide
# ---------------------------------------------------------------------------

class TestAddStatSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_stat_slide_uses_blank_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape
        tb = MagicMock()
        tb.text_frame.paragraphs = [MagicMock()]
        tb.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = tb
        mock_prs.slides.add_slide.return_value = slide

        gen.add_stat_slide("67%", "de ROI")
        # LAYOUT_BLANK = 16
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[16])

    def test_stat_slide_with_subtitle_and_context(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []

        shape = MagicMock()
        shape.text_frame.paragraphs = [MagicMock()]
        shape.text_frame.paragraphs[0].add_run.return_value = MagicMock()
        slide.shapes.add_shape.return_value = shape

        tb = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tb.text_frame.paragraphs = [para]
        slide.shapes.add_textbox.return_value = tb

        mock_prs.slides.add_slide.return_value = slide

        gen.add_stat_slide("100M€", "CA", subtitle="KPI Principal", context="Contexte additionnel")
        # Both subtitle and context boxes should be added
        assert slide.shapes.add_textbox.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_quote_slide
# ---------------------------------------------------------------------------

class TestAddQuoteSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_quote_slide_uses_blank_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        shape = MagicMock()
        slide.shapes.add_shape.return_value = shape
        tb = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tb.text_frame.paragraphs = [para]
        slide.shapes.add_textbox.return_value = tb
        mock_prs.slides.add_slide.return_value = slide

        gen.add_quote_slide("Life is what happens")
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[16])

    def test_quote_slide_with_author(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_shape.return_value = MagicMock()
        tb = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tb.text_frame.paragraphs = [para]
        slide.shapes.add_textbox.return_value = tb
        mock_prs.slides.add_slide.return_value = slide

        gen.add_quote_slide("To be or not to be", author="Shakespeare")
        assert slide.shapes.add_textbox.called

    def test_quote_slide_with_title(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = []
        slide.shapes.add_shape.return_value = MagicMock()
        tb = MagicMock()
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        tb.text_frame.paragraphs = [para]
        slide.shapes.add_textbox.return_value = tb
        mock_prs.slides.add_slide.return_value = slide

        gen.add_quote_slide("Quote text", title="Vision strategique")
        assert slide.shapes.add_textbox.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_highlight_slide
# ---------------------------------------------------------------------------

class TestAddHighlightSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_highlight_slide_with_2_points(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]
        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_highlight_slide("T", ["Point A", "Point B"])
        assert slide.shapes.add_shape.called

    def test_highlight_slide_with_3_points(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]
        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_highlight_slide("T", ["P1", "P2", "P3"])
        assert slide.shapes.add_shape.called

    def test_highlight_slide_with_4_points(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]
        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_highlight_slide("T", ["P1", "P2", "P3", "P4"])
        assert slide.shapes.add_shape.called

    def test_highlight_slide_custom_color(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]
        slide.shapes.add_shape.return_value = MagicMock()
        slide.shapes.add_textbox.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_highlight_slide("T", ["P1"], highlight_color="terracotta")
        assert slide.shapes.add_shape.called


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.add_closing_slide
# ---------------------------------------------------------------------------

class TestAddClosingSlide:
    def _make_gen(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            return ProposalPPTXGenerator(str(tmp_path / "t.pptx")), mock_prs

    def test_closing_slide_uses_section_layout(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)
        slide = MagicMock()
        slide.placeholders = [_make_placeholder(0), _make_placeholder(1)]
        para = MagicMock()
        run = MagicMock()
        para.add_run.return_value = run
        slide.placeholders[0].text_frame.paragraphs = [para]
        slide.placeholders[0].text_frame.add_paragraph.return_value = MagicMock()
        mock_prs.slides.add_slide.return_value = slide

        gen.add_closing_slide("Jane", "Consulting Tools")
        mock_prs.slides.add_slide.assert_called_with(mock_prs.slide_layouts[15])

    def test_closing_slide_sets_merci(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph0 = _make_placeholder(0)
        ph1 = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph0, ph1]

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph0.text_frame.paragraphs = [para0]

        para1 = MagicMock()
        run1 = MagicMock()
        para1.add_run.return_value = run1
        ph0.text_frame.add_paragraph.return_value = para1

        mock_prs.slides.add_slide.return_value = slide

        gen.add_closing_slide("Consultant Name", "Company")
        assert run0.text == "MERCI"

    def test_closing_slide_with_contact_info(self, tmp_path):
        gen, mock_prs = self._make_gen(tmp_path)

        ph0 = _make_placeholder(0)
        ph1 = _make_placeholder(1)
        slide = MagicMock()
        slide.placeholders = [ph0, ph1]

        para0 = MagicMock()
        run0 = MagicMock()
        para0.add_run.return_value = run0
        ph0.text_frame.paragraphs = [para0]

        para1 = MagicMock()
        run1 = MagicMock()
        para1.add_run.return_value = run1

        para2 = MagicMock()
        run2 = MagicMock()
        para2.add_run.return_value = run2

        ph0.text_frame.add_paragraph.side_effect = [para1, para2]
        mock_prs.slides.add_slide.return_value = slide

        gen.add_closing_slide("Consultant", "Corp", contact_info="contact@corp.com")
        # contact_info paragraph added
        assert ph0.text_frame.add_paragraph.call_count >= 2


# ---------------------------------------------------------------------------
# ProposalPPTXGenerator.save
# ---------------------------------------------------------------------------

class TestSave:
    def test_save_calls_prs_save(self, tmp_path):
        mock_prs = _make_presentation_mock()
        with patch("utils.pptx_generator.Presentation", return_value=mock_prs):
            gen = ProposalPPTXGenerator(str(tmp_path / "t.pptx"))

        output = str(tmp_path / "output.pptx")
        gen.save(output)
        mock_prs.save.assert_called_once_with(output)


# ---------------------------------------------------------------------------
# build_proposal_pptx (main orchestration function)
# ---------------------------------------------------------------------------

class TestBuildProposalPptx:
    def _make_gen_mock(self):
        """Return a mock ProposalPPTXGenerator."""
        gen = MagicMock(spec=ProposalPPTXGenerator)
        return gen

    def test_returns_output_path(self, tmp_path):
        """build_proposal_pptx retourne le chemin de sortie."""
        output = str(tmp_path / "output.pptx")

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            result = build_proposal_pptx(
                template_path="fake_template.pptx",
                slides_data=[],
                output_path=output,
                consultant_info={"name": "Jane", "company": "Corp"},
            )

        assert result == output

    def test_calls_save_at_end(self, tmp_path):
        output = str(tmp_path / "out.pptx")

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", [], output, {"name": "N"})
            mock_gen.save.assert_called_once_with(output)

    def test_cover_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "cover", "client": "Acme", "project": "Transform", "date": "2026"}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_cover_slide.assert_called_once_with(
                client_name="Acme",
                project_title="Transform",
                date="2026",
                consultant_name="Jane",
            )

    def test_section_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "section", "title": "Context", "number": 1}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_section_slide.assert_called_once_with(
                section_title="Context",
                section_number=1,
            )

    def test_content_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "content", "title": "Overview", "bullets": ["A", "B"], "subtitle": ""}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_content_slide.assert_called_once_with(
                title="Overview",
                bullet_points=["A", "B"],
                subtitle="",
            )

    def test_two_column_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "two_column",
            "title": "Comparison",
            "left_title": "Left",
            "left_points": ["L1"],
            "right_title": "Right",
            "right_points": ["R1"],
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_two_column_slide.assert_called_once()

    def test_table_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "table", "title": "Budget", "headers": ["Phase", "Cost"], "rows": []}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_table_slide.assert_called_once_with(
                title="Budget",
                headers=["Phase", "Cost"],
                rows=[],
            )

    def test_cv_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "cv",
            "name": "Alice",
            "title": "Engineer",
            "experiences": ["Exp A"],
            "skills": ["Python"],
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_cv_slide.assert_called_once_with(
                name="Alice",
                title="Engineer",
                experiences=["Exp A"],
                skills=["Python"],
            )

    def test_image_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "image",
            "title": "Architecture",
            "image_path": "/img.png",
            "caption": "Cap",
            "bullets": ["B1"],
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_image_slide.assert_called_once_with(
                title="Architecture",
                image_path="/img.png",
                caption="Cap",
                bullets=["B1"],
            )

    def test_diagram_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "diagram",
            "title": "Process",
            "diagram_type": "flow",
            "elements": ["A", "B"],
            "description": "desc",
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_diagram_slide.assert_called_once_with(
                title="Process",
                diagram_type="flow",
                elements=["A", "B"],
                description="desc",
            )

    def test_stat_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "stat",
            "stat_value": "67%",
            "stat_label": "ROI",
            "context": "ctx",
            "subtitle": "sub",
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_stat_slide.assert_called_once_with(
                stat_value="67%",
                stat_label="ROI",
                context="ctx",
                subtitle="sub",
            )

    def test_quote_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "quote", "quote_text": "To be...", "author": "W.S.", "title": "Quote"}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_quote_slide.assert_called_once_with(
                quote_text="To be...",
                author="W.S.",
                title="Quote",
            )

    def test_highlight_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{
            "type": "highlight",
            "title": "Key Points",
            "key_points": ["P1", "P2"],
            "highlight_color": "corail",
        }]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_highlight_slide.assert_called_once_with(
                title="Key Points",
                key_points=["P1", "P2"],
                highlight_color="corail",
            )

    def test_closing_slide_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "closing"}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane", "company": "Corp"})
            mock_gen.add_closing_slide.assert_called_once_with(
                consultant_name="Jane",
                company="Corp",
            )

    def test_unknown_slide_type_skipped(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [{"type": "UNKNOWN_TYPE"}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            result = build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})

        assert result == output
        mock_gen.save.assert_called_once()

    def test_multiple_slides_dispatched(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        slides = [
            {"type": "cover", "client": "X", "project": "Y", "date": "D"},
            {"type": "section", "title": "S1", "number": 1},
            {"type": "content", "title": "C1", "bullets": []},
            {"type": "closing"},
        ]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane", "company": "Corp"})

            mock_gen.add_cover_slide.assert_called_once()
            mock_gen.add_section_slide.assert_called_once()
            mock_gen.add_content_slide.assert_called_once()
            mock_gen.add_closing_slide.assert_called_once()

    def test_empty_slides_data(self, tmp_path):
        output = str(tmp_path / "out.pptx")

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            result = build_proposal_pptx("t.pptx", [], output, {})

        assert result == output
        mock_gen.save.assert_called_once_with(output)

    def test_slide_defaults_used_when_keys_missing(self, tmp_path):
        output = str(tmp_path / "out.pptx")
        # A content slide with no 'bullets' or 'subtitle' keys
        slides = [{"type": "content"}]

        with patch("utils.pptx_generator.ProposalPPTXGenerator") as MockGen:
            mock_gen = self._make_gen_mock()
            MockGen.return_value = mock_gen

            build_proposal_pptx("t.pptx", slides, output, {"name": "Jane"})
            mock_gen.add_content_slide.assert_called_once_with(
                title="",
                bullet_points=[],
                subtitle="",
            )
