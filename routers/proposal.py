"""Router: Proposal Generator — Génération, modulation et export de propositions commerciales"""

import asyncio
import io
import json
import os
import re
import threading
import uuid
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, File, Form, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, StreamingResponse
from PyPDF2 import PdfReader

from agents.proposal_generator import ProposalGeneratorAgent
from utils.validation import sanitize_text_input
from routers.shared import (
    BASE_DIR,
    jobs,
    limiter,
    templates,
    CONSULTANT_NAME,
    COMPANY_NAME,
    safe_error_message,
    safe_traceback,
    send_sse,
)

router = APIRouter()


# === API: PROPOSAL ===


@router.post("/api/proposal/generate")
@limiter.limit("5/minute")
async def api_proposal_generate(
    request: Request,
    tender_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Lance la generation d'une proposition commerciale"""
    text = ""

    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            # Extraire le texte du PDF
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif tender_text:
        text = tender_text
    else:
        return JSONResponse({"error": "Aucun appel d'offre fourni."}, status_code=400)

    if len(text.strip()) < 50:
        return JSONResponse(
            {"error": "L'appel d'offre semble trop court (min 50 caracteres)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "tender_text": text,
    }

    # Lancer en background
    thread = threading.Thread(target=_run_proposal, args=(job_id, text), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_proposal(job_id: str, tender_text: str):
    """Execute la generation de proposition en background"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Step 1: Analyse
        job["steps"].append({"step": "analyze", "status": "active", "progress": 5})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 15})

        # Step 2: Template
        job["steps"].append({"step": "template", "status": "active", "progress": 18})
        template = agent.load_template()
        job["steps"].append({"step": "template", "status": "done", "progress": 25})

        # Step 3: References
        job["steps"].append({"step": "references", "status": "active", "progress": 28})
        references = agent.match_references(tender_analysis)
        job["steps"].append({"step": "references", "status": "done", "progress": 45})

        # Step 4: CVs
        job["steps"].append({"step": "generate", "status": "active", "progress": 48})
        cvs = agent.load_cvs()
        adapted_cvs = agent.adapt_cvs(cvs, tender_analysis) if cvs else []

        # Step 5: Generate slides structure + markdown content
        slides = agent.generate_slides_structure(tender_analysis, references, template, adapted_cvs)
        proposal = agent.generate_proposal_content(tender_analysis, references, template)
        job["steps"].append({"step": "generate", "status": "done", "progress": 85})

        # Sauvegarder les fichiers
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        json_path = output_dir / f"proposal_{client_name}_{timestamp}.json"
        md_path = output_dir / f"proposal_{client_name}_{timestamp}.md"
        pptx_path = output_dir / f"proposal_{client_name}_{timestamp}.pptx"

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(proposal, f, indent=2, ensure_ascii=False)

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(
                f"# Proposition Commerciale - {tender_analysis.get('project_title', 'Projet')}\n\n"
            )
            f.write(f"**Client:** {tender_analysis.get('client_name', 'N/A')}\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y')}\n\n---\n\n")
            f.write(proposal["content"])

        # Generer le PPTX
        try:
            from utils.pptx_generator import build_proposal_pptx

            build_proposal_pptx(
                template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                slides_data=slides,
                output_path=str(pptx_path),
                consultant_info={
                    "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                    "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
                },
            )
        except Exception as e:
            print(f"   ⚠️  Erreur PPTX: {e}")
            pptx_path = None

        job["status"] = "done"
        job["result"] = {
            "content": proposal["content"],
            "client_name": tender_analysis.get("client_name", "Projet"),
            "project_title": tender_analysis.get("project_title", ""),
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "json_path": str(json_path.relative_to(BASE_DIR)),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)) if pptx_path else None,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@router.get("/api/proposal/stream/{job_id}")
async def api_proposal_stream(job_id: str):
    """SSE stream pour la progression de la proposition"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            # Envoyer les nouvelles etapes
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: PROPOSAL FEEDBACK ===


@router.post("/api/proposal/regenerate")
async def api_proposal_regenerate(request: Request):
    """Regenere la proposition avec le feedback utilisateur"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_content = body.get("previous_content", "")
    job_id_ref = body.get("job_id", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    # Recuperer le tender_text du job original si possible
    tender_text = ""
    if job_id_ref and job_id_ref in jobs:
        tender_text = jobs[job_id_ref].get("tender_text", "")

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "tender_text": tender_text,
    }

    thread = threading.Thread(
        target=_run_proposal_feedback,
        args=(job_id, previous_content, feedback, tender_text),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


# === API: PROPOSAL MODULAR GENERATION ===


@router.post("/api/proposal/generate-section")
@limiter.limit("10/minute")
async def api_proposal_generate_section(
    request: Request,
    section: str = Form(...),
    tender_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Génère une section spécifique de la proposition"""
    # Récupérer le texte de l'appel d'offre
    text = ""
    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif tender_text:
        text = tender_text
    else:
        return JSONResponse({"error": "Aucun appel d'offre fourni."}, status_code=400)

    if len(text.strip()) < 50:
        return JSONResponse({"error": "L'appel d'offre semble trop court."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal_section",
        "section": section,
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_proposal_section,
        args=(job_id, text, section),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_proposal_section(job_id: str, tender_text: str, section: str):
    """Génère une section spécifique de la proposition"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Analyser l'appel d'offre
        job["steps"].append({"step": "analyze", "status": "active", "progress": 10})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 30})

        # Générer la section demandée
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        slides = []
        section_name = ""

        if section == "agenda":
            slides = agent.generate_agenda_slide(tender_analysis)
            section_name = "Agenda"
        elif section == "context":
            slides = agent.generate_context_slides(tender_analysis)
            section_name = "Contexte"
        elif section == "approach":
            references = agent.match_references(tender_analysis)
            slides = agent.generate_approach_slides(tender_analysis, references)
            section_name = "Approche"
        elif section == "planning":
            slides = agent.generate_planning_slide(tender_analysis)
            section_name = "Planning"
        elif section == "budget":
            slides = agent.generate_budget_slide(tender_analysis)
            section_name = "Chiffrage"
        elif section == "references":
            references = agent.load_references()
            slides = agent.generate_references_slides(
                tender_analysis, {"all_references": references}
            )
            section_name = "Références"
        elif section == "cvs":
            slides = agent.generate_cv_slides(tender_analysis)
            section_name = "CVs"
        else:
            raise ValueError(f"Section inconnue: {section}")

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Générer le PPTX SANS cover ni closing (mode modulaire)
        job["steps"].append({"step": "pptx", "status": "active", "progress": 85})

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}.pptx"

        # Mode modulaire : UNIQUEMENT les slides de la section (sans cover ni closing)
        # L'utilisateur assemble ensuite les sections manuellement
        full_slides = slides

        from utils.pptx_generator import build_proposal_pptx

        build_proposal_pptx(
            template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
            slides_data=full_slides,
            output_path=str(pptx_path),
            consultant_info={
                "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
            },
        )

        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "section": section,
            "section_name": section_name,
            "slides_count": len(slides),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)),
            "tender_text": tender_text,  # Inclure le texte pour la régénération
            "slides_data": slides,  # Inclure les données des slides
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@router.post("/api/proposal/preview")
@limiter.limit("20/minute")
async def api_proposal_preview(
    request: Request,
    pptx_path: str = Form(...),
):
    """Génère des images de prévisualisation des slides"""
    try:
        import subprocess
        import tempfile

        from PIL import Image

        # Résoudre le chemin absolu du PPTX
        file_path = BASE_DIR / pptx_path

        if not file_path.exists():
            return JSONResponse({"error": "Fichier PPTX introuvable"}, status_code=404)

        # Créer un répertoire temporaire pour les images
        temp_dir = Path(tempfile.mkdtemp(dir=BASE_DIR / "output"))

        try:
            # Convertir PPTX en PNG avec LibreOffice (méthode la plus fiable)
            # Note: LibreOffice doit être installé sur le système
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "png",
                    "--outdir",
                    str(temp_dir),
                    str(file_path),
                ],
                check=True,
                timeout=30,
                capture_output=True,
            )

            # Récupérer toutes les images générées
            preview_images = sorted(temp_dir.glob("*.png"))

            if not preview_images:
                # Fallback: essayer avec python-pptx + Pillow (moins fiable
                # mais sans dépendance externe)
                from pptx import Presentation

                prs = Presentation(str(file_path))
                preview_images = []

                for i, slide in enumerate(prs.slides):
                    # Cette approche est limitée car python-pptx ne rend pas visuellement
                    # On crée juste un placeholder pour indiquer que la
                    # conversion a échoué
                    pass

                if not preview_images:
                    return JSONResponse(
                        {
                            "error": "Impossible de générer la prévisualisation. LibreOffice n'est pas installé.",
                            "preview_images": [],
                        },
                        status_code=500,
                    )

            # Créer des versions miniatures pour la galerie
            thumbnails = []
            for img_path in preview_images:
                # Copier vers le dossier output avec un nom unique
                thumb_name = f"preview_{uuid.uuid4().hex[:8]}_{img_path.name}"
                thumb_path = BASE_DIR / "output" / thumb_name

                # Créer une miniature (optionnel, pour optimiser le chargement)
                with Image.open(img_path) as img:
                    img.thumbnail((800, 600))
                    img.save(thumb_path, "PNG")

                thumbnails.append(str(Path("output") / thumb_name))

                # Nettoyer le fichier temporaire
                img_path.unlink()

            # Nettoyer le répertoire temporaire
            temp_dir.rmdir()

            return JSONResponse({"preview_images": thumbnails, "slide_count": len(thumbnails)})

        except subprocess.CalledProcessError:
            # LibreOffice n'est pas disponible ou a échoué
            return JSONResponse(
                {
                    "error": "La génération de prévisualisation nécessite LibreOffice. Installez-le avec: brew install libreoffice (Mac) ou apt-get install libreoffice (Linux)",
                    "preview_images": [],
                },
                status_code=500,
            )
        except Exception as e:
            return JSONResponse(
                {
                    "error": f"Erreur lors de la génération de prévisualisation: {str(e)}",
                    "preview_images": [],
                },
                status_code=500,
            )

    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@router.post("/api/proposal/regenerate-section")
@limiter.limit("10/minute")
async def api_proposal_regenerate_section(
    request: Request,
    section: str = Form(...),
    feedback: str = Form(...),
    tender_text: Optional[str] = Form(None),
    # JSON string des slides précédentes
    previous_slides: Optional[str] = Form(None),
):
    """Régénère une section avec feedback en langage naturel"""

    if not feedback.strip():
        return JSONResponse({"error": "Le feedback ne peut pas être vide."}, status_code=400)

    if not tender_text or len(tender_text.strip()) < 50:
        return JSONResponse(
            {"error": "Le texte de l'appel d'offre est requis pour la régénération."},
            status_code=400,
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal_section",
        "section": section,
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    # Parse previous_slides si fourni
    prev_slides = []
    if previous_slides:
        try:
            prev_slides = json.loads(previous_slides)
        except BaseException:
            pass

    thread = threading.Thread(
        target=_run_proposal_section_with_feedback,
        args=(job_id, tender_text, section, feedback, prev_slides),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_proposal_section_with_feedback(
    job_id: str, tender_text: str, section: str, feedback: str, previous_slides: List[dict]
):
    """Régénère une section en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Analyser l'appel d'offre
        job["steps"].append({"step": "analyze", "status": "active", "progress": 10})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 30})

        # Interpréter le feedback avec le LLM
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        # Construire le contexte des slides précédentes
        previous_context = ""
        if previous_slides:
            previous_context = """
SLIDES PRÉCÉDENTES:
{json.dumps(previous_slides, indent=2, ensure_ascii=False)[:2000]}
"""

        # Utiliser le LLM pour interpréter le feedback et générer des
        # instructions structurées
        interpretation = agent.llm_client.generate(
            prompt="""L'utilisateur veut modifier la section "{section}" de sa proposition commerciale.

{previous_context}

FEEDBACK UTILISATEUR (en langage naturel):
{feedback}

Analyse ce feedback et génère des instructions STRUCTURÉES pour régénérer la section.
Réponds en JSON avec ces champs:
{
                "changes_requested": ["Liste des modifications demandées"],
  "visual_changes": ["Modifications visuelles (diagrammes, mise en page)"],
  "content_changes": ["Modifications de contenu (texte, chiffres, structure)"],
  "tone_changes": "Changement de ton si demandé",
  "specific_instructions": "Instructions précises pour la régénération"
}

Exemples:
- "Planning trop court" → ajouter des phases, augmenter les durées
- "Moins de bullets" → réduire à 2-3 bullets par slide
- "Diagramme flow → cycle" → remplacer le type de diagramme
- "Augmenter budget 20%" → multiplier les montants par 1.2""",
            system_prompt="Tu es un expert en interprétation de feedback pour la génération de documents. Sois précis et actionnable.",
            temperature=0.4,
            max_tokens=1500,
        )

        # Parser l'interprétation
        try:
            feedback_data = json.loads(interpretation.strip())
        except BaseException:
            feedback_data = {"specific_instructions": feedback, "changes_requested": [feedback]}

        # Générer la section avec les instructions du feedback
        slides = []
        section_name = ""
        references = None

        # Enrichir tender_analysis avec les instructions du feedback
        tender_analysis["feedback_instructions"] = feedback_data.get(
            "specific_instructions", feedback
        )
        tender_analysis["feedback_data"] = feedback_data

        if section == "agenda":
            slides = agent.generate_agenda_slide(tender_analysis)
            section_name = "Agenda"
        elif section == "context":
            slides = agent.generate_context_slides(tender_analysis)
            section_name = "Contexte"
        elif section == "approach":
            references = agent.match_references(tender_analysis)
            slides = agent.generate_approach_slides(tender_analysis, references)
            section_name = "Approche"
        elif section == "planning":
            slides = agent.generate_planning_slide(tender_analysis)
            section_name = "Planning"
        elif section == "budget":
            slides = agent.generate_budget_slide(tender_analysis)
            section_name = "Chiffrage"
        elif section == "references":
            references = agent.load_references()
            slides = agent.generate_references_slides(
                tender_analysis, {"all_references": references}
            )
            section_name = "Références"
        elif section == "cvs":
            slides = agent.generate_cv_slides(tender_analysis)
            section_name = "CVs"
        else:
            raise ValueError(f"Section inconnue: {section}")

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Générer le PPTX
        job["steps"].append({"step": "pptx", "status": "active", "progress": 85})

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}_v2.pptx"

        full_slides = slides

        from utils.pptx_generator import build_proposal_pptx

        build_proposal_pptx(
            template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
            slides_data=full_slides,
            output_path=str(pptx_path),
            consultant_info={
                "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
            },
        )

        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "section": section,
            "section_name": section_name,
            "slides_count": len(slides),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)),
            "feedback_applied": feedback_data.get("changes_requested", [feedback]),
            "slides_data": slides,  # Inclure les données des slides pour la prochaine itération
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)
        print(f"Error in regenerate section: {safe_traceback()}")


@router.post("/api/proposal/regenerate-slide")
@limiter.limit("20/minute")
async def api_proposal_regenerate_slide(
    request: Request,
    section: str = Form(...),
    slide_index: int = Form(...),
    slide_title: str = Form(...),
    slide_bullets: str = Form(...),  # JSON array
    diagram_type: Optional[str] = Form(None),
    tender_text: Optional[str] = Form(None),
):
    """Régénère une seule slide avec le nouveau contenu"""

    if not tender_text or len(tender_text.strip()) < 50:
        return JSONResponse({"error": "Le texte de l'appel d'offre est requis."}, status_code=400)

    try:
        # Parse bullets
        bullets = json.loads(slide_bullets)

        # Charger le PPTX existant pour récupérer toutes les slides
        # On va chercher le fichier le plus récent pour cette section
        output_dir = BASE_DIR / "output"
        matching_files = sorted(
            output_dir.glob(f"proposal_{section}_*.pptx"),
            key=lambda p: p.stat().st_mtime,
            reverse=True,
        )

        if not matching_files:
            return JSONResponse(
                {"error": "Aucun PPTX existant trouvé pour cette section"}, status_code=404
            )

        latest_pptx = matching_files[0]

        # Charger les slides depuis le PPTX existant
        from pptx import Presentation

        prs = Presentation(str(latest_pptx))

        # Vérifier que l'index est valide (en tenant compte de la slide de
        # couverture)
        if slide_index + 1 >= len(prs.slides):
            return JSONResponse(
                {"error": f"Index de slide invalide: {slide_index}"}, status_code=400
            )

        # Modifier la slide ciblée
        # +1 car slide 0 = couverture
        target_slide = prs.slides[slide_index + 1]

        # Trouver les textframes et mettre à jour
        for shape in target_slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                # Le premier textframe est généralement le titre
                if shape.name.startswith("Title") or "titre" in shape.name.lower():
                    text_frame.text = slide_title
                    break

        # Mettre à jour le contenu (bullets)
        for shape in target_slide.shapes:
            if shape.has_text_frame and not (
                shape.name.startswith("Title") or "titre" in shape.name.lower()
            ):
                text_frame = shape.text_frame
                text_frame.clear()

                # Ajouter les nouveaux bullets
                for i, bullet in enumerate(bullets):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

                    # Styling (utiliser les styles existants si possible)
                    from pptx.util import Pt

                    if p.runs:
                        p.runs[0].font.size = Pt(11)

                break

        # Sauvegarder le PPTX modifié
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        agent = ProposalGeneratorAgent()
        tender_analysis = agent.analyze_tender(tender_text)
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")

        new_pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}_edited.pptx"
        prs.save(str(new_pptx_path))

        # Extraire les données des slides pour le frontend
        slides_data = []
        for i, slide in enumerate(prs.slides):
            if i == 0:  # Skip cover
                continue

            slide_data = {"title": "", "bullets": [], "visual": None}

            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.name.startswith("Title") or "titre" in shape.name.lower():
                        slide_data["title"] = shape.text_frame.text
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_data["bullets"].append(paragraph.text.strip())

            slides_data.append(slide_data)

        return JSONResponse(
            {
                "pptx_path": str(new_pptx_path.relative_to(BASE_DIR)),
                "slides_data": slides_data,
                "message": "Slide régénérée avec succès",
            }
        )

    except Exception as e:
        print(f"Error in regenerate slide: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@router.post("/api/proposal/export-to-slides")
@limiter.limit("5/minute")
async def api_export_to_slides(
    request: Request,
    section: str = Form(...),
    tender_text: Optional[str] = Form(None),
    slides_data: str = Form(...),  # JSON string
):
    """Exporte une section vers Google Slides"""

    try:
        # Parser les données des slides avec sanitization robuste
        # Fix: le JSON envoyé depuis le frontend peut contenir des caractères de contrôle
        # ou des chaînes mal échappées, causant "Unterminated string"
        import re

        sanitized = slides_data.strip()
        # Supprimer les caractères de contrôle (sauf \n, \t, \r qui sont gérés
        # par JSON)
        sanitized = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", sanitized)
        # Remplacer les vrais retours à la ligne dans les valeurs string par \\n
        # (ceux qui ne sont pas déjà échappés)
        sanitized = sanitized.replace("\r\n", "\\n").replace("\r", "\\n")

        try:
            slides = json.loads(sanitized)
        except json.JSONDecodeError:
            # Fallback: essayer de nettoyer plus agressivement
            # Parfois les guillemets internes ne sont pas échappés
            sanitized = sanitized.replace("\t", "\\t")
            # Tenter un nettoyage ligne par ligne
            lines = sanitized.split("\n")
            cleaned_lines = []
            for line in lines:
                # Ne pas modifier les lignes structurelles JSON
                cleaned_lines.append(line)
            sanitized = "\n".join(cleaned_lines)
            slides = json.loads(sanitized)

        # Analyser le tender pour obtenir le titre
        if tender_text and len(tender_text.strip()) >= 50:
            agent = ProposalGeneratorAgent()
            tender_analysis = agent.analyze_tender(tender_text)
            client_name = tender_analysis.get("client_name", "Client")
            project_title = tender_analysis.get("project_title", "Projet")
        else:
            client_name = "Client"
            project_title = "Projet"

        # Générer le titre de la présentation
        section_names = {
            "agenda": "Agenda",
            "context": "Contexte",
            "approach": "Approche",
            "planning": "Planning",
            "budget": "Chiffrage",
            "references": "Références",
            "cvs": "CVs",
        }

        section_name = section_names.get(section, section.capitalize())
        presentation_title = f"{client_name} - {project_title} - {section_name}"

        # Exporter vers Google Slides
        from utils.google_api import GoogleAPIClient

        try:
            google_client = GoogleAPIClient()
            presentation_id = google_client.export_pptx_to_slides(
                slides_data=slides, title=presentation_title
            )

            if presentation_id:
                presentation_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"

                return JSONResponse(
                    {
                        "presentation_id": presentation_id,
                        "presentation_url": presentation_url,
                        "message": "Présentation Google Slides créée avec succès",
                    }
                )
            else:
                return JSONResponse(
                    {"error": "Échec de la création de la présentation Google Slides"},
                    status_code=500,
                )

        except Exception as e:
            # Si l'authentification Google échoue ou n'est pas configurée
            if "credentials" in str(e).lower() or "token" in str(e).lower():
                return JSONResponse(
                    {
                        "error": "Google API non configurée. Veuillez configurer vos credentials Google dans config/google_credentials.json",
                        "setup_required": True,
                    },
                    status_code=400,
                )
            else:
                raise

    except Exception as e:
        print(f"Error in export to slides: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


def _run_proposal_feedback(job_id: str, previous_content: str, feedback: str, tender_text: str):
    """Regenere la proposition en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        job["steps"].append({"step": "analyze", "status": "done", "progress": 10})
        job["steps"].append({"step": "template", "status": "done", "progress": 20})
        job["steps"].append({"step": "references", "status": "done", "progress": 30})
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        # Analyser le tender_text si disponible pour avoir plus de contexte
        tender_context = ""
        if tender_text:
            try:
                agent.analyze_tender(tender_text)
                tender_context = """
Contexte de l'appel d'offre:
Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {', '.join(tender_analysis.get('objectives', [])[:3])}
"""
            except Exception:
                tender_context = f"Appel d'offre: {tender_text[:1000]}"

        # Regenerer le contenu markdown avec feedback
        revised_content = agent.llm_client.generate(
            prompt="""Voici une proposition commerciale generee precedemment:

{previous_content[:6000]}

FEEDBACK DE L'UTILISATEUR:
{feedback}

{tender_context}

Reecris la proposition en integrant les corrections demandees dans le feedback.
Conserve la structure et le style professionnel. Ne modifie que ce qui est demande dans le feedback.
Sois concret et precis. Evite les majuscules inutiles et les emojis.
Fournis la proposition complete revisee.""",
            system_prompt="""Tu es {agent.consultant_info['name']}, {agent.consultant_info['title']} chez {agent.consultant_info['company']}.
Tu corriges une proposition commerciale selon le retour du consultant.
Reponds en francais de maniere professionnelle.""",
            temperature=0.6,
            max_tokens=8000,
        )

        job["steps"].append({"step": "generate", "status": "done", "progress": 60})
        job["steps"].append({"step": "slides", "status": "active", "progress": 65})

        # Regenerer les slides avec feedback - avec meilleur prompt
        revised_slides = None
        try:
            slides_response = agent.llm_client.generate(
                prompt="""Genere une structure complete de slides PowerPoint pour cette proposition commerciale revisee.

PROPOSITION REVISEE:
{revised_content[:5000]}

FEEDBACK APPLIQUE:
{feedback}

{tender_context}

Genere la structure COMPLETE avec tout le contenu. Reponds UNIQUEMENT en JSON valide:

[
  {{"type":"cover","client":"...","project":"...","date":"{datetime.now().strftime('%d/%m/%Y')}"}},
  {{"type":"section","title":"Notre comprehension du contexte","number":1}},
  {{"type":"content","title":"Contexte et enjeux","bullets":["Bullet 1","Bullet 2","Bullet 3","Bullet 4"],"subtitle":""}},
  {{"type":"content","title":"Objectifs","bullets":["Objectif 1","Objectif 2","Objectif 3"]}},
  {{"type":"section","title":"Notre expertise","number":2}},
  {{"type":"content","title":"Notre valeur ajoutee","bullets":["Valeur 1","Valeur 2","Valeur 3","Valeur 4"]}},
  {{"type":"section","title":"Notre demarche","number":3}},
  {{"type":"content","title":"Phases du projet","bullets":["Phase 1: Description","Phase 2: Description","Phase 3: Description"]}},
  {{"type":"section","title":"Planning","number":4}},
  {{"type":"table","title":"Planning previsionnel","headers":["Phase","Duree","Livrables"],"rows":[["Phase 1","2 sem","Livrable"],["Phase 2","3 sem","Livrable"]]}},
  {{"type":"section","title":"Budget","number":5}},
  {{"type":"table","title":"Chiffrage","headers":["Prestation","Jours","TJM","Total"],"rows":[["Phase 1","10j","1000€","10000€"],["Total","","","XX XXX€"]]}},
  {{"type":"closing"}}
]

IMPORTANT:
- Chaque slide "content" DOIT avoir un champ "bullets" avec 4-6 points concrets
- Chaque slide "table" DOIT avoir "headers" et "rows"
- Total: 10-16 slides
- Integre les modifications demandees dans le feedback""",
                system_prompt="Tu generes des structures de slides PowerPoint completes en JSON. Sois concret et professionnel.",
                temperature=0.5,
                max_tokens=8000,
            )

            json_start = slides_response.find("[")
            json_end = slides_response.rfind("]") + 1
            if json_start >= 0 and json_end > json_start:
                revised_slides = json.loads(slides_response[json_start:json_end])

                # Valider que les slides ont du contenu
                for slide in revised_slides:
                    if slide.get("type") == "content" and "bullets" not in slide:
                        slide["bullets"] = [
                            "Point a developper",
                            "Deuxieme point",
                            "Troisieme point",
                        ]
                    elif slide.get("type") == "table" and "rows" not in slide:
                        slide["headers"] = slide.get("headers", ["Colonne 1", "Colonne 2"])
                        slide["rows"] = [["Donnee 1", "Donnee 2"]]

                print(f"   Slides revisees generees: {len(revised_slides)}")
        except Exception as e:
            print(f"   Erreur generation slides: {e}")
            revised_slides = None

        job["steps"].append({"step": "slides", "status": "done", "progress": 85})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"proposal_revised_{timestamp}.md"
        json_path = output_dir / f"proposal_revised_{timestamp}.json"

        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Proposition commerciale (revisee)\n\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y')}\n\n---\n\n")
            f.write(revised_content)

        proposal_data = {
            "content": revised_content,
            "feedback_applied": feedback,
            "slides": revised_slides,
        }
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(proposal_data, f, indent=2, ensure_ascii=False)

        # Generer le PPTX revise
        pptx_path = None
        if revised_slides:
            try:
                from utils.pptx_generator import build_proposal_pptx

                pptx_path = output_dir / f"proposal_revised_{timestamp}.pptx"
                build_proposal_pptx(
                    template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                    slides_data=revised_slides,
                    output_path=str(pptx_path),
                    consultant_info={
                        "name": CONSULTANT_NAME,
                        "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
                    },
                )
                print(f"   PPTX revise genere: {pptx_path}")
            except Exception as e:
                print(f"   Erreur PPTX revise: {e}")
                pptx_path = None

        job["status"] = "done"
        job["result"] = {
            "content": revised_content,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "json_path": str(json_path.relative_to(BASE_DIR)),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)) if pptx_path else None,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === API: DOWNLOAD ===


@router.get("/api/download")
async def api_download(path: str):
    """Telecharge un fichier genere"""
    try:
        # Nettoyer le path pour éviter les problèmes d'encodage
        path = path.strip()
        file_path = BASE_DIR / path

        # Vérifier que le fichier existe
        if not file_path.exists():
            print(f"Fichier non trouve: {file_path}")
            return JSONResponse({"error": "Fichier non trouve"}, status_code=404)

        # Vérifier que le fichier est bien dans le répertoire BASE_DIR
        # (sécurité)
        if not file_path.is_relative_to(BASE_DIR):
            print(f"Tentative d'accès à un fichier hors du répertoire autorisé: {file_path}")
            return JSONResponse({"error": "Accès non autorisé"}, status_code=403)

        # Déterminer le type MIME en fonction de l'extension
        media_type = None
        if file_path.suffix == ".pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_path.suffix == ".pdf":
            media_type = "application/pdf"
        elif file_path.suffix == ".md":
            media_type = "text/markdown"
        elif file_path.suffix == ".json":
            media_type = "application/json"
        elif file_path.suffix == ".png":
            media_type = "image/png"
        elif file_path.suffix == ".jpg" or file_path.suffix == ".jpeg":
            media_type = "image/jpeg"

        # Retourner le fichier avec les bons headers
        # Pour les images (PNG, JPG), utiliser inline pour affichage dans le
        # navigateur
        disposition = "inline" if file_path.suffix in [".png", ".jpg", ".jpeg"] else "attachment"

        return FileResponse(
            file_path,
            filename=file_path.name,
            media_type=media_type,
            headers={
                "Content-Disposition": f'{disposition}; filename="{file_path.name}"',
                "Cache-Control": "no-cache",
            },
        )

    except Exception as e:
        print(f"Erreur lors du téléchargement: {safe_error_message(e)}")
        return JSONResponse(
            {"error": f"Erreur lors du téléchargement: {safe_error_message(e)}"}, status_code=500
        )


@router.post("/api/convert-to-pdf")
async def api_convert_to_pdf(request: Request):
    """
    Convertit un fichier PPTX ou Markdown en PDF

    Body JSON:
        - file_path: Chemin relatif du fichier à convertir (PPTX ou MD)
        - file_type: Type de fichier ('pptx' ou 'markdown')
    """
    try:
        from utils.pdf_converter import pdf_converter

        body = await request.json()
        file_path = body.get("file_path", "").strip()
        file_type = body.get("file_type", "pptx")

        if not file_path:
            return JSONResponse({"error": "Chemin de fichier manquant"}, status_code=400)

        # Chemin absolu
        full_path = BASE_DIR / file_path

        if not full_path.exists():
            return JSONResponse({"error": "Fichier non trouvé"}, status_code=404)

        # Vérifier la sécurité
        if not full_path.is_relative_to(BASE_DIR):
            return JSONResponse({"error": "Accès non autorisé"}, status_code=403)

        # Convertir selon le type
        pdf_path = None
        if file_type == "pptx" and full_path.suffix == ".pptx":
            pdf_path = pdf_converter.pptx_to_pdf(str(full_path))
        elif file_type == "markdown" and full_path.suffix in [".md", ".markdown"]:
            pdf_path = pdf_converter.markdown_to_pdf(str(full_path))
        else:
            return JSONResponse({"error": "Type de fichier non supporté"}, status_code=400)

        if pdf_path:
            # Retourner le chemin relatif du PDF
            pdf_rel_path = str(Path(pdf_path).relative_to(BASE_DIR))
            return {"success": True, "pdf_path": pdf_rel_path, "message": "Conversion réussie"}
        else:
            return JSONResponse(
                {
                    "error": "Conversion PDF échouée",
                    "message": "Vérifiez que LibreOffice est installé (pour PPTX) ou pandoc/weasyprint (pour Markdown)",
                },
                status_code=500,
            )

    except Exception as e:
        print(f"Erreur lors de la conversion PDF: {safe_error_message(e)}")
        return JSONResponse({"error": f"Erreur: {safe_error_message(e)}"}, status_code=500)


@router.get("/api/pdf-capabilities")
async def api_pdf_capabilities():
    """Retourne les capacités de conversion PDF disponibles"""
    try:
        from utils.pdf_converter import pdf_converter

        capabilities = pdf_converter.is_pdf_conversion_available()
        return {
            "capabilities": capabilities,
            "message": "Pour activer la conversion PDF, installez LibreOffice (PPTX→PDF) ou pandoc/weasyprint (MD→PDF)",
        }
    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# ===================================
# PROPOSAL CANVA (Mode Gemini)
# ===================================


@router.get("/proposal-canva")
async def proposal_canva_page(request: Request):
    """Page de génération de propositions en mode Canva (conversationnel avec Gemini)"""
    return templates.TemplateResponse(
        "proposal-canva.html",
        {"request": request, "consultant_name": CONSULTANT_NAME, "active": "proposal"},
    )


@router.post("/api/proposal-canva/generate")
async def api_proposal_canva_generate(request: Request):
    """Génère des propositions en mode conversationnel avec Gemini"""
    try:
        data = await request.json()
        user_message = data.get("message", "")
        conversation_history = data.get("conversation_history", [])
        data.get("current_proposal")

        if not user_message:
            return JSONResponse({"error": "Message manquant"}, status_code=400)

        # Initialiser le client Gemini
        from utils.llm_client import LLMClient

        llm = LLMClient(provider="gemini", model="gemini-3-flash-preview")

        # System prompt pour le mode conversationnel
        system_prompt = """Tu es un assistant expert en propositions commerciales pour {COMPANY_NAME}.
Tu aides les consultants à créer des propositions professionnelles de manière conversationnelle.

INSTRUCTIONS:
1. Analyse la demande de l'utilisateur
2. Si c'est la première fois, demande le contexte de l'appel d'offres
3. Génère des slides structurées en JSON selon le besoin
4. Réponds de manière naturelle et conversationnelle
5. Propose des améliorations si nécessaire

FORMAT DE SORTIE:
- Si tu génères des slides, utilise ce format JSON à la fin de ta réponse:
```json
{"slides": [
    {"type": "cover", "title": "...", "subtitle": "..."} ,
    {"type": "content", "title": "...", "bullets": ["...", "..."]} ,
    {"type": "diagram", "title": "...", "diagram_type": "flow", "elements": ["..."]} ,
    {"type": "stat", "stat_value": "67%", "stat_label": "de ROI", "context": "..."} ,
    {"type": "quote", "quote_text": "...", "author": "..."} ,
    {"type": "highlight", "title": "...", "key_points": ["...", "..."]}
  ]
}
```

TYPES DE SLIDES DISPONIBLES:
- cover: Slide de couverture (title, subtitle)
- section: Séparateur de section (title)
- content: Contenu avec bullets (title, bullets)
- diagram: Diagramme visuel (title, diagram_type, elements, description)
- stat: Statistique impactante (stat_value, stat_label, context, subtitle)
- quote: Citation/message clé (quote_text, author, title)
- highlight: Points clés en encadrés (title, key_points, highlight_color)
- closing: Slide de clôture

Consultant: {CONSULTANT_NAME}
Société: {COMPANY_NAME} """

        # Construire les messages pour l'API
        messages = []
        for msg in conversation_history:
            if msg["role"] in ["user", "assistant"]:
                messages.append({"role": msg["role"], "content": msg["content"]})

        # Ajouter le nouveau message
        messages.append({"role": "user", "content": user_message})

        # Générer la réponse
        response_text = llm.generate_with_context(
            messages=messages, system_prompt=system_prompt, temperature=0.7, max_tokens=8000
        )

        # Extraire le JSON s'il existe
        slides = []
        pptx_path = None

        if "```json" in response_text:
            # Extraire le JSON
            json_start = response_text.find("```json") + 7
            json_end = response_text.find("```", json_start)
            json_str = response_text[json_start:json_end].strip()

            try:
                slides_data = json.loads(json_str)
                slides = slides_data.get("slides", [])

                # Générer le PPTX si des slides sont présentes
                if slides:
                    from utils.pptx_generator import build_proposal_pptx

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_dir = BASE_DIR / "output"
                    output_dir.mkdir(exist_ok=True)

                    output_path = output_dir / f"proposal_canva_{timestamp}.pptx"

                    consultant_info = {"name": CONSULTANT_NAME, "company": COMPANY_NAME}

                    pptx_path = build_proposal_pptx(
                        template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                        slides_data=slides,
                        output_path=str(output_path),
                        consultant_info=consultant_info,
                    )

                    pptx_path = str(Path(pptx_path).relative_to(BASE_DIR))

            except json.JSONDecodeError as e:
                print(f"Erreur JSON: {e}")
                # Continue sans slides

        # Nettoyer la réponse (enlever le JSON)
        response_text_clean = response_text.split("```json")[0].strip()

        # Mettre à jour l'historique
        conversation_history.append({"role": "user", "content": user_message})
        conversation_history.append({"role": "assistant", "content": response_text_clean})

        return {
            "response": response_text_clean,
            "slides": slides,
            "pptx_path": pptx_path,
            "conversation_history": conversation_history,
        }

    except Exception as e:
        print(f"Erreur génération Canva: {e}")
        import traceback

        traceback.print_exc()
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@router.post("/api/proposal-canva/generate-section")
async def api_proposal_canva_generate_section(request: Request):
    """Génère une section spécifique de proposition (mode modulaire)"""
    try:
        data = await request.json()
        section = data.get("section", "")
        tender_text = data.get("tender_text", "")
        data.get("current_proposal")

        if not section or not tender_text:
            return JSONResponse({"error": "Section et tender_text requis"}, status_code=400)

        # Initialiser le client Gemini
        from utils.llm_client import LLMClient

        llm = LLMClient(provider="gemini", model="gemini-3-flash-preview")

        # Mapping des sections
        section_prompts = {
            "cover": "Génère une slide de couverture professionnelle",
            "agenda": "Génère une slide d'agenda avec la structure de la proposition",
            "context": "Génère 2-3 slides de contexte avec les enjeux et objectifs",
            "approach": "Génère 3-4 slides d'approche méthodologique",
            "planning": "Génère 2-3 slides de planning avec les phases",
            "budget": "Génère 1-2 slides de chiffrage avec estimation budgétaire",
            "references": "Génère 2-3 slides de références projets similaires",
            "cvs": "Génère 2-3 slides de CVs adaptés au projet",
        }

        section_instruction = section_prompts.get(
            section, f"Génère des slides pour la section {section}"
        )

        # System prompt
        system_prompt = """Tu es un expert en propositions commerciales pour {COMPANY_NAME}.

Génère des slides pour la section "{section}" basées sur cet appel d'offres :

{tender_text[:3000]}

INSTRUCTION: {section_instruction}

FORMAT DE SORTIE - JSON uniquement:
```json
{{
  "slides": [
    {{"type": "content", "title": "...", "bullets": ["...", "..."]}},
    {{"type": "diagram", "title": "...", "diagram_type": "flow", "elements": ["..."]}},
    {{"type": "stat", "stat_value": "67%", "stat_label": "...", "context": "..."}},
    {{"type": "quote", "quote_text": "...", "author": "..."}},
    {{"type": "highlight", "title": "...", "key_points": ["...", "..."]}}
  ]
}}
```

TYPES DE SLIDES:
- cover: couverture (title, subtitle, client_name, project_name)
- section: séparateur (title, section_number)
- content: contenu avec bullets (title, bullets)
- diagram: diagramme (title, diagram_type, elements, description)
- stat: statistique impactante (stat_value, stat_label, context)
- quote: citation (quote_text, author, title)
- highlight: points clés (title, key_points, highlight_color)

Consultant: {CONSULTANT_NAME}
Société: {COMPANY_NAME}
"""

        # Générer la réponse
        response_text = llm.generate(
            prompt=f"Génère les slides pour la section {section}",
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=4000,
        )

        # Extraire le JSON
        slides = []
        pptx_path = None

        if "```json" in response_text:
            json_start = response_text.find("```json") + 7
            json_end = response_text.find("```", json_start)
            json_str = response_text[json_start:json_end].strip()

            try:
                slides_data = json.loads(json_str)
                slides = slides_data.get("slides", [])

                # Générer le PPTX si des slides sont présentes
                if slides:
                    from utils.pptx_generator import build_proposal_pptx

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_dir = BASE_DIR / "output"
                    output_dir.mkdir(exist_ok=True)

                    output_path = output_dir / f"proposal_section_{section}_{timestamp}.pptx"

                    consultant_info = {"name": CONSULTANT_NAME, "company": COMPANY_NAME}

                    pptx_path = build_proposal_pptx(
                        template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                        slides_data=slides,
                        output_path=str(output_path),
                        consultant_info=consultant_info,
                    )

                    pptx_path = str(Path(pptx_path).relative_to(BASE_DIR))

            except json.JSONDecodeError as e:
                print(f"Erreur JSON section: {e}")
                return JSONResponse({"error": "Erreur parsing JSON"}, status_code=500)

        return {
            "response": f"Section '{section}' générée avec succès ({len(slides)} slides)",
            "slides": slides,
            "pptx_path": pptx_path,
            "section": section,
        }

    except Exception as e:
        print(f"Erreur génération section: {e}")
        import traceback

        traceback.print_exc()
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)
