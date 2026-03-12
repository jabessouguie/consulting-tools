"""Router: E-Learning Adaptatif — Cours, Quiz, Parcours et Simulation d'entretien"""

import asyncio
import json
import threading
import uuid
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, File, Form, Request, UploadFile
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse

from agents.elearning_agent import ElearningAgent
from utils.consultant_db import ConsultantDatabase
from utils.validation import sanitize_text_input, sanitize_filename, validate_file_upload
from routers.shared import (
    BASE_DIR,
    jobs,
    limiter,
    templates,
    CONSULTANT_NAME,
    safe_error_message,
    elearning_db,
    send_sse,
)

router = APIRouter()


@router.get("/elearning", response_class=HTMLResponse)
async def elearning_page(request: Request):
    return templates.TemplateResponse(
        "elearning.html",
        {
            "request": request,
            "active": "elearning",
            "consultant_name": CONSULTANT_NAME,
        },
    )


# --- Sessions ---


@router.post("/api/elearning/session/init")
async def api_elearning_session_init(request: Request):
    """Initialise ou recupere une session etudiant"""
    body = await request.json()
    identifier = body.get("session_identifier")
    session = elearning_db.init_session(identifier)
    return session


# --- Cours: Generation ---


def _run_course_generator(
    job_id,
    topic,
    audience,
    difficulty,
    duration,
    document_content=None,
    mode="free",
    interview_type="",
    consultant_id=0,
):
    """Background job pour generer un cours"""
    try:
        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        # Charger le profil consultant si fourni
        consultant_profile = None
        if consultant_id and int(consultant_id) > 0:
            try:
                cdb = ConsultantDatabase()
                consultant_profile = cdb.get_consultant(int(consultant_id))
            except Exception as e:
                print(f"Erreur chargement consultant {consultant_id}: {e}")

        if document_content:
            result = agent.generate_course_from_document(
                document_content=document_content,
                target_audience=audience,
                difficulty=difficulty,
                duration_hours=duration,
                progress_callback=progress,
                mode=mode,
            )
        else:
            result = agent.generate_course(
                topic=topic,
                target_audience=audience,
                difficulty=difficulty,
                duration_hours=duration,
                progress_callback=progress,
                mode=mode,
                interview_type=interview_type,
                consultant_profile=consultant_profile,
            )

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            course_id = elearning_db.save_course(result)
            result["id"] = course_id
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "course_id": course_id,
                "title": result.get("title", ""),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du cours")


@router.post("/api/elearning/course/generate")
@limiter.limit("3/minute")
async def api_elearning_generate_course(
    request: Request,
    topic: str = Form(""),
    target_audience: str = Form("Professionnels"),
    difficulty: str = Form("intermediate"),
    duration_hours: int = Form(3),
    document_content: str = Form(""),
    mode: str = Form("free"),
    interview_type: str = Form(""),
    consultant_id: int = Form(0),
    auto_duration: str = Form("false"),
):
    """Lance la generation d'un cours (depuis sujet ou document)"""
    topic = sanitize_text_input(topic, max_length=1000)
    target_audience = sanitize_text_input(target_audience)
    document_content = sanitize_text_input(document_content, max_length=50000)
    interview_type = sanitize_text_input(interview_type)

    if not topic and not document_content:
        return JSONResponse(
            {"error": "Sujet ou document requis"},
            status_code=400,
        )

    if difficulty not in ("beginner", "intermediate", "advanced"):
        return JSONResponse({"error": "Difficulte invalide"}, status_code=400)

    valid_modes = ("free", "interview", "certification", "training")
    if mode not in valid_modes:
        mode = "free"

    valid_interview_types = ("rh", "technique", "cas", "fit", "")
    if interview_type not in valid_interview_types:
        interview_type = ""

    # Durée auto : passer 0 à l'agent pour qu'il la détermine lui-même
    if auto_duration.lower() in ("true", "1", "on"):
        duration_hours = 0

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    t = threading.Thread(
        target=_run_course_generator,
        args=(
            job_id,
            topic,
            target_audience,
            difficulty,
            duration_hours,
            document_content or None,
            mode,
            interview_type,
            consultant_id,
        ),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@router.post("/api/elearning/course/upload-document")
@limiter.limit("10/minute")
async def api_elearning_upload_document(
    request: Request,
    file: UploadFile = File(...),
):
    """Parse un document uploade pour e-learning"""
    try:
        content = await validate_file_upload(file)
        filename = sanitize_filename(file.filename or "document.txt")
        ext = Path(filename).suffix.lower()

        if ext in (".md", ".txt", ".markdown"):
            text = content.decode("utf-8")
        elif ext == ".pdf":
            import io

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext in (".html", ".htm"):
            text = content.decode("utf-8")
        elif ext == ".pptx":
            import io

            from pptx import Presentation as PptxPres

            prs = PptxPres(io.BytesIO(content))
            slide_texts = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                texts.append(p.text.strip())
                if texts:
                    slide_texts.append("\n".join(texts))
            text = "\n\n---\n\n".join(slide_texts)
        else:
            return JSONResponse(
                {"error": f"Format non supporte: {ext}"},
                status_code=400,
            )

        return {
            "text": text,
            "filename": filename,
            "length": len(text),
        }
    except Exception as e:
        return JSONResponse(
            {"error": safe_error_message(e)},
            status_code=500,
        )


@router.get("/api/elearning/course/stream/{job_id}")
async def api_elearning_course_stream(job_id: str):
    """SSE stream pour la generation de cours"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


# --- Cours: Regeneration ---


def _run_course_regenerator(job_id, course_id, feedback):
    """Background job pour regenerer un cours"""
    try:
        course = elearning_db.get_course(course_id)
        if not course:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Cours non trouve"
            return

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        result = agent.regenerate_with_feedback(course, feedback, progress_callback=progress)

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            elearning_db.delete_course(course_id)
            new_id = elearning_db.save_course(result)
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "course_id": new_id,
                "title": result.get("title", ""),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Regeneration du cours")


@router.post("/api/elearning/course/{course_id}/regenerate")
@limiter.limit("3/minute")
async def api_elearning_regenerate_course(
    request: Request,
    course_id: int,
    feedback: str = Form(...),
):
    """Regenere un cours avec feedback"""
    feedback = sanitize_text_input(feedback)
    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    t = threading.Thread(
        target=_run_course_regenerator,
        args=(job_id, course_id, feedback),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


# --- Cours: CRUD ---


@router.get("/api/elearning/courses")
async def api_elearning_list_courses():
    """Liste tous les cours"""
    courses = elearning_db.get_all_courses()
    return {"courses": courses}


@router.get("/api/elearning/course/{course_id}")
async def api_elearning_get_course(course_id: int):
    """Recupere un cours complet"""
    course = elearning_db.get_course(course_id)
    if not course:
        return JSONResponse({"error": "Cours non trouve"}, status_code=404)
    return {"course": course}


@router.delete("/api/elearning/course/{course_id}")
async def api_elearning_delete_course(course_id: int):
    """Supprime un cours"""
    deleted = elearning_db.delete_course(course_id)
    if not deleted:
        return JSONResponse({"error": "Cours non trouve"}, status_code=404)
    return {"ok": True}


@router.post("/api/elearning/lesson/{lesson_id}/solutions")
@limiter.limit("10/minute")
async def api_elearning_lesson_solutions(request: Request, lesson_id: int):
    """Génère les solutions pour les exercices d'une leçon existante."""
    body = await request.json()
    consultant_id = int(body.get("consultant_id", 0))

    # Trouver la leçon dans la DB via le cours
    course_id = int(body.get("course_id", 0))
    topic = sanitize_text_input(body.get("topic", ""), max_length=500)
    lesson_title = sanitize_text_input(body.get("lesson_title", ""), max_length=500)
    exercises = body.get("exercises", [])

    if not exercises:
        return JSONResponse({"error": "Aucun exercice fourni"}, status_code=400)

    consultant_profile = None
    if consultant_id > 0:
        try:
            consultant_db = ConsultantDatabase()
            consultant_profile = consultant_db.get_consultant(consultant_id)
        except Exception:
            pass

    try:
        agent = ElearningAgent()
        updated_exercises = agent.generate_exercise_solutions(
            exercises=exercises,
            topic=topic,
            lesson_title=lesson_title,
            consultant_profile=consultant_profile,
        )
        # Persister en DB si lesson_id fourni
        if lesson_id > 0:
            elearning_db.update_lesson_exercises(lesson_id, updated_exercises)
        return {"exercises": updated_exercises}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# --- Quiz: Generation ---


def _run_quiz_generator(job_id, course_id, lesson_id, difficulty, mode="free"):
    """Background job pour generer un quiz"""
    try:
        if lesson_id:
            course = elearning_db.get_course(course_id)
            if not course:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Cours non trouve"
                return

            lesson = None
            for mod in course.get("modules", []):
                for les in mod.get("lessons", []):
                    if les["id"] == lesson_id:
                        lesson = les
                        break

            if not lesson:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Lecon non trouvee"
                return

            lesson_title = lesson.get("title", "")
            lesson_content = lesson.get("content_markdown", "")
        else:
            course = elearning_db.get_course(course_id)
            if not course:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Cours non trouve"
                return

            lesson_title = course.get("title", "")
            all_content = []
            for mod in course.get("modules", []):
                for les in mod.get("lessons", []):
                    all_content.append(les.get("content_markdown", ""))
            lesson_content = "\n\n".join(all_content)

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        result = agent.generate_quiz(
            lesson_title=lesson_title,
            lesson_content=lesson_content,
            difficulty=difficulty,
            progress_callback=progress,
            mode=mode,
        )

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            result["course_id"] = course_id
            result["lesson_id"] = lesson_id
            result["difficulty_level"] = difficulty
            quiz_id = elearning_db.save_quiz(result)
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "quiz_id": quiz_id,
                "questions_count": len(result.get("questions", [])),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du quiz")


@router.post("/api/elearning/quiz/generate")
@limiter.limit("5/minute")
async def api_elearning_generate_quiz(
    request: Request,
    course_id: int = Form(...),
    lesson_id: Optional[int] = Form(None),
    difficulty: str = Form("medium"),
    mode: str = Form("free"),
):
    """Lance la generation d'un quiz"""
    if difficulty not in ("easy", "medium", "hard"):
        return JSONResponse({"error": "Difficulte invalide"}, status_code=400)

    valid_modes = ("free", "interview", "certification", "training")
    if mode not in valid_modes:
        mode = "free"

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    t = threading.Thread(
        target=_run_quiz_generator,
        args=(job_id, course_id, lesson_id, difficulty, mode),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@router.get("/api/elearning/quiz/stream/{job_id}")
async def api_elearning_quiz_stream(job_id: str):
    """SSE stream pour la generation de quiz"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


# --- Quiz: Tentatives ---


@router.post("/api/elearning/quiz/start")
async def api_elearning_quiz_start(
    request: Request,
    quiz_id: int = Form(...),
    session_identifier: str = Form(...),
):
    """Demarre une tentative de quiz"""
    session = elearning_db.get_session(session_identifier)
    if not session:
        session = elearning_db.init_session(session_identifier)

    quiz = elearning_db.get_quiz(quiz_id)
    if not quiz:
        return JSONResponse({"error": "Quiz non trouve"}, status_code=404)

    attempt_id = elearning_db.create_attempt(session["id"], quiz_id)

    # Get first question at medium difficulty
    questions = quiz.get("questions", [])
    first_question = questions[0] if questions else None

    if first_question:
        # Remove correct_answer from response
        q = {**first_question}
        q.pop("correct_answer", None)
        q.pop("explanation", None)
    else:
        q = None

    return {
        "attempt_id": attempt_id,
        "total_questions": len(questions),
        "first_question": q,
    }


@router.post("/api/elearning/quiz/answer")
async def api_elearning_quiz_answer(
    request: Request,
    attempt_id: int = Form(...),
    question_id: int = Form(...),
    answer: str = Form(...),
    time_spent: int = Form(0),
):
    """Soumet une reponse et obtient la suivante (adaptatif)"""
    results = elearning_db.get_attempt_results(attempt_id)
    if not results:
        return JSONResponse({"error": "Tentative non trouvee"}, status_code=404)

    quiz = elearning_db.get_quiz(results["quiz_id"])
    if not quiz:
        return JSONResponse({"error": "Quiz non trouve"}, status_code=404)

    # Find the question
    question = None
    for q in quiz.get("questions", []):
        if q["id"] == question_id:
            question = q
            break

    if not question:
        return JSONResponse({"error": "Question non trouvee"}, status_code=404)

    # Evaluate answer
    agent = ElearningAgent()
    evaluation = agent.evaluate_answer(question, answer)

    # Record answer
    elearning_db.record_answer(
        attempt_id,
        question_id,
        answer,
        evaluation["is_correct"],
        time_spent,
    )

    # Adaptive difficulty
    recent = elearning_db.get_recent_answers(attempt_id, 3)
    new_difficulty = agent.adapt_difficulty(recent)
    current_difficulty = results.get("current_difficulty", "medium")

    if new_difficulty:
        elearning_db.update_attempt_difficulty(attempt_id, new_difficulty)
        current_difficulty = new_difficulty

    # Find next unanswered question
    answered_ids = {a["question_id"] for a in results.get("answers", [])}
    answered_ids.add(question_id)

    next_question = None
    for q in quiz.get("questions", []):
        if q["id"] not in answered_ids:
            next_question = {**q}
            next_question.pop("correct_answer", None)
            next_question.pop("explanation", None)
            break

    return {
        "is_correct": evaluation["is_correct"],
        "explanation": evaluation.get("explanation", ""),
        "feedback": evaluation.get("feedback", ""),
        "current_difficulty": current_difficulty,
        "difficulty_changed": new_difficulty is not None,
        "next_question": next_question,
        "questions_remaining": len(quiz.get("questions", [])) - len(answered_ids),
    }


@router.get("/api/elearning/quiz/results/{attempt_id}")
async def api_elearning_quiz_results(attempt_id: int):
    """Recupere les resultats d'une tentative"""
    # Complete the attempt first
    elearning_db.complete_attempt(attempt_id)

    results = elearning_db.get_attempt_results(attempt_id)
    if not results:
        return JSONResponse({"error": "Tentative non trouvee"}, status_code=404)
    return {"results": results}


@router.get("/api/elearning/quizzes/{course_id}")
async def api_elearning_list_quizzes(course_id: int):
    """Liste les quiz d'un cours"""
    quizzes = elearning_db.get_quizzes_for_course(course_id)
    return {"quizzes": quizzes}


# --- Parcours d'apprentissage ---


def _run_learning_path_generator(job_id, session_id, course_id, goals):
    """Background job pour generer un parcours"""
    try:
        course = elearning_db.get_course(course_id)
        if not course:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Cours non trouve"
            return

        session = elearning_db.get_session(session_id)
        if not session:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Session non trouvee"
            return

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        # Get quiz results for gap analysis
        quiz_results = elearning_db.get_session_quiz_results(session["id"], course_id)

        progress("analyzing", "Analyse des performances...")

        gaps = agent.analyze_knowledge_gaps(quiz_results, course.get("modules", []))

        progress("generating", "Generation du parcours...")

        path_data = agent.generate_learning_path(
            gaps=gaps,
            goals=goals,
            course_modules=course.get("modules", []),
            progress_callback=progress,
        )

        # Save learning path
        path_id = elearning_db.save_learning_path(
            {
                "session_id": session["id"],
                "course_id": course_id,
                "stated_goals": goals,
                "knowledge_gaps": path_data.get("knowledge_gaps", []),
                "recommendations": path_data.get("recommendations", []),
            }
        )

        jobs[job_id]["status"] = "done"
        jobs[job_id]["result"] = {
            "path_id": path_id,
            "steps_count": len(path_data.get("path_steps", [])),
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du parcours")


@router.post("/api/elearning/learning-path/create")
@limiter.limit("3/minute")
async def api_elearning_create_learning_path(
    request: Request,
    session_identifier: str = Form(...),
    course_id: int = Form(...),
    goals: str = Form(""),
):
    """Cree un parcours d'apprentissage personnalise"""
    try:
        goals_list = json.loads(goals) if goals else []
    except json.JSONDecodeError:
        goals_list = [g.strip() for g in goals.split(",") if g.strip()]

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    t = threading.Thread(
        target=_run_learning_path_generator,
        args=(job_id, session_identifier, course_id, goals_list),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@router.get("/api/elearning/learning-path/stream/{job_id}")
async def api_elearning_learning_path_stream(job_id: str):
    """SSE stream pour la generation de parcours"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


@router.get("/api/elearning/learning-path/{session_identifier}/{course_id}")
async def api_elearning_get_learning_path(session_identifier: str, course_id: int):
    """Recupere un parcours d'apprentissage"""
    session = elearning_db.get_session(session_identifier)
    if not session:
        return JSONResponse({"error": "Session non trouvee"}, status_code=404)

    path = elearning_db.get_learning_path(session["id"], course_id)
    if not path:
        return JSONResponse({"error": "Parcours non trouve"}, status_code=404)
    return {"path": path}


@router.post("/api/elearning/learning-path/update-progress")
async def api_elearning_update_progress(
    request: Request,
    learning_path_id: int = Form(...),
    module_id: int = Form(...),
    completion_pct: float = Form(...),
    mastery_level: str = Form("none"),
):
    """Met a jour la progression d'un module"""
    valid_levels = (
        "none",
        "beginner",
        "intermediate",
        "proficient",
        "expert",
    )
    if mastery_level not in valid_levels:
        return JSONResponse(
            {"error": "Niveau de maitrise invalide"},
            status_code=400,
        )

    elearning_db.update_module_progress(
        learning_path_id,
        module_id,
        completion_pct,
        mastery_level,
    )
    return {"ok": True}


# === INTERVIEW CHAT ===


@router.post("/api/elearning/interview/chat")
@limiter.limit("30/minute")
async def api_elearning_interview_chat(request: Request):
    """
    Tour de conversation dans un entretien simulé interactif.

    Body JSON :
      {
        "messages": [{"role": "user"|"assistant", "content": "..."}],
        "topic": "Python",          # optionnel
        "role": "Data Engineer"     # optionnel
      }

    Retourne la réponse de l'interviewer.
    """
    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"error": "Corps JSON invalide"}, status_code=400)

    messages = body.get("messages")
    if not messages or not isinstance(messages, list):
        return JSONResponse({"error": "Champ 'messages' requis (liste)"}, status_code=400)

    # Valider les messages
    for msg in messages:
        if not isinstance(msg, dict) or msg.get("role") not in ("user", "assistant"):
            return JSONResponse(
                {"error": "Chaque message doit avoir 'role' (user|assistant) et 'content'"},
                status_code=400,
            )

    topic = str(body.get("topic", ""))[:200]
    role = str(body.get("role", ""))[:200]
    interviewer_name = str(body.get("interviewer_name", ""))[:200]
    interviewer_linkedin = str(body.get("interviewer_linkedin", ""))[:500]
    interview_type = str(body.get("interview_type", ""))[:20]

    try:
        agent = ElearningAgent()
        reply = agent.interview_chat(
            messages,
            topic=topic,
            role=role,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
            interview_type=interview_type,
        )
    except Exception as e:
        return JSONResponse({"error": f"Erreur entretien : {str(e)}"}, status_code=500)

    return {"reply": reply}


@router.post("/api/elearning/interview/analyze")
@limiter.limit("10/minute")
async def api_elearning_interview_analyze(request: Request):
    """
    Analyse les performances d'un entretien simulé complet.

    Body JSON :
      {
        "messages": [{"role": "user"|"assistant", "content": "..."}],
        "topic": "Python",
        "role": "Data Engineer",
        "interview_type": "rh|technique|cas|fit"
      }

    Retourne {overall_score, grade, strengths, improvements, detailed_feedback, ...}
    """
    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"error": "Corps JSON invalide"}, status_code=400)

    messages = body.get("messages")
    if not messages or not isinstance(messages, list):
        return JSONResponse({"error": "Champ 'messages' requis (liste)"}, status_code=400)

    topic = str(body.get("topic", ""))[:200]
    role = str(body.get("role", ""))[:200]
    interviewer_name = str(body.get("interviewer_name", ""))[:200]
    interviewer_linkedin = str(body.get("interviewer_linkedin", ""))[:500]
    interview_type = str(body.get("interview_type", ""))[:20]

    try:
        agent = ElearningAgent()
        analysis = agent.analyze_interview_performance(
            messages,
            topic=topic,
            role=role,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
            interview_type=interview_type,
        )
    except Exception as e:
        return JSONResponse({"error": f"Erreur analyse : {str(e)}"}, status_code=500)

    return {"analysis": analysis}


# === PREMIUM RESOURCES ===


@router.post("/api/elearning/course/{course_id}/premium-resources")
@limiter.limit("5/minute")
async def api_elearning_premium_resources(
    request: Request,
    course_id: int,
):
    """
    Génère des ressources premium pour un cours :
    fiches mémo, résumés condensés, cheatsheet globale.

    Body JSON optionnel :
      { "resource_types": ["memo", "summary", "cheatsheet"] }
    """
    course = elearning_db.get_course(course_id)
    if not course:
        return JSONResponse({"error": "Cours non trouvé"}, status_code=404)

    try:
        body = await request.json()
    except Exception:
        body = {}

    resource_types = body.get("resource_types") if isinstance(body, dict) else None
    valid_types = {"memo", "summary", "cheatsheet"}
    if resource_types is not None:
        resource_types = [t for t in resource_types if t in valid_types] or None

    try:
        agent = ElearningAgent()
        resources = agent.generate_premium_resources(course, resource_types)
    except Exception as e:
        return JSONResponse(
            {"error": f"Erreur génération ressources premium : {str(e)}"},
            status_code=500,
        )

    return {"resources": resources}
