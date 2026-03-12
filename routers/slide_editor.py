"""Router: Slide Editor — Génération HTML-First de présentations et doc-to-presentation"""

import asyncio
import json
import os
import threading
import uuid
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, File, Form, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse

from utils.validation import sanitize_text_input, sanitize_filename, validate_file_upload
from routers.shared import (
    BASE_DIR,
    jobs,
    limiter,
    templates,
    CONSULTANT_NAME,
    COMPANY_NAME,
    safe_error_message,
    safe_traceback,
    send_sse,
)

router = APIRouter()


# ===================================
# DOCUMENT TO PRESENTATION
# ===================================


@router.get("/doc-to-presentation", response_class=HTMLResponse)
async def doc_to_presentation_page(request: Request):
    return templates.TemplateResponse(
        "doc-to-presentation.html",
        {
            "request": request,
            "active": "doc-to-presentation",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@router.post("/api/doc-to-presentation/generate")
async def api_doc_to_presentation_generate(
    request: Request,
    target_audience: str = Form(...),
    objective: str = Form(...),
    files: List[UploadFile] = File(...),
):
    """Lance la generation de presentation a partir de documents"""
    if not files:
        return JSONResponse({"error": "Aucun fichier fourni."}, status_code=400)

    # Lire les fichiers uploades
    documents = []
    for f in files:
        content = await f.read()
        documents.append({"filename": f.filename, "content": content})

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "doc-to-presentation",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_doc_to_presentation,
        args=(job_id, documents, target_audience, objective),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_doc_to_presentation(job_id: str, documents: list, target_audience: str, objective: str):
    """Execute la generation de presentation en background"""
    job = jobs[job_id]
    try:
        from agents.doc_to_presentation import DocToPresentationAgent

        agent = DocToPresentationAgent()

        job["steps"].append({"step": "parse", "status": "active", "progress": 10})
        result = agent.run(documents, target_audience, objective)
        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        if result.get("error"):
            job["status"] = "error"
            job["error"] = result["error"]
        else:
            job["status"] = "done"
            job["result"] = result

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@router.get("/api/doc-to-presentation/stream/{job_id}")
async def api_doc_to_presentation_stream(job_id: str):
    """SSE stream pour la progression"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(event_generator(), media_type="text/event-stream")


# ===================================
# HTML SLIDES PREMIUM (Gemini 3.1 Pro)
# ===================================


@router.get("/html-slides")
async def html_slides_page():
    """Redirige vers le Slide Editor (HTML Premium integre)"""
    from starlette.responses import RedirectResponse

    return RedirectResponse(url="/slide-editor", status_code=301)


@router.post("/api/html-slides/generate")
@limiter.limit("3/minute")
async def api_html_slides_generate(
    request: Request,
    topic: str = Form(...),
    num_slides: int = Form(10),
    language: str = Form("fr"),
):
    """Lance la generation de slides HTML premium via Gemini 3.1 Pro"""
    topic = sanitize_text_input(topic)
    if not topic or len(topic.strip()) < 5:
        return JSONResponse(
            {"error": "Le sujet est trop court (minimum 5 caracteres)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "html-slides",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_html_slides_generator,
        args=(job_id, topic, num_slides, language),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_html_slides_generator(job_id: str, topic: str, num_slides: int, language: str):
    """Generate HTML slides in background"""
    job = jobs[job_id]
    try:
        from agents.html_slide_generator import HtmlSlideGeneratorAgent

        agent = HtmlSlideGeneratorAgent()

        # Step 1: Extract design system
        job["steps"].append({"step": "design", "status": "active", "progress": 10})
        design_system = agent.extract_design_system()
        job["steps"].append({"step": "design", "status": "done", "progress": 30})

        # Step 2: Generate HTML via Gemini 3.1 Pro
        job["steps"].append({"step": "generate", "status": "active", "progress": 35})

        system_instruction, user_prompt = agent.build_prompt(topic, num_slides, design_system)

        lang_prefix = ""
        if language == "en":
            lang_prefix = "Generate the presentation content in English. "
        user_prompt = lang_prefix + user_prompt

        messages = [{"role": "user", "content": user_prompt}]
        raw_html = agent.llm.generate_with_context(
            messages=messages,
            system_prompt=system_instruction,
            temperature=0.2,
            max_tokens=16384,
        )
        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Step 3: Clean and save
        job["steps"].append({"step": "save", "status": "active", "progress": 85})
        html_content = agent.clean_html_response(raw_html)
        html_path = agent.save_html(html_content)
        job["steps"].append({"step": "save", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "html_path": html_path,
            "html_content": html_content,
            "topic": topic,
            "num_slides": num_slides,
        }

    except Exception as e:
        print(f"Error in HTML slides: {safe_traceback()}")
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@router.get("/api/html-slides/stream/{job_id}")
async def api_html_slides_stream(job_id: str):
    """SSE stream pour la progression de generation de slides HTML"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === SLIDE EDITOR (HTML-First) ===


@router.get("/slide-editor", response_class=HTMLResponse)
async def slide_editor_page(request: Request):
    return templates.TemplateResponse(
        "slide-editor.html",
        {
            "request": request,
            "active": "slide-editor",
            "consultant_name": CONSULTANT_NAME,
        },
    )


def _extract_images_from_pdf(content_bytes):
    """Extract images from PDF as base64 data URIs"""
    images = []
    try:
        import base64
        import io

        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(content_bytes))
        for page_num, page in enumerate(reader.pages):
            for img_key in page.images:
                img_data = img_key.data
                ext = img_key.name.rsplit(".", 1)[-1].lower()
                mime = {
                    "jpg": "image/jpeg",
                    "jpeg": "image/jpeg",
                    "png": "image/png",
                    "gif": "image/gif",
                }.get(ext, "image/png")
                b64 = base64.b64encode(img_data).decode()
                images.append(
                    {
                        "data_uri": f"data:{mime};base64,{b64}",
                        "name": img_key.name,
                        "page": page_num + 1,
                    }
                )
                if len(images) >= 20:
                    return images
    except Exception as e:
        print(f"Erreur extraction images PDF: {e}")
    return images


def _extract_images_from_pptx(content_bytes):
    """Extract images from PPTX as base64 data URIs"""
    images = []
    try:
        import base64
        import io

        from pptx import Presentation as PptxPres

        prs = PptxPres(io.BytesIO(content_bytes))
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    img = shape.image
                    ext = img.content_type.split("/")[-1]
                    mime = img.content_type
                    b64 = base64.b64encode(img.blob).decode()
                    images.append(
                        {
                            "data_uri": f"data:{mime};base64,{b64}",
                            "name": f"slide_{slide_num + 1}_{shape.name}.{ext}",
                            "slide": slide_num + 1,
                        }
                    )
                    if len(images) >= 20:
                        return images
    except Exception as e:
        print(f"Erreur extraction images PPTX: {e}")
    return images


@router.post("/api/slide-editor/parse-document")
@limiter.limit("20/minute")
async def api_slide_editor_parse_document(request: Request, file: UploadFile = File(...)):
    """Parse un document uploade et retourne son contenu texte + images extraites"""
    try:
        # Valider le fichier uploade (taille, type)
        content = await validate_file_upload(file)
        filename = sanitize_filename(file.filename or "document.txt")
        ext = Path(filename).suffix.lower()
        images = []

        if ext in (".md", ".txt"):
            text = content.decode("utf-8")
        elif ext == ".pdf":
            import io

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            images = _extract_images_from_pdf(content)
        elif ext == ".docx":
            import io

            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            images = _extract_images_from_pptx(content)
            # Pour PPTX: sauvegarder temporairement le fichier et retourner le
            # chemin
            import tempfile

            temp_dir = Path(tempfile.gettempdir()) / "consulting-tools_uploads"
            temp_dir.mkdir(exist_ok=True)

            temp_path = temp_dir / f"upload_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
            with open(temp_path, "wb") as f:
                f.write(content)

            # Retourner le chemin au lieu du texte
            return {
                "text": str(temp_path),
                "filename": filename,
                "length": len(content),
                "is_pptx": True,
                "images": images,
                "image_count": len(images),
            }
        elif ext == ".html":
            text = content.decode("utf-8")
        else:
            return JSONResponse({"error": f"Format non supporte: {ext}"}, status_code=400)

        return {
            "text": text,
            "filename": filename,
            "length": len(text),
            "images": images,
            "image_count": len(images),
        }
    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


def _generate_slide_illustrations(slides, job, topic, gen_type="presentation"):
    """Add Nano Banana image generation prompts to relevant slides (no actual image generation)"""
    try:
        # Slide types that benefit from illustrations
        visual_types = ["content", "highlight", "stat", "diagram", "image", "two_column"]

        # Context adapte au type de generation
        context_map = {
            "formation": "training presentation",
            "proposal": "business proposal",
            "rex": "experience feedback presentation",
            "presentation": "business presentation",
        }
        context_map.get(gen_type, "business presentation")

        for idx, slide in enumerate(slides):
            slide_type = slide.get("type", "")

            # Skip non-visual slides
            if slide_type not in visual_types:
                continue

            # Skip if already has an image or image_prompt
            if slide.get("image") or slide.get("image_prompt"):
                continue

            # Generate prompt for the slide (Nano Banana format)
            slide.get("title", "")
            slide.get("content", "")
            slide.get("bullets", [])
            slide.get("key_points", [])

            # Prompt adapte au type
            if gen_type == "formation":
                prompt = """Create a premium, professional illustration for a training presentation slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: Educational yet professional, Unreal Engine 5 render, engaging and clear.
Colors: Cool blues, warm amber/gold accents, approachable palette.
Mood: Pedagogical, inspiring, professional learning environment.
Format: Wide 16:9, 1792x1024px."""
            elif gen_type == "proposal":
                prompt = """Create a premium, professional illustration for a business proposal slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: High-end corporate, Unreal Engine 5 render, sophisticated and impactful.
Colors: Premium blues, gold/amber accents, executive palette.
Mood: Professional, trustworthy, results-oriented, winning proposal aesthetic.
Format: Wide 16:9, 1792x1024px."""
            else:
                prompt = """Create a premium, professional illustration for a {context} slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: Corporate tech aesthetic, Unreal Engine 5 render, clean and modern.
Colors: Cool blues, warm amber/gold accents, professional palette.
Mood: Sophisticated, futuristic but grounded, business presentation quality.
Format: Wide 16:9, 1792x1024px."""

            # Add prompt and dimensions to slide (for Nano Banana generation)
            slide["image_prompt"] = prompt.strip()
            slide["image_dimensions"] = "1792x1024px (16:9)"
            job["steps"].append(
                {"message": f"Prompt Nano Banana ajoute pour slide {idx + 1}", "step": 4}
            )

    except Exception as e:
        print(f"  Erreur ajout prompts illustrations: {e}")
        # Non-blocking - continue sans prompts


def _extract_slides_from_buffer(buffer, job):
    """Parse complete JSON slide objects from a partial LLM response buffer.
    Adds newly found slides to job['slides'] incrementally."""
    # Find the slides array content
    content = buffer
    # Strip markdown fences
    if "```json" in content:
        content = content.split("```json", 1)[1]
        if "```" in content:
            content = content.rsplit("```", 1)[0]
    elif "```" in content:
        parts = content.split("```")
        if len(parts) >= 2:
            content = parts[1]

    # Find the "slides" array or top-level array
    start = -1
    for marker in ['"slides"', "'slides'"]:
        idx = content.find(marker)
        if idx != -1:
            bracket = content.find("[", idx)
            if bracket != -1:
                start = bracket
                break

    if start == -1:
        # Maybe it's a top-level array
        bracket = content.find("[")
        if bracket != -1:
            start = bracket
        else:
            return

    # Extract individual slide objects by tracking brace depth
    existing_count = len(job["slides"])
    obj_start = -1
    depth = 0
    slide_count = 0

    for i in range(start, len(content)):
        c = content[i]
        if c == "{":
            if depth == 0:
                obj_start = i
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0 and obj_start != -1:
                slide_count += 1
                if slide_count > existing_count:
                    # New slide found - try to parse it
                    obj_str = content[obj_start : i + 1]
                    try:
                        slide = json.loads(obj_str)
                        if isinstance(slide, dict) and ("type" in slide or "title" in slide):
                            job["slides"].append(slide)
                    except json.JSONDecodeError:
                        pass
                obj_start = -1


@router.post("/api/slide-editor/start-generate")
@limiter.limit("5/minute")
async def api_slide_editor_start_generate(request: Request):
    """Demarre la generation de slides en arriere-plan et retourne un job_id"""
    data = await request.json()
    topic = data.get("topic", "")
    audience = data.get("audience", "")
    slide_count = min(int(data.get("slide_count", 10)), 40)
    gen_type = data.get("type", "presentation")
    model = data.get("model", "gemini-2.0-flash")
    document_text = data.get("document_text", "")
    feedback = data.get("feedback", "")
    previous_slides = data.get("previous_slides", "")
    provider = "claude" if model.startswith("claude") else "gemini"

    _ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    _uid = id(data) % 10000
    job_id = f"slide_{_ts}_{_uid}"

    jobs[job_id] = {
        "type": "slide-editor",
        "status": "running",
        "steps": [],
        "slides": [],
        "result": None,
        "error": None,
    }

    def run_generation():
        job = jobs[job_id]
        try:
            job["steps"].append(
                {
                    "message": f"Initialisation {model}...",
                    "substep": "Chargement du modele",
                    "step": 0,
                }
            )

            from agents.html_slide_generator import HtmlSlideGeneratorAgent

            agent = HtmlSlideGeneratorAgent(model=model, provider=provider)

            # Load extra context for proposals (references + CVs)
            extra_context = ""
            if gen_type == "proposal":
                try:
                    references_path = Path(BASE_DIR) / "data" / "notebooklm" / "references.json"
                    if references_path.exists():
                        with open(references_path, "r", encoding="utf-8") as f:
                            refs = json.load(f)
                        projects = json.dumps(
                            refs.get("projects", []), indent=2, ensure_ascii=False
                        )[:3000]
                        expertise = json.dumps(refs.get("expertise", []), ensure_ascii=False)[:1000]
                        methodologies = json.dumps(
                            refs.get("methodologies", []), ensure_ascii=False
                        )[:1000]
                        extra_context += (
                            f"\nREFERENCES Consulting Tools :\n{projects}"
                            f"\nEXPERTISE : {expertise}"
                            f"\nMETHODOLOGIES : {methodologies}"
                        )
                except Exception as e:
                    print(f"  Erreur chargement references: {e}")

                try:
                    bio_path = Path(BASE_DIR) / "Biographies - CV All Consulting Tools.pptx"
                    if bio_path.exists():
                        from pptx import Presentation as PptxPres

                        prs = PptxPres(str(bio_path))
                        cvs = []
                        for slide in prs.slides:
                            texts = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for p in shape.text_frame.paragraphs:
                                        if p.text.strip():
                                            texts.append(p.text.strip())
                            full = "\n".join(texts)
                            if len(full) > 50:
                                cvs.append(full)
                        if cvs:
                            extra_context += (
                                "\nCVs EQUIPE Consulting Tools :\n" + "\n---\n".join(cvs[:5])[:3000]
                            )
                except Exception as e:
                    print(f"  Erreur chargement CVs: {e}")

            # Handle feedback/regeneration
            if feedback and previous_slides:
                extra_context += (
                    "\n\nSLIDES ACTUELLES (a modifier selon le feedback) :\n"
                    + previous_slides[:6000]
                    + "\n\nFEEDBACK UTILISATEUR :\n"
                    + feedback
                    + "\n\nModifie les slides en tenant compte du feedback."
                )

            job["steps"].append(
                {
                    "message": "Generation HTML Premium...",
                    "substep": "Analyse du sujet et creation des slides",
                    "step": 1,
                }
            )
            job["head_html"] = ""
            job["html_sections"] = []

            # Stream HTML sections via Gemini 3.1 Pro
            for result in agent.run_streaming(
                topic=topic,
                num_slides=slide_count,
                gen_type=gen_type,
                audience=audience,
                document_text=document_text,
                extra_context=extra_context,
            ):
                if result["head_html"]:
                    job["head_html"] = result["head_html"]

                job["html_sections"].append(result["html_section"])
                slide_data = result["slide_json"]
                slide_data["html_section"] = result["html_section"]
                job["slides"].append(slide_data)

                slide_title = (
                    result.get("slide_json", {}).get("title", "") or f"Slide {result['index'] + 1}"
                )
                job["steps"].append(
                    {
                        "message": f"Slide {result['index'] + 1} generee",
                        "substep": slide_title,
                        "step": 2,
                    }
                )

            if not job["slides"]:
                raise ValueError("Aucune slide generee - reponse LLM invalide")

            job["result"] = {
                "slides": job["slides"],
                "total": len(job["slides"]),
                "head_html": job["head_html"],
            }
            job["status"] = "done"

            # Add image prompts in background (non-blocking)
            _generate_slide_illustrations(job["slides"], job, topic, gen_type)

        except Exception as e:
            print(f"Erreur slide-editor generate: {e}")
            import traceback

            traceback.print_exc()
            job["error"] = safe_error_message(e)
            job["status"] = "error"

    import threading

    threading.Thread(target=run_generation, daemon=True).start()

    return {"job_id": job_id}


@router.get("/api/slide-editor/stream/{job_id}")
async def api_slide_editor_stream(job_id: str):
    """SSE stream pour la progression de generation de slides"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        last_slide_idx = 0

        while True:
            # Send new steps
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("status", step)
                last_step_idx += 1

            # Send slides progressively as they appear
            while last_slide_idx < len(job["slides"]):
                slide = job["slides"][last_slide_idx]
                yield send_sse(
                    "slide",
                    {
                        "index": last_slide_idx,
                        "total": len(job["slides"]),
                        "slide": slide,
                        "head_html": job.get("head_html", "") if last_slide_idx == 0 else "",
                    },
                )
                last_slide_idx += 1
                await asyncio.sleep(0.1)

            if job["status"] == "done":
                # Send any remaining slides
                while last_slide_idx < len(job["slides"]):
                    slide = job["slides"][last_slide_idx]
                    yield send_sse(
                        "slide",
                        {
                            "index": last_slide_idx,
                            "total": len(job["slides"]),
                            "slide": slide,
                            "head_html": job.get("head_html", "") if last_slide_idx == 0 else "",
                        },
                    )
                    last_slide_idx += 1
                yield send_sse(
                    "done",
                    {
                        "total": len(job["slides"]),
                        "head_html": job.get("head_html", ""),
                    },
                )
                return

            if job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.3)

    return StreamingResponse(event_generator(), media_type="text/event-stream")


def _get_slide_json_format():
    return """
IMPORTANT : Tu dois retourner UNIQUEMENT un JSON valide avec cette structure :
{{"slides": [...]}}

Types de slides disponibles et leur format JSON :
- "cover" : {{"type":"cover", "title":"...", "subtitle":"...", "meta":"..."}}
- "section" : {{"type":"section", "title":"...", "number":"01"}}
- "content" : {{"type":"content", "title":"...", "bullets":["...","..."]}}
- "stat" : {{"type":"stat", "stat_value":"73%", "stat_label":"...", "context":"...", "subtitle":"..."}}
- "quote" : {{"type":"quote", "quote_text":"...", "author":"...", "title":"..."}}
- "highlight" : {{"type":"highlight", "title":"...", "key_points":["...","...","..."]}}
- "diagram" : {{"type":"diagram", "title":"...", "diagram_type":"flow|pyramid|timeline|cycle|matrix", "elements":["...","..."]}}
- "two_column" : {{"type":"two_column", "title":"...", "left_title":"...", "left_points":["..."], "right_title":"...", "right_points":["..."]}}
- "table" : {{"type":"table", "title":"...", "headers":["...","..."], "rows":[["...","..."]]}}
- "closing" : {{"type":"closing", "title":"Merci", "subtitle":"..."}}

Optionnel : chaque slide peut avoir un champ "tags": ["TAG1", "TAG2"] pour la categoriser.
"""


def _get_proposal_system_prompt(company_name):
    return """Tu agis en tant que Consultant Senior chez {company_name}. Ta mission est de rediger le contenu et la structure d'une proposition commerciale au format "Slides".

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.
- Style : Tags en en-tete (#TECH #DATA), marqueur vertical "Consulting Tools | 2025-2026", mise en page aeree.

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Utilise des types varies : stat pour les chiffres, diagram pour les processus, highlight pour les points cles, two_column pour les comparaisons
- Le contenu doit etre redige sur un ton expert, convaincant et direct
- Maximum 4 bullet points par slide, phrases courtes et percutantes
- Tu determines toi-meme le nombre de slides adapte au contenu et a la complexite du besoin client
- Ajoute des slides CV equipe et references missions si ces informations sont fournies"""


def _get_proposal_prompt(
    topic,
    audience,
    company_name,
    consultant_name,
    json_format,
    document_text="",
    references="",
    cvs="",
):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT FOURNI (utilise ce contenu comme base) :
{document_text[:6000]}
"""

    ref_section = ""
    if references:
        ref_section = """

REFERENCES DE MISSIONS SIMILAIRES (integre-les dans la proposition) :
{references}
"""

    cv_section = ""
    if cvs:
        cv_section = """

EQUIPE DISPONIBLE (integre les profils pertinents dans la proposition) :
{cvs}
"""

    return """Genere une proposition commerciale pour le sujet suivant : {topic}

PUBLIC CIBLE : {audience or 'Direction et decideurs'}
CONSULTANT : {consultant_name} - {company_name}
{doc_section}{ref_section}{cv_section}
Determine toi-meme le nombre de slides et la structure les plus adaptes au besoin client.
Adapte la profondeur et le nombre de slides a la complexite du sujet.

### ELEMENTS A COUVRIR (ordre et decoupage libres) :
- Couverture (cover) : titre impactant, meta "{company_name} | 2025-2026"
- Comprehension du besoin : montrer que tu as compris la realite du client
- Approche et methodologie : comment tu vas repondre au besoin
- Demarche / Phases : diagram flow ou timeline
- Planning / Feuille de route
- Dispositif et gouvernance : equipe, roles, rituels
- Facteurs cles de succes
- Pourquoi {company_name} : valeur ajoutee, differenciants
- Presentation {company_name} : Cabinet de conseil en transformation digitale (Strategy & Vision, Data & Insights, Product Management, Scalability Platform)
- Si des CVs sont fournis : slides CV equipe (type "cv") avec experiences REFORMULEES pour le besoin client
- Si des references sont fournies : slides references (type "highlight") avec resultats concrets et paralleles avec la mission
- Cloture (closing) : contact {consultant_name}

{json_format}

Genere maintenant la proposition pour : {topic}"""


def _get_presentation_system_prompt(company_name):
    return """Tu es un expert en creation de presentations professionnelles pour {company_name}.
Tu crees des presentations visuelles, impactantes et variees.

REGLES STRICTES :
- VARIETE : Alterne les types de slides. Ne jamais avoir 2 slides "content" consecutives.
- VISUELS : Au moins 40% des slides doivent etre visuelles (stat, diagram, highlight, quote)
- CONCIS : Maximum 3-4 bullet points par slide, phrases courtes
- IMPACT : Chiffres-cles en "stat", principes en "quote", comparaisons en "two_column"

Tu retournes UNIQUEMENT un JSON valide."""


def _get_presentation_prompt(topic, audience, slide_count, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme base) :
{document_text[:6000]}
"""

    return """Genere une presentation de {slide_count} slides sur le sujet suivant :

SUJET : {topic}
PUBLIC CIBLE : {audience or 'Managers et decideurs'}
NOMBRE DE SLIDES : {slide_count}
{doc_section}
Structure attendue : commence par cover, puis alterne sections et contenu, termine par closing.

{json_format}

Reponds UNIQUEMENT avec le JSON."""


def _get_formation_system_prompt(company_name):
    return """Tu es un formateur expert chez {company_name}, specialise en data, IA et transformation digitale.
Tu crees des supports de formation pedagogiques, progressifs et engageants.

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.

### REGLES PEDAGOGIQUES
- Progression : du concept a la pratique (theorie > exemples > exercices)
- 1 idee par slide maximum
- Alterner theorie (content), visualisation (diagram, stat), et recapitulatif (highlight)
- Inclure des slides quiz/exercice avec le type "highlight" (key_points comme options)
- Ajouter un champ "tags" (tableau de strings) pour categoriser chaque slide : ex ["THEORIE"], ["PRATIQUE"], ["QUIZ"], ["DEMO"]
- Maximum 3 bullet points par slide, phrases courtes et pedagogiques

### NOMBRE DE SLIDES
- Tu determines toi-meme le nombre de slides adapte au sujet, a la duree et au contenu du brief
- Chaque module doit avoir : intro, theorie, exercice/quiz, recap

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Au moins 40% de slides visuelles (stat, diagram, highlight, quote)
- Ne jamais avoir 2 slides "content" consecutives"""


def _get_formation_prompt(topic, audience, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme matiere pedagogique) :
{document_text[:6000]}
"""

    return """Genere un support de formation complet sur le sujet suivant :

SUJET : {topic or 'Formation basee sur le document fourni'}
PUBLIC CIBLE : {audience or 'Equipes techniques et managers'}
{doc_section}
Determine le nombre de slides adapte au contenu et a la duree mentionnee dans le brief.
Si aucune duree n est precisee, genere 15-20 slides (1 journee).

### STRUCTURE ATTENDUE :
1. Cover : titre de la formation + sous-titre
2. Section "Objectifs" : ce que les apprenants sauront faire
3. Sections thematiques : alterner content, diagram, stat, highlight
4. Pour chaque section : concept > illustration > exercice/quiz
5. Recapitulatif final (highlight avec les points cles)
6. Closing : prochaines etapes + contact

IMPORTANT : Ajoute un champ "tags" (tableau) a chaque slide pour la categoriser.
Exemples de tags : "THEORIE", "PRATIQUE", "QUIZ", "DEMO", "RECAP", "OBJECTIFS"

{json_format}

Reponds UNIQUEMENT avec le JSON."""


def _get_rex_system_prompt(company_name):
    return """Tu es un consultant senior chez {company_name} qui redige un Retour d Experience (REX) de mission.

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.

### REGLES DU REX
- Ton factuel et professionnel, oriente resultats
- Chiffres concrets (KPIs, delais, gains, ROI)
- Structure : Contexte > Enjeux > Demarche > Realisations > Resultats > Apprentissages > Recommandations
- Utilise des types visuels varies : stat pour les resultats, diagram pour la demarche, timeline pour le planning
- Maximum 3-4 bullet points par slide
- Ajoute un champ "tags" a chaque slide : "CONTEXTE", "ENJEUX", "DEMARCHE", "RESULTATS", "APPRENTISSAGES"

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Au moins 40% de slides visuelles (stat, diagram, highlight, timeline)"""


def _get_rex_prompt(topic, client, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme base pour le REX) :
{document_text[:6000]}
"""

    return """Genere un Retour d Experience (REX) de mission pour :

MISSION : {topic or 'Mission basee sur le document fourni'}
CLIENT / CONTEXTE : {client or 'Client entreprise'}
{doc_section}
### STRUCTURE OBLIGATOIRE (10 slides) :

**Slide 1 : Couverture (cover)**
- Titre : "REX Mission - [Nom de la mission]"
- Sous-titre avec le client et la periode

**Slide 2 : Contexte Client (content)**
- Qui est le client, son secteur, sa taille
- Contexte de la mission (appel d offre, besoin identifie)
- Tags: ["CONTEXTE"]

**Slide 3 : Enjeux & Objectifs (highlight)**
- 3-4 enjeux business et objectifs mesurables
- Tags: ["ENJEUX"]

**Slide 4 : Notre Demarche (diagram)**
- Type "flow" : phases de la mission (Cadrage > Audit > Design > Build > Run)
- Tags: ["DEMARCHE"]

**Slide 5 : Planning & Jalons (diagram)**
- Type "timeline" : jalons cles de la mission
- Tags: ["DEMARCHE"]

**Slide 6 : Realisations Cles (highlight)**
- 3-4 livrables ou realisations majeures
- Tags: ["RESULTATS"]

**Slide 7 : Resultats Chiffres (stat)**
- KPI principal avec valeur impactante (ex: +35 pourcent, ROI x3, etc.)
- Tags: ["RESULTATS"]

**Slide 8 : Apprentissages (two_column)**
- Colonne gauche : "Ce qui a bien marche" (facteurs de succes)
- Colonne droite : "Points d'amelioration" (lessons learned)
- Tags: ["APPRENTISSAGES"]

**Slide 9 : Recommandations (content)**
- 3-4 recommandations pour la suite ou pour de futures missions similaires
- Tags: ["APPRENTISSAGES"]

**Slide 10 : Cloture (closing)**
- "Merci" + contact

{json_format}

Reponds UNIQUEMENT avec le JSON."""


@router.post("/api/slide-editor/generate")
@limiter.limit("5/minute")
async def api_slide_editor_generate(request: Request):
    """Genere des slides JSON via LLM pour le slide editor (non-streaming fallback)"""
    try:
        data = await request.json()

        # Valider et sanitizer les inputs texte
        topic = sanitize_text_input(data.get("topic", ""), max_length=1000, field_name="topic")
        audience = sanitize_text_input(
            data.get("audience", ""), max_length=500, field_name="audience"
        )
        document_text = sanitize_text_input(
            data.get("document_text", ""), max_length=50000, field_name="document_text"
        )

        slide_count = data.get("slide_count", 10)
        gen_type = data.get("type", "presentation")
        model = data.get("model", "")

        from utils.llm_client import LLMClient

        llm_kwargs = {"max_tokens": 8192}
        if model:
            llm_kwargs["model"] = model
            llm_kwargs["provider"] = "gemini"
        llm = LLMClient(**llm_kwargs)

        consultant_name = os.getenv("CONSULTANT_NAME", "Jean-Sebastien Abessouguie Bayiha")
        company_name = os.getenv("COMPANY_NAME", "Consulting Tools")
        json_format = _get_slide_json_format()

        if gen_type == "proposal":
            system_prompt = _get_proposal_system_prompt(company_name)
            prompt = _get_proposal_prompt(
                topic, audience, company_name, consultant_name, json_format, document_text
            )
        elif gen_type == "formation":
            system_prompt = _get_formation_system_prompt(company_name)
            prompt = _get_formation_prompt(topic, audience, json_format, document_text)
        elif gen_type == "rex":
            system_prompt = _get_rex_system_prompt(company_name)
            prompt = _get_rex_prompt(topic, audience, json_format, document_text)
        else:
            system_prompt = _get_presentation_system_prompt(company_name)
            prompt = _get_presentation_prompt(
                topic, audience, slide_count, json_format, document_text
            )

        response = llm.generate(prompt=prompt, system_prompt=system_prompt, temperature=0.7)

        if "```json" in response:
            json_str = response.split("```json")[1].split("```")[0].strip()
        elif "```" in response:
            json_str = response.split("```")[1].split("```")[0].strip()
        else:
            json_str = response.strip()

        result = json.loads(json_str)
        return JSONResponse(result)

    except json.JSONDecodeError as e:
        print(f"Erreur JSON slide-editor: {e}")
        return JSONResponse({"error": "Erreur de parsing JSON", "slides": []}, status_code=422)
    except Exception as e:
        print(f"Erreur slide-editor generate: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)
