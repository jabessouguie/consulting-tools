"""
ThemeManager - Gestion de la charte graphique et des settings de l'application
"""
import json
import os
from pathlib import Path
from typing import Any, Dict

BASE_DIR = Path(__file__).parent.parent

DEFAULT_SETTINGS: Dict[str, Any] = {
    "app_name": os.getenv("APP_NAME", "Consulting Tools"),
    "app_tagline": os.getenv("APP_TAGLINE", "Agents IA"),
    "llm_provider": os.getenv("LLM_PROVIDER", "gemini"),
    "gemini_model": os.getenv("GEMINI_MODEL", "gemini-2.5-flash"),
    "openai_model": os.getenv("OPENAI_MODEL", "gpt-4o"),
    "claude_model": os.getenv("CLAUDE_MODEL", "claude-sonnet-4-6"),
    "microsoft_tenant_id": os.getenv("MICROSOFT_TENANT_ID", ""),
    "microsoft_client_id": os.getenv("MICROSOFT_CLIENT_ID", ""),
    "microsoft_client_secret": os.getenv("MICROSOFT_CLIENT_SECRET", ""),
    "teams_channel_id": os.getenv("TEAMS_CHANNEL_ID", ""),
    "onedrive_folder": os.getenv("ONEDRIVE_FOLDER", ""),
    "theme": {
        "title_font": "Chakra Petch",
        "body_font": "Inter",
        "primary_color": "#FF6B58",
        "secondary_color": "#FBF0F4",
        "background_color": "#1F1F1F",
        "text_color": "#474747",
        "accent_color": "#3A3A3B",
    },
}

SETTINGS_PATH = BASE_DIR / "data" / "settings.json"


class ThemeManager:
    """Gestion centralisee des settings et de la charte graphique."""

    @staticmethod
    def load() -> Dict[str, Any]:
        """
        Charge les settings depuis settings.json.
        Retourne les defaults si le fichier est absent ou invalide.

        Returns:
            Dict avec app_name, app_tagline, llm_provider, theme, etc.
        """
        if not SETTINGS_PATH.exists():
            return dict(DEFAULT_SETTINGS)

        try:
            with open(SETTINGS_PATH, "r", encoding="utf-8") as f:
                data = json.load(f)
            # Merge avec defaults pour les cles manquantes
            merged = dict(DEFAULT_SETTINGS)
            merged.update(data)
            if "theme" in data:
                merged["theme"] = dict(DEFAULT_SETTINGS["theme"])
                merged["theme"].update(data["theme"])
            return merged
        except (json.JSONDecodeError, OSError):
            return dict(DEFAULT_SETTINGS)

    @staticmethod
    def save(settings: Dict[str, Any]) -> None:
        """
        Sauvegarde les settings dans settings.json.

        Args:
            settings: Dict des settings a sauvegarder
        """
        SETTINGS_PATH.parent.mkdir(parents=True, exist_ok=True)
        with open(SETTINGS_PATH, "w", encoding="utf-8") as f:
            json.dump(settings, f, indent=2, ensure_ascii=False)

    @staticmethod
    def get_css_vars() -> Dict[str, str]:
        """
        Retourne les variables CSS a partir du theme actuel.

        Returns:
            Dict mapping nom_variable → valeur CSS
        """
        settings = ThemeManager.load()
        theme = settings.get("theme", DEFAULT_SETTINGS["theme"])
        return {
            "primary_color": theme.get("primary_color", "#FF6B58"),
            "secondary_color": theme.get("secondary_color", "#FBF0F4"),
            "background_color": theme.get("background_color", "#1F1F1F"),
            "text_color": theme.get("text_color", "#474747"),
            "accent_color": theme.get("accent_color", "#3A3A3B"),
            "title_font": theme.get("title_font", "Chakra Petch"),
            "body_font": theme.get("body_font", "Inter"),
        }

    @staticmethod
    def import_from_pptx(file_path: str) -> Dict[str, str]:
        """
        Extrait les couleurs et polices depuis un fichier PPTX.

        Args:
            file_path: Chemin vers le fichier PPTX

        Returns:
            Dict theme extrait {primary_color, title_font, body_font, ...}
        """
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor

            prs = Presentation(file_path)
            theme_data: Dict[str, str] = {}

            # Extraire les couleurs du theme PPTX
            slide_master = prs.slide_master
            theme_element = slide_master.theme_color_map

            # Extraire depuis les slides directement si le theme n'est pas accessible
            accent_colors = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.type:
                                    try:
                                        rgb = run.font.color.rgb
                                        hex_color = str(rgb)
                                        if hex_color and hex_color != "000000":
                                            accent_colors.append("#" + hex_color)
                                    except Exception:
                                        pass
                    if shape.fill.type is not None:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            hex_color = str(rgb)
                            if hex_color not in ("000000", "FFFFFF", "ffffff"):
                                accent_colors.append("#" + hex_color)
                        except Exception:
                            pass

            if accent_colors:
                theme_data["primary_color"] = accent_colors[0]
                if len(accent_colors) > 1:
                    theme_data["accent_color"] = accent_colors[1]

            # Extraire les polices
            fonts = set()
            for layout in list(slide_master.slide_layouts)[:3]:
                for ph in layout.placeholders:
                    if ph.has_text_frame:
                        for para in ph.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)

            if fonts:
                font_list = list(fonts)
                theme_data["title_font"] = font_list[0]
                if len(font_list) > 1:
                    theme_data["body_font"] = font_list[1]

            return theme_data

        except ImportError:
            return {}
        except Exception:
            return {}

    @staticmethod
    def import_from_pdf(file_path: str) -> Dict[str, str]:
        """
        Extrait les couleurs dominantes depuis la premiere page d'un PDF.

        Args:
            file_path: Chemin vers le fichier PDF

        Returns:
            Dict theme extrait {primary_color, secondary_color, ...}
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {}

        try:
            doc = fitz.open(file_path)
            if len(doc) == 0:
                return {}

            page = doc[0]
            # Rendu haute resolution pour meilleure extraction des couleurs
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")

            # Extraire les couleurs dominantes
            try:
                from PIL import Image
                import io

                img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                img_small = img.resize((100, 100))
                pixels = list(img_small.getdata())

                # Compter les couleurs
                color_counts: Dict[tuple, int] = {}
                for r, g, b in pixels:
                    # Ignorer le blanc et le noir
                    if (r, g, b) in ((255, 255, 255), (0, 0, 0)):
                        continue
                    if r > 240 and g > 240 and b > 240:
                        continue
                    if r < 15 and g < 15 and b < 15:
                        continue
                    key = (r // 32 * 32, g // 32 * 32, b // 32 * 32)
                    color_counts[key] = color_counts.get(key, 0) + 1

                if color_counts:
                    sorted_colors = sorted(color_counts.items(), key=lambda x: x[1], reverse=True)
                    dominant = sorted_colors[0][0]
                    primary = "#{:02X}{:02X}{:02X}".format(*dominant)

                    theme_data = {"primary_color": primary}
                    if len(sorted_colors) > 1:
                        second = sorted_colors[1][0]
                        theme_data["accent_color"] = "#{:02X}{:02X}{:02X}".format(*second)

                    return theme_data

            except ImportError:
                pass

            return {}

        except Exception:
            return {}
