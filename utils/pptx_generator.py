"""
Generateur de presentations PPTX
Utilise le template PPTX pour creer des propositions commerciales esthetiques
"""

import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Chemin par defaut du branding
BRANDING_FILE = Path(__file__).parent.parent / "data" / "branding.json"


def _load_branding() -> Dict[str, Any]:
    """Charge le branding depuis branding.json"""
    try:
        if BRANDING_FILE.exists():
            with open(BRANDING_FILE, "r") as f:
                return json.load(f)
    except Exception:
        pass
    # Fallback theme defaults
    return {
        "colors": {
            "anthracite": "3A3A3B",
            "noir_profond": "1F1F1F",
            "gris_fonce": "2D2D2D",
            "gris_moyen": "474747",
            "gris_clair": "EEEEEE",
            "rose_poudre": "FBF0F4",
            "terracotta": "C0504D",
            "corail": "FF6B58",
            "blanc": "FFFFFF",
        },
        "fonts": {"title": "Chakra Petch", "body": "Inter"},
    }


BRANDING = _load_branding()


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convertit hex (ex: 'FF6B58') en RGBColor"""
    hex_str = hex_str.replace("#", "").replace(",", "")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# Palette dynamique
COLORS = {k: _hex_to_rgb(v) for k, v in BRANDING["colors"].items()}

# Polices dynamiques
FONT_TITLE = BRANDING["fonts"].get("title", "Chakra Petch")
FONT_BODY = BRANDING["fonts"].get("body", "Inter")

# Dimensions de la slide (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layouts du template PPTX
LAYOUT_CONTENT = 0  # TITLE_ONLY : titre + sous-titre + corps
LAYOUT_CONTENT_IMAGE = 1  # TITLE_ONLY_4 : titre + image + corps
LAYOUT_TWO_COL = 4  # TITLE_ONLY_3_1_1 : titre + 2 colonnes corps
LAYOUT_SECTION = 15  # Section - Couverture : titre section
LAYOUT_BLANK = 16  # BLANK
LAYOUT_TITLE_ONLY = 17  # Corps - Titre seul


class ProposalPPTXGenerator:
    """Generateur de PPTX pour propositions commerciales"""

    def __init__(self, template_path: str):
        """
        Initialise le generateur avec le template Consulting Tools

        Args:
            template_path: Chemin vers le fichier PPTX template
        """
        self.template_path = template_path
        self.prs = Presentation(template_path)

        # Supprimer les slides existantes du template (on garde juste les
        # layouts)
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def _set_text_style(
        self, run, font_name=FONT_BODY, size=Pt(14), color=None, bold=False, italic=False
    ):
        """Applique un style a un run de texte"""
        run.font.name = font_name
        run.font.size = size
        if color:
            run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic

    def _add_paragraph(
        self,
        text_frame,
        text,
        font_name=FONT_BODY,
        size=Pt(14),
        color=None,
        bold=False,
        alignment=PP_ALIGN.LEFT,
        space_after=Pt(6),
    ):
        """Ajoute un paragraphe stylise a un text frame"""
        p = text_frame.add_paragraph()
        p.alignment = alignment
        p.space_after = space_after
        run = p.add_run()
        run.text = text
        self._set_text_style(run, font_name, size, color, bold)
        return p

    def add_cover_slide(
        self, client_name: str, project_title: str, date: str, consultant_name: str
    ):
        """Ajoute la slide de couverture"""
        layout = self.prs.slide_layouts[LAYOUT_SECTION]
        slide = self.prs.slides.add_slide(layout)

        # Titre
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # TITLE
                ph.text = ""
                tf = ph.text_frame
                tf.clear()

                # Ligne 1 : "PROPOSITION COMMERCIALE"
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = "PROPOSITION COMMERCIALE"
                self._set_text_style(run, FONT_TITLE, Pt(28), COLORS["corail"], bold=True)

                # Ligne 2 : nom du projet
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.LEFT
                p2.space_before = Pt(12)
                run2 = p2.add_run()
                run2.text = project_title
                self._set_text_style(run2, FONT_TITLE, Pt(22), COLORS["blanc"], bold=True)

            elif ph.placeholder_format.idx == 1:  # SUBTITLE
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"{client_name}  |  {date}  |  {consultant_name}"
                self._set_text_style(run, FONT_BODY, Pt(14), COLORS["blanc"])

    def add_section_slide(self, section_title: str, section_number: int):
        """Ajoute une slide de section (separateur)"""
        layout = self.prs.slide_layouts[LAYOUT_SECTION]
        slide = self.prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # TITLE
                ph.text = ""
                tf = ph.text_frame
                tf.clear()

                # Numero
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"0{section_number}" if section_number < 10 else str(section_number)
                self._set_text_style(run, FONT_TITLE, Pt(48), COLORS["corail"], bold=True)

                # Titre
                p2 = tf.add_paragraph()
                p2.space_before = Pt(8)
                run2 = p2.add_run()
                run2.text = section_title.upper()
                self._set_text_style(run2, FONT_TITLE, Pt(28), COLORS["blanc"], bold=True)

            elif ph.placeholder_format.idx == 1:  # SUBTITLE
                ph.text = ""

    def add_content_slide(self, title: str, bullet_points: List[str], subtitle: str = ""):
        """Ajoute une slide de contenu avec puces"""
        layout = self.prs.slide_layouts[LAYOUT_CONTENT]
        slide = self.prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # TITLE
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

            elif ph.placeholder_format.idx == 1:  # SUBTITLE
                if subtitle:
                    ph.text = ""
                    tf = ph.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = subtitle
                    self._set_text_style(run, FONT_BODY, Pt(12), COLORS["gris_moyen"])
                else:
                    ph.text = ""

            elif ph.placeholder_format.idx == 2:  # BODY
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                tf.word_wrap = True

                for i, point in enumerate(bullet_points):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.space_after = Pt(10)  # Augmenté pour plus d'aération
                    p.level = 0

                    # Puce coloree
                    run_bullet = p.add_run()
                    run_bullet.text = "  "
                    self._set_text_style(run_bullet, FONT_BODY, Pt(11), COLORS["corail"], bold=True)

                    # Texte - Réduit à Pt(11) pour éviter débordements
                    run_text = p.add_run()
                    run_text.text = point
                    self._set_text_style(run_text, FONT_BODY, Pt(11), COLORS["gris_moyen"])

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_points: List[str],
        right_title: str,
        right_points: List[str],
    ):
        """Ajoute une slide avec 2 colonnes"""
        layout = self.prs.slide_layouts[LAYOUT_TWO_COL]
        slide = self.prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # TITLE
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

            elif ph.placeholder_format.idx == 1:  # SUBTITLE
                ph.text = ""

            elif ph.placeholder_format.idx == 2:  # Colonne gauche
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                tf.word_wrap = True

                # Titre colonne
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = left_title
                self._set_text_style(run, FONT_TITLE, Pt(16), COLORS["corail"], bold=True)

                for point in left_points:
                    p = tf.add_paragraph()
                    p.space_after = Pt(6)
                    run = p.add_run()
                    run.text = f"  {point}"
                    self._set_text_style(run, FONT_BODY, Pt(12), COLORS["gris_moyen"])

            elif ph.placeholder_format.idx == 3:  # Colonne droite
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                tf.word_wrap = True

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = right_title
                self._set_text_style(run, FONT_TITLE, Pt(16), COLORS["corail"], bold=True)

                for point in right_points:
                    p = tf.add_paragraph()
                    p.space_after = Pt(6)
                    run = p.add_run()
                    run.text = f"  {point}"
                    self._set_text_style(run, FONT_BODY, Pt(12), COLORS["gris_moyen"])

    def add_table_slide(self, title: str, headers: List[str], rows: List[List[str]]):
        """Ajoute une slide avec un tableau"""
        layout = self.prs.slide_layouts[LAYOUT_TITLE_ONLY]
        slide = self.prs.slides.add_slide(layout)

        # Titre
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

        # Tableau
        num_rows = len(rows) + 1  # +1 pour header
        num_cols = len(headers)
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.5)
        height = Inches(0.4) * num_rows

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # Style header
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT_TITLE
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = COLORS["blanc"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["anthracite"]

        # Style lignes
        for i, row in enumerate(rows):
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONT_BODY
                        run.font.size = Pt(11)
                        run.font.color.rgb = COLORS["gris_moyen"]
                # Alternance de couleur
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["rose_poudre"]

    def add_cv_slide(self, name: str, title: str, experiences: List[str], skills: List[str]):
        """Ajoute une slide CV one-page"""
        layout = self.prs.slide_layouts[LAYOUT_TITLE_ONLY]
        slide = self.prs.slides.add_slide(layout)

        # Titre de la slide
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"CV  |  {name}"
                self._set_text_style(run, FONT_TITLE, Pt(22), COLORS["noir_profond"], bold=True)

        # Barre de titre coloree
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["corail"]
        bar.line.fill.background()
        tf_bar = bar.text_frame
        tf_bar.word_wrap = True
        p = tf_bar.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"   {name}  —  {title}"
        self._set_text_style(run, FONT_TITLE, Pt(16), COLORS["blanc"], bold=True)

        # Zone experiences (colonne gauche)
        exp_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(6.5), Inches(4.5))
        tf_exp = exp_box.text_frame
        tf_exp.word_wrap = True

        p = tf_exp.paragraphs[0]
        run = p.add_run()
        run.text = "EXPERIENCES CLES"
        self._set_text_style(run, FONT_TITLE, Pt(14), COLORS["terracotta"], bold=True)

        for exp in experiences:
            p = tf_exp.add_paragraph()
            p.space_after = Pt(6)
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = f"  {exp}"
            self._set_text_style(run, FONT_BODY, Pt(11), COLORS["gris_moyen"])

        # Zone competences (colonne droite)
        skills_box = slide.shapes.add_textbox(Inches(7.8), Inches(2.3), Inches(4.5), Inches(4.5))
        tf_skills = skills_box.text_frame
        tf_skills.word_wrap = True

        p = tf_skills.paragraphs[0]
        run = p.add_run()
        run.text = "COMPETENCES"
        self._set_text_style(run, FONT_TITLE, Pt(14), COLORS["terracotta"], bold=True)

        for skill in skills:
            p = tf_skills.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = f"  {skill}"
            self._set_text_style(run, FONT_BODY, Pt(11), COLORS["gris_moyen"])

    def add_image_slide(
        self, title: str, image_path: str, caption: str = "", bullets: List[str] = None
    ):
        """Ajoute une slide avec une image/schema"""
        layout = self.prs.slide_layouts[LAYOUT_CONTENT_IMAGE]
        slide = self.prs.slides.add_slide(layout)

        # Titre
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

        # Ajouter l'image si le fichier existe
        import os

        if os.path.exists(image_path):
            try:
                # Position de l'image (centre-gauche)
                left = Inches(0.8)
                top = Inches(1.8)
                width = Inches(6.5)
                slide.shapes.add_picture(image_path, left, top, width=width)
            except Exception as e:
                print(f"   Erreur ajout image: {e}")

        # Zone de texte pour caption et bullets (droite)
        text_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.5), Inches(4.8))
        tf = text_box.text_frame
        tf.word_wrap = True

        if caption:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = caption
            self._set_text_style(run, FONT_BODY, Pt(11), COLORS["gris_moyen"], italic=True)

        if bullets:
            for bullet in bullets:
                p = tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = f"• {bullet}"
                self._set_text_style(run, FONT_BODY, Pt(12), COLORS["gris_moyen"])

    def add_diagram_slide(
        self, title: str, diagram_type: str, elements: List[str], description: str = ""
    ):
        """
        Ajoute une slide avec un diagramme simple (texte uniquement pour l'instant)
        Les vrais diagrammes nécessiteraient une bibliothèque de génération de schémas

        Args:
            title: Titre de la slide
            diagram_type: Type de diagramme ('flow', 'cycle', 'pyramid', 'timeline')
            elements: Liste des éléments du diagramme
            description: Description optionnelle
        """
        layout = self.prs.slide_layouts[LAYOUT_TITLE_ONLY]
        slide = self.prs.slides.add_slide(layout)

        # Titre
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

        # Représentation textuelle du diagramme
        y_pos = Inches(1.8)
        Inches(2.5)

        if diagram_type == "flow":
            # Diagramme en flux horizontal - VERSION AMÉLIORÉE
            num_elements = len(elements)
            # Espacement adaptatif selon le nombre d'éléments
            box_width = Inches(2.2)
            arrow_width = Inches(0.6)
            total_width = num_elements * box_width + (num_elements - 1) * arrow_width
            start_x = Inches((13.33 - total_width / Inches(1)) / 2)  # Centrer

            for i, element in enumerate(elements):
                left = start_x + i * (box_width + arrow_width)

                # Boîte avec ombre et couleurs vives
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left, y_pos, box_width, Inches(1.0)
                )
                shape.fill.solid()
                # Couleurs alternées pour plus de dynamisme
                if i % 2 == 0:
                    shape.fill.fore_color.rgb = COLORS["corail"]
                else:
                    shape.fill.fore_color.rgb = COLORS["terracotta"]
                shape.line.color.rgb = COLORS["anthracite"]
                shape.line.width = Pt(2.5)
                shape.shadow.inherit = False

                tf = shape.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = element
                self._set_text_style(run, FONT_BODY, Pt(13), COLORS["blanc"], bold=True)

                # Flèche plus visible entre les éléments
                if i < num_elements - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        left + box_width,
                        y_pos + Inches(0.35),
                        arrow_width,
                        Inches(0.3),
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = COLORS["anthracite"]
                    arrow.line.fill.background()

        elif diagram_type == "pyramid":
            # Diagramme pyramidal - VERSION AMÉLIORÉE (plus beau, plus grand)
            pyramid_levels = len(elements)
            start_y = Inches(1.8)
            level_height = Inches(0.9)  # Plus haut pour plus de lisibilité

            for i, element in enumerate(elements):
                # Plus large en bas, plus étroit en haut
                level_index = pyramid_levels - i - 1
                width_ratio = (level_index + 1) / pyramid_levels
                width = Inches(9.5 * width_ratio)  # Plus large
                left = Inches(6.66 - 4.75 * width_ratio)  # Centre
                top = start_y + i * level_height

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.8)
                )
                shape.fill.solid()
                # Gradient de couleur progressif du haut vers le bas
                if i == 0:
                    shape.fill.fore_color.rgb = COLORS["terracotta"]
                elif i == pyramid_levels - 1:
                    shape.fill.fore_color.rgb = RGBColor(230, 200, 200)  # Rose très clair
                else:
                    # Interpolation entre terracotta et rose clair
                    progress = i / (pyramid_levels - 1)
                    r = int(192 * (1 - progress) + 230 * progress)
                    g = int(80 * (1 - progress) + 200 * progress)
                    b = int(77 * (1 - progress) + 200 * progress)
                    shape.fill.fore_color.rgb = RGBColor(r, g, b)

                shape.line.color.rgb = COLORS["anthracite"]
                shape.line.width = Pt(2.5)
                shape.shadow.inherit = False

                tf = shape.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = element
                # Texte blanc en haut, foncé en bas
                text_color = COLORS["blanc"] if i < pyramid_levels // 2 else COLORS["noir_profond"]
                self._set_text_style(run, FONT_BODY, Pt(13), text_color, bold=True)

        elif diagram_type == "timeline":
            # Diagramme timeline horizontal - VERSION AMÉLIORÉE
            total_width = Inches(10.0)
            start_x = Inches(1.66)
            y_pos = Inches(3.5)
            num_elements = len(elements)

            # Ligne de temps plus épaisse
            timeline_line = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                start_x,
                y_pos - Inches(0.025),
                total_width,
                Inches(0.1),
            )
            timeline_line.fill.solid()
            timeline_line.fill.fore_color.rgb = COLORS["anthracite"]
            timeline_line.line.fill.background()

            # Points sur la timeline
            spacing = total_width / (num_elements - 1) if num_elements > 1 else total_width
            for i, element in enumerate(elements):
                x_pos = start_x + (spacing * i if num_elements > 1 else total_width / 2)

                # Point plus gros avec effet
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x_pos - Inches(0.22),
                    y_pos - Inches(0.22),
                    Inches(0.44),
                    Inches(0.44),
                )
                circle.fill.solid()
                # Couleur progressive
                if i == 0:
                    circle.fill.fore_color.rgb = COLORS["terracotta"]
                elif i == num_elements - 1:
                    circle.fill.fore_color.rgb = COLORS["corail"]
                else:
                    circle.fill.fore_color.rgb = COLORS["corail"]
                circle.line.color.rgb = COLORS["anthracite"]
                circle.line.width = Pt(3)
                circle.shadow.inherit = False

                # Label dans une boîte colorée alternée haut/bas
                is_top = i % 2 == 0
                label_y = y_pos - Inches(1.3) if is_top else y_pos + Inches(0.5)

                # Boîte pour le label
                label_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x_pos - Inches(0.9),
                    label_y,
                    Inches(1.8),
                    Inches(0.65),
                )
                label_shape.fill.solid()
                label_shape.fill.fore_color.rgb = COLORS["rose_poudre"]
                label_shape.line.color.rgb = COLORS["terracotta"]
                label_shape.line.width = Pt(2)

                # Texte du label
                tf = label_shape.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = element
                self._set_text_style(run, FONT_BODY, Pt(11), COLORS["noir_profond"], bold=True)

        elif diagram_type == "matrix":
            # Diagramme matrice 2x2 (ou 2x3)
            # Diviser les elements en quadrants
            matrix_size = min(len(elements), 6)  # Max 6 elements (2x3)
            cols = (matrix_size + 1) // 2

            cell_width = Inches(4.0)
            cell_height = Inches(2.0)
            start_x = Inches(2.5)
            start_y = Inches(2.0)

            for i, element in enumerate(elements[:matrix_size]):
                row = i // cols
                col = i % cols

                x = start_x + col * cell_width
                y = start_y + row * cell_height

                cell = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x,
                    y,
                    cell_width - Inches(0.2),
                    cell_height - Inches(0.2),
                )
                cell.fill.solid()
                # Couleur alternee
                if (row + col) % 2 == 0:
                    cell.fill.fore_color.rgb = COLORS["rose_poudre"]
                else:
                    cell.fill.fore_color.rgb = COLORS["gris_clair"]
                cell.line.color.rgb = COLORS["terracotta"]
                cell.line.width = Pt(2)

                tf = cell.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = element
                self._set_text_style(run, FONT_BODY, Pt(11), COLORS["gris_moyen"], bold=True)

        elif diagram_type == "cycle":
            # Diagramme circulaire - VERSION AMÉLIORÉE (vrai cycle avec
            # flèches)
            import math

            num_elements = len(elements)
            center_x = Inches(6.66)
            center_y = Inches(3.8)
            radius = Inches(2.2)

            # Positionner les boîtes en cercle
            for i, element in enumerate(elements):
                # Angle pour positionner en cercle (commencer en haut, sens
                # horaire)
                angle = (i / num_elements) * 2 * math.pi - math.pi / 2

                # Position de la boîte
                x = center_x + radius * math.cos(angle) - Inches(1.0)
                y = center_y + radius * math.sin(angle) - Inches(0.45)

                # Boîte du cycle
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(0.9)
                )
                shape.fill.solid()
                # Couleur dégradée selon la position
                colors = [
                    COLORS["corail"],
                    COLORS["terracotta"],
                    COLORS["rose_poudre"],
                    COLORS["corail"],
                ]
                shape.fill.fore_color.rgb = colors[i % len(colors)]
                shape.line.color.rgb = COLORS["anthracite"]
                shape.line.width = Pt(2.5)

                tf = shape.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = f"{i + 1}. {element}"
                self._set_text_style(run, FONT_BODY, Pt(11), COLORS["blanc"], bold=True)

            # Ajouter des flèches circulaires entre les éléments
            for i in range(num_elements):
                angle_start = (i / num_elements) * 2 * math.pi - math.pi / 2
                angle_end = ((i + 1) / num_elements) * 2 * math.pi - math.pi / 2

                # Position de la flèche (entre deux boîtes)
                mid_angle = (angle_start + angle_end) / 2
                arrow_x = center_x + (radius + Inches(0.3)) * math.cos(mid_angle) - Inches(0.15)
                arrow_y = center_y + (radius + Inches(0.3)) * math.sin(mid_angle) - Inches(0.15)

                # Flèche courbée (simulée par une flèche tournée)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.4), Inches(0.25)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = COLORS["anthracite"]
                arrow.line.fill.background()
                # Rotation de la flèche selon l'angle
                arrow.rotation = math.degrees(mid_angle + math.pi / 2)

        else:
            # Par défaut, liste simple
            y = Inches(1.8)
            for i, element in enumerate(elements):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(2.0),
                    y + i * Inches(0.9),
                    Inches(8.0),
                    Inches(0.7),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLORS["rose_poudre"]
                shape.line.color.rgb = COLORS["terracotta"]

                tf = shape.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"  {element}"
                self._set_text_style(run, FONT_BODY, Pt(12), COLORS["gris_moyen"])

        # Description en bas
        if description:
            desc_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(10.0), Inches(0.8))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = description
            self._set_text_style(run, FONT_BODY, Pt(11), COLORS["gris_moyen"], italic=True)

    def add_stat_slide(
        self, stat_value: str, stat_label: str, context: str = "", subtitle: str = ""
    ):
        """
        Ajoute une slide de stat/chiffre clé impactante (style Veolia)
        Met en avant un chiffre important de manière visuelle

        Args:
            stat_value: Le chiffre/stat à mettre en avant (ex: "67%", "15M€", "3 mois")
            stat_label: Label du chiffre (ex: "de ROI", "de chiffre d'affaires")
            context: Contexte additionnel (1-2 lignes max)
            subtitle: Sous-titre optionnel en haut
        """
        layout = self.prs.slide_layouts[LAYOUT_BLANK]
        slide = self.prs.slides.add_slide(layout)

        # Fond subtil
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLORS["blanc"]
        bg_shape.line.fill.background()

        # Barre décorative gauche
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_HEIGHT)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS["corail"]
        accent_bar.line.fill.background()

        # Sous-titre en haut si présent
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(0.8), Inches(10.0), Inches(0.6)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = subtitle
            self._set_text_style(run, FONT_BODY, Pt(14), COLORS["gris_moyen"])

        # Chiffre principal - TRÈS GRAND
        stat_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(2.5))
        tf = stat_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat_value
        self._set_text_style(run, FONT_TITLE, Pt(120), COLORS["corail"], bold=True)

        # Label du chiffre
        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.0), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat_label
        self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

        # Contexte en bas
        if context:
            context_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(5.8), Inches(8.0), Inches(1.0)
            )
            tf = context_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = context
            self._set_text_style(run, FONT_BODY, Pt(13), COLORS["gris_moyen"])

    def add_quote_slide(self, quote_text: str, author: str = "", title: str = ""):
        """
        Ajoute une slide de citation/key message impactante

        Args:
            quote_text: Le texte de la citation ou message clé
            author: Auteur de la citation (optionnel)
            title: Titre de la slide (optionnel)
        """
        layout = self.prs.slide_layouts[LAYOUT_BLANK]
        slide = self.prs.slides.add_slide(layout)

        # Fond
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLORS["rose_poudre"]
        bg_shape.line.fill.background()

        # Titre si présent
        y_offset = Inches(1.5)
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(0.8), Inches(10.0), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            self._set_text_style(run, FONT_TITLE, Pt(20), COLORS["terracotta"], bold=True)
            y_offset = Inches(2.0)

        # Guillemets décoratifs
        quote_mark = slide.shapes.add_textbox(
            Inches(2.0), y_offset - Inches(0.3), Inches(1.0), Inches(1.0)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '"'
        self._set_text_style(run, FONT_TITLE, Pt(80), COLORS["corail"], bold=True)

        # Citation principale
        quote_box = slide.shapes.add_textbox(Inches(2.5), y_offset, Inches(8.0), Inches(3.0))
        tf = quote_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = quote_text
        self._set_text_style(run, FONT_BODY, Pt(20), COLORS["noir_profond"], italic=True)

        # Auteur
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(2.5), y_offset + Inches(3.2), Inches(8.0), Inches(0.6)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"— {author}"
            self._set_text_style(run, FONT_BODY, Pt(14), COLORS["gris_fonce"], bold=True)

    def add_highlight_slide(
        self, title: str, key_points: List[str], highlight_color: str = "corail"
    ):
        """
        Ajoute une slide avec points clés mis en évidence dans des encadrés colorés
        Plus visuel et impactant qu'une slide de contenu classique

        Args:
            title: Titre de la slide
            key_points: Liste de 2-4 points clés à mettre en avant
            highlight_color: Couleur des encadrés ('corail', 'terracotta', 'rose_poudre')
        """
        layout = self.prs.slide_layouts[LAYOUT_TITLE_ONLY]
        slide = self.prs.slides.add_slide(layout)

        # Titre
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                self._set_text_style(run, FONT_TITLE, Pt(24), COLORS["noir_profond"], bold=True)

        # Couleur des encadrés
        box_color = COLORS.get(highlight_color, COLORS["corail"])

        # Disposition des points clés en grille
        num_points = min(len(key_points), 4)

        if num_points <= 2:
            # 2 colonnes
            cols = 2
            box_width = Inches(5.0)
            box_height = Inches(3.5)
        elif num_points == 3:
            # 3 colonnes
            cols = 3
            box_width = Inches(3.5)
            box_height = Inches(3.5)
        else:
            # 2x2
            cols = 2
            box_width = Inches(5.0)
            box_height = Inches(2.2)

        start_x = Inches(1.0)
        start_y = Inches(2.0)
        h_spacing = Inches(0.3)
        v_spacing = Inches(0.3)

        for i, point in enumerate(key_points[:4]):
            row = i // cols
            col = i % cols

            x = start_x + col * (box_width + h_spacing)
            y = start_y + row * (box_height + v_spacing)

            # Encadré coloré
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = box_color
            box.line.color.rgb = COLORS["anthracite"]
            box.line.width = Pt(2)
            box.shadow.inherit = False

            # Numéro en grand
            num_tf = slide.shapes.add_textbox(
                x + Inches(0.3), y + Inches(0.2), Inches(1.0), Inches(0.8)
            )
            tf = num_tf.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(i + 1)
            self._set_text_style(run, FONT_TITLE, Pt(48), COLORS["blanc"], bold=True)

            # Texte du point
            text_tf = slide.shapes.add_textbox(
                x + Inches(0.4), y + Inches(1.0), box_width - Inches(0.8), box_height - Inches(1.2)
            )
            tf = text_tf.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = point
            self._set_text_style(run, FONT_BODY, Pt(13), COLORS["blanc"], bold=True)

    def add_closing_slide(self, consultant_name: str, company: str, contact_info: str = ""):
        """Ajoute la slide de cloture"""
        layout = self.prs.slide_layouts[LAYOUT_SECTION]
        slide = self.prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = ""
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = "MERCI"
                self._set_text_style(run, FONT_TITLE, Pt(48), COLORS["corail"], bold=True)

                p2 = tf.add_paragraph()
                p2.space_before = Pt(16)
                run2 = p2.add_run()
                run2.text = f"{consultant_name}\n{company}"
                self._set_text_style(run2, FONT_BODY, Pt(18), COLORS["blanc"])

                if contact_info:
                    p3 = tf.add_paragraph()
                    p3.space_before = Pt(8)
                    run3 = p3.add_run()
                    run3.text = contact_info
                    self._set_text_style(run3, FONT_BODY, Pt(14), COLORS["gris_clair"])

            elif ph.placeholder_format.idx == 1:
                ph.text = ""

    def save(self, output_path: str):
        """Sauvegarde la presentation"""
        self.prs.save(output_path)
        print(f"   ✅ PPTX sauvegarde: {output_path}")


def build_proposal_pptx(
    template_path: str,
    slides_data: List[Dict[str, Any]],
    output_path: str,
    consultant_info: Dict[str, str],
) -> str:
    """
    Construit une presentation PPTX a partir de donnees structurees

    Args:
        template_path: Chemin vers le template Consulting Tools
        slides_data: Liste de dictionnaires decrivant chaque slide
        output_path: Chemin de sortie
        consultant_info: Infos du consultant

    Returns:
        Chemin du fichier genere
    """
    gen = ProposalPPTXGenerator(template_path)

    for slide in slides_data:
        slide_type = slide.get("type", "content")

        if slide_type == "cover":
            gen.add_cover_slide(
                client_name=slide.get("client", ""),
                project_title=slide.get("project", ""),
                date=slide.get("date", ""),
                consultant_name=consultant_info.get("name", ""),
            )

        elif slide_type == "section":
            gen.add_section_slide(
                section_title=slide.get("title", ""),
                section_number=slide.get("number", 1),
            )

        elif slide_type == "content":
            gen.add_content_slide(
                title=slide.get("title", ""),
                bullet_points=slide.get("bullets", []),
                subtitle=slide.get("subtitle", ""),
            )

        elif slide_type == "two_column":
            gen.add_two_column_slide(
                title=slide.get("title", ""),
                left_title=slide.get("left_title", ""),
                left_points=slide.get("left_points", []),
                right_title=slide.get("right_title", ""),
                right_points=slide.get("right_points", []),
            )

        elif slide_type == "table":
            gen.add_table_slide(
                title=slide.get("title", ""),
                headers=slide.get("headers", []),
                rows=slide.get("rows", []),
            )

        elif slide_type == "cv":
            gen.add_cv_slide(
                name=slide.get("name", ""),
                title=slide.get("title", ""),
                experiences=slide.get("experiences", []),
                skills=slide.get("skills", []),
            )

        elif slide_type == "image":
            gen.add_image_slide(
                title=slide.get("title", ""),
                image_path=slide.get("image_path", ""),
                caption=slide.get("caption", ""),
                bullets=slide.get("bullets", []),
            )

        elif slide_type == "diagram":
            gen.add_diagram_slide(
                title=slide.get("title", ""),
                diagram_type=slide.get("diagram_type", "flow"),
                elements=slide.get("elements", []),
                description=slide.get("description", ""),
            )

        elif slide_type == "stat":
            gen.add_stat_slide(
                stat_value=slide.get("stat_value", ""),
                stat_label=slide.get("stat_label", ""),
                context=slide.get("context", ""),
                subtitle=slide.get("subtitle", ""),
            )

        elif slide_type == "quote":
            gen.add_quote_slide(
                quote_text=slide.get("quote_text", ""),
                author=slide.get("author", ""),
                title=slide.get("title", ""),
            )

        elif slide_type == "highlight":
            gen.add_highlight_slide(
                title=slide.get("title", ""),
                key_points=slide.get("key_points", []),
                highlight_color=slide.get("highlight_color", "corail"),
            )

        elif slide_type == "closing":
            gen.add_closing_slide(
                consultant_name=consultant_info.get("name", ""),
                company=consultant_info.get("company", ""),
            )

    gen.save(output_path)
    return output_path
