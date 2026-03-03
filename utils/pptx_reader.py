"""
Utilitaire pour lire et extraire le contenu de fichiers PowerPoint (PPTX)
"""

import os
from typing import Any, Dict, List

from pptx import Presentation


def extract_text_from_shape(shape, depth: int = 0) -> List[str]:
    """
    Extrait recursivement le texte d'un shape (y compris les groupes)

    Args:
        shape: Shape PowerPoint
        depth: Profondeur de recursion

    Returns:
        Liste de textes extraits
    """
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
    # Recursion dans les groupes
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child, depth + 1))
    return texts


def read_pptx_template(pptx_path: str) -> Dict[str, Any]:
    """
    Lit un template PowerPoint et extrait sa structure

    Args:
        pptx_path: Chemin vers le fichier PPTX

    Returns:
        Dictionnaire avec la structure du template
    """
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "N/A"
        texts = []

        for shape in slide.shapes:
            extracted = extract_text_from_shape(shape)
            texts.extend(extracted)

        slides_data.append(
            {
                "slide_number": i,
                "layout": layout_name,
                "content": texts,
            }
        )

    return {
        "title": os.path.basename(pptx_path),
        "total_slides": len(slides_data),
        "slides": slides_data,
    }


def read_pptx_reference(pptx_path: str) -> Dict[str, Any]:
    """
    Lit un fichier de reference PPTX et extrait les informations de projet

    Args:
        pptx_path: Chemin vers le fichier PPTX de reference

    Returns:
        Dictionnaire avec les informations du projet
    """
    prs = Presentation(pptx_path)
    all_texts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            all_texts.extend(extract_text_from_shape(shape))

    return {
        "filename": os.path.basename(pptx_path),
        "content": all_texts,
        "full_text": "\n".join(all_texts),
    }


def read_all_references(references_dir: str) -> List[Dict[str, Any]]:
    """
    Lit tous les fichiers PPTX de reference dans un repertoire

    Args:
        references_dir: Chemin vers le repertoire des references

    Returns:
        Liste de dictionnaires avec les informations de chaque reference
    """
    references = []

    if not os.path.exists(references_dir):
        print(f"Repertoire non trouve: {references_dir}")
        return references

    for filename in os.listdir(references_dir):
        if filename.endswith(".pptx"):
            filepath = os.path.join(references_dir, filename)
            try:
                ref = read_pptx_reference(filepath)
                references.append(ref)
                print(f"  Reference chargee: {filename}")
            except Exception as e:
                print(f"  Erreur lors de la lecture de {filename}: {e}")

    return references


def extract_template_structure(pptx_path: str) -> str:
    """
    Extrait la structure du template sous forme de texte lisible
    pour le LLM

    Args:
        pptx_path: Chemin vers le template PPTX

    Returns:
        Texte decrivant la structure du template
    """
    template = read_pptx_template(pptx_path)
    sections = []

    for slide in template["slides"]:
        content = slide["content"]
        if content:
            slide_num = slide['slide_number']
            slide_layout = slide['layout']
            slide_text = f"### Slide {slide_num} ({slide_layout})\n"
            slide_text += "\n".join(f"- {t}" for t in content[:10])
            sections.append(slide_text)

    return f"""# Template: {template['title']}
Nombre total de slides: {template['total_slides']}

{chr(10).join(sections)}"""
