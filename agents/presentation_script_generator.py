"""
Agent de generation de script de presentation
Prend un PPTX et genere le discours pour chaque slide
"""

from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

from config import get_consultant_info
from utils.llm_client import LLMClient


class PresentationScriptGenerator:
    """Genere un script de presentation a partir d un PPTX"""

    def __init__(self):
        self.llm = LLMClient(max_tokens=8192)
        # Informations consultant (depuis config centralisee)
        self.consultant_info = get_consultant_info()

    def extract_slides_content(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Extrait le contenu de chaque slide du PPTX

        Args:
            pptx_path: Chemin vers le fichier PPTX

        Returns:
            Liste de slides avec titre, texte, notes
        """
        prs = Presentation(pptx_path)
        slides_content = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_data = {"slide_number": i, "title": "", "content": [], "notes": ""}

            # Extraire titre
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # Extraire contenu textuel
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Eviter doublons avec le titre
                    if shape != slide.shapes.title:
                        slide_data["content"].append(shape.text.strip())

            # Extraire notes du presentateur
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_data["notes"] = notes_slide.notes_text_frame.text.strip()

            slides_content.append(slide_data)

        return slides_content

    def generate_script_for_slide(self, slide_data: Dict[str, Any], context: str = "") -> str:
        """
        Genere le script de presentation pour une slide

        Args:
            slide_data: Donnees de la slide (titre, contenu, notes)
            context: Contexte de la presentation (theme general)

        Returns:
            Script markdown pour cette slide
        """
        self.consultant_info["name"]
        self.consultant_info["company"]

        system_prompt = """Tu es {consultant_name}, consultant chez {company}.

TON ROLE : Generer un script de presentation oral pour une slide PowerPoint.

OBJECTIF : Aider le presentateur a delivrer un discours fluide, naturel et impactant.

STRUCTURE DU SCRIPT PAR SLIDE :

## Slide {slide_data['slide_number']} : [Titre de la slide]

### 🎯 Message cle
[En 1 phrase : le takeaway de cette slide]

### 🗣️ Script oral (2-3 min)

[Discours a prononcer, repartit en paragraphes courts]

**Points cles a mentionner** :
- Point 1
- Point 2
- Point 3

**Transition vers slide suivante** :
[Phrase de transition naturelle]

### 💡 Conseils de presentation
- Conseil pratique 1 (ex: pointer un element visuel)
- Conseil pratique 2 (ex: faire une pause apres une stat)

---

REGLES D ECRITURE :
- Ton **conversationnel** : ecris comme on parle (pas comme on lit)
- **Phrases courtes** : faciles a prononcer
- **Exemples concrets** quand possible
- **Chiffres cles** a mentionner explicitement
- **Pauses** suggerees avec [...] pour respirer
- Evite le jargon technique sauf si absolument necessaire
- Si slide technique : vulgarise avec une metaphore

TIMING :
- Slide classique : 2-3 min de discours
- Slide de transition : 30s-1min
- Slide avec graphique/stat : 3-4 min (temps pour pointer et expliquer)

STYLE :
- Naturel, comme si tu parlais a un collegue
- Engage l audience avec des questions rhetoriques
- Cree du rythme : varie longueur de phrases
- Humanise les chiffres (ex: "1 million, c est l equivalent de...")

IMPORTANT :
- NE PAS lire betement le texte de la slide (l audience sait lire)
- APPORTER de la valeur ajoutee par le discours oral
- CREER du lien entre les slides pour fluidite
"""

        # Construire le contenu de la slide pour le prompt
        slide_content = """SLIDE {slide_data['slide_number']}

**Titre** : {slide_data['title']}

**Contenu de la slide** :
{chr(10).join(slide_data['content']) if slide_data['content'] else "(slide visuelle sans texte)"}

**Notes presentateur** :
{slide_data['notes'] if slide_data['notes'] else "(aucune note)"}

**Contexte de la presentation** :
{context if context else "(contexte general)"}
"""

        prompt = """{slide_content}

Genere le script de presentation pour cette slide.

Retourne UNIQUEMENT le script formate en markdown, sans preambule."""

        result = self.llm.generate(prompt=prompt, system_prompt=system_prompt, temperature=0.7)

        # Nettoyer markdown fences
        result = result.strip()
        if result.startswith("```markdown"):
            result = result[len("```markdown") :].strip()
        if result.startswith("```"):
            result = result[3:].strip()
        if result.endswith("```"):
            result = result[:-3].strip()

        return result

    def run(self, pptx_path: str, presentation_context: str = "") -> Dict[str, Any]:
        """
        Genere le script complet de presentation

        Args:
            pptx_path: Chemin vers le PPTX
            presentation_context: Contexte general (audience, objectif)

        Returns:
            Dict avec script complet, timing estime, conseils generaux
        """
        print("\n🎬 Generation du script de presentation")
        print(f"   📄 Fichier : {Path(pptx_path).name}")

        # 1. Extraire contenu des slides
        print("\n   📊 Extraction du contenu des slides...")
        slides = self.extract_slides_content(pptx_path)
        print(f"   ✓ {len(slides)} slides extraites")

        # 2. Generer script pour chaque slide
        print("\n   ✍️  Generation du script pour chaque slide...")
        scripts = []

        for i, slide_data in enumerate(slides, start=1):
            print(f"      → Slide {i}/{len(slides)} : {slide_data['title'][:50]}...")

            script = self.generate_script_for_slide(slide_data, presentation_context)
            scripts.append({"slide_number": i, "title": slide_data["title"], "script": script})

        # 3. Assembler le document final
        print("\n   📝 Assemblage du document final...")

        full_script = """# Script de Presentation

**Presentation** : {Path(pptx_path).stem}
**Nombre de slides** : {len(slides)}
**Timing estime** : {len(slides) * 2.5:.0f}-{len(slides) * 3.5:.0f} minutes

{f"**Contexte** : {presentation_context}" if presentation_context else ""}

---

"""

        # Ajouter conseils generaux
        full_script += """## 🎯 Conseils generaux avant de commencer

**Preparation** :
- Repetez au moins 1 fois a voix haute (seul ou devant collegue)
- Chronometrez-vous : adaptez le rythme si besoin
- Anticipez les questions probables

**Pendant la presentation** :
- Gardez le contact visuel avec l audience (ne lisez PAS vos slides)
- Utilisez vos mains pour illustrer (graphiques, concepts)
- Faites des pauses apres les points importants
- Adaptez le rythme selon les reactions de l audience

**Gestion du temps** :
- Si en retard : sautez les details, gardez l essentiel
- Si en avance : approfondissez avec exemples concrets

---

"""

        # Ajouter scripts de chaque slide
        for script_data in scripts:
            full_script += f"\n{script_data['script']}\n\n---\n"

        # Conclusion generale
        full_script += """
## 🎤 Conclusion de la presentation

**Derniere slide (recapitulatif)** :

Terminez par :
1. **Recap en 30s** : les 3 points cles de la presentation
2. **Call to action** : qu attendez-vous de l audience ?
3. **Ouverture Q&A** : "Je suis dispo pour vos questions"

**Apres la presentation** :
- Notez les questions recurrentes → ameliorer slides
- Demandez du feedback a 2-3 personnes
- Ajustez pour prochaine fois

---

*Script genere par WEnvision Agents - Bonne presentation !* 🚀
"""

        result = {
            "markdown": full_script,
            "num_slides": len(slides),
            "estimated_duration": f"{len(slides) * 2.5:.0f}-{len(slides) * 3.5:.0f} min",
            "slides": slides,
        }

        print("\n✅ Script genere avec succes")
        print(f"   📊 {len(slides)} slides traitees")
        print(f"   ⏱️  Duree estimee : {result['estimated_duration']}")

        return result


def main():
    """Test de l agent"""
    import argparse

    parser = argparse.ArgumentParser(description="Genere un script de presentation")
    parser.add_argument("pptx_path", help="Chemin vers le fichier PPTX")
    parser.add_argument("--context", "-c", default="", help="Contexte de la presentation")
    parser.add_argument("--output", "-o", help="Fichier de sortie (optionnel)")

    args = parser.parse_args()

    agent = PresentationScriptGenerator()
    result = agent.run(args.pptx_path, args.context)

    # Sauvegarder
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = Path(args.pptx_path).with_suffix(".script.md")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(result["markdown"])

    print(f"\n📄 Script sauvegarde : {output_path}")


if __name__ == "__main__":
    main()
