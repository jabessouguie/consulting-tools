"""
Agent de génération de propositions commerciales
Analyse un appel d'offre et génère une proposition commerciale en utilisant
les références Wenvision (fichiers PPTX locaux) et un template PowerPoint
"""
import os
import json
from typing import Dict, Any, Optional
from datetime import datetime
import sys

# Ajouter le répertoire parent au path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from dotenv import load_dotenv
load_dotenv()

from typing import List
from utils.llm_client import LLMClient
from utils.pptx_reader import (
    read_pptx_template,
    read_pptx_reference,
    read_all_references,
    extract_template_structure,
)
from utils.pptx_generator import build_proposal_pptx
from utils.image_generator import ImageGenerator, ImageLibrary, DiagramGenerator

# Chemin de base du projet
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


class ProposalGeneratorAgent:
    """Agent pour générer des propositions commerciales"""

    def __init__(self):
        """Initialise l'agent"""
        self.llm_client = LLMClient()

        # Chemins locaux
        self.template_path = os.path.join(
            BASE_DIR, "WENVISION_Template_Palette 2026.pptx"
        )
        self.references_dir = os.path.join(BASE_DIR, "data", "references")
        self.notebooklm_data_path = os.path.join(
            BASE_DIR, "data", "notebooklm", "references.json"
        )
        self.biographies_path = os.path.join(
            BASE_DIR, "Biographies - CV All WEnvision.pptx"
        )

        # Informations consultant
        self.consultant_info = {
            'name': os.getenv('CONSULTANT_NAME', 'Jean-Sébastien Abessouguie Bayiha'),
            'title': os.getenv('CONSULTANT_TITLE', 'Consultant en stratégie data et IA'),
            'company': os.getenv('COMPANY_NAME', 'Wenvision'),
            'profile': os.getenv('CONSULTANT_PROFILE', '')
        }

        # Diagram generation (Claude + Mermaid) et images
        self.diagram_generator = DiagramGenerator(llm_client=self.llm_client)
        self.image_generator = ImageGenerator()  # DALL-E optionnel
        self.image_library = ImageLibrary()

    def load_references(self) -> Dict[str, Any]:
        """
        Charge les références depuis les fichiers locaux:
        1. Fichiers PPTX dans data/references/
        2. Fichier JSON dans data/notebooklm/references.json

        Returns:
            Dictionnaire contenant toutes les références
        """
        result = {
            'pptx_references': [],
            'projects': [],
            'expertise': [],
            'methodologies': [],
            'differentiators': [],
        }

        # 1. Charger les références PPTX
        print("Chargement des références PPTX...")
        pptx_refs = read_all_references(self.references_dir)
        if pptx_refs:
            result['pptx_references'] = pptx_refs
            print(f"   {len(pptx_refs)} référence(s) PPTX chargée(s)")
        else:
            print("   Aucune référence PPTX trouvée dans data/references/")

        # 2. Charger le fichier JSON de références
        if os.path.exists(self.notebooklm_data_path):
            print("Chargement des références JSON (NotebookLM)...")
            with open(self.notebooklm_data_path, 'r', encoding='utf-8') as f:
                json_refs = json.load(f)
            result['projects'] = json_refs.get('projects', [])
            result['expertise'] = json_refs.get('expertise', [])
            result['methodologies'] = json_refs.get('methodologies', [])
            result['differentiators'] = json_refs.get('differentiators', [])
            print(f"   {len(result['projects'])} projet(s) JSON chargé(s)")
        else:
            print(f"   Fichier JSON non trouvé: {self.notebooklm_data_path}")

        return result

    def load_template(self) -> str:
        """
        Charge le template de présentation depuis le fichier PPTX local

        Returns:
            Structure du template sous forme de texte
        """
        print("Chargement du template PowerPoint local...")

        if not os.path.exists(self.template_path):
            print(f"   Template non trouvé: {self.template_path}")
            return None

        template_text = extract_template_structure(self.template_path)
        template_data = read_pptx_template(self.template_path)
        print(f"   Template chargé: {template_data['total_slides']} slides")

        return template_text

    def analyze_tender(self, tender_text: str) -> Dict[str, Any]:
        """
        Analyse un appel d'offre pour extraire les informations clés

        Args:
            tender_text: Texte de l'appel d'offre

        Returns:
            Dictionnaire avec les informations extraites
        """
        print("Analyse de l'appel d'offre...")

        schema = {
            "client_name": "string",
            "project_title": "string",
            "objectives": ["string"],
            "requirements": {
                "technical": ["string"],
                "functional": ["string"],
                "constraints": ["string"]
            },
            "budget": "string or null",
            "timeline": "string or null",
            "deliverables": ["string"],
            "evaluation_criteria": ["string"],
            "keywords": ["string"]
        }

        prompt = f"""Analyse cet appel d'offre et extrais les informations clés de manière structurée:

{tender_text}

Identifie particulièrement:
- Le nom du client et le titre du projet
- Les objectifs principaux
- Les exigences techniques et fonctionnelles
- Les contraintes (budget, délais)
- Les livrables attendus
- Les critères d'évaluation
- Les mots-clés techniques importants"""

        analysis = self.llm_client.extract_structured_data(prompt, schema)
        print(f"   Analyse terminée: {analysis.get('project_title', 'Sans titre')}")

        return analysis

    def match_references(self, tender_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """
        Trouve les références pertinentes par rapport à l'appel d'offre

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Références sélectionnées
        """
        print("Sélection des références pertinentes...")

        references = self.load_references()

        # Construire le contexte des references PPTX
        pptx_context = ""
        for ref in references.get('pptx_references', []):
            pptx_context += f"\n--- Référence: {ref['filename']} ---\n"
            pptx_context += ref['full_text'] + "\n"

        # Construire le contexte des references JSON
        json_context = ""
        for project in references.get('projects', []):
            json_context += f"\n--- Projet: {project.get('title', 'N/A')} ---\n"
            json_context += f"Client: {project.get('client', 'N/A')}\n"
            json_context += f"Secteur: {project.get('sector', 'N/A')}\n"
            json_context += f"Description: {project.get('description', 'N/A')}\n"
            json_context += f"Challenge: {project.get('challenge', 'N/A')}\n"
            json_context += f"Solution: {project.get('solution', 'N/A')}\n"
            json_context += f"Technologies: {', '.join(project.get('technologies', []))}\n"
            json_context += f"Résultats: {', '.join(project.get('results', []))}\n"

        prompt = f"""À partir de cet appel d'offre analysé:
{json.dumps(tender_analysis, indent=2, ensure_ascii=False)}

Et de ces références disponibles:

RÉFÉRENCES PROJETS (PPTX):
{pptx_context}

RÉFÉRENCES PROJETS (JSON):
{json_context}

EXPERTISE WENVISION:
{json.dumps(references.get('expertise', []), indent=2, ensure_ascii=False)}

MÉTHODOLOGIES:
{json.dumps(references.get('methodologies', []), indent=2, ensure_ascii=False)}

DIFFÉRENCIATEURS:
{json.dumps(references.get('differentiators', []), indent=2, ensure_ascii=False)}

Sélectionne et justifie les références les plus pertinentes pour cet appel d'offre.
Pour chaque référence, explique pourquoi elle est pertinente."""

        response = self.llm_client.generate(prompt, temperature=0.5)

        return {
            'selected_references': response,
            'all_references': references
        }

    def load_cvs(self) -> List[Dict[str, Any]]:
        """
        Charge les CVs depuis le fichier Biographies PPTX.
        Chaque slide = un CV.

        Returns:
            Liste de CVs avec texte brut par slide
        """
        print("Chargement des biographies...")

        if not os.path.exists(self.biographies_path):
            print(f"   Fichier biographies non trouvé: {self.biographies_path}")
            return []

        from pptx import Presentation
        prs = Presentation(self.biographies_path)

        cvs = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            full_text = "\n".join(texts)
            if len(full_text) > 50:  # Ignorer les slides quasi vides
                cvs.append({
                    'slide_index': i,
                    'raw_text': full_text,
                })

        print(f"   {len(cvs)} CV(s) chargé(s)")
        return cvs

    def adapt_cvs(self, cvs: List[Dict[str, Any]], tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Reecrit les CVs pour les adapter au contexte de l'appel d'offre

        Args:
            cvs: CVs bruts extraits du PPTX
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            CVs adaptes au format slides
        """
        print("Adaptation des CVs au contexte de la mission...")

        if not cvs:
            return []

        # Envoyer tous les CVs au LLM pour selection et adaptation
        cvs_text = ""
        for i, cv in enumerate(cvs):
            cvs_text += f"\n--- CV {i+1} ---\n{cv['raw_text'][:2000]}\n"

        prompt = f"""A partir de cet appel d'offre:
Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {json.dumps(tender_analysis.get('objectives', []), ensure_ascii=False)}
Mots-cles: {json.dumps(tender_analysis.get('keywords', []), ensure_ascii=False)}

Et de ces CVs de l'equipe Wenvision:
{cvs_text}

Selectionne les 2-3 profils les plus pertinents pour cette mission et pour chacun, genere un CV one-page adapte.

Reponds UNIQUEMENT en JSON valide avec cette structure:
[
  {{
    "name": "Prenom Nom",
    "title": "Titre/Poste",
    "experiences": [
      "Experience 1 reformulee pour mettre en avant la pertinence pour ce projet",
      "Experience 2...",
      "Experience 3...",
      "Experience 4..."
    ],
    "skills": [
      "Competence 1",
      "Competence 2",
      "Competence 3",
      "Competence 4",
      "Competence 5"
    ]
  }}
]

IMPORTANT:
- Reformule les experiences pour mettre en avant ce qui est pertinent pour cette mission
- Ne fabrique PAS d'experiences ou competences inexistantes
- Mets en avant les technologies et methodologies mentionnees dans l'appel d'offre
- Maximum 5 experiences et 6 competences par CV"""

        response = self.llm_client.generate(
            prompt=prompt,
            temperature=0.5,
            max_tokens=3000
        )

        # Parser le JSON
        try:
            # Extraire le JSON de la reponse
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start >= 0 and json_end > json_start:
                adapted_cvs = json.loads(response[json_start:json_end])
                print(f"   {len(adapted_cvs)} CV(s) adapté(s)")
                return adapted_cvs
        except json.JSONDecodeError:
            print("   Erreur de parsing des CVs, utilisation du format brut")

        return []

    def _suggest_diagram_for_slide(self, slide_title: str, context: str) -> Optional[Dict[str, Any]]:
        """
        Suggere un diagramme pertinent pour une slide donnee

        Args:
            slide_title: Titre de la slide
            context: Contexte de la slide (contenu)

        Returns:
            Dict avec les parametres du diagramme ou None
        """
        title_lower = slide_title.lower()
        context_lower = context.lower()

        # Detecter si un diagramme flow serait pertinent
        if any(word in title_lower for word in ['demarche', 'methodologie', 'processus', 'etapes', 'phases']):
            return {
                'type': 'flow',
                'trigger': 'methodologie',
                'description': 'Processus etape par etape'
            }

        # Detecter si un diagramme cycle serait pertinent
        if any(word in title_lower for word in ['succes', 'facteurs', 'cles', 'iteration', 'amelioration']):
            return {
                'type': 'cycle',
                'trigger': 'facteurs_cles',
                'description': 'Elements interdependants'
            }

        # Detecter si une architecture serait pertinente
        if any(word in context_lower for word in ['architecture', 'infrastructure', 'composants', 'systeme']):
            return {
                'type': 'flow',
                'trigger': 'architecture',
                'description': 'Architecture de la solution'
            }

        return None

    def generate_slides_structure(
        self,
        tender_analysis: Dict[str, Any],
        references: Dict[str, Any],
        template_structure: str,
        adapted_cvs: List[Dict[str, Any]],
    ) -> List[Dict[str, Any]]:
        """
        Genere la structure complete des slides en JSON pour le PPTX

        Returns:
            Liste de dictionnaires decrivant chaque slide avec leur contenu
        """
        print("Generation de la structure des slides...")

        system_prompt = f"""Tu es {self.consultant_info['name']}, {self.consultant_info['title']} chez {self.consultant_info['company']}.
Tu generes une proposition commerciale structuree en slides PowerPoint.
IMPORTANT: Genere du contenu concret et professionnel, sans majuscules inutiles ni emojis."""

        # Resumer les informations pour le contexte
        context_summary = f"""Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {', '.join(tender_analysis.get('objectives', [])[:4])}
Mots-cles: {', '.join(tender_analysis.get('keywords', [])[:6])}
Exigences techniques: {', '.join(tender_analysis.get('requirements', {}).get('technical', [])[:3])}"""

        # References resume
        refs_summary = ""
        if references and references.get('all_references'):
            projects = references['all_references'].get('projects', [])[:2]
            for p in projects:
                refs_summary += f"\n- {p.get('title', 'Projet')}: {p.get('description', '')[:100]}"

        prompt = f"""Genere une proposition commerciale VISUELLE et MODERNE en slides PowerPoint.

CONTEXTE:
{context_summary}

REFERENCES WENVISION:{refs_summary}

IMPORTANT - APPROCHE VISUELLE:
- PRIVILEGIE les diagrammes et visuels (60% des slides doivent etre visuelles)
- LIMITE le texte: 2-3 bullets COURTS par slide maximum
- EVITE les longs paragraphes et les listes a rallonge
- STRUCTURE visuelle: 1 idee = 1 slide = 1 visuel

Types de slides disponibles:
- "content": slide avec titre et 2-3 bullets COURTS (pas plus!)
- "diagram": slide avec diagramme visuel (A PRIVILEGIER)
- "table": slide avec tableau compact
- "image": slide avec image + 1-2 bullets

Types de diagrammes disponibles:
- "flow": processus sequentiels (Etape 1 → Etape 2 → Etape 3)
- "cycle": processus cycliques/iteratifs
- "pyramid": hierarchies et niveaux
- "timeline": chronologie et jalons
- "matrix": categorisation 2x2 ou 2x3

REGLES DE DESIGN:
1. Chaque section DOIT commencer par un diagramme (pas de slide section)
2. Maximum 2-3 bullets par slide content (phrases courtes, <10 mots)
3. Alterner slides textuelles et visuelles (jamais 2 slides content consecutives)
4. Utiliser des diagrammes pour TOUT ce qui peut etre visualise

Exemple de structure VISUELLE (pas de slides section, debut direct par diagramme):

[
  {{"type":"cover","client":"{tender_analysis.get('client_name', '')}","project":"{tender_analysis.get('project_title', '')}","date":"{datetime.now().strftime('%d/%m/%Y')}"}},
  {{"type":"diagram","title":"Contexte du projet","diagram_type":"pyramid","elements":["Vision strategique","Objectifs metier","Actions prioritaires"],"description":"Pyramide des enjeux"}},
  {{"type":"content","title":"Nos objectifs","bullets":["Objectif 1 (court)","Objectif 2 (court)","Objectif 3 (court)"]}},
  {{"type":"diagram","title":"Notre demarche","diagram_type":"flow","elements":["Cadrage","Conception","Pilote","Deploiement"],"description":"Methodologie en 4 phases"}},
  {{"type":"content","title":"Phase de cadrage","bullets":["Diagnostic rapide","Ateliers collaboratifs"]}},
  {{"type":"diagram","title":"Architecture proposee","diagram_type":"flow","elements":["Frontend","API","Base de donnees","Analytics"],"description":"Architecture cloud-native"}},
  {{"type":"timeline","title":"Planning","diagram_type":"timeline","elements":["J0: Kick-off","M1: POC","M3: Pilote","M6: Production"],"description":"Roadmap 6 mois"}},
  {{"type":"table","title":"Budget","headers":["Phase","Jours","Total"],"rows":[["Cadrage","10j","10k€"],["Pilote","20j","20k€"],["Total","30j","30k€"]]}},
  {{"type":"diagram","title":"Facteurs de succes","diagram_type":"cycle","elements":["Co-construction","Iterations","Mesure","Ajustement"],"description":"Cycle d'amelioration"}},
  {{"type":"closing"}}
]

IMPORTANT - STRUCTURE VISUELLE:
- PAS de slides "section" (debut direct par diagramme)
- Chaque slide "content" = 2-3 bullets COURTS maximum (phrases <10 mots)
- 60% de slides visuelles (diagram/table/image) vs 40% textuelles
- Alterner: diagram → content → diagram → content
- Total: 10-12 slides (plus compact, plus visuel)
- Diagrammes: 4-5 par proposition (au moins 40% des slides)
- Chaque concept = 1 diagramme si possible"""

        response = self.llm_client.generate(
            prompt=prompt,
            system_prompt=system_prompt,
            temperature=0.6,
            max_tokens=8000
        )

        # Parser le JSON
        try:
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start >= 0 and json_end > json_start:
                slides = json.loads(response[json_start:json_end])

                # Valider et completer les slides
                for slide in slides:
                    if slide.get('type') == 'content' and 'bullets' not in slide:
                        # Ajouter des bullets par defaut si manquants
                        slide['bullets'] = [
                            "Point a developper selon le contexte",
                            "Deuxieme point a completer",
                            "Troisieme point a personnaliser"
                        ]
                    elif slide.get('type') == 'table' and 'rows' not in slide:
                        # Ajouter des donnees par defaut
                        slide['headers'] = slide.get('headers', ['Colonne 1', 'Colonne 2'])
                        slide['rows'] = [['Donnee 1', 'Donnee 2'], ['Total', 'A definir']]

                # Ajouter les slides CV (sans section, juste les CVs)
                if adapted_cvs:
                    # Supprimer closing si presente
                    if slides and slides[-1].get('type') == 'closing':
                        slides.pop()

                    # Ajouter les CVs directement
                    for cv in adapted_cvs:
                        slides.append({
                            "type": "cv",
                            "name": cv.get('name', ''),
                            "title": cv.get('title', ''),
                            "experiences": cv.get('experiences', []),
                            "skills": cv.get('skills', []),
                        })

                    # Pas de closing slide (on termine sur les CVs ou derniere slide utile)

                print(f"   {len(slides)} slides completes generees")
                return slides

        except (json.JSONDecodeError, ValueError) as e:
            print(f"   Erreur parsing JSON: {e}")
            print(f"   Utilisation d'une structure par defaut...")

            # Structure de fallback VISUELLE (pas de slides section, focus diagrammes)
            fallback_slides = [
                {"type": "cover", "client": tender_analysis.get('client_name', ''), "project": tender_analysis.get('project_title', ''), "date": datetime.now().strftime('%d/%m/%Y')},
                # Debut direct par un diagramme
                {"type": "diagram", "title": "Vision du projet", "diagram_type": "pyramid", "elements": [
                    tender_analysis.get('project_title', 'Projet'),
                    "Objectifs strategiques",
                    "Actions prioritaires"
                ], "description": "Pyramide des enjeux"},
                {"type": "content", "title": "Nos objectifs", "bullets": tender_analysis.get('objectives', ['Objectif 1', 'Objectif 2'])[:3]},
                # Diagramme methodologie
                {"type": "diagram", "title": "Notre demarche", "diagram_type": "flow", "elements": [
                    "Cadrage",
                    "Conception",
                    "Pilote",
                    "Deploiement"
                ], "description": "Methodologie iterative en 4 phases"},
                {"type": "content", "title": "Approche collaborative", "bullets": [
                    "Co-construction avec vos equipes",
                    "Iterations courtes et validation continue"
                ]},
                # Timeline planning
                {"type": "diagram", "title": "Planning", "diagram_type": "timeline", "elements": [
                    "J0: Kick-off",
                    "M1: POC",
                    "M3: Pilote",
                    "M6: Production"
                ], "description": "Roadmap de deploiement sur 6 mois"},
                # Budget compact
                {"type": "table", "title": "Budget", "headers": ["Phase", "Duree", "Total"], "rows": [
                    ["Cadrage", "2 sem", "10k€"],
                    ["Pilote", "8 sem", "40k€"],
                    ["Production", "4 sem", "20k€"],
                    ["Total", "14 sem", "70k€"]
                ]},
                # Facteurs de succes en cycle
                {"type": "diagram", "title": "Facteurs de succes", "diagram_type": "cycle", "elements": [
                    "Co-construction",
                    "Iterations rapides",
                    "Mesure impact",
                    "Ajustement continu"
                ], "description": "Cycle d'amelioration continue"},
                {"type": "content", "title": "Notre valeur ajoutee", "bullets": [
                    "Expertise data & IA eprouvee",
                    "Accompagnement pragmatique"
                ]},
                {"type": "closing"}
            ]

            # Ajouter CVs si disponibles (sans section, sans closing)
            if adapted_cvs:
                # Supprimer closing
                if fallback_slides[-1].get('type') == 'closing':
                    fallback_slides.pop()

                # Ajouter CVs directement
                for cv in adapted_cvs:
                    fallback_slides.append({
                        "type": "cv",
                        "name": cv.get('name', ''),
                        "title": cv.get('title', ''),
                        "experiences": cv.get('experiences', []),
                        "skills": cv.get('skills', [])
                    })
            else:
                # Pas de CVs: supprimer closing quand meme pour finir sur du contenu utile
                if fallback_slides[-1].get('type') == 'closing':
                    fallback_slides.pop()

            return fallback_slides

    def enhance_slides_with_images(
        self,
        slides: List[Dict[str, Any]],
        tender_analysis: Dict[str, Any],
        generate_images: bool = False
    ) -> List[Dict[str, Any]]:
        """
        Ameliore les slides avec des diagrammes/images generes

        Args:
            slides: Liste des slides a ameliorer
            tender_analysis: Analyse de l'appel d'offre pour le contexte
            generate_images: Si True, genere des diagrammes via Claude/Mermaid et optionnellement DALL-E

        Returns:
            Liste des slides ameliorees avec diagrammes/images
        """
        if not generate_images:
            return slides

        context = {
            'client_name': tender_analysis.get('client_name', 'Client'),
            'project_title': tender_analysis.get('project_title', 'Projet')
        }

        enhanced_slides = []
        for slide in slides:
            enhanced_slides.append(slide)

            # Detecter les slides qui pourraient beneficier de diagrammes generes
            slide_title = slide.get('title', '').lower()

            # Architecture/Infrastructure → generer diagramme Mermaid via Claude
            if any(keyword in slide_title for keyword in ['architecture', 'infrastructure', 'technique', 'solution', 'composants']):
                if slide.get('type') == 'content':
                    # Construire description depuis les bullets
                    bullets = slide.get('bullets', [])
                    components = [b.split(':')[0].strip() for b in bullets[:6]]

                    print(f"   Generation diagramme architecture pour: {slide.get('title', 'Architecture')}")

                    # Methode 1: Claude + Mermaid (privilegie, gratuit)
                    image_path = self.diagram_generator.generate_architecture_diagram(
                        components=components,
                        context=context,
                        description=slide.get('title', 'Architecture')
                    )

                    # Methode 2: DALL-E en fallback si Mermaid echoue et OPENAI_API_KEY disponible
                    if not image_path and os.getenv('OPENAI_API_KEY'):
                        print("   Fallback DALL-E pour architecture...")
                        image_path = self.image_generator.generate_architecture_diagram(
                            components=components,
                            context=context
                        )

                    if image_path:
                        # Ajouter une slide image apres la slide content
                        enhanced_slides.append({
                            "type": "image",
                            "title": f"Vue d'ensemble - {slide.get('title', '')}",
                            "image_path": image_path,
                            "caption": "Architecture proposee pour le projet",
                            "bullets": components[:3]
                        })

            # Flux/Processus → possiblement generer flowchart Mermaid
            elif any(keyword in slide_title for keyword in ['flux', 'processus', 'workflow']) and slide.get('type') == 'content':
                bullets = slide.get('bullets', [])
                steps = [b.split(':')[0].strip() for b in bullets[:6]]

                if len(steps) >= 3:
                    print(f"   Generation diagramme flux pour: {slide.get('title', 'Processus')}")
                    image_path = self.diagram_generator.generate_flow_diagram(
                        steps=steps,
                        context=context,
                        description=slide.get('title', 'Processus')
                    )

                    if image_path:
                        enhanced_slides.append({
                            "type": "image",
                            "title": f"Flux - {slide.get('title', '')}",
                            "image_path": image_path,
                            "caption": "Processus detaille",
                            "bullets": steps[:3]
                        })

        return enhanced_slides

    def suggest_library_image(self, slide_title: str, slide_content: str) -> Optional[str]:
        """
        Suggere une image de la bibliotheque pour une slide donnee

        Args:
            slide_title: Titre de la slide
            slide_content: Contenu de la slide

        Returns:
            Chemin vers l'image suggeree ou None
        """
        title_lower = slide_title.lower()
        content_lower = slide_content.lower()

        # Mapping de mots-cles vers categories
        keyword_to_category = {
            'architecture': 'architecture',
            'infrastructure': 'infrastructure',
            'processus': 'process',
            'methodologie': 'methodology',
            'demarche': 'methodology',
            'equipe': 'team',
            'resultats': 'success',
            'reussite': 'success',
            'dashboard': 'dashboard',
            'tableau de bord': 'dashboard',
            'donnees': 'data',
            'data': 'data',
            'technologie': 'technology',
            'mockup': 'mockup',
            'interface': 'mockup'
        }

        # Chercher une categorie correspondante
        for keyword, category in keyword_to_category.items():
            if keyword in title_lower or keyword in content_lower:
                image_path = self.image_library.get_image_by_category(category)
                if image_path:
                    return image_path

        return None

    def generate_proposal_content(
        self,
        tender_analysis: Dict[str, Any],
        references: Dict[str, Any],
        template_structure: str
    ) -> Dict[str, Any]:
        """
        Génère le contenu de la proposition commerciale

        Args:
            tender_analysis: Analyse de l'appel d'offre
            references: Références sélectionnées
            template_structure: Structure du template (texte)

        Returns:
            Contenu de la proposition
        """
        print("Génération de la proposition commerciale...")

        system_prompt = f"""Tu es {self.consultant_info['name']}, {self.consultant_info['title']} chez {self.consultant_info['company']}.

Tu dois générer une proposition commerciale professionnelle et convaincante en français.

Style WEnvision:
- Professionnel mais accessible
- Axé sur la valeur pour le client
- Concret avec des exemples et résultats chiffrés
- Démontrer l'expertise sans jargon excessif
- Philosophie: Co-construction, Alignement, Pragmatisme"""

        prompt = f"""Génère une proposition commerciale complète basée sur:

APPEL D'OFFRE ANALYSÉ:
{json.dumps(tender_analysis, indent=2, ensure_ascii=False)}

RÉFÉRENCES WENVISION PERTINENTES:
{references.get('selected_references')}

STRUCTURE DU TEMPLATE WENVISION:
{template_structure}

La proposition DOIT suivre la structure du template Wenvision 2026:

1. NOTRE COMPRÉHENSION DU CONTEXTE
   - Contexte et objectifs de la mission
   - Analyse du contexte: Qui ? (parties prenantes), Quoi ? (périmètre), Pourquoi ? (motivations)
   - Philosophie d'intervention: Alignement, Co-construction, Pragmatisme

2. NOTRE EXPERTISE ET DÉMARCHE À VOTRE SERVICE
   - Références similaires avec résultats concrets
   - Expertise de l'équipe Wenvision
   - Valeur ajoutée pour ce projet

3. NOTRE DÉMARCHE D'ACCOMPAGNEMENT
   - Phases détaillées avec activités, ateliers et livrables
   - Gouvernance et points de synchronisation (COPIL, Comités de suivi, Points d'équipe)
   - Hypothèses et pré-requis

4. NOTRE FEUILLE DE ROUTE
   - Planning semaine par semaine
   - Ateliers et entretiens planifiés
   - Livrables par phase
   - Durée totale

5. FACTEURS CLÉS DE SUCCÈS
   - 5 facteurs déterminants pour garantir la réussite

6. OFFRE FINANCIÈRE
   - Détail des postes et tarification
   - Modalités de facturation (ex: 30% à la commande, 70% à la fin)

Pour chaque section, fournis du contenu détaillé et directement exploitable.
Sois concret, spécifique au client et à ses enjeux."""

        proposal = self.llm_client.generate(
            prompt=prompt,
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=8000
        )

        return {
            'content': proposal,
            'tender_analysis': tender_analysis,
            'references_used': {
                'selected_references': references.get('selected_references')
            },
            'generated_at': datetime.now().isoformat(),
            'consultant': self.consultant_info
        }

    def generate_proposal(
        self,
        tender_path: str,
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Génère une proposition commerciale complète

        Args:
            tender_path: Chemin vers le fichier d'appel d'offre
            output_path: Chemin de sortie (optionnel)

        Returns:
            Proposition générée
        """
        print(f"\n{'='*60}")
        print("GÉNÉRATION DE PROPOSITION COMMERCIALE - WENVISION")
        print(f"{'='*60}\n")

        # 1. Charger l'appel d'offre
        print(f"Chargement de l'appel d'offre: {tender_path}")
        with open(tender_path, 'r', encoding='utf-8') as f:
            tender_text = f.read()

        # 2. Analyser l'appel d'offre
        tender_analysis = self.analyze_tender(tender_text)

        # 3. Charger le template local
        template = self.load_template()

        # 4. Sélectionner les références
        references = self.match_references(tender_analysis)

        # 5. Charger les CVs
        cvs = self.load_cvs()

        # 6. Adapter les CVs au contexte
        adapted_cvs = self.adapt_cvs(cvs, tender_analysis) if cvs else []

        # 7. Générer la structure des slides PPTX
        slides_structure = self.generate_slides_structure(
            tender_analysis,
            references,
            template,
            adapted_cvs
        )

        # 7b. Ameliorer avec images generees (optionnel, si OPENAI_API_KEY configuree)
        use_dalle = os.getenv('USE_DALLE_IMAGES', 'false').lower() == 'true'
        if use_dalle and os.getenv('OPENAI_API_KEY'):
            print("Amelioration des slides avec images DALL-E...")
            slides_structure = self.enhance_slides_with_images(
                slides_structure,
                tender_analysis,
                generate_images=True
            )

        # 8. Générer la proposition Markdown
        proposal = self.generate_proposal_content(
            tender_analysis,
            references,
            template
        )

        # 9. Sauvegarder
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            client_name = tender_analysis.get('client_name', 'client').replace(' ', '_')
            output_path = os.path.join(
                BASE_DIR, "output", f"proposal_{client_name}_{timestamp}.json"
            )

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(proposal, f, indent=2, ensure_ascii=False)

        # Générer aussi une version markdown
        md_path = output_path.replace('.json', '.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(f"# Proposition Commerciale - {tender_analysis.get('project_title', 'Projet')}\n\n")
            f.write(f"**Client:** {tender_analysis.get('client_name', 'N/A')}\n")
            f.write(f"**Consultant:** {self.consultant_info['name']}\n")
            f.write(f"**Entreprise:** {self.consultant_info['company']}\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y')}\n\n")
            f.write("---\n\n")
            f.write(proposal['content'])

        # 10. Générer le PPTX
        pptx_path = output_path.replace('.json', '.pptx')
        try:
            build_proposal_pptx(
                template_path=self.template_path,
                slides_data=slides_structure,
                output_path=pptx_path,
                consultant_info=self.consultant_info
            )
            print(f"   PPTX: {pptx_path}")
        except Exception as e:
            print(f"   Erreur PPTX: {e}")
            pptx_path = None

        print(f"\nProposition générée avec succès!")
        print(f"   JSON: {output_path}")
        print(f"   Markdown: {md_path}")
        if pptx_path:
            print(f"   PowerPoint: {pptx_path}")
        print(f"\n{'='*60}\n")

        return {
            **proposal,
            'pptx_path': pptx_path,
            'slides_structure': slides_structure,
        }

    # === GENERATION MODULAIRE PAR SECTIONS ===

    def generate_agenda_slide(self, tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère une slide d'agenda

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Liste de slides (agenda)
        """
        print("Génération de l'agenda...")

        prompt = f"""Génère une slide d'agenda pour cette proposition commerciale:

Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}

Crée un agenda structuré avec ces sections (5-7 points max):
1. Notre compréhension du contexte
2. Notre vision et approche
3. Planning de l'intervention
4. Chiffrage
5. Références pertinentes
6. L'équipe proposée

Retourne UNIQUEMENT en JSON:
{{
  "title": "Agenda",
  "bullets": ["Section 1", "Section 2", ...]
}}"""

        response = self.llm_client.generate(prompt, temperature=0.5, max_tokens=500)

        # Parser le JSON
        try:
            json_start = response.find('{')
            json_end = response.rfind('}') + 1
            if json_start >= 0 and json_end > json_start:
                agenda_data = json.loads(response[json_start:json_end])
                print("   Agenda généré")
                return [{
                    "type": "content",
                    "title": agenda_data.get("title", "Agenda"),
                    "bullets": agenda_data.get("bullets", [])
                }]
        except Exception as e:
            print(f"   Erreur: {e}")

        return [{"type": "content", "title": "Agenda", "bullets": ["Erreur de génération"]}]

    def generate_context_slides(self, tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère les slides de contexte (enjeux et objectifs) - APPROCHE VISUELLE

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Liste de slides de contexte (diagramme + texte court)
        """
        print("Génération du contexte (approche visuelle)...")

        # Vérifier s'il y a des instructions de feedback
        feedback_instructions = ""
        if 'feedback_instructions' in tender_analysis:
            feedback_instructions = f"\n\nINSTRUCTIONS DE MODIFICATION:\n{tender_analysis['feedback_instructions']}\n\nAPPLIQUE CES MODIFICATIONS PRÉCISÉMENT.\n"

        prompt = f"""Génère des slides de contexte VISUELLES et IMPACTANTES (style Veolia) pour cette mission:

Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {json.dumps(tender_analysis.get('objectives', []), ensure_ascii=False)}
Contraintes: {json.dumps(tender_analysis.get('requirements', {}), ensure_ascii=False)}
{feedback_instructions}
APPROCHE VISUELLE MODERNE (style Veolia) - 3-4 slides :
1. Slide STAT impactante (si chiffre clé disponible : budget, délai, ROI attendu)
   OU Quote/Key message si c'est une phrase forte du client
2. Diagramme pyramidal (pyramid) ou cycle des enjeux (4-5 éléments)
3. Slide HIGHLIGHT avec 3-4 objectifs clés dans des encadrés colorés (plus visuel que bullets)
4. [Optionnel] Slide de contenu classique si nécessaire

TYPES DE SLIDES DISPONIBLES:
- "stat": chiffre clé en très grand (stat_value, stat_label, context, subtitle)
- "quote": citation/message clé impactant (quote_text, author, title)
- "highlight": points clés dans encadrés colorés (title, key_points [2-4], highlight_color)
- "diagram": diagramme visuel (title, diagram_type, elements, description)
- "content": slide classique avec bullets (title, bullets)

Retourne UNIQUEMENT en JSON (3-4 slides, privilégie stat/quote/highlight):
[
  {{
    "type": "stat",
    "stat_value": "18 mois",
    "stat_label": "pour transformer la data",
    "context": "Durée cible de la mission",
    "subtitle": "Timeline du projet"
  }},
  {{
    "type": "diagram",
    "title": "Enjeux stratégiques",
    "diagram_type": "pyramid",
    "elements": ["Vision stratégique", "Objectifs métier", "Quick wins", "Fondations"],
    "description": "Hiérarchie des priorités"
  }},
  {{
    "type": "highlight",
    "title": "Nos 3 objectifs clés",
    "key_points": ["Objectif 1 (15 mots max)", "Objectif 2", "Objectif 3"],
    "highlight_color": "corail"
  }}
]"""

        response = self.llm_client.generate(prompt, temperature=0.6, max_tokens=1500)

        try:
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start >= 0 and json_end > json_start:
                slides_data = json.loads(response[json_start:json_end])
                print(f"   {len(slides_data)} slides de contexte générées (approche visuelle)")
                return slides_data
        except Exception as e:
            print(f"   Erreur: {e}")

        # Fallback visuel
        return [
            {
                "type": "diagram",
                "title": "Enjeux du projet",
                "diagram_type": "pyramid",
                "elements": ["Vision stratégique", "Objectifs métier", "Fondations techniques"],
                "description": "Hiérarchie des enjeux"
            },
            {
                "type": "content",
                "title": "Objectifs",
                "bullets": ["Transformer la donnée", "Créer de la valeur"]
            }
        ]

    def generate_approach_slides(self, tender_analysis: Dict[str, Any], references: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère les slides de vision et approche - APPROCHE VISUELLE

        Args:
            tender_analysis: Analyse de l'appel d'offre
            references: Références sélectionnées

        Returns:
            Liste de slides d'approche (diagrammes + texte court)
        """
        print("Génération de la vision et approche (approche visuelle)...")

        # Résumer les références pour réduire les tokens
        refs_summary = ""
        if references and references.get('selected_references'):
            refs_summary = str(references.get('selected_references', ''))[:1500]
        elif references and references.get('all_references'):
            projects = references['all_references'].get('projects', [])[:2]
            for p in projects:
                refs_summary += f"\n- {p.get('title', 'Projet')}: {p.get('description', '')[:100]}"

        # Extraire les objectifs pour contexte
        objectives = ', '.join(tender_analysis.get('objectives', [])[:3])

        # Vérifier s'il y a des instructions de feedback
        feedback_instructions = ""
        if 'feedback_instructions' in tender_analysis:
            feedback_instructions = f"\n\nINSTRUCTIONS DE MODIFICATION:\n{tender_analysis['feedback_instructions']}\n\nAPPLIQUE CES MODIFICATIONS PRÉCISÉMENT.\n"

        prompt = f"""Génère des slides VISUELLES et IMPACTANTES (style Veolia) pour notre vision et approche:

Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {objectives}

Références Wenvision: {refs_summary}
{feedback_instructions}
APPROCHE VISUELLE MODERNE (style Veolia) - 4-5 slides:
1. Slide QUOTE avec notre vision/philosophie (phrase impactante de 20-30 mots)
   OU Slide HIGHLIGHT avec 3 piliers de notre approche (encadrés colorés)
2. Diagramme FLOW de la méthodologie (5-6 étapes claires)
3. Diagramme CYCLE des facteurs de succès (4-5 éléments)
4. [Optionnel] Slide avec nos différenciateurs vs concurrence

TYPES DE SLIDES DISPONIBLES:
- "quote": citation/message clé impactant (quote_text, author optional, title optional)
- "highlight": points clés dans encadrés (title, key_points [2-4], highlight_color)
- "diagram": flow, cycle, pyramid, timeline, matrix (title, diagram_type, elements, description)
- "content": slide classique avec bullets (title, bullets)

Retourne UNIQUEMENT en JSON (4-5 slides, privilégie quote/highlight/diagram):
[
  {{
    "type": "quote",
    "quote_text": "Transformer la data en valeur métier nécessite une approche pragmatique, itérative et centrée sur l'humain",
    "title": "Notre philosophie"
  }},
  {{
    "type": "highlight",
    "title": "Nos 3 piliers",
    "key_points": ["Pilier 1 (concis)", "Pilier 2", "Pilier 3"],
    "highlight_color": "terracotta"
  }},
  {{
    "type": "diagram",
    "title": "Notre démarche",
    "diagram_type": "flow",
    "elements": ["Cadrage", "Conception", "POC", "Pilote", "Déploiement"],
    "description": "Méthodologie itérative éprouvée"
  }},
  {{
    "type": "diagram",
    "title": "Facteurs clés de succès",
    "diagram_type": "cycle",
    "elements": ["Implication métier", "Itérations courtes", "Transfert compétences", "Quick wins"],
    "description": "Cercle vertueux de la réussite"
  }}
]"""

        response = self.llm_client.generate(prompt, temperature=0.6, max_tokens=2000)

        try:
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start >= 0 and json_end > json_start:
                slides_data = json.loads(response[json_start:json_end])
                print(f"   {len(slides_data)} slides d'approche générées (approche visuelle)")
                return slides_data
        except (json.JSONDecodeError, ValueError) as e:
            print(f"   Erreur parsing JSON: {e}")

        # Fallback visuel
        print("   Utilisation de slides visuelles par défaut...")
        return [
            {
                "type": "content",
                "title": "Notre vision",
                "bullets": [
                    f"Accompagner {tender_analysis.get('client_name', 'le client')} dans sa transformation",
                    "Approche pragmatique orientée résultats"
                ]
            },
            {
                "type": "diagram",
                "title": "Notre démarche",
                "diagram_type": "flow",
                "elements": ["Cadrage", "Conception", "Pilote", "Déploiement"],
                "description": "Méthodologie itérative"
            },
            {
                "type": "diagram",
                "title": "Facteurs clés de succès",
                "diagram_type": "cycle",
                "elements": [
                    "Implication parties prenantes",
                    "Iterations courtes",
                    "Alignement métier",
                    "Transfert de compétences continu"
                ]
            }
        ]

    def generate_planning_slide(self, tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère la slide de planning - APPROCHE VISUELLE

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Liste avec timeline + tableau (approche visuelle)
        """
        print("Génération du planning (approche visuelle)...")

        # Vérifier s'il y a des instructions de feedback
        feedback_instructions = ""
        if 'feedback_instructions' in tender_analysis:
            feedback_instructions = f"\n\nINSTRUCTIONS DE MODIFICATION:\n{tender_analysis['feedback_instructions']}\n\nAPPLIQUE CES MODIFICATIONS PRÉCISÉMENT.\n"

        prompt = f"""Génère un planning VISUEL pour cette mission:

Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Délais: {tender_analysis.get('timeline', 'Non spécifié')}
{feedback_instructions}
IMPORTANT - APPROCHE VISUELLE:
- Slide 1: Diagramme TIMELINE (4-5 jalons clés avec dates relatives)
- Slide 2: Tableau détaillé des phases (optionnel)
- Privilégie la TIMELINE visuelle

Retourne UNIQUEMENT en JSON (1 ou 2 slides):
[
  {{
    "type": "diagram",
    "title": "Roadmap de déploiement",
    "diagram_type": "timeline",
    "elements": ["T0: Cadrage", "T+1M: POC", "T+3M: Pilote", "T+6M: Déploiement"],
    "description": "Jalons clés du projet"
  }}
]

OU si besoin d'un tableau détaillé en plus:
[
  {{
    "type": "diagram",
    "title": "Roadmap",
    "diagram_type": "timeline",
    "elements": ["T0: Cadrage", "T+1M: Conception", "T+2M: Pilote", "T+4M: Déploiement"],
    "description": "Timeline du projet"
  }},
  {{
    "type": "table",
    "title": "Planning détaillé",
    "headers": ["Phase", "Durée", "Livrables"],
    "rows": [["Cadrage", "2 sem", "Roadmap"], ["Conception", "4 sem", "Architecture"]]
  }}
]"""

        response = self.llm_client.generate(prompt, temperature=0.5, max_tokens=1500)

        try:
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start >= 0 and json_end > json_start:
                planning_data = json.loads(response[json_start:json_end])
                print(f"   {len(planning_data)} slide(s) de planning générée(s) (approche visuelle)")
                return planning_data
        except Exception as e:
            print(f"   Erreur: {e}")

        # Fallback visuel
        return [{
            "type": "diagram",
            "title": "Roadmap de déploiement",
            "diagram_type": "timeline",
            "elements": ["T0: Cadrage", "T+1M: Conception", "T+3M: Pilote", "T+6M: Déploiement"],
            "description": "Timeline du projet sur 6 mois"
        }]

    def generate_budget_slide(self, tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère la slide de chiffrage - FORMAT SIMPLIFIÉ ET VISUEL

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Liste avec slide de budget (tableau simplifié)
        """
        print("Génération du chiffrage (format simplifié)...")

        # Vérifier s'il y a des instructions de feedback
        feedback_instructions = ""
        if 'feedback_instructions' in tender_analysis:
            feedback_instructions = f"\n\nINSTRUCTIONS DE MODIFICATION:\n{tender_analysis['feedback_instructions']}\n\nAPPLIQUE CES MODIFICATIONS PRÉCISÉMENT (ex: augmenter budget 20% = multiplier par 1.2).\n"

        prompt = f"""Génère un chiffrage SIMPLIFIÉ et VISUEL pour cette mission:

Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Budget indicatif: {tender_analysis.get('budget', 'Non spécifié')}
{feedback_instructions}
IMPORTANT - FORMAT SIMPLIFIÉ:
- Maximum 3-4 lignes de prestations + ligne TOTAL
- Tableau épuré et facile à lire
- Éviter les détails excessifs

Retourne UNIQUEMENT en JSON avec 3-4 lignes + total:
{{
  "title": "Chiffrage de l'intervention",
  "headers": ["Prestation", "Jours", "Total"],
  "rows": [
    ["Cadrage et conception", "15 j", "15 000€"],
    ["Mise en œuvre et pilote", "25 j", "25 000€"],
    ["Déploiement et transfert", "10 j", "10 000€"],
    ["TOTAL", "50 j", "50 000€"]
  ]
}}"""

        response = self.llm_client.generate(prompt, temperature=0.5, max_tokens=800)

        try:
            json_start = response.find('{')
            json_end = response.rfind('}') + 1
            if json_start >= 0 and json_end > json_start:
                budget_data = json.loads(response[json_start:json_end])
                print("   Chiffrage généré (format simplifié)")
                return [{
                    "type": "table",
                    "title": budget_data.get("title", "Chiffrage"),
                    "headers": budget_data.get("headers", []),
                    "rows": budget_data.get("rows", [])
                }]
        except Exception as e:
            print(f"   Erreur: {e}")

        # Fallback simplifié
        return [{
            "type": "table",
            "title": "Chiffrage de l'intervention",
            "headers": ["Prestation", "Jours", "Total"],
            "rows": [
                ["Cadrage et conception", "15 j", "15 000€"],
                ["Mise en œuvre", "25 j", "25 000€"],
                ["Déploiement", "10 j", "10 000€"],
                ["TOTAL", "50 j", "50 000€"]
            ]
        }]

    def generate_references_slides(self, tender_analysis: Dict[str, Any], references: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère les slides de références pertinentes

        Args:
            tender_analysis: Analyse de l'appel d'offre
            references: Références chargées

        Returns:
            Liste de slides de références
        """
        print("Génération des références...")

        # Extraire les projets pertinents
        projects = references.get('all_references', {}).get('projects', [])[:3]

        if not projects:
            return [{"type": "content", "title": "Références", "bullets": ["Aucune référence disponible"]}]

        slides = []
        for i, project in enumerate(projects, 1):
            slides.append({
                "type": "content",
                "title": f"Référence {i}: {project.get('title', 'Projet')}",
                "bullets": [
                    f"Client: {project.get('client', 'N/A')}",
                    f"Secteur: {project.get('sector', 'N/A')}",
                    f"Challenge: {project.get('challenge', 'N/A')[:150]}",
                    f"Solution: {project.get('solution', 'N/A')[:150]}",
                    f"Technologies: {', '.join(project.get('technologies', [])[:4])}"
                ]
            })

        print(f"   {len(slides)} slides de références générées")
        return slides

    def generate_cv_slides(self, tender_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Génère les slides de CVs adaptés à la mission

        Args:
            tender_analysis: Analyse de l'appel d'offre

        Returns:
            Liste de slides CV
        """
        print("Génération des CVs...")

        # Charger et adapter les CVs
        cvs = self.load_cvs()
        if not cvs:
            return [{"type": "content", "title": "Équipe", "bullets": ["Aucun CV disponible"]}]

        adapted_cvs = self.adapt_cvs(cvs, tender_analysis)
        if not adapted_cvs:
            return [{"type": "content", "title": "Équipe", "bullets": ["Erreur lors de l'adaptation des CVs"]}]

        # Convertir en slides
        cv_slides = []
        for cv in adapted_cvs:
            cv_slides.append({
                "type": "cv",
                "name": cv.get('name', ''),
                "title": cv.get('title', ''),
                "experiences": cv.get('experiences', []),
                "skills": cv.get('skills', [])
            })

        print(f"   {len(cv_slides)} CVs générés")
        return cv_slides


def main():
    """Point d'entrée principal"""
    import argparse

    parser = argparse.ArgumentParser(
        description='Agent de génération de propositions commerciales Wenvision'
    )
    parser.add_argument(
        'tender_file',
        help='Chemin vers le fichier d\'appel d\'offre (txt, md, pdf)'
    )
    parser.add_argument(
        '-o', '--output',
        help='Chemin de sortie pour la proposition',
        default=None
    )

    args = parser.parse_args()

    # Créer l'agent
    agent = ProposalGeneratorAgent()

    # Générer la proposition
    agent.generate_proposal(
        tender_path=args.tender_file,
        output_path=args.output
    )


if __name__ == '__main__':
    main()
