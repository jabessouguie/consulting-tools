"""
Consulting Tools Agents - Application Web
FastAPI server exposant les agents de propositions commerciales et veille LinkedIn
"""

import asyncio
import io
import json
import os
import sys
import threading
import uuid
from datetime import datetime
from pathlib import Path
from typing import List, Optional
from urllib.parse import urlparse

from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, Request, UploadFile
from fastapi.responses import (
    FileResponse,
    HTMLResponse,
    JSONResponse,
    RedirectResponse,
    StreamingResponse,
)
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from PyPDF2 import PdfReader
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.errors import RateLimitExceeded
from slowapi.util import get_remote_address
from starlette.middleware.base import BaseHTTPMiddleware
from starlette.middleware.sessions import SessionMiddleware

from agents.article_to_post import ArticleToPostAgent
from agents.dataset_analyzer import DatasetAnalyzerAgent
from agents.elearning_agent import ElearningAgent
from agents.linkedin_commenter import LinkedInCommenterAgent
from agents.linkedin_monitor import LinkedInMonitorAgent
from agents.meeting_capture_agent import GEMINI_MODELS as MC_GEMINI_MODELS
from agents.meeting_capture_agent import APIKeyError as MCAPIKeyError
from agents.meeting_capture_agent import (
    DraftCreationError,
    MeetingCaptureAgent,
    MeetingGmailClient,
    VideoProcessingError,
)
from agents.meeting_summarizer import MeetingSummarizerAgent
from agents.proposal_generator import ProposalGeneratorAgent
from agents.rfp_responder import RFPResponderAgent
from agents.skills_market import SkillsMarketAgent
from agents.tech_monitor import TechMonitorAgent
from agents.tender_scout_agent import AnalysisError, ScrapingError, TenderScoutAgent
from agents.workshop_planner import WorkshopPlannerAgent
from config import get_consultant_info
from utils.auth import authenticate_user, get_current_user, get_session_secret
from utils.consultant_db import ConsultantDatabase
from utils.elearning_db import ElearningDatabase
from utils.tender_db import TenderDatabase
from utils.validation import (
    sanitize_error_message,
    sanitize_filename,
    sanitize_text_input,
    validate_file_upload,
)

# Charger l'environnement
BASE_DIR = Path(__file__).parent
load_dotenv(BASE_DIR / ".env")


# Ajouter le repertoire au path pour les imports
sys.path.insert(0, str(BASE_DIR))


# Rate limiting

# === APP SETUP ===

app = FastAPI(title="Consulting Tools Agents", version="1.0.0")

# Rate limiter configuration
limiter = Limiter(key_func=get_remote_address, default_limits=["100 per minute"])
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# Templates and static files
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")
templates = Jinja2Templates(directory=str(BASE_DIR / "templates"))


# === CSRF PROTECTION MIDDLEWARE ===


class CSRFProtectionMiddleware(BaseHTTPMiddleware):
    """
    Middleware CSRF basé sur Origin/Referer checking
    Plus simple que token-based CSRF, adapté pour une app locale
    """

    async def dispatch(self, request: Request, call_next):
        # Méthodes qui nécessitent protection CSRF
        unsafe_methods = {"POST", "PUT", "DELETE", "PATCH"}

        # Routes exemptées (par exemple les webhooks externes)
        exempt_paths = []

        # Si méthode safe ou route exemptée, passer
        if request.method not in unsafe_methods:
            return await call_next(request)

        if any(request.url.path.startswith(path) for path in exempt_paths):
            return await call_next(request)

        # Vérifier Origin ou Referer
        origin = request.headers.get("origin")
        referer = request.headers.get("referer")

        # Obtenir l'hôte autorisé
        host = request.headers.get("host", "localhost:8000")
        allowed_origins = [
            f"http://{host}",
            f"https://{host}",
            "http://localhost:8000",
            "https://localhost:8000",
            "http://127.0.0.1:8000",
            "https://127.0.0.1:8000",
        ]

        # Vérifier l'origine
        origin_valid = False

        if origin:
            origin_valid = origin in allowed_origins
        elif referer:
            # Si pas d'Origin, vérifier Referer
            try:
                parsed = urlparse(referer)
                referer_origin = f"{parsed.scheme}://{parsed.netloc}"
                origin_valid = referer_origin in allowed_origins
            except BaseException:
                origin_valid = False

        if not origin_valid:
            print(f"⚠️  CSRF blocked: {request.method} {request.url.path} from {origin or referer}")
            return JSONResponse(
                {"detail": "CSRF validation failed. Request origin not allowed."}, status_code=403
            )

        return await call_next(request)


# === AUTHENTICATION MIDDLEWARE ===


class AuthMiddleware(BaseHTTPMiddleware):
    """Middleware pour vérifier l'authentification"""

    async def dispatch(self, request: Request, call_next):
        # Routes publiques (pas besoin d'auth)
        public_paths = ["/login", "/static"]

        # Si la route est publique, passer
        if any(request.url.path.startswith(path) for path in public_paths):
            return await call_next(request)

        # Vérifier l'authentification
        if not get_current_user(request):
            # Si c'est une requête API, renvoyer 401
            if request.url.path.startswith("/api/"):
                return JSONResponse({"detail": "Non authentifié"}, status_code=401)
            # Sinon, rediriger vers /login
            return RedirectResponse(url="/login", status_code=302)

        return await call_next(request)


# Ajouter les middlewares dans l'ordre correct
# L'ordre d'ajout est inversé : le dernier ajouté s'exécute en premier
# Protection CSRF (origin checking)
app.add_middleware(CSRFProtectionMiddleware)
# app.add_middleware(AuthMiddleware)  # DESACTIVE - Pas de login requis
app.add_middleware(SessionMiddleware, secret_key=get_session_secret())  # S'exécute en premier

# Consultant info (depuis config centralisee)
try:
    _consultant_config = get_consultant_info()
    CONSULTANT_NAME = _consultant_config["name"]
    COMPANY_NAME = _consultant_config["company"]
except ValueError as e:
    # Si config non chargee, on laisse l'app demarrer mais avec avertissement
    print(f"\n⚠️  {e}")
    CONSULTANT_NAME = "CONFIGURE_CONSULTANT_NAME"
    COMPANY_NAME = "CONFIGURE_COMPANY_NAME"

# Job store (in-memory)
jobs = {}


# === SECURITY: Safe Error Handling (Phase 3) ===


def safe_error_message(error: Exception, context: str = "") -> str:
    """
    Convertit une exception en message d erreur safe pour les logs et reponses
    Masque automatiquement les secrets (API keys, passwords, tokens)

    Args:
        error: L exception
        context: Contexte optionnel pour le message

    Returns:
        Message d erreur nettoye sans secrets
    """
    error_msg = str(error)
    sanitized = sanitize_error_message(error_msg)

    if context:
        return f"{context}: {sanitized}"
    return sanitized


def safe_traceback() -> str:
    """
    Retourne un traceback nettoye sans secrets

    Returns:
        Traceback sanitize
    """
    import traceback

    tb = traceback.format_exc()
    return sanitize_error_message(tb)


# Global model settings (in-memory, persists per session)
AVAILABLE_GEMINI_MODELS = {
    # Gemini 3.x (Preview)
    "gemini-3.1-pro-preview": "Gemini 3.1 Pro Preview 🔥",
    "gemini-3-flash-preview": "Gemini 3 Flash Preview ⚡",
    "gemini-3-pro-preview": "Gemini 3 Pro Preview",
    "gemini-3.1-flash-lite-preview": "Gemini 3.1 Flash Lite",
    "deep-research-pro-preview-12-2025": "Deep Research Pro 🔬",
    # Gemini 2.5
    "gemini-2.5-pro": "Gemini 2.5 Pro",
    "gemini-2.5-flash": "Gemini 2.5 Flash",
    "gemini-2.5-flash-lite": "Gemini 2.5 Flash Lite",
    # Gemini 2.0
    "gemini-2.0-flash": "Gemini 2.0 Flash",
    "gemini-2.0-flash-001": "Gemini 2.0 Flash 001",
    "gemini-2.0-flash-lite": "Gemini 2.0 Flash Lite",
    "gemini-2.0-flash-lite-001": "Gemini 2.0 Flash Lite 001",
    # Latest aliases
    "gemini-flash-latest": "Gemini Flash (Latest)",
    "gemini-pro-latest": "Gemini Pro (Latest)",
    # Gemini 1.5
    "gemini-1.5-pro": "Gemini 1.5 Pro",
    "gemini-1.5-flash": "Gemini 1.5 Flash",
}
SELECTED_GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.0-flash")
IMAGE_MODEL = "gemini-3-pro-image-preview"  # Nano Banana Pro for images

# Settings persistence
SETTINGS_FILE = BASE_DIR / "data" / "settings.json"


def load_settings():
    """Charge les settings depuis le fichier JSON"""
    global SELECTED_GEMINI_MODEL
    if SETTINGS_FILE.exists():
        try:
            with open(SETTINGS_FILE, "r") as f:
                settings = json.load(f)
            SELECTED_GEMINI_MODEL = settings.get("gemini_model", SELECTED_GEMINI_MODEL)
        except Exception:
            pass


def save_settings():
    """Sauvegarde les settings dans le fichier JSON"""
    SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
    settings = {"gemini_model": SELECTED_GEMINI_MODEL}
    with open(SETTINGS_FILE, "w") as f:
        json.dump(settings, f, indent=2)


load_settings()


# === HELPERS ===


def get_stats():
    """Calcule les statistiques du dashboard"""
    output_dir = BASE_DIR / "output"
    monitoring_dir = BASE_DIR / "data" / "monitoring"

    proposals = len(list(output_dir.glob("proposal_*.md"))) if output_dir.exists() else 0
    posts = len(list(output_dir.glob("linkedin_post_*.md"))) if output_dir.exists() else 0

    articles = 0
    if monitoring_dir.exists():
        for f in monitoring_dir.glob("veille_*.json"):
            try:
                data = json.loads(f.read_text())
                articles += data.get("total_relevant", 0)
            except Exception:
                pass

    return {"proposals": proposals, "posts": posts, "articles": articles}


def get_recent_files(limit=8):
    """Liste les fichiers recemment generes"""
    output_dir = BASE_DIR / "output"
    if not output_dir.exists():
        return []

    files = []
    for f in sorted(output_dir.glob("*.md"), key=lambda x: x.stat().st_mtime, reverse=True)[:limit]:
        files.append(
            {
                "name": f.name,
                "date": datetime.fromtimestamp(f.stat().st_mtime).strftime("%d/%m/%Y %H:%M"),
                "path": str(f),
            }
        )
    return files


def send_sse(event: str, data: dict) -> str:
    """Formate un message SSE"""
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


# === AUTHENTICATION ROUTES ===


@app.get("/login", response_class=HTMLResponse)
async def login_page(request: Request):
    """Page de connexion"""
    # Si déjà connecté, rediriger vers le dashboard
    if get_current_user(request):
        return RedirectResponse(url="/", status_code=302)

    return templates.TemplateResponse("login.html", {"request": request})


@app.post("/login")
@limiter.limit("10/minute")
async def login(request: Request, username: str = Form(...), password: str = Form(...)):
    """Authentification"""
    if authenticate_user(username, password):
        request.session["user"] = username
        return {"success": True, "redirect": "/"}
    else:
        return JSONResponse(
            {"detail": "Nom d'utilisateur ou mot de passe incorrect"}, status_code=401
        )


@app.get("/logout")
async def logout(request: Request):
    """Déconnexion"""
    request.session.clear()
    return RedirectResponse(url="/login", status_code=302)


# === LINKEDIN OAUTH ===


@app.get("/auth/linkedin")
async def linkedin_auth_start():
    """Demarre le flux OAuth LinkedIn"""
    from utils.linkedin_client import LinkedInClient, is_linkedin_configured

    if not is_linkedin_configured():
        return HTMLResponse(
            """
            <h1>LinkedIn OAuth non configure</h1>
            <p>Veuillez configurer les variables suivantes dans votre .env :</p>
            <ul>
                <li>LINKEDIN_CLIENT_ID</li>
                <li>LINKEDIN_CLIENT_SECRET</li>
                <li>LINKEDIN_REDIRECT_URI</li>
            </ul>
            <p>Consultez <a href="https://www.linkedin.com/developers/apps">LinkedIn Developers</a> pour créer une app.</p>
        """,
            status_code=400,
        )

    try:
        client = LinkedInClient()
        auth_url = client.get_auth_url()
        return RedirectResponse(auth_url)
    except Exception as e:
        return HTMLResponse(f"<h1>Erreur</h1><p>{safe_error_message(e)}</p>", status_code=500)


@app.get("/auth/linkedin/callback")
async def linkedin_auth_callback(
    code: str = None, error: str = None, error_description: str = None
):
    """Gere le callback OAuth LinkedIn"""
    from utils.linkedin_client import LinkedInClient

    # Check for errors
    if error:
        return HTMLResponse(
            """
            <h1>Erreur OAuth LinkedIn</h1>
            <p><strong>Erreur :</strong> {error}</p>
            <p><strong>Description :</strong> {error_description or 'N/A'}</p>
            <p><a href="/">Retour au dashboard</a></p>
        """,
            status_code=400,
        )

    if not code:
        return HTMLResponse(
            """
            <h1>Erreur OAuth LinkedIn</h1>
            <p>Code d autorisation manquant</p>
            <p><a href="/">Retour au dashboard</a></p>
        """,
            status_code=400,
        )

    try:
        client = LinkedInClient()
        token_data = client.exchange_code_for_token(code)
        token_data.get("access_token")
        token_data.get("expires_in", "N/A")

        return HTMLResponse(
            """
            <!DOCTYPE html>
            <html>
            <head>
                <title>LinkedIn Connected!</title>
                <style>
                    body {{
                        font-family: 'Inter', sans-serif;
                        max-width: 800px;
                        margin: 50px auto;
                        padding: 20px;
                        background: #f8f9fa;
                    }}
                    .success {{
                        background: white;
                        padding: 30px;
                        border-radius: 12px;
                        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
                    }}
                    h1 {{ color: #0a66c2; }}
                    code {{
                        background: #f0f0f0;
                        padding: 15px;
                        display: block;
                        border-radius: 8px;
                        margin: 15px 0;
                        font-family: monospace;
                        word-break: break-all;
                    }}
                    .btn {{
                        background: #0a66c2;
                        color: white;
                        padding: 12px 24px;
                        text-decoration: none;
                        border-radius: 8px;
                        display: inline-block;
                        margin-top: 20px;
                    }}
                </style>
            </head>
            <body>
                <div class="success">
                    <h1>✅ LinkedIn connecté avec succès !</h1>
                    <p>Votre application est maintenant autorisée à publier sur LinkedIn.</p>
                    <p><strong>Expire dans :</strong> {expires_in} secondes</p>

                    <h2>Configuration</h2>
                    <p>Ajoutez cette ligne à votre fichier <code>.env</code> :</p>
                    <code>LINKEDIN_ACCESS_TOKEN={access_token}</code>

                    <p><strong>Important :</strong> Ce token expire. Vous devrez répéter le processus OAuth lorsqu il expirera.</p>

                    <a href="/" class="btn">Retour au dashboard</a>
                </div>
            </body>
            </html>
        """
        )

    except Exception:
        return HTMLResponse(
            """
            <h1>Erreur lors de l echange du code</h1>
            <p>{safe_error_message(e)}</p>
            <p><a href="/">Retour au dashboard</a></p>
        """,
            status_code=500,
        )


# === SETTINGS API ===


@app.get("/api/settings/model")
async def get_model_settings():
    """Retourne le modele Gemini selectionne et la liste des modeles disponibles"""
    return {
        "current_model": SELECTED_GEMINI_MODEL,
        "available_models": AVAILABLE_GEMINI_MODELS,
        "image_model": IMAGE_MODEL,
    }


@app.post("/api/settings/model")
async def set_model_settings(request: Request):
    """Met a jour le modele Gemini selectionne"""
    global SELECTED_GEMINI_MODEL
    data = await request.json()
    model = data.get("model")
    if model not in AVAILABLE_GEMINI_MODELS:
        return JSONResponse({"error": "Modele inconnu"}, status_code=400)
    SELECTED_GEMINI_MODEL = model
    save_settings()
    return {"success": True, "model": model}


# === PAGE ROUTES ===


@app.get("/", response_class=HTMLResponse)
async def dashboard(request: Request):
    return templates.TemplateResponse(
        "index.html",
        {
            "request": request,
            "active": "dashboard",
            "consultant_name": CONSULTANT_NAME,
            "stats": get_stats(),
            "recent_files": get_recent_files(),
        },
    )


@app.get("/proposal", response_class=HTMLResponse)
async def proposal_page(request: Request):
    return templates.TemplateResponse(
        "proposal.html",
        {
            "request": request,
            "active": "proposal",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.get("/proposal-modular", response_class=HTMLResponse)
async def proposal_modular_page(request: Request):
    return templates.TemplateResponse(
        "proposal-modular.html",
        {
            "request": request,
            "active": "proposal",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.get("/linkedin", response_class=HTMLResponse)
async def linkedin_page(request: Request):
    return templates.TemplateResponse(
        "linkedin.html",
        {
            "request": request,
            "active": "linkedin",
            "consultant_name": CONSULTANT_NAME,
        },
    )


# === API: PROPOSAL ===


@app.post("/api/proposal/generate")
@limiter.limit("5/minute")
async def api_proposal_generate(
    request: Request,
    tender_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Lance la generation d'une proposition commerciale"""
    text = ""

    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            # Extraire le texte du PDF
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif tender_text:
        text = tender_text
    else:
        return JSONResponse({"error": "Aucun appel d'offre fourni."}, status_code=400)

    if len(text.strip()) < 50:
        return JSONResponse(
            {"error": "L'appel d'offre semble trop court (min 50 caracteres)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "tender_text": text,
    }

    # Lancer en background
    thread = threading.Thread(target=_run_proposal, args=(job_id, text), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_proposal(job_id: str, tender_text: str):
    """Execute la generation de proposition en background"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Step 1: Analyse
        job["steps"].append({"step": "analyze", "status": "active", "progress": 5})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 15})

        # Step 2: Template
        job["steps"].append({"step": "template", "status": "active", "progress": 18})
        template = agent.load_template()
        job["steps"].append({"step": "template", "status": "done", "progress": 25})

        # Step 3: References
        job["steps"].append({"step": "references", "status": "active", "progress": 28})
        references = agent.match_references(tender_analysis)
        job["steps"].append({"step": "references", "status": "done", "progress": 45})

        # Step 4: CVs
        job["steps"].append({"step": "generate", "status": "active", "progress": 48})
        cvs = agent.load_cvs()
        adapted_cvs = agent.adapt_cvs(cvs, tender_analysis) if cvs else []

        # Step 5: Generate slides structure + markdown content
        slides = agent.generate_slides_structure(tender_analysis, references, template, adapted_cvs)
        proposal = agent.generate_proposal_content(tender_analysis, references, template)
        job["steps"].append({"step": "generate", "status": "done", "progress": 85})

        # Sauvegarder les fichiers
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        json_path = output_dir / f"proposal_{client_name}_{timestamp}.json"
        md_path = output_dir / f"proposal_{client_name}_{timestamp}.md"
        pptx_path = output_dir / f"proposal_{client_name}_{timestamp}.pptx"

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(proposal, f, indent=2, ensure_ascii=False)

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(
                f"# Proposition Commerciale - {tender_analysis.get('project_title', 'Projet')}\n\n"
            )
            f.write(f"**Client:** {tender_analysis.get('client_name', 'N/A')}\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y')}\n\n---\n\n")
            f.write(proposal["content"])

        # Generer le PPTX
        try:
            from utils.pptx_generator import build_proposal_pptx

            build_proposal_pptx(
                template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                slides_data=slides,
                output_path=str(pptx_path),
                consultant_info={
                    "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                    "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
                },
            )
        except Exception as e:
            print(f"   ⚠️  Erreur PPTX: {e}")
            pptx_path = None

        job["status"] = "done"
        job["result"] = {
            "content": proposal["content"],
            "client_name": tender_analysis.get("client_name", "Projet"),
            "project_title": tender_analysis.get("project_title", ""),
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "json_path": str(json_path.relative_to(BASE_DIR)),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)) if pptx_path else None,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/proposal/stream/{job_id}")
async def api_proposal_stream(job_id: str):
    """SSE stream pour la progression de la proposition"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            # Envoyer les nouvelles etapes
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: PROPOSAL FEEDBACK ===


@app.post("/api/proposal/regenerate")
async def api_proposal_regenerate(request: Request):
    """Regenere la proposition avec le feedback utilisateur"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_content = body.get("previous_content", "")
    job_id_ref = body.get("job_id", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    # Recuperer le tender_text du job original si possible
    tender_text = ""
    if job_id_ref and job_id_ref in jobs:
        tender_text = jobs[job_id_ref].get("tender_text", "")

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "tender_text": tender_text,
    }

    thread = threading.Thread(
        target=_run_proposal_feedback,
        args=(job_id, previous_content, feedback, tender_text),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


# === API: PROPOSAL MODULAR GENERATION ===


@app.post("/api/proposal/generate-section")
@limiter.limit("10/minute")
async def api_proposal_generate_section(
    request: Request,
    section: str = Form(...),
    tender_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Génère une section spécifique de la proposition"""
    # Récupérer le texte de l'appel d'offre
    text = ""
    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif tender_text:
        text = tender_text
    else:
        return JSONResponse({"error": "Aucun appel d'offre fourni."}, status_code=400)

    if len(text.strip()) < 50:
        return JSONResponse({"error": "L'appel d'offre semble trop court."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal_section",
        "section": section,
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_proposal_section,
        args=(job_id, text, section),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_proposal_section(job_id: str, tender_text: str, section: str):
    """Génère une section spécifique de la proposition"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Analyser l'appel d'offre
        job["steps"].append({"step": "analyze", "status": "active", "progress": 10})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 30})

        # Générer la section demandée
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        slides = []
        section_name = ""

        if section == "agenda":
            slides = agent.generate_agenda_slide(tender_analysis)
            section_name = "Agenda"
        elif section == "context":
            slides = agent.generate_context_slides(tender_analysis)
            section_name = "Contexte"
        elif section == "approach":
            references = agent.match_references(tender_analysis)
            slides = agent.generate_approach_slides(tender_analysis, references)
            section_name = "Approche"
        elif section == "planning":
            slides = agent.generate_planning_slide(tender_analysis)
            section_name = "Planning"
        elif section == "budget":
            slides = agent.generate_budget_slide(tender_analysis)
            section_name = "Chiffrage"
        elif section == "references":
            references = agent.load_references()
            slides = agent.generate_references_slides(
                tender_analysis, {"all_references": references}
            )
            section_name = "Références"
        elif section == "cvs":
            slides = agent.generate_cv_slides(tender_analysis)
            section_name = "CVs"
        else:
            raise ValueError(f"Section inconnue: {section}")

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Générer le PPTX SANS cover ni closing (mode modulaire)
        job["steps"].append({"step": "pptx", "status": "active", "progress": 85})

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}.pptx"

        # Mode modulaire : UNIQUEMENT les slides de la section (sans cover ni closing)
        # L'utilisateur assemble ensuite les sections manuellement
        full_slides = slides

        from utils.pptx_generator import build_proposal_pptx

        build_proposal_pptx(
            template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
            slides_data=full_slides,
            output_path=str(pptx_path),
            consultant_info={
                "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
            },
        )

        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "section": section,
            "section_name": section_name,
            "slides_count": len(slides),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)),
            "tender_text": tender_text,  # Inclure le texte pour la régénération
            "slides_data": slides,  # Inclure les données des slides
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/proposal/preview")
@limiter.limit("20/minute")
async def api_proposal_preview(
    request: Request,
    pptx_path: str = Form(...),
):
    """Génère des images de prévisualisation des slides"""
    try:
        import subprocess
        import tempfile

        from PIL import Image

        # Résoudre le chemin absolu du PPTX
        file_path = BASE_DIR / pptx_path

        if not file_path.exists():
            return JSONResponse({"error": "Fichier PPTX introuvable"}, status_code=404)

        # Créer un répertoire temporaire pour les images
        temp_dir = Path(tempfile.mkdtemp(dir=BASE_DIR / "output"))

        try:
            # Convertir PPTX en PNG avec LibreOffice (méthode la plus fiable)
            # Note: LibreOffice doit être installé sur le système
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "png",
                    "--outdir",
                    str(temp_dir),
                    str(file_path),
                ],
                check=True,
                timeout=30,
                capture_output=True,
            )

            # Récupérer toutes les images générées
            preview_images = sorted(temp_dir.glob("*.png"))

            if not preview_images:
                # Fallback: essayer avec python-pptx + Pillow (moins fiable
                # mais sans dépendance externe)
                from pptx import Presentation

                prs = Presentation(str(file_path))
                preview_images = []

                for i, slide in enumerate(prs.slides):
                    # Cette approche est limitée car python-pptx ne rend pas visuellement
                    # On crée juste un placeholder pour indiquer que la
                    # conversion a échoué
                    pass

                if not preview_images:
                    return JSONResponse(
                        {
                            "error": "Impossible de générer la prévisualisation. LibreOffice n'est pas installé.",
                            "preview_images": [],
                        },
                        status_code=500,
                    )

            # Créer des versions miniatures pour la galerie
            thumbnails = []
            for img_path in preview_images:
                # Copier vers le dossier output avec un nom unique
                thumb_name = f"preview_{uuid.uuid4().hex[:8]}_{img_path.name}"
                thumb_path = BASE_DIR / "output" / thumb_name

                # Créer une miniature (optionnel, pour optimiser le chargement)
                with Image.open(img_path) as img:
                    img.thumbnail((800, 600))
                    img.save(thumb_path, "PNG")

                thumbnails.append(str(Path("output") / thumb_name))

                # Nettoyer le fichier temporaire
                img_path.unlink()

            # Nettoyer le répertoire temporaire
            temp_dir.rmdir()

            return JSONResponse({"preview_images": thumbnails, "slide_count": len(thumbnails)})

        except subprocess.CalledProcessError:
            # LibreOffice n'est pas disponible ou a échoué
            return JSONResponse(
                {
                    "error": "La génération de prévisualisation nécessite LibreOffice. Installez-le avec: brew install libreoffice (Mac) ou apt-get install libreoffice (Linux)",
                    "preview_images": [],
                },
                status_code=500,
            )
        except Exception as e:
            return JSONResponse(
                {
                    "error": f"Erreur lors de la génération de prévisualisation: {str(e)}",
                    "preview_images": [],
                },
                status_code=500,
            )

    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/proposal/regenerate-section")
@limiter.limit("10/minute")
async def api_proposal_regenerate_section(
    request: Request,
    section: str = Form(...),
    feedback: str = Form(...),
    tender_text: Optional[str] = Form(None),
    # JSON string des slides précédentes
    previous_slides: Optional[str] = Form(None),
):
    """Régénère une section avec feedback en langage naturel"""

    if not feedback.strip():
        return JSONResponse({"error": "Le feedback ne peut pas être vide."}, status_code=400)

    if not tender_text or len(tender_text.strip()) < 50:
        return JSONResponse(
            {"error": "Le texte de l'appel d'offre est requis pour la régénération."},
            status_code=400,
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "proposal_section",
        "section": section,
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    # Parse previous_slides si fourni
    prev_slides = []
    if previous_slides:
        try:
            prev_slides = json.loads(previous_slides)
        except BaseException:
            pass

    thread = threading.Thread(
        target=_run_proposal_section_with_feedback,
        args=(job_id, tender_text, section, feedback, prev_slides),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_proposal_section_with_feedback(
    job_id: str, tender_text: str, section: str, feedback: str, previous_slides: List[dict]
):
    """Régénère une section en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        # Analyser l'appel d'offre
        job["steps"].append({"step": "analyze", "status": "active", "progress": 10})
        tender_analysis = agent.analyze_tender(tender_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 30})

        # Interpréter le feedback avec le LLM
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        # Construire le contexte des slides précédentes
        previous_context = ""
        if previous_slides:
            previous_context = """
SLIDES PRÉCÉDENTES:
{json.dumps(previous_slides, indent=2, ensure_ascii=False)[:2000]}
"""

        # Utiliser le LLM pour interpréter le feedback et générer des
        # instructions structurées
        interpretation = agent.llm_client.generate(
            prompt="""L'utilisateur veut modifier la section "{section}" de sa proposition commerciale.

{previous_context}

FEEDBACK UTILISATEUR (en langage naturel):
{feedback}

Analyse ce feedback et génère des instructions STRUCTURÉES pour régénérer la section.
Réponds en JSON avec ces champs:
{
                "changes_requested": ["Liste des modifications demandées"],
  "visual_changes": ["Modifications visuelles (diagrammes, mise en page)"],
  "content_changes": ["Modifications de contenu (texte, chiffres, structure)"],
  "tone_changes": "Changement de ton si demandé",
  "specific_instructions": "Instructions précises pour la régénération"
}

Exemples:
- "Planning trop court" → ajouter des phases, augmenter les durées
- "Moins de bullets" → réduire à 2-3 bullets par slide
- "Diagramme flow → cycle" → remplacer le type de diagramme
- "Augmenter budget 20%" → multiplier les montants par 1.2""",
            system_prompt="Tu es un expert en interprétation de feedback pour la génération de documents. Sois précis et actionnable.",
            temperature=0.4,
            max_tokens=1500,
        )

        # Parser l'interprétation
        try:
            feedback_data = json.loads(interpretation.strip())
        except BaseException:
            feedback_data = {"specific_instructions": feedback, "changes_requested": [feedback]}

        # Générer la section avec les instructions du feedback
        slides = []
        section_name = ""
        references = None

        # Enrichir tender_analysis avec les instructions du feedback
        tender_analysis["feedback_instructions"] = feedback_data.get(
            "specific_instructions", feedback
        )
        tender_analysis["feedback_data"] = feedback_data

        if section == "agenda":
            slides = agent.generate_agenda_slide(tender_analysis)
            section_name = "Agenda"
        elif section == "context":
            slides = agent.generate_context_slides(tender_analysis)
            section_name = "Contexte"
        elif section == "approach":
            references = agent.match_references(tender_analysis)
            slides = agent.generate_approach_slides(tender_analysis, references)
            section_name = "Approche"
        elif section == "planning":
            slides = agent.generate_planning_slide(tender_analysis)
            section_name = "Planning"
        elif section == "budget":
            slides = agent.generate_budget_slide(tender_analysis)
            section_name = "Chiffrage"
        elif section == "references":
            references = agent.load_references()
            slides = agent.generate_references_slides(
                tender_analysis, {"all_references": references}
            )
            section_name = "Références"
        elif section == "cvs":
            slides = agent.generate_cv_slides(tender_analysis)
            section_name = "CVs"
        else:
            raise ValueError(f"Section inconnue: {section}")

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Générer le PPTX
        job["steps"].append({"step": "pptx", "status": "active", "progress": 85})

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}_v2.pptx"

        full_slides = slides

        from utils.pptx_generator import build_proposal_pptx

        build_proposal_pptx(
            template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
            slides_data=full_slides,
            output_path=str(pptx_path),
            consultant_info={
                "name": os.getenv("CONSULTANT_NAME", "Jean-Sébastien Abessouguie Bayiha"),
                "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
            },
        )

        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "section": section,
            "section_name": section_name,
            "slides_count": len(slides),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)),
            "feedback_applied": feedback_data.get("changes_requested", [feedback]),
            "slides_data": slides,  # Inclure les données des slides pour la prochaine itération
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)
        print(f"Error in regenerate section: {safe_traceback()}")


@app.post("/api/proposal/regenerate-slide")
@limiter.limit("20/minute")
async def api_proposal_regenerate_slide(
    request: Request,
    section: str = Form(...),
    slide_index: int = Form(...),
    slide_title: str = Form(...),
    slide_bullets: str = Form(...),  # JSON array
    diagram_type: Optional[str] = Form(None),
    tender_text: Optional[str] = Form(None),
):
    """Régénère une seule slide avec le nouveau contenu"""

    if not tender_text or len(tender_text.strip()) < 50:
        return JSONResponse({"error": "Le texte de l'appel d'offre est requis."}, status_code=400)

    try:
        # Parse bullets
        bullets = json.loads(slide_bullets)

        # Charger le PPTX existant pour récupérer toutes les slides
        # On va chercher le fichier le plus récent pour cette section
        output_dir = BASE_DIR / "output"
        matching_files = sorted(
            output_dir.glob(f"proposal_{section}_*.pptx"),
            key=lambda p: p.stat().st_mtime,
            reverse=True,
        )

        if not matching_files:
            return JSONResponse(
                {"error": "Aucun PPTX existant trouvé pour cette section"}, status_code=404
            )

        latest_pptx = matching_files[0]

        # Charger les slides depuis le PPTX existant
        from pptx import Presentation

        prs = Presentation(str(latest_pptx))

        # Vérifier que l'index est valide (en tenant compte de la slide de
        # couverture)
        if slide_index + 1 >= len(prs.slides):
            return JSONResponse(
                {"error": f"Index de slide invalide: {slide_index}"}, status_code=400
            )

        # Modifier la slide ciblée
        # +1 car slide 0 = couverture
        target_slide = prs.slides[slide_index + 1]

        # Trouver les textframes et mettre à jour
        for shape in target_slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                # Le premier textframe est généralement le titre
                if shape.name.startswith("Title") or "titre" in shape.name.lower():
                    text_frame.text = slide_title
                    break

        # Mettre à jour le contenu (bullets)
        for shape in target_slide.shapes:
            if shape.has_text_frame and not (
                shape.name.startswith("Title") or "titre" in shape.name.lower()
            ):
                text_frame = shape.text_frame
                text_frame.clear()

                # Ajouter les nouveaux bullets
                for i, bullet in enumerate(bullets):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

                    # Styling (utiliser les styles existants si possible)
                    from pptx.util import Pt

                    if p.runs:
                        p.runs[0].font.size = Pt(11)

                break

        # Sauvegarder le PPTX modifié
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        agent = ProposalGeneratorAgent()
        tender_analysis = agent.analyze_tender(tender_text)
        client_name = tender_analysis.get("client_name", "client").replace(" ", "_")

        new_pptx_path = output_dir / f"proposal_{section}_{client_name}_{timestamp}_edited.pptx"
        prs.save(str(new_pptx_path))

        # Extraire les données des slides pour le frontend
        slides_data = []
        for i, slide in enumerate(prs.slides):
            if i == 0:  # Skip cover
                continue

            slide_data = {"title": "", "bullets": [], "visual": None}

            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.name.startswith("Title") or "titre" in shape.name.lower():
                        slide_data["title"] = shape.text_frame.text
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_data["bullets"].append(paragraph.text.strip())

            slides_data.append(slide_data)

        return JSONResponse(
            {
                "pptx_path": str(new_pptx_path.relative_to(BASE_DIR)),
                "slides_data": slides_data,
                "message": "Slide régénérée avec succès",
            }
        )

    except Exception as e:
        print(f"Error in regenerate slide: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/proposal/export-to-slides")
@limiter.limit("5/minute")
async def api_export_to_slides(
    request: Request,
    section: str = Form(...),
    tender_text: Optional[str] = Form(None),
    slides_data: str = Form(...),  # JSON string
):
    """Exporte une section vers Google Slides"""

    try:
        # Parser les données des slides avec sanitization robuste
        # Fix: le JSON envoyé depuis le frontend peut contenir des caractères de contrôle
        # ou des chaînes mal échappées, causant "Unterminated string"
        import re

        sanitized = slides_data.strip()
        # Supprimer les caractères de contrôle (sauf \n, \t, \r qui sont gérés
        # par JSON)
        sanitized = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", sanitized)
        # Remplacer les vrais retours à la ligne dans les valeurs string par \\n
        # (ceux qui ne sont pas déjà échappés)
        sanitized = sanitized.replace("\r\n", "\\n").replace("\r", "\\n")

        try:
            slides = json.loads(sanitized)
        except json.JSONDecodeError:
            # Fallback: essayer de nettoyer plus agressivement
            # Parfois les guillemets internes ne sont pas échappés
            sanitized = sanitized.replace("\t", "\\t")
            # Tenter un nettoyage ligne par ligne
            lines = sanitized.split("\n")
            cleaned_lines = []
            for line in lines:
                # Ne pas modifier les lignes structurelles JSON
                cleaned_lines.append(line)
            sanitized = "\n".join(cleaned_lines)
            slides = json.loads(sanitized)

        # Analyser le tender pour obtenir le titre
        if tender_text and len(tender_text.strip()) >= 50:
            agent = ProposalGeneratorAgent()
            tender_analysis = agent.analyze_tender(tender_text)
            client_name = tender_analysis.get("client_name", "Client")
            project_title = tender_analysis.get("project_title", "Projet")
        else:
            client_name = "Client"
            project_title = "Projet"

        # Générer le titre de la présentation
        section_names = {
            "agenda": "Agenda",
            "context": "Contexte",
            "approach": "Approche",
            "planning": "Planning",
            "budget": "Chiffrage",
            "references": "Références",
            "cvs": "CVs",
        }

        section_name = section_names.get(section, section.capitalize())
        presentation_title = f"{client_name} - {project_title} - {section_name}"

        # Exporter vers Google Slides
        from utils.google_api import GoogleAPIClient

        try:
            google_client = GoogleAPIClient()
            presentation_id = google_client.export_pptx_to_slides(
                slides_data=slides, title=presentation_title
            )

            if presentation_id:
                presentation_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"

                return JSONResponse(
                    {
                        "presentation_id": presentation_id,
                        "presentation_url": presentation_url,
                        "message": "Présentation Google Slides créée avec succès",
                    }
                )
            else:
                return JSONResponse(
                    {"error": "Échec de la création de la présentation Google Slides"},
                    status_code=500,
                )

        except Exception as e:
            # Si l'authentification Google échoue ou n'est pas configurée
            if "credentials" in str(e).lower() or "token" in str(e).lower():
                return JSONResponse(
                    {
                        "error": "Google API non configurée. Veuillez configurer vos credentials Google dans config/google_credentials.json",
                        "setup_required": True,
                    },
                    status_code=400,
                )
            else:
                raise

    except Exception as e:
        print(f"Error in export to slides: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


def _run_proposal_feedback(job_id: str, previous_content: str, feedback: str, tender_text: str):
    """Regenere la proposition en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = ProposalGeneratorAgent()

        job["steps"].append({"step": "analyze", "status": "done", "progress": 10})
        job["steps"].append({"step": "template", "status": "done", "progress": 20})
        job["steps"].append({"step": "references", "status": "done", "progress": 30})
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})

        # Analyser le tender_text si disponible pour avoir plus de contexte
        tender_context = ""
        if tender_text:
            try:
                agent.analyze_tender(tender_text)
                tender_context = """
Contexte de l'appel d'offre:
Client: {tender_analysis.get('client_name', 'N/A')}
Projet: {tender_analysis.get('project_title', 'N/A')}
Objectifs: {', '.join(tender_analysis.get('objectives', [])[:3])}
"""
            except Exception:
                tender_context = f"Appel d'offre: {tender_text[:1000]}"

        # Regenerer le contenu markdown avec feedback
        revised_content = agent.llm_client.generate(
            prompt="""Voici une proposition commerciale generee precedemment:

{previous_content[:6000]}

FEEDBACK DE L'UTILISATEUR:
{feedback}

{tender_context}

Reecris la proposition en integrant les corrections demandees dans le feedback.
Conserve la structure et le style professionnel. Ne modifie que ce qui est demande dans le feedback.
Sois concret et precis. Evite les majuscules inutiles et les emojis.
Fournis la proposition complete revisee.""",
            system_prompt="""Tu es {agent.consultant_info['name']}, {agent.consultant_info['title']} chez {agent.consultant_info['company']}.
Tu corriges une proposition commerciale selon le retour du consultant.
Reponds en francais de maniere professionnelle.""",
            temperature=0.6,
            max_tokens=8000,
        )

        job["steps"].append({"step": "generate", "status": "done", "progress": 60})
        job["steps"].append({"step": "slides", "status": "active", "progress": 65})

        # Regenerer les slides avec feedback - avec meilleur prompt
        revised_slides = None
        try:
            slides_response = agent.llm_client.generate(
                prompt="""Genere une structure complete de slides PowerPoint pour cette proposition commerciale revisee.

PROPOSITION REVISEE:
{revised_content[:5000]}

FEEDBACK APPLIQUE:
{feedback}

{tender_context}

Genere la structure COMPLETE avec tout le contenu. Reponds UNIQUEMENT en JSON valide:

[
  {{"type":"cover","client":"...","project":"...","date":"{datetime.now().strftime('%d/%m/%Y')}"}},
  {{"type":"section","title":"Notre comprehension du contexte","number":1}},
  {{"type":"content","title":"Contexte et enjeux","bullets":["Bullet 1","Bullet 2","Bullet 3","Bullet 4"],"subtitle":""}},
  {{"type":"content","title":"Objectifs","bullets":["Objectif 1","Objectif 2","Objectif 3"]}},
  {{"type":"section","title":"Notre expertise","number":2}},
  {{"type":"content","title":"Notre valeur ajoutee","bullets":["Valeur 1","Valeur 2","Valeur 3","Valeur 4"]}},
  {{"type":"section","title":"Notre demarche","number":3}},
  {{"type":"content","title":"Phases du projet","bullets":["Phase 1: Description","Phase 2: Description","Phase 3: Description"]}},
  {{"type":"section","title":"Planning","number":4}},
  {{"type":"table","title":"Planning previsionnel","headers":["Phase","Duree","Livrables"],"rows":[["Phase 1","2 sem","Livrable"],["Phase 2","3 sem","Livrable"]]}},
  {{"type":"section","title":"Budget","number":5}},
  {{"type":"table","title":"Chiffrage","headers":["Prestation","Jours","TJM","Total"],"rows":[["Phase 1","10j","1000€","10000€"],["Total","","","XX XXX€"]]}},
  {{"type":"closing"}}
]

IMPORTANT:
- Chaque slide "content" DOIT avoir un champ "bullets" avec 4-6 points concrets
- Chaque slide "table" DOIT avoir "headers" et "rows"
- Total: 10-16 slides
- Integre les modifications demandees dans le feedback""",
                system_prompt="Tu generes des structures de slides PowerPoint completes en JSON. Sois concret et professionnel.",
                temperature=0.5,
                max_tokens=8000,
            )

            json_start = slides_response.find("[")
            json_end = slides_response.rfind("]") + 1
            if json_start >= 0 and json_end > json_start:
                revised_slides = json.loads(slides_response[json_start:json_end])

                # Valider que les slides ont du contenu
                for slide in revised_slides:
                    if slide.get("type") == "content" and "bullets" not in slide:
                        slide["bullets"] = [
                            "Point a developper",
                            "Deuxieme point",
                            "Troisieme point",
                        ]
                    elif slide.get("type") == "table" and "rows" not in slide:
                        slide["headers"] = slide.get("headers", ["Colonne 1", "Colonne 2"])
                        slide["rows"] = [["Donnee 1", "Donnee 2"]]

                print(f"   Slides revisees generees: {len(revised_slides)}")
        except Exception as e:
            print(f"   Erreur generation slides: {e}")
            revised_slides = None

        job["steps"].append({"step": "slides", "status": "done", "progress": 85})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"proposal_revised_{timestamp}.md"
        json_path = output_dir / f"proposal_revised_{timestamp}.json"

        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Proposition commerciale (revisee)\n\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y')}\n\n---\n\n")
            f.write(revised_content)

        proposal_data = {
            "content": revised_content,
            "feedback_applied": feedback,
            "slides": revised_slides,
        }
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(proposal_data, f, indent=2, ensure_ascii=False)

        # Generer le PPTX revise
        pptx_path = None
        if revised_slides:
            try:
                from utils.pptx_generator import build_proposal_pptx

                pptx_path = output_dir / f"proposal_revised_{timestamp}.pptx"
                build_proposal_pptx(
                    template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                    slides_data=revised_slides,
                    output_path=str(pptx_path),
                    consultant_info={
                        "name": CONSULTANT_NAME,
                        "company": os.getenv("COMPANY_NAME", "Consulting Tools"),
                    },
                )
                print(f"   PPTX revise genere: {pptx_path}")
            except Exception as e:
                print(f"   Erreur PPTX revise: {e}")
                pptx_path = None

        job["status"] = "done"
        job["result"] = {
            "content": revised_content,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "json_path": str(json_path.relative_to(BASE_DIR)),
            "pptx_path": str(pptx_path.relative_to(BASE_DIR)) if pptx_path else None,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === API: LINKEDIN ===


@app.post("/api/linkedin/generate")
@limiter.limit("5/minute")
async def api_linkedin_generate(request: Request):
    """Lance la veille et generation de posts LinkedIn"""
    body = await request.json()
    post_type = body.get("post_type", "insight")
    num_posts = min(body.get("num_posts", 3), 5)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "linkedin",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_linkedin, args=(job_id, post_type, num_posts), daemon=True
    )
    thread.start()

    return {"job_id": job_id}


def _run_linkedin(job_id: str, post_type: str, num_posts: int):
    """Execute la veille LinkedIn en background"""
    job = jobs[job_id]

    try:
        agent = LinkedInMonitorAgent()

        # Step 1: RSS
        job["steps"].append({"step": "rss", "status": "active", "progress": 5})
        monitoring_data = agent.collect_monitoring_data()
        job["steps"].append({"step": "rss", "status": "done", "progress": 30})
        job["steps"].append({"step": "web", "status": "done", "progress": 40})

        # Step 2: Trends
        job["steps"].append({"step": "trends", "status": "active", "progress": 45})
        trends = agent.analyze_trends(monitoring_data["articles"])
        job["steps"].append({"step": "trends", "status": "done", "progress": 60})

        # Step 3: Posts
        job["steps"].append({"step": "posts", "status": "active", "progress": 65})

        if num_posts == 1:
            posts = [agent.generate_linkedin_post(trends, monitoring_data["articles"], post_type)]
        else:
            posts = agent.generate_multiple_posts(trends, monitoring_data["articles"], num_posts)

        job["steps"].append({"step": "posts", "status": "done", "progress": 100})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        for i, post in enumerate(posts, 1):
            md_path = output_dir / f"linkedin_post_{timestamp}_{i}.md"
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(f"# Post LinkedIn - {post['post_type'].title()}\n\n")
                f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
                f.write("## Post Principal\n\n")
                f.write(post["main_post"])

        job["status"] = "done"
        job["result"] = {
            "posts": [
                {
                    "main_post": p["main_post"],
                    "post_type": p["post_type"],
                    "source_articles": p.get("source_articles", []),
                }
                for p in posts
            ],
            "articles": [
                {
                    "title": a.get("title", ""),
                    "link": a.get("link", ""),
                    "source": a.get("source", a.get("source_type", "")),
                    "relevance_score": a.get("relevance_score", 0),
                }
                for a in monitoring_data["articles"][:20]
            ],
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/linkedin/stream/{job_id}")
async def api_linkedin_stream(job_id: str):
    """SSE stream pour la progression LinkedIn"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: LINKEDIN FEEDBACK ===


@app.post("/api/linkedin/regenerate")
async def api_linkedin_regenerate(request: Request):
    """Regenere les posts LinkedIn avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_posts = body.get("previous_posts", [])

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "linkedin",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_linkedin_feedback,
        args=(job_id, previous_posts, feedback),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_linkedin_feedback(job_id: str, previous_posts: list, feedback: str):
    """Regenere les posts LinkedIn en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = LinkedInMonitorAgent()

        job["steps"].append({"step": "rss", "status": "done", "progress": 10})
        job["steps"].append({"step": "web", "status": "done", "progress": 20})
        job["steps"].append({"step": "trends", "status": "done", "progress": 30})
        job["steps"].append({"step": "posts", "status": "active", "progress": 40})

        revised_posts = []
        for i, post in enumerate(previous_posts):
            revised = agent.llm_client.generate(
                prompt="""Voici un post LinkedIn genere precedemment:

{
                    post.get(
                        'main_post',
                        '')}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce post en integrant les corrections demandees.
Conserve le format LinkedIn (hook, contenu, question d'engagement, hashtags).
Ne modifie que ce qui est demande dans le feedback.""",
                system_prompt="""Tu es {
                    agent.consultant_info['name']}, {
                    agent.consultant_info['title']} chez {
                        agent.consultant_info['company']}.
Tu corriges un post LinkedIn selon le retour du consultant.

REGLES IMPERATIVES:
- Ne JAMAIS inventer d'anecdotes, d'exemples fictifs ou d'experiences personnelles
- Ne JAMAIS fabriquer de chiffres, statistiques ou citations""",
                temperature=0.6,
                max_tokens=1500,
            )
            revised_posts.append(
                {
                    "main_post": revised,
                    "post_type": post.get("post_type", "insight"),
                    "source_articles": post.get("source_articles", []),
                }
            )

        job["steps"].append({"step": "posts", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "posts": revised_posts,
            "articles": [],
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/linkedin/publish")
@limiter.limit("5/minute")
async def api_linkedin_publish(request: Request):
    """Publie un post directement sur LinkedIn"""
    from utils.linkedin_client import LinkedInClient, has_linkedin_access_token
    from utils.validation import sanitize_text_input

    # Check if LinkedIn is configured
    if not has_linkedin_access_token():
        return JSONResponse(
            {"error": "LinkedIn non configure. Completez le flux OAuth via /auth/linkedin"},
            status_code=400,
        )

    body = await request.json()
    text = body.get("text", "").strip()
    visibility = body.get("visibility", "PUBLIC")

    # Validate
    if not text:
        return JSONResponse({"error": "Texte du post manquant"}, status_code=400)

    # Sanitize and validate length
    try:
        text = sanitize_text_input(text, max_length=3000, field_name="post")
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=400)

    if visibility not in ["PUBLIC", "CONNECTIONS"]:
        visibility = "PUBLIC"

    try:
        client = LinkedInClient()
        result = client.publish_post(text=text, visibility=visibility)

        return JSONResponse(
            {
                "message": "Post publie avec succes sur LinkedIn",
                "id": result.get("id"),
                "url": result.get("url"),
                "status": result.get("status"),
            }
        )

    except ValueError as e:
        # Configuration or validation errors
        return JSONResponse(
            {"error": f"Erreur de configuration: {safe_error_message(e)}"}, status_code=400
        )
    except Exception as e:
        # API errors
        return JSONResponse(
            {"error": f"Erreur lors de la publication: {safe_error_message(e)}"}, status_code=500
        )


@app.get("/api/linkedin/status")
async def api_linkedin_status():
    """Verifie le status de connexion LinkedIn"""
    from utils.linkedin_client import has_linkedin_access_token

    return JSONResponse({"connected": has_linkedin_access_token()})


# === PAGE & API: ARTICLE TO POST ===


@app.get("/article", response_class=HTMLResponse)
async def article_page(request: Request):
    return templates.TemplateResponse(
        "article.html",
        {
            "request": request,
            "active": "article",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/article/generate")
async def api_article_generate(request: Request):
    """Lance la generation d'un post a partir d'un article"""
    body = await request.json()
    url = body.get("url", "").strip()
    tone = body.get("tone", "expert")

    if not url:
        return JSONResponse({"error": "URL manquante."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "article",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(target=_run_article_to_post, args=(job_id, url, tone), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_article_to_post(job_id: str, url: str, tone: str):
    """Execute la generation de post depuis un article en background"""
    job = jobs[job_id]

    try:
        agent = ArticleToPostAgent()

        # Step 1: Fetch
        job["steps"].append({"step": "fetch", "status": "active", "progress": 10})
        article = agent.fetch_article(url)
        job["steps"].append({"step": "fetch", "status": "done", "progress": 30})

        # Step 2: Generate main post
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})
        result = agent.generate_post(article, tone=tone)
        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Step 3: Done (variant already generated inside generate_post)
        job["steps"].append({"step": "variant", "status": "done", "progress": 100})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"article_post_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Post LinkedIn - Partage d'article\n\n")
            f.write(f"**Article:** [{article['title']}]({url})\n")
            f.write(f"**Ton:** {tone}\n\n## Post Principal\n\n")
            f.write(result["main_post"])
            f.write("\n\n---\n\n## Version Courte\n\n")
            f.write(result["short_version"])

        job["status"] = "done"
        job["result"] = {
            "main_post": result["main_post"],
            "short_version": result["short_version"],
            "article_title": article["title"],
            "article_url": url,
            "tone": tone,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/article/stream/{job_id}")
async def api_article_stream(job_id: str):
    """SSE stream pour la progression article-to-post"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: ARTICLE FEEDBACK ===


@app.post("/api/article/regenerate")
async def api_article_regenerate(request: Request):
    """Regenere le post article avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_main = body.get("previous_main", "")
    previous_short = body.get("previous_short", "")
    article_url = body.get("article_url", "")
    tone = body.get("tone", "expert")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "article",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_article_feedback,
        args=(job_id, previous_main, previous_short, feedback, article_url, tone),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_article_feedback(
    job_id: str, previous_main: str, previous_short: str, feedback: str, article_url: str, tone: str
):
    """Regenere le post article en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = ArticleToPostAgent()

        job["steps"].append({"step": "fetch", "status": "done", "progress": 10})
        job["steps"].append({"step": "generate", "status": "active", "progress": 30})

        # Regenerer le post principal
        revised_main = agent.llm_client.generate(
            prompt="""Voici un post LinkedIn genere precedemment pour partager un article:

{previous_main}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce post en integrant les corrections demandees.
Conserve le format LinkedIn (hook, perspective, points cles, appel a l'action, hashtags).
Ne modifie que ce qui est demande dans le feedback.""",
            system_prompt="""Tu es {
                agent.consultant_info['name']}, {
                agent.consultant_info['title']} chez {
                agent.consultant_info['company']}.
Tu corriges un post LinkedIn selon le retour du consultant.

REGLES IMPERATIVES:
- Base-toi UNIQUEMENT sur le contenu de l'article. Ne fabrique aucun fait, chiffre ou exemple.
- Ne JAMAIS inventer d'anecdotes personnelles ou d'experiences fictives""",
            temperature=0.6,
            max_tokens=1500,
        )

        job["steps"].append({"step": "generate", "status": "done", "progress": 70})
        job["steps"].append({"step": "variant", "status": "active", "progress": 75})

        # Regenerer la version courte
        revised_short = agent.llm_client.generate(
            prompt="""A partir de ce post LinkedIn revise, cree une version courte (400-600 caracteres max) qui va droit au point.

Post revise:
{revised_main}

Garde le hook et la question finale, compresse le milieu.
NE PAS inventer d'anecdotes ou d'exemples.""",
            system_prompt="""Tu es {
                agent.consultant_info['name']}, {
                agent.consultant_info['title']} chez {
                agent.consultant_info['company']}.""",
            temperature=0.6,
            max_tokens=600,
        )

        job["steps"].append({"step": "variant", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "main_post": revised_main,
            "short_version": revised_short,
            "article_url": article_url,
            "tone": tone,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: MEETING SUMMARIZER ===


@app.get("/meeting", response_class=HTMLResponse)
async def meeting_page(request: Request):
    return templates.TemplateResponse(
        "meeting.html",
        {
            "request": request,
            "active": "meeting",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/meeting/generate")
async def api_meeting_generate(
    transcript_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Lance la generation d'un compte rendu de reunion"""
    text = ""

    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif transcript_text:
        text = transcript_text
    else:
        return JSONResponse({"error": "Aucun transcript fourni."}, status_code=400)

    if len(text.strip()) < 30:
        return JSONResponse({"error": "Le transcript semble trop court."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "meeting",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(target=_run_meeting_summarizer, args=(job_id, text), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_meeting_summarizer(job_id: str, transcript: str):
    """Execute la generation de compte rendu en background"""
    job = jobs[job_id]

    try:
        from agents.meeting_summarizer import MeetingSummarizerAgent

        agent = MeetingSummarizerAgent()

        # Step 1: Extraction
        job["steps"].append({"step": "extract", "status": "active", "progress": 10})
        extracted_result = agent.extract_key_info(transcript)
        extracted_info = extracted_result["extracted_info"]
        job["steps"].append({"step": "extract", "status": "done", "progress": 35})

        # Step 2: Compte rendu
        job["steps"].append({"step": "minutes", "status": "active", "progress": 40})
        minutes = agent.generate_minutes(transcript, extracted_info)
        job["steps"].append({"step": "minutes", "status": "done", "progress": 70})

        # Step 3: Email
        job["steps"].append({"step": "email", "status": "active", "progress": 75})
        email_result = agent.generate_email(extracted_info, minutes)
        job["steps"].append({"step": "email", "status": "done", "progress": 100})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"meeting_summary_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Compte Rendu de Reunion\n\n")
            f.write(f"**Genere le:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n---\n\n")
            f.write(minutes)
            f.write("\n\n---\n\n# Mail de Partage\n\n")
            f.write(f"**Objet:** {email_result['subject']}\n\n")
            f.write(email_result["body"])

        job["status"] = "done"
        job["result"] = {
            "minutes": minutes,
            "email": email_result,
            "md_path": str(md_path.relative_to(BASE_DIR)),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)
        print(f"Error in meeting summarizer: {safe_traceback()}")


@app.get("/api/meeting/stream/{job_id}")
async def api_meeting_stream(job_id: str):
    """SSE stream pour la progression du compte rendu"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: MEETING FEEDBACK ===


@app.post("/api/meeting/regenerate")
async def api_meeting_regenerate(request: Request):
    """Regenere le compte rendu avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_minutes = body.get("previous_minutes", "")
    previous_email = body.get("previous_email", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "meeting",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_meeting_feedback,
        args=(job_id, previous_minutes, previous_email, feedback),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_meeting_feedback(job_id: str, previous_minutes: str, previous_email: str, feedback: str):
    """Regenere le compte rendu en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = MeetingSummarizerAgent()

        job["steps"].append({"step": "extract", "status": "done", "progress": 10})
        job["steps"].append({"step": "minutes", "status": "active", "progress": 30})

        # Regenerer le compte rendu
        revised_minutes = agent.llm_client.generate(
            prompt="""Voici un compte rendu de reunion genere precedemment:

{previous_minutes}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce compte rendu en integrant les corrections demandees.
Conserve la structure professionnelle (contexte, points abordes, decisions, plan d'actions, points en suspens, prochaines etapes).
Ne modifie que ce qui est demande dans le feedback.""",
            system_prompt="""Tu es {
                agent.consultant_info['name']}, {
                agent.consultant_info['title']} chez {
                agent.consultant_info['company']}.
Tu corriges un compte rendu de reunion selon le retour du consultant.""",
            temperature=0.5,
            max_tokens=3000,
        )

        job["steps"].append({"step": "minutes", "status": "done", "progress": 70})
        job["steps"].append({"step": "email", "status": "active", "progress": 75})

        # Regenerer le mail
        revised_email = agent.llm_client.generate(
            prompt="""Voici un mail de partage de compte rendu genere precedemment:

{previous_email}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce mail en integrant les corrections demandees.
Conserve le format professionnel (objet, resume executif, decisions cles, actions, signature).
Ne modifie que ce qui est demande dans le feedback.""",
            system_prompt="""Tu es {
                agent.consultant_info['name']}, {
                agent.consultant_info['title']} chez {
                agent.consultant_info['company']}.
Tu corriges un mail professionnel selon le retour du consultant.""",
            temperature=0.5,
            max_tokens=1500,
        )

        job["steps"].append({"step": "email", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "minutes": revised_minutes,
            "email": revised_email,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/meeting/share-email")
@limiter.limit("10/minute")
async def api_meeting_share_email(request: Request):
    """Partage le compte rendu de reunion par email avec piece jointe"""
    from utils.gmail_client import GmailClient
    from utils.validation import validate_email

    body = await request.json()
    to_email = body.get("to_email", "").strip()
    meeting_summary = body.get("meeting_summary", "")
    meeting_title = body.get("meeting_title", "Sans titre")

    # Valider email
    if not to_email:
        return JSONResponse({"error": "Email destinataire manquant"}, status_code=400)

    if not validate_email(to_email):
        return JSONResponse({"error": "Email destinataire invalide"}, status_code=400)

    if not meeting_summary:
        return JSONResponse({"error": "Compte rendu manquant"}, status_code=400)

    try:
        # Get consultant info
        get_consultant_info()

        # Creer fichier temporaire pour le compte rendu
        import tempfile
        from datetime import datetime

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        temp_file = tempfile.NamedTemporaryFile(
            mode="w",
            suffix=".md",
            prefix=f"compte_rendu_{timestamp}_",
            delete=False,
            encoding="utf-8",
        )

        temp_file.write(meeting_summary)
        temp_file.close()

        # Construire email
        subject = f"Compte rendu de reunion - {meeting_title}"
        body = """Bonjour,

Veuillez trouver ci-joint le compte rendu de notre reunion.

Cordialement,
{consultant_info['name']}
{consultant_info['title']}
{consultant_info['company']}
"""

        # Envoyer email
        gmail = GmailClient()
        result = gmail.send_email(
            to=to_email, subject=subject, body=body, attachments=[temp_file.name]
        )

        # Supprimer fichier temporaire
        import os

        os.unlink(temp_file.name)

        return JSONResponse({"message": "Email envoye avec succes", "id": result["id"]})

    except FileNotFoundError as e:
        return JSONResponse(
            {"error": f"Fichier non trouve: {safe_error_message(e)}"}, status_code=404
        )
    except Exception as e:
        return JSONResponse(
            {"error": f"Erreur lors de l envoi: {safe_error_message(e)}"}, status_code=500
        )


# === PAGE & API: LINKEDIN COMMENT ===


@app.get("/comment", response_class=HTMLResponse)
async def comment_page(request: Request):
    return templates.TemplateResponse(
        "comment.html",
        {
            "request": request,
            "active": "comment",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/comment/generate")
@limiter.limit("10/minute")
async def api_comment_generate(request: Request):
    """Lance la generation de commentaires LinkedIn"""
    body = await request.json()
    post_input = body.get("post_input", "").strip()
    style = body.get("style", "insightful")

    if not post_input:
        return JSONResponse({"error": "Post LinkedIn manquant."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "comment",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_comment_generator, args=(job_id, post_input, style), daemon=True
    )
    thread.start()

    return {"job_id": job_id}


def _run_comment_generator(job_id: str, post_input: str, style: str):
    """Execute la generation de commentaires en background"""
    job = jobs[job_id]

    try:
        agent = LinkedInCommenterAgent()

        # Step 1: Extraction
        job["steps"].append({"step": "extract", "status": "active", "progress": 10})
        post_content = agent.extract_post_content(post_input)
        job["steps"].append({"step": "extract", "status": "done", "progress": 35})

        # Step 2: Generation
        job["steps"].append({"step": "generate", "status": "active", "progress": 40})
        result = agent.generate_comments(post_content, style=style)
        job["steps"].append({"step": "generate", "status": "done", "progress": 90})

        # Step 3: Sauvegarde
        job["steps"].append({"step": "save", "status": "active", "progress": 95})
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"linkedin_comment_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Commentaires LinkedIn\n\n")
            f.write(f"**Style:** {style}\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
            f.write("## Post original (extrait)\n\n")
            f.write(f"> {result['post_preview']}\n\n")
            f.write("---\n\n")
            f.write("## Commentaire Court\n\n")
            f.write(result["short"])
            f.write("\n\n---\n\n## Commentaire Moyen\n\n")
            f.write(result["medium"])
            f.write("\n\n---\n\n## Commentaire Long\n\n")
            f.write(result["long"])
            f.write("\n")

        job["steps"].append({"step": "save", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "short": result["short"],
            "medium": result["medium"],
            "long": result["long"],
            "post_preview": result["post_preview"],
            "style": style,
            "md_path": str(md_path.relative_to(BASE_DIR)),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/comment/stream/{job_id}")
async def api_comment_stream(job_id: str):
    """Streame le progres de la generation de commentaires"""

    async def event_generator():
        while True:
            if job_id not in jobs:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break

            job = jobs[job_id]
            yield f"data: {json.dumps(job)}\n\n"

            if job["status"] in ["done", "error"]:
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: COMMENT FEEDBACK ===


@app.post("/api/comment/regenerate")
async def api_comment_regenerate(request: Request):
    """Regenere les commentaires avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_short = body.get("previous_short", "")
    previous_medium = body.get("previous_medium", "")
    previous_long = body.get("previous_long", "")
    post_preview = body.get("post_preview", "")
    style = body.get("style", "insightful")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "comment",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_comment_feedback,
        args=(
            job_id,
            previous_short,
            previous_medium,
            previous_long,
            feedback,
            post_preview,
            style,
        ),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_comment_feedback(
    job_id: str,
    previous_short: str,
    previous_medium: str,
    previous_long: str,
    feedback: str,
    post_preview: str,
    style: str,
):
    """Regenere les commentaires en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = LinkedInCommenterAgent()

        job["steps"].append({"step": "extract", "status": "done", "progress": 10})
        job["steps"].append({"step": "generate", "status": "active", "progress": 30})

        # Charger le persona
        persona_path = BASE_DIR / "data" / "linkedin_persona.md"
        persona_style = ""
        if persona_path.exists():
            persona_content = persona_path.read_text(encoding="utf-8")
            start_idx = persona_content.find("### ✨ Ton & Style")
            end_idx = persona_content.find("### 📝 Structure de posts")
            if start_idx != -1 and end_idx != -1:
                persona_style = (
                    "\n\nSTYLE 'PARISIEN GENZ' A APPLIQUER:\n" + persona_content[start_idx:end_idx]
                )

        system_prompt = """Tu es {
            agent.consultant_info['name']}, {
            agent.consultant_info['title']} chez {
            agent.consultant_info['company']}.
Base : Paris | Génération Z assumée
{persona_style}

Tu corriges des commentaires LinkedIn selon le feedback de l'utilisateur."""

        # Regenerer commentaire court
        revised_short = agent.llm_client.generate(
            prompt="""Voici un commentaire court genere precedemment pour ce post:

POST ORIGINAL:
{post_preview}

COMMENTAIRE COURT PRECEDENT:
{previous_short}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce commentaire court (50-150 caracteres) en integrant les corrections demandees.
Reste specifique au post. Apporte de la valeur. Ton Parisien GenZ.""",
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=200,
        )

        # Regenerer commentaire moyen
        revised_medium = agent.llm_client.generate(
            prompt="""Voici un commentaire moyen genere precedemment pour ce post:

POST ORIGINAL:
{post_preview}

COMMENTAIRE MOYEN PRECEDENT:
{previous_medium}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce commentaire moyen (150-300 caracteres) en integrant les corrections demandees.
Reste specifique au post. Apporte de la valeur. Ton Parisien GenZ.""",
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=400,
        )

        # Regenerer commentaire long
        revised_long = agent.llm_client.generate(
            prompt="""Voici un commentaire long genere precedemment pour ce post:

POST ORIGINAL:
{post_preview}

COMMENTAIRE LONG PRECEDENT:
{previous_long}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reecris ce commentaire long (300-500 caracteres) en integrant les corrections demandees.
Reste specifique au post. Apporte de la valeur. Ton Parisien GenZ.""",
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=600,
        )

        job["steps"].append({"step": "generate", "status": "done", "progress": 90})

        # Sauvegarde
        job["steps"].append({"step": "save", "status": "active", "progress": 95})
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"linkedin_comment_{timestamp}_revised.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Commentaires LinkedIn (Révisés)\n\n")
            f.write(f"**Style:** {style}\n")
            f.write(f"**Date:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
            f.write("## Post original (extrait)\n\n")
            f.write(f"> {post_preview}\n\n")
            f.write("---\n\n")
            f.write("## Commentaire Court\n\n")
            f.write(revised_short)
            f.write("\n\n---\n\n## Commentaire Moyen\n\n")
            f.write(revised_medium)
            f.write("\n\n---\n\n## Commentaire Long\n\n")
            f.write(revised_long)
            f.write("\n")

        job["steps"].append({"step": "save", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "short": revised_short,
            "medium": revised_medium,
            "long": revised_long,
            "post_preview": post_preview,
            "style": style,
            "md_path": str(md_path.relative_to(BASE_DIR)),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: TECH WATCH ===


@app.get("/techwatch")
async def techwatch_page():
    """Redirection vers la nouvelle page veille"""
    from starlette.responses import RedirectResponse

    return RedirectResponse(url="/veille", status_code=301)


@app.post("/api/techwatch/generate")
@limiter.limit("3/minute")
async def api_techwatch_generate(request: Request):
    """Lance la génération d'un digest de veille tech"""
    body = await request.json()
    keywords_str = body.get("keywords", "").strip()
    keywords = [k.strip() for k in keywords_str.split(",")] if keywords_str else None
    days = int(body.get("days", 7))
    period_type = body.get("period_type", "weekly")

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "techwatch",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_tech_monitor, args=(job_id, keywords, days, period_type), daemon=True
    )
    thread.start()

    return {"job_id": job_id}


def _run_tech_monitor(job_id: str, keywords: List[str], days: int, period_type: str):
    """Execute la génération de digest en background"""
    job = jobs[job_id]

    try:
        agent = TechMonitorAgent()

        # Step 1: Collecte
        job["steps"].append({"step": "collect", "status": "active", "progress": 10})
        articles = agent.collect_articles(keywords=keywords, days=days)
        job["steps"].append({"step": "collect", "status": "done", "progress": 40})

        if not articles:
            job["status"] = "done"
            job["result"] = {
                "digest": "# ⚠️ Aucun article trouvé\n\nAucun article correspondant aux critères n'a été trouvé pour cette période.",
                "num_articles": 0,
                "period": period_type,
                "md_path": None,
            }
            job["steps"].append({"step": "analyze", "status": "done", "progress": 60})
            job["steps"].append({"step": "generate", "status": "done", "progress": 100})
            return

        # Step 2: Analyse
        job["steps"].append({"step": "analyze", "status": "active", "progress": 45})
        trends = agent.analyze_trends(articles)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 60})

        # Step 3: Generation
        job["steps"].append({"step": "generate", "status": "active", "progress": 65})
        result = agent.generate_digest(articles, trends, period=period_type)
        job["steps"].append({"step": "generate", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"tech_digest_{period_type}_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(f"# Veille Technologique {period_type.capitalize()}\n\n")
            f.write(f"**Généré le:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n")
            f.write(f"**Période:** {days} derniers jours\n")
            f.write(f"**Articles:** {len(articles)}\n\n")
            f.write("---\n\n")
            f.write(result["content"])

        job["status"] = "done"
        job["result"] = {
            "digest": result["content"],
            "num_articles": len(articles),
            "period": period_type,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/techwatch/stream/{job_id}")
async def api_techwatch_stream(job_id: str):
    """Streame le progres de la génération de digest"""

    async def event_generator():
        while True:
            if job_id not in jobs:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break

            job = jobs[job_id]
            yield f"data: {json.dumps(job)}\n\n"

            if job["status"] in ["done", "error"]:
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: TECH WATCH FEEDBACK ===


@app.post("/api/techwatch/regenerate")
async def api_techwatch_regenerate(request: Request):
    """Regenere le digest avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_digest = body.get("previous_digest", "")
    num_articles = body.get("num_articles", 0)
    period = body.get("period", "weekly")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "techwatch",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_techwatch_feedback,
        args=(job_id, previous_digest, feedback, num_articles, period),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_techwatch_feedback(
    job_id: str, previous_digest: str, feedback: str, num_articles: int, period: str
):
    """Regenere le digest en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = TechMonitorAgent()

        job["steps"].append({"step": "collect", "status": "done", "progress": 20})
        job["steps"].append({"step": "analyze", "status": "done", "progress": 40})
        job["steps"].append({"step": "generate", "status": "active", "progress": 50})

        system_prompt = """Tu es {
            agent.consultant_info['name']}, {
            agent.consultant_info['title']} chez {
            agent.consultant_info['company']}.
Tu corriges un digest de veille technologique selon le retour du consultant."""

        # Regenerer le digest
        revised_digest = agent.llm_client.generate(
            prompt="""Voici un digest de veille technologique généré précédemment:

{previous_digest}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reécris ce digest en intégrant les corrections demandées.
Conserve le format Markdown avec emojis et la structure (Tendances, Articles clés, Insights, Sources).
Ne modifie que ce qui est demandé dans le feedback.""",
            system_prompt=system_prompt,
            temperature=0.6,
            max_tokens=3000,
        )

        job["steps"].append({"step": "generate", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"tech_digest_{period}_{timestamp}_revised.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(f"# Veille Technologique {period.capitalize()} (Révisé)\n\n")
            f.write(f"**Généré le:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
            f.write("---\n\n")
            f.write(revised_digest)

        job["status"] = "done"
        job["result"] = {
            "digest": revised_digest,
            "num_articles": num_articles,
            "period": period,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: DATASET ANALYZER ===


@app.get("/dataset", response_class=HTMLResponse)
async def dataset_page(request: Request):
    return templates.TemplateResponse(
        "dataset.html",
        {
            "request": request,
            "active": "dataset",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/dataset/analyze")
@limiter.limit("5/minute")
async def api_dataset_analyze(request: Request, file: UploadFile = File(...)):
    """Lance l'analyse d'un dataset"""

    # Valider le fichier (taille, type)
    allowed_exts = {".csv", ".xlsx", ".xls"}
    content = await validate_file_upload(file, allowed_extensions=allowed_exts)
    filename = sanitize_filename(file.filename or "dataset.csv")

    # Sauvegarder temporairement le fichier
    temp_dir = BASE_DIR / "temp"
    temp_dir.mkdir(exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    temp_path = temp_dir / f"upload_{timestamp}_{filename}"

    with open(temp_path, "wb") as f:
        f.write(content)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "dataset",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "temp_file": str(temp_path),
    }

    thread = threading.Thread(
        target=_run_dataset_analyzer, args=(job_id, str(temp_path), filename), daemon=True
    )
    thread.start()

    return {"job_id": job_id}


def _run_dataset_analyzer(job_id: str, file_path: str, filename: str):
    """Execute l'analyse de dataset en background"""
    job = jobs[job_id]

    try:
        agent = DatasetAnalyzerAgent()

        # Step 1: Load
        job["steps"].append({"step": "load", "status": "active", "progress": 10})
        df = agent.load_dataset(file_path)
        job["steps"].append({"step": "load", "status": "done", "progress": 25})

        # Step 2: Structure
        job["steps"].append({"step": "structure", "status": "active", "progress": 30})
        structure = agent.analyze_structure(df)
        job["steps"].append({"step": "structure", "status": "done", "progress": 45})

        # Step 3: Quality
        job["steps"].append({"step": "quality", "status": "active", "progress": 50})
        quality = agent.analyze_quality(df)
        job["steps"].append({"step": "quality", "status": "done", "progress": 65})

        # Step 4: Stats
        job["steps"].append({"step": "stats", "status": "active", "progress": 70})
        stats = agent.analyze_statistics(df, structure)
        job["steps"].append({"step": "stats", "status": "done", "progress": 85})

        # Step 5: Report
        job["steps"].append({"step": "report", "status": "active", "progress": 90})
        result = agent.generate_report(df, structure, quality, stats, filename)
        job["steps"].append({"step": "report", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"dataset_analysis_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(result["report"])

        job["status"] = "done"
        job["result"] = {
            "report": result["report"],
            "filename": filename,
            "num_rows": structure["num_rows"],
            "num_columns": structure["num_columns"],
            "memory_mb": round(structure["memory_usage"], 2),
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

        # Nettoyer le fichier temporaire
        try:
            os.remove(file_path)
        except BaseException:
            pass

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)
        # Nettoyer le fichier temporaire en cas d'erreur
        try:
            if "temp_file" in job:
                os.remove(job["temp_file"])
        except BaseException:
            pass


@app.get("/api/dataset/stream/{job_id}")
async def api_dataset_stream(job_id: str):
    """Streame le progres de l'analyse"""

    async def event_generator():
        while True:
            if job_id not in jobs:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break

            job = jobs[job_id]
            yield f"data: {json.dumps(job)}\n\n"

            if job["status"] in ["done", "error"]:
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# === API: DATASET FEEDBACK ===


@app.post("/api/dataset/regenerate")
async def api_dataset_regenerate(request: Request):
    """Regenere le rapport avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_report = body.get("previous_report", "")
    filename = body.get("filename", "dataset")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "dataset",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_dataset_feedback,
        args=(job_id, previous_report, feedback, filename),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_dataset_feedback(job_id: str, previous_report: str, feedback: str, filename: str):
    """Regenere le rapport en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        agent = DatasetAnalyzerAgent()

        # Skip aux étapes finales
        job["steps"].append({"step": "load", "status": "done", "progress": 20})
        job["steps"].append({"step": "structure", "status": "done", "progress": 40})
        job["steps"].append({"step": "quality", "status": "done", "progress": 60})
        job["steps"].append({"step": "stats", "status": "done", "progress": 70})
        job["steps"].append({"step": "report", "status": "active", "progress": 80})

        system_prompt = """Tu es {
            agent.consultant_info['name']}, {
            agent.consultant_info['title']} chez {
            agent.consultant_info['company']}.
Tu corriges un rapport d'analyse de dataset selon le retour du consultant."""

        # Regenerer le rapport
        revised_report = agent.llm_client.generate(
            prompt="""Voici un rapport d'analyse de dataset généré précédemment:

{previous_report}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reécris ce rapport en intégrant les corrections demandées.
Conserve le format Markdown avec emojis et la structure (Vue d'ensemble, Structure, Qualité, Insights, Recommandations).
Ne modifie que ce qui est demandé dans le feedback.""",
            system_prompt=system_prompt,
            temperature=0.5,
            max_tokens=2000,
        )

        job["steps"].append({"step": "report", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"dataset_analysis_{timestamp}_revised.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(revised_report)

        job["status"] = "done"
        job["result"] = {
            "report": revised_report,
            "filename": filename,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: WORKSHOP PLANNER ===


@app.get("/workshop", response_class=HTMLResponse)
async def workshop_page(request: Request):
    return templates.TemplateResponse(
        "workshop.html",
        {
            "request": request,
            "active": "workshop",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/workshop/generate")
@limiter.limit("5/minute")
async def api_workshop_generate(request: Request):
    """Lance la génération d'un plan de workshop"""
    body = await request.json()
    topic = body.get("topic", "").strip()
    duration = body.get("duration", "full_day")
    audience = body.get("audience", "Professionnels").strip()
    objectives = body.get("objectives", "").strip()

    if not topic:
        return JSONResponse({"error": "Le sujet est obligatoire."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "workshop",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_workshop_planner,
        args=(job_id, topic, duration, audience, objectives),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_workshop_planner(job_id: str, topic: str, duration: str, audience: str, objectives: str):
    """Execute la génération de plan en background"""
    job = jobs[job_id]

    try:
        agent = WorkshopPlannerAgent()

        job["steps"].append({"step": "plan", "status": "active", "progress": 30})
        result = agent.generate_workshop_plan(topic, duration, audience, objectives)
        job["steps"].append({"step": "plan", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        safe_topic = "".join(c for c in topic if c.isalnum() or c in (" ", "-", "_")).strip()[:50]
        safe_topic = safe_topic.replace(" ", "_")

        md_path = output_dir / f"workshop_{safe_topic}_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(result["plan"])

        job["status"] = "done"
        job["result"] = {
            "plan": result["plan"],
            "topic": topic,
            "duration": result["duration"],
            "audience": audience,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/workshop/stream/{job_id}")
async def api_workshop_stream(job_id: str):
    """Streame le progres"""

    async def event_generator():
        while True:
            if job_id not in jobs:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break

            job = jobs[job_id]
            yield f"data: {json.dumps(job)}\n\n"

            if job["status"] in ["done", "error"]:
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/workshop/regenerate")
async def api_workshop_regenerate(request: Request):
    """Regenere avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_plan = body.get("previous_plan", "")
    topic = body.get("topic", "Workshop")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "workshop",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_workshop_feedback,
        args=(job_id, previous_plan, feedback, topic),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_workshop_feedback(job_id: str, previous_plan: str, feedback: str, topic: str):
    """Regenere avec feedback"""
    job = jobs[job_id]

    try:
        agent = WorkshopPlannerAgent()

        job["steps"].append({"step": "plan", "status": "active", "progress": 50})

        system_prompt = """Tu es {
            agent.consultant_info['name']}, {
            agent.consultant_info['title']} chez {
            agent.consultant_info['company']}.
Tu corriges un plan de formation selon le retour du consultant."""

        revised_plan = agent.llm_client.generate(
            prompt="""Voici un plan de workshop généré précédemment:

{previous_plan}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reécris ce plan en intégrant les corrections demandées.
Conserve le format Markdown et la structure.
Ne modifie que ce qui est demandé dans le feedback.""",
            system_prompt=system_prompt,
            temperature=0.6,
            max_tokens=3500,
        )

        job["steps"].append({"step": "plan", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        safe_topic = "".join(c for c in topic if c.isalnum() or c in (" ", "-", "_")).strip()[:50]
        safe_topic = safe_topic.replace(" ", "_")

        md_path = output_dir / f"workshop_{safe_topic}_{timestamp}_revised.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(revised_plan)

        job["status"] = "done"
        job["result"] = {
            "plan": revised_plan,
            "topic": topic,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: RFP RESPONDER ===


@app.get("/rfp", response_class=HTMLResponse)
async def rfp_page(request: Request):
    return templates.TemplateResponse(
        "rfp.html",
        {
            "request": request,
            "active": "rfp",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/rfp/generate")
@limiter.limit("3/minute")
async def api_rfp_generate(
    request: Request,
    rfp_text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
):
    """Lance la génération de réponse au RFP"""
    text = ""

    if file:
        content = await file.read()
        filename = file.filename or ""
        if filename.lower().endswith(".pdf"):
            pdf_reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        else:
            text = content.decode("utf-8")
    elif rfp_text:
        text = rfp_text
    else:
        return JSONResponse({"error": "Aucun RFP fourni."}, status_code=400)

    if len(text.strip()) < 100:
        return JSONResponse({"error": "Le RFP semble trop court."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "rfp",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(target=_run_rfp_responder, args=(job_id, text), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_rfp_responder(job_id: str, rfp_text: str):
    """Execute la génération en background"""
    job = jobs[job_id]

    try:
        agent = RFPResponderAgent()

        # Step 1: Analyze
        job["steps"].append({"step": "analyze", "status": "active", "progress": 20})
        analysis = agent.analyze_rfp(rfp_text)
        job["steps"].append({"step": "analyze", "status": "done", "progress": 40})

        # Step 2: Generate response
        job["steps"].append({"step": "response", "status": "active", "progress": 50})
        result = agent.generate_response(rfp_text, analysis)
        job["steps"].append({"step": "response", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"rfp_response_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# Réponse à l'appel d'offres\n\n")
            f.write(f"**Généré le:** {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
            f.write("---\n\n")
            f.write(result["response"])

        job["status"] = "done"
        job["result"] = {
            "response": result["response"],
            "analysis": analysis,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/rfp/stream/{job_id}")
async def api_rfp_stream(job_id: str):
    """Streame le progres"""

    async def event_generator():
        while True:
            if job_id not in jobs:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break

            job = jobs[job_id]
            yield f"data: {json.dumps(job)}\n\n"

            if job["status"] in ["done", "error"]:
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/rfp/regenerate")
async def api_rfp_regenerate(request: Request):
    """Regenere avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_response = body.get("previous_response", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "rfp",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_rfp_feedback,
        args=(job_id, previous_response, feedback),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_rfp_feedback(job_id: str, previous_response: str, feedback: str):
    """Regenere avec feedback"""
    job = jobs[job_id]

    try:
        agent = RFPResponderAgent()

        job["steps"].append({"step": "analyze", "status": "done", "progress": 30})
        job["steps"].append({"step": "response", "status": "active", "progress": 50})

        system_prompt = """Tu es {
            agent.consultant_info['name']}, {
            agent.consultant_info['title']} chez {
            agent.consultant_info['company']}.
Tu corriges une réponse à un appel d'offres selon le retour du consultant."""

        revised_response = agent.llm_client.generate(
            prompt="""Voici une réponse à un RFP générée précédemment:

{previous_response}

FEEDBACK DE L'UTILISATEUR:
{feedback}

Reécris cette réponse en intégrant les corrections demandées.
Conserve le format Markdown et la structure en 10 sections.
Ne modifie que ce qui est demandé dans le feedback.""",
            system_prompt=system_prompt,
            temperature=0.5,
            max_tokens=4000,
        )

        job["steps"].append({"step": "response", "status": "done", "progress": 100})

        # Sauvegarde
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"rfp_response_{timestamp}_revised.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(revised_response)

        job["status"] = "done"
        job["result"] = {
            "response": revised_response,
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "generated_at": datetime.now().isoformat(),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


# === PAGE & API: ARTICLE GENERATOR ===


@app.get("/article-generator", response_class=HTMLResponse)
async def article_generator_page(request: Request):
    return templates.TemplateResponse(
        "article-generator.html",
        {
            "request": request,
            "active": "article-generator",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/article-generator/generate")
@limiter.limit("5/minute")
async def api_article_generator_generate(
    request: Request,
    idea_text: str = Form(...),
    use_context: bool = Form(False),
):
    """Lance la génération d'un article de blog"""
    if len(idea_text.strip()) < 20:
        return JSONResponse(
            {"error": "L'idée est trop courte (minimum 20 caractères)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "article-generator",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_article_generator,
        args=(job_id, idea_text, use_context),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_article_generator(job_id: str, idea: str, use_context: bool = False):
    """Execute la generation d'article en background (4 etapes)"""
    job = jobs[job_id]

    try:
        from agents.article_generator import ArticleGeneratorAgent

        agent = ArticleGeneratorAgent()

        job["steps"].append({"step": "article", "status": "active", "progress": 10})

        # Pipeline complet (article + linkedin + image + sources)
        result = agent.run(idea, target_length="medium", use_context=use_context)

        job["steps"].append({"step": "article", "status": "done", "progress": 100})

        word_count = len(result["article"].split())

        job["status"] = "done"
        job["result"] = {
            "content": result["article"],
            "word_count": word_count,
            "md_path": result["article_path"],
            "illustration_prompt": result.get("illustration_prompt", ""),
            "linkedin_post": result.get("linkedin_post", ""),
            "image_path": result.get("image_path"),
            "sources": result.get("sources", []),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/article-generator/stream/{job_id}")
async def api_article_generator_stream(job_id: str):
    """SSE stream pour la progression de la génération d'article"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouvé"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/article-generator/regenerate")
async def api_article_generator_regenerate(request: Request):
    """Régénère l'article avec le feedback utilisateur"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_content = body.get("previous_content", "")
    previous_idea = body.get("previous_idea", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "article-generator",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_article_generator_feedback,
        args=(job_id, previous_content, feedback, previous_idea),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_article_generator_feedback(
    job_id: str, previous_content: str, feedback: str, previous_idea: str
):
    """Régénère l'article en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        from agents.article_generator import ArticleGeneratorAgent

        agent = ArticleGeneratorAgent()

        job["steps"].append({"step": "generate", "status": "active", "progress": 20})

        revised_prompt = """{previous_idea}

CONTEXTE : Voici l'article précédent que tu as généré:

{previous_content}

FEEDBACK de l'utilisateur:
{feedback}

Régénère l'article en tenant compte du feedback. Applique les modifications demandées tout en gardant le style et la qualité Consulting Tools."""

        revised_content = agent.generate_article(revised_prompt, target_length="medium")

        # Générer le nouveau prompt d'illustration
        illustration_prompt = agent.generate_illustration_prompt(revised_content)

        job["steps"].append({"step": "generate", "status": "done", "progress": 100})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        # Créer un slug pour le nom de fichier
        import re

        slug = re.sub(r"[^\w\s-]", "", previous_idea.lower())
        slug = re.sub(r"[-\s]+", "-", slug)[:50]

        md_path = output_dir / f"article_{slug}_revised_{timestamp}.md"

        # Ajouter les métadonnées
        full_article = """---
title: {previous_idea}
author: {CONSULTANT_NAME}
company: {COMPANY_NAME}
date: {datetime.now().strftime('%Y-%m-%d')}
revised: true
illustration_prompt: |
  {illustration_prompt}
---

{revised_content}
"""

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(full_article)

        job["status"] = "done"
        job["result"] = {
            "content": revised_content,
            "word_count": len(revised_content.split()),
            "md_path": str(md_path.relative_to(BASE_DIR)),
            "illustration_prompt": illustration_prompt,
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/article-generator/export-gdocs")
@limiter.limit("5/minute")
async def api_article_generator_export_gdocs(
    request: Request,
    content: str = Form(...),
    title: str = Form("Article de Blog"),
):
    """Exporte un article genere vers Google Docs avec formatage markdown"""
    try:
        from utils.google_api import GoogleAPIClient

        try:
            google_client = GoogleAPIClient()
            doc_url = google_client.export_markdown_to_docs(content, title)

            if doc_url:
                return JSONResponse(
                    {"doc_url": doc_url, "message": "Article exporte vers Google Docs"}
                )
            else:
                return JSONResponse(
                    {"error": "Echec de la creation du document Google Docs"},
                    status_code=500,
                )

        except Exception as e:
            if "credentials" in str(e).lower() or "token" in str(e).lower():
                return JSONResponse(
                    {
                        "error": "Google API non configuree. Configurez vos credentials.",
                        "setup_required": True,
                    },
                    status_code=400,
                )
            else:
                raise

    except Exception as e:
        print(f"Erreur export article GDocs: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# ====================================
# FORMATION GENERATOR
# ====================================


@app.get("/formation", response_class=HTMLResponse)
async def formation_page(request: Request):
    return templates.TemplateResponse(
        "formation.html",
        {
            "request": request,
            "active": "formation",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/formation/generate")
@limiter.limit("5/minute")
async def api_formation_generate(
    request: Request,
    client_needs: str = Form(...),
):
    """Lance la génération d'un programme de formation"""
    if len(client_needs.strip()) < 20:
        return JSONResponse(
            {"error": "Le besoin est trop court (minimum 20 caractères)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "formation",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_formation_generator,
        args=(job_id, client_needs),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_formation_generator(job_id: str, client_needs: str):
    """Génère le programme de formation en background"""
    job = jobs[job_id]

    try:
        from agents.formation_generator import FormationGeneratorAgent

        agent = FormationGeneratorAgent()

        job["steps"].append({"step": "analyze", "status": "active", "progress": 10})
        job["steps"].append({"step": "analyze", "status": "done", "progress": 20})
        job["steps"].append({"step": "generate", "status": "active", "progress": 30})

        result = agent.generate_programme(client_needs)

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})
        job["steps"].append({"step": "format", "status": "active", "progress": 85})

        # Sauvegarder le markdown
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        title = result["metadata"].get("title", "Programme_Formation")
        safe_title = title.replace(" ", "_").replace("/", "-")[:50]
        md_path = output_dir / f"formation_{safe_title}_{timestamp}.md"

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(result["markdown"])

        job["steps"].append({"step": "format", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "content": result["markdown"],
            "metadata": result["metadata"],
            "md_path": str(md_path.relative_to(BASE_DIR)),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/formation/stream/{job_id}")
async def api_formation_stream(job_id: str):
    """SSE stream pour la progression de la génération de formation"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouvé"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/formation/regenerate")
async def api_formation_regenerate(request: Request):
    """Régénère le programme avec feedback"""
    body = await request.json()
    feedback = body.get("feedback", "").strip()
    previous_content = body.get("previous_content", "")

    if not feedback:
        return JSONResponse({"error": "Aucun feedback fourni."}, status_code=400)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "formation",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_formation_regenerate,
        args=(job_id, previous_content, feedback),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_formation_regenerate(job_id: str, previous_content: str, feedback: str):
    """Régénère le programme en tenant compte du feedback"""
    job = jobs[job_id]

    try:
        from agents.formation_generator import FormationGeneratorAgent

        agent = FormationGeneratorAgent()

        job["steps"].append({"step": "generate", "status": "active", "progress": 20})

        result = agent.regenerate_with_feedback(previous_content, feedback)

        job["steps"].append({"step": "generate", "status": "done", "progress": 80})
        job["steps"].append({"step": "format", "status": "active", "progress": 85})

        # Sauvegarder
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = BASE_DIR / "output"
        output_dir.mkdir(exist_ok=True)

        md_path = output_dir / f"formation_revised_{timestamp}.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(result["markdown"])

        job["steps"].append({"step": "format", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "content": result["markdown"],
            "metadata": result["metadata"],
            "md_path": str(md_path.relative_to(BASE_DIR)),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/formation/export-gdocs")
@limiter.limit("5/minute")
async def api_formation_export_gdocs(
    request: Request,
    content: str = Form(...),
    title: str = Form("Programme de Formation"),
):
    """Exporte le programme de formation vers Google Docs"""
    try:
        from utils.google_api import GoogleAPIClient

        try:
            google_client = GoogleAPIClient()
            doc_url = google_client.export_markdown_to_docs(content, title)

            if doc_url:
                return JSONResponse(
                    {"doc_url": doc_url, "message": "Document Google Docs créé avec succès"}
                )
            else:
                return JSONResponse(
                    {"error": "Échec de la création du document Google Docs"}, status_code=500
                )

        except Exception as e:
            if "credentials" in str(e).lower() or "token" in str(e).lower():
                return JSONResponse(
                    {
                        "error": "Google API non configurée. Configurez vos credentials dans config/google_credentials.json",
                        "setup_required": True,
                    },
                    status_code=400,
                )
            else:
                raise

    except Exception as e:
        print(f"Error in export to Google Docs: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# ====================================
# TRAINING SLIDES GENERATOR
# ====================================


@app.get("/training-slides", response_class=HTMLResponse)
async def training_slides_page(request: Request):
    return templates.TemplateResponse(
        "training-slides.html",
        {
            "request": request,
            "active": "training-slides",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/training-slides/generate")
@limiter.limit("3/minute")
async def api_training_slides_generate(
    request: Request,
    programme_text: str = Form(None),
    file: Optional[UploadFile] = File(None),
    files: List[UploadFile] = File(default=[]),
):
    """
    Lance la génération des slides de formation.

    Entrées acceptées (par priorité) :
    - files : plusieurs fichiers (TXT, MD, DOCX, PDF) — leur contenu est concaténé
    - file  : un seul fichier (compatibilité ascendante)
    - programme_text : texte brut du programme
    """
    from utils.document_parser import document_parser

    # Fusionner file (ancien) et files (nouveau) en une seule liste
    all_files: List[UploadFile] = []
    if files:
        all_files.extend(files)
    if file:
        all_files.append(file)

    text_parts: List[str] = []

    for upload in all_files:
        temp_file = BASE_DIR / "output" / f"temp_{upload.filename}"
        temp_file.parent.mkdir(exist_ok=True)
        try:
            content = await upload.read()
            with open(temp_file, "wb") as fp:
                fp.write(content)
            extracted = document_parser.parse_file(str(temp_file))
            temp_file.unlink()
            if extracted:
                text_parts.append(extracted)
            else:
                return JSONResponse(
                    {
                        "error": (
                            f"Impossible d'extraire le texte du fichier {upload.filename}. "
                            "Formats supportés : TXT, MD, DOCX, PDF"
                        )
                    },
                    status_code=400,
                )
        except Exception as e:
            if temp_file.exists():
                temp_file.unlink()
            return JSONResponse(
                {"error": f"Erreur lecture fichier '{upload.filename}' : {safe_error_message(e)}"},
                status_code=400,
            )

    if text_parts:
        # Concaténer les textes avec un séparateur clair
        text = "\n\n---\n\n".join(text_parts)
    elif programme_text:
        text = programme_text.strip()
    else:
        text = ""

    if not text or len(text) < 50:
        return JSONResponse(
            {"error": "Le programme est trop court (minimum 50 caractères)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "training-slides",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_training_slides_generator,
        args=(job_id, text),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_training_slides_generator(job_id: str, programme_text: str):
    """Génère les slides de formation en background"""
    job = jobs[job_id]

    try:
        from agents.training_slides_generator import TrainingSlidesGeneratorAgent

        agent = TrainingSlidesGeneratorAgent()

        job["steps"].append({"step": "parse", "status": "active", "progress": 10})

        result = agent.generate_all_slides(programme_text)

        job["steps"].append({"step": "parse", "status": "done", "progress": 30})
        job["steps"].append({"step": "generate", "status": "done", "progress": 80})
        job["steps"].append({"step": "pptx", "status": "active", "progress": 85})

        # Générer les PPTX par module
        pptx_paths = {}
        for module_name, module_slides in result["modules_slides"].items():
            if module_slides:
                pptx_path = agent.generate_module_pptx(module_slides, module_name)
                pptx_paths[module_name] = pptx_path

        # Générer le PPTX complet
        all_pptx_path = agent.generate_module_pptx(result["all_slides"], "Complet")

        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "programme_data": result["programme_data"],
            "modules_slides": {k: v for k, v in result["modules_slides"].items()},
            "all_slides": result["all_slides"],
            "total_slides": result["total_slides"],
            "pptx_paths": pptx_paths,
            "all_pptx_path": all_pptx_path,
        }

    except Exception as e:
        print(f"Error in training slides: {safe_traceback()}")
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/training-slides/stream/{job_id}")
async def api_training_slides_stream(job_id: str):
    """SSE stream pour la progression des slides de formation"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouvé"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/training-slides/export-slides")
@limiter.limit("3/minute")
async def api_training_slides_export(
    request: Request,
    slides_data: str = Form(...),
    title: str = Form("Support de Formation"),
):
    """Exporte les slides de formation vers Google Slides"""
    try:
        import re

        # Sanitize pour éviter les erreurs JSON (supprime les caractères de
        # contrôle)
        sanitized = slides_data.strip()
        sanitized = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", sanitized)

        slides = json.loads(sanitized)

        from utils.google_api import GoogleAPIClient

        try:
            google_client = GoogleAPIClient()
            presentation_id = google_client.export_pptx_to_slides(slides_data=slides, title=title)

            if presentation_id:
                presentation_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"
                return JSONResponse(
                    {
                        "presentation_id": presentation_id,
                        "presentation_url": presentation_url,
                        "message": "Présentation Google Slides créée avec succès",
                    }
                )
            else:
                return JSONResponse(
                    {"error": "Échec de la création de la présentation"}, status_code=500
                )

        except Exception as e:
            if "credentials" in str(e).lower() or "token" in str(e).lower():
                return JSONResponse(
                    {"error": "Google API non configurée.", "setup_required": True}, status_code=400
                )
            else:
                raise

    except Exception as e:
        print(f"Error in training slides export: {safe_traceback()}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# === API: DOWNLOAD ===


@app.get("/api/download")
async def api_download(path: str):
    """Telecharge un fichier genere"""
    try:
        # Nettoyer le path pour éviter les problèmes d'encodage
        path = path.strip()
        file_path = BASE_DIR / path

        # Vérifier que le fichier existe
        if not file_path.exists():
            print(f"Fichier non trouve: {file_path}")
            return JSONResponse({"error": "Fichier non trouve"}, status_code=404)

        # Vérifier que le fichier est bien dans le répertoire BASE_DIR
        # (sécurité)
        if not file_path.is_relative_to(BASE_DIR):
            print(f"Tentative d'accès à un fichier hors du répertoire autorisé: {file_path}")
            return JSONResponse({"error": "Accès non autorisé"}, status_code=403)

        # Déterminer le type MIME en fonction de l'extension
        media_type = None
        if file_path.suffix == ".pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_path.suffix == ".pdf":
            media_type = "application/pdf"
        elif file_path.suffix == ".md":
            media_type = "text/markdown"
        elif file_path.suffix == ".json":
            media_type = "application/json"
        elif file_path.suffix == ".png":
            media_type = "image/png"
        elif file_path.suffix == ".jpg" or file_path.suffix == ".jpeg":
            media_type = "image/jpeg"

        # Retourner le fichier avec les bons headers
        # Pour les images (PNG, JPG), utiliser inline pour affichage dans le
        # navigateur
        disposition = "inline" if file_path.suffix in [".png", ".jpg", ".jpeg"] else "attachment"

        return FileResponse(
            file_path,
            filename=file_path.name,
            media_type=media_type,
            headers={
                "Content-Disposition": f'{disposition}; filename="{file_path.name}"',
                "Cache-Control": "no-cache",
            },
        )

    except Exception as e:
        print(f"Erreur lors du téléchargement: {safe_error_message(e)}")
        return JSONResponse(
            {"error": f"Erreur lors du téléchargement: {safe_error_message(e)}"}, status_code=500
        )


@app.post("/api/convert-to-pdf")
async def api_convert_to_pdf(request: Request):
    """
    Convertit un fichier PPTX ou Markdown en PDF

    Body JSON:
        - file_path: Chemin relatif du fichier à convertir (PPTX ou MD)
        - file_type: Type de fichier ('pptx' ou 'markdown')
    """
    try:
        from utils.pdf_converter import pdf_converter

        body = await request.json()
        file_path = body.get("file_path", "").strip()
        file_type = body.get("file_type", "pptx")

        if not file_path:
            return JSONResponse({"error": "Chemin de fichier manquant"}, status_code=400)

        # Chemin absolu
        full_path = BASE_DIR / file_path

        if not full_path.exists():
            return JSONResponse({"error": "Fichier non trouvé"}, status_code=404)

        # Vérifier la sécurité
        if not full_path.is_relative_to(BASE_DIR):
            return JSONResponse({"error": "Accès non autorisé"}, status_code=403)

        # Convertir selon le type
        pdf_path = None
        if file_type == "pptx" and full_path.suffix == ".pptx":
            pdf_path = pdf_converter.pptx_to_pdf(str(full_path))
        elif file_type == "markdown" and full_path.suffix in [".md", ".markdown"]:
            pdf_path = pdf_converter.markdown_to_pdf(str(full_path))
        else:
            return JSONResponse({"error": "Type de fichier non supporté"}, status_code=400)

        if pdf_path:
            # Retourner le chemin relatif du PDF
            pdf_rel_path = str(Path(pdf_path).relative_to(BASE_DIR))
            return {"success": True, "pdf_path": pdf_rel_path, "message": "Conversion réussie"}
        else:
            return JSONResponse(
                {
                    "error": "Conversion PDF échouée",
                    "message": "Vérifiez que LibreOffice est installé (pour PPTX) ou pandoc/weasyprint (pour Markdown)",
                },
                status_code=500,
            )

    except Exception as e:
        print(f"Erreur lors de la conversion PDF: {safe_error_message(e)}")
        return JSONResponse({"error": f"Erreur: {safe_error_message(e)}"}, status_code=500)


@app.get("/api/pdf-capabilities")
async def api_pdf_capabilities():
    """Retourne les capacités de conversion PDF disponibles"""
    try:
        from utils.pdf_converter import pdf_converter

        capabilities = pdf_converter.is_pdf_conversion_available()
        return {
            "capabilities": capabilities,
            "message": "Pour activer la conversion PDF, installez LibreOffice (PPTX→PDF) ou pandoc/weasyprint (MD→PDF)",
        }
    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# ===================================
# DOCUMENT TO PRESENTATION
# ===================================


@app.get("/doc-to-presentation", response_class=HTMLResponse)
async def doc_to_presentation_page(request: Request):
    return templates.TemplateResponse(
        "doc-to-presentation.html",
        {
            "request": request,
            "active": "doc-to-presentation",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/doc-to-presentation/generate")
async def api_doc_to_presentation_generate(
    request: Request,
    target_audience: str = Form(...),
    objective: str = Form(...),
    files: List[UploadFile] = File(...),
):
    """Lance la generation de presentation a partir de documents"""
    if not files:
        return JSONResponse({"error": "Aucun fichier fourni."}, status_code=400)

    # Lire les fichiers uploades
    documents = []
    for f in files:
        content = await f.read()
        documents.append({"filename": f.filename, "content": content})

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "doc-to-presentation",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_doc_to_presentation,
        args=(job_id, documents, target_audience, objective),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_doc_to_presentation(job_id: str, documents: list, target_audience: str, objective: str):
    """Execute la generation de presentation en background"""
    job = jobs[job_id]
    try:
        from agents.doc_to_presentation import DocToPresentationAgent

        agent = DocToPresentationAgent()

        job["steps"].append({"step": "parse", "status": "active", "progress": 10})
        result = agent.run(documents, target_audience, objective)
        job["steps"].append({"step": "pptx", "status": "done", "progress": 100})

        if result.get("error"):
            job["status"] = "error"
            job["error"] = result["error"]
        else:
            job["status"] = "done"
            job["result"] = result

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/doc-to-presentation/stream/{job_id}")
async def api_doc_to_presentation_stream(job_id: str):
    """SSE stream pour la progression"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(event_generator(), media_type="text/event-stream")


# ===================================
# HTML SLIDES PREMIUM (Gemini 3.1 Pro)
# ===================================


@app.get("/html-slides")
async def html_slides_page():
    """Redirige vers le Slide Editor (HTML Premium integre)"""
    from starlette.responses import RedirectResponse

    return RedirectResponse(url="/slide-editor", status_code=301)


@app.post("/api/html-slides/generate")
@limiter.limit("3/minute")
async def api_html_slides_generate(
    request: Request,
    topic: str = Form(...),
    num_slides: int = Form(10),
    language: str = Form("fr"),
):
    """Lance la generation de slides HTML premium via Gemini 3.1 Pro"""
    topic = sanitize_text_input(topic)
    if not topic or len(topic.strip()) < 5:
        return JSONResponse(
            {"error": "Le sujet est trop court (minimum 5 caracteres)."}, status_code=400
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "html-slides",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_html_slides_generator,
        args=(job_id, topic, num_slides, language),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_html_slides_generator(job_id: str, topic: str, num_slides: int, language: str):
    """Generate HTML slides in background"""
    job = jobs[job_id]
    try:
        from agents.html_slide_generator import HtmlSlideGeneratorAgent

        agent = HtmlSlideGeneratorAgent()

        # Step 1: Extract design system
        job["steps"].append({"step": "design", "status": "active", "progress": 10})
        design_system = agent.extract_design_system()
        job["steps"].append({"step": "design", "status": "done", "progress": 30})

        # Step 2: Generate HTML via Gemini 3.1 Pro
        job["steps"].append({"step": "generate", "status": "active", "progress": 35})

        system_instruction, user_prompt = agent.build_prompt(topic, num_slides, design_system)

        lang_prefix = ""
        if language == "en":
            lang_prefix = "Generate the presentation content in English. "
        user_prompt = lang_prefix + user_prompt

        messages = [{"role": "user", "content": user_prompt}]
        raw_html = agent.llm.generate_with_context(
            messages=messages,
            system_prompt=system_instruction,
            temperature=0.2,
            max_tokens=16384,
        )
        job["steps"].append({"step": "generate", "status": "done", "progress": 80})

        # Step 3: Clean and save
        job["steps"].append({"step": "save", "status": "active", "progress": 85})
        html_content = agent.clean_html_response(raw_html)
        html_path = agent.save_html(html_content)
        job["steps"].append({"step": "save", "status": "done", "progress": 100})

        job["status"] = "done"
        job["result"] = {
            "html_path": html_path,
            "html_content": html_content,
            "topic": topic,
            "num_slides": num_slides,
        }

    except Exception as e:
        print(f"Error in HTML slides: {safe_traceback()}")
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.get("/api/html-slides/stream/{job_id}")
async def api_html_slides_stream(job_id: str):
    """SSE stream pour la progression de generation de slides HTML"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


# ===================================
# PROPOSAL CANVA (Mode Gemini)
# ===================================


@app.get("/proposal-canva")
async def proposal_canva_page(request: Request):
    """Page de génération de propositions en mode Canva (conversationnel avec Gemini)"""
    return templates.TemplateResponse(
        "proposal-canva.html",
        {"request": request, "consultant_name": CONSULTANT_NAME, "active": "proposal"},
    )


@app.post("/api/proposal-canva/generate")
async def api_proposal_canva_generate(request: Request):
    """Génère des propositions en mode conversationnel avec Gemini"""
    try:
        data = await request.json()
        user_message = data.get("message", "")
        conversation_history = data.get("conversation_history", [])
        data.get("current_proposal")

        if not user_message:
            return JSONResponse({"error": "Message manquant"}, status_code=400)

        # Initialiser le client Gemini
        from utils.llm_client import LLMClient

        llm = LLMClient(provider="gemini", model="gemini-3-flash-preview")

        # System prompt pour le mode conversationnel
        system_prompt = """Tu es un assistant expert en propositions commerciales pour {COMPANY_NAME}.
Tu aides les consultants à créer des propositions professionnelles de manière conversationnelle.

INSTRUCTIONS:
1. Analyse la demande de l'utilisateur
2. Si c'est la première fois, demande le contexte de l'appel d'offres
3. Génère des slides structurées en JSON selon le besoin
4. Réponds de manière naturelle et conversationnelle
5. Propose des améliorations si nécessaire

FORMAT DE SORTIE:
- Si tu génères des slides, utilise ce format JSON à la fin de ta réponse:
```json
{"slides": [
    {"type": "cover", "title": "...", "subtitle": "..."} ,
    {"type": "content", "title": "...", "bullets": ["...", "..."]} ,
    {"type": "diagram", "title": "...", "diagram_type": "flow", "elements": ["..."]} ,
    {"type": "stat", "stat_value": "67%", "stat_label": "de ROI", "context": "..."} ,
    {"type": "quote", "quote_text": "...", "author": "..."} ,
    {"type": "highlight", "title": "...", "key_points": ["...", "..."]}
  ]
}
```

TYPES DE SLIDES DISPONIBLES:
- cover: Slide de couverture (title, subtitle)
- section: Séparateur de section (title)
- content: Contenu avec bullets (title, bullets)
- diagram: Diagramme visuel (title, diagram_type, elements, description)
- stat: Statistique impactante (stat_value, stat_label, context, subtitle)
- quote: Citation/message clé (quote_text, author, title)
- highlight: Points clés en encadrés (title, key_points, highlight_color)
- closing: Slide de clôture

Consultant: {CONSULTANT_NAME}
Société: {COMPANY_NAME} """

        # Construire les messages pour l'API
        messages = []
        for msg in conversation_history:
            if msg["role"] in ["user", "assistant"]:
                messages.append({"role": msg["role"], "content": msg["content"]})

        # Ajouter le nouveau message
        messages.append({"role": "user", "content": user_message})

        # Générer la réponse
        response_text = llm.generate_with_context(
            messages=messages, system_prompt=system_prompt, temperature=0.7, max_tokens=8000
        )

        # Extraire le JSON s'il existe
        slides = []
        pptx_path = None

        if "```json" in response_text:
            # Extraire le JSON
            json_start = response_text.find("```json") + 7
            json_end = response_text.find("```", json_start)
            json_str = response_text[json_start:json_end].strip()

            try:
                slides_data = json.loads(json_str)
                slides = slides_data.get("slides", [])

                # Générer le PPTX si des slides sont présentes
                if slides:
                    from utils.pptx_generator import build_proposal_pptx

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_dir = BASE_DIR / "output"
                    output_dir.mkdir(exist_ok=True)

                    output_path = output_dir / f"proposal_canva_{timestamp}.pptx"

                    consultant_info = {"name": CONSULTANT_NAME, "company": COMPANY_NAME}

                    pptx_path = build_proposal_pptx(
                        template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                        slides_data=slides,
                        output_path=str(output_path),
                        consultant_info=consultant_info,
                    )

                    pptx_path = str(Path(pptx_path).relative_to(BASE_DIR))

            except json.JSONDecodeError as e:
                print(f"❌ Erreur JSON: {e}")
                # Continue sans slides

        # Nettoyer la réponse (enlever le JSON)
        response_text_clean = response_text.split("```json")[0].strip()

        # Mettre à jour l'historique
        conversation_history.append({"role": "user", "content": user_message})
        conversation_history.append({"role": "assistant", "content": response_text_clean})

        return {
            "response": response_text_clean,
            "slides": slides,
            "pptx_path": pptx_path,
            "conversation_history": conversation_history,
        }

    except Exception as e:
        print(f"❌ Erreur génération Canva: {e}")
        import traceback

        traceback.print_exc()
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/proposal-canva/generate-section")
async def api_proposal_canva_generate_section(request: Request):
    """Génère une section spécifique de proposition (mode modulaire)"""
    try:
        data = await request.json()
        section = data.get("section", "")
        tender_text = data.get("tender_text", "")
        data.get("current_proposal")

        if not section or not tender_text:
            return JSONResponse({"error": "Section et tender_text requis"}, status_code=400)

        # Initialiser le client Gemini
        from utils.llm_client import LLMClient

        llm = LLMClient(provider="gemini", model="gemini-3-flash-preview")

        # Mapping des sections
        section_prompts = {
            "cover": "Génère une slide de couverture professionnelle",
            "agenda": "Génère une slide d'agenda avec la structure de la proposition",
            "context": "Génère 2-3 slides de contexte avec les enjeux et objectifs",
            "approach": "Génère 3-4 slides d'approche méthodologique",
            "planning": "Génère 2-3 slides de planning avec les phases",
            "budget": "Génère 1-2 slides de chiffrage avec estimation budgétaire",
            "references": "Génère 2-3 slides de références projets similaires",
            "cvs": "Génère 2-3 slides de CVs adaptés au projet",
        }

        section_instruction = section_prompts.get(
            section, f"Génère des slides pour la section {section}"
        )

        # System prompt
        system_prompt = """Tu es un expert en propositions commerciales pour {COMPANY_NAME}.

Génère des slides pour la section "{section}" basées sur cet appel d'offres :

{tender_text[:3000]}

INSTRUCTION: {section_instruction}

FORMAT DE SORTIE - JSON uniquement:
```json
{{
  "slides": [
    {{"type": "content", "title": "...", "bullets": ["...", "..."]}},
    {{"type": "diagram", "title": "...", "diagram_type": "flow", "elements": ["..."]}},
    {{"type": "stat", "stat_value": "67%", "stat_label": "...", "context": "..."}},
    {{"type": "quote", "quote_text": "...", "author": "..."}},
    {{"type": "highlight", "title": "...", "key_points": ["...", "..."]}}
  ]
}}
```

TYPES DE SLIDES:
- cover: couverture (title, subtitle, client_name, project_name)
- section: séparateur (title, section_number)
- content: contenu avec bullets (title, bullets)
- diagram: diagramme (title, diagram_type, elements, description)
- stat: statistique impactante (stat_value, stat_label, context)
- quote: citation (quote_text, author, title)
- highlight: points clés (title, key_points, highlight_color)

Consultant: {CONSULTANT_NAME}
Société: {COMPANY_NAME}
"""

        # Générer la réponse
        response_text = llm.generate(
            prompt=f"Génère les slides pour la section {section}",
            system_prompt=system_prompt,
            temperature=0.7,
            max_tokens=4000,
        )

        # Extraire le JSON
        slides = []
        pptx_path = None

        if "```json" in response_text:
            json_start = response_text.find("```json") + 7
            json_end = response_text.find("```", json_start)
            json_str = response_text[json_start:json_end].strip()

            try:
                slides_data = json.loads(json_str)
                slides = slides_data.get("slides", [])

                # Générer le PPTX si des slides sont présentes
                if slides:
                    from utils.pptx_generator import build_proposal_pptx

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_dir = BASE_DIR / "output"
                    output_dir.mkdir(exist_ok=True)

                    output_path = output_dir / f"proposal_section_{section}_{timestamp}.pptx"

                    consultant_info = {"name": CONSULTANT_NAME, "company": COMPANY_NAME}

                    pptx_path = build_proposal_pptx(
                        template_path=str(BASE_DIR / "Consulting Tools_Template_Palette 2026.pptx"),
                        slides_data=slides,
                        output_path=str(output_path),
                        consultant_info=consultant_info,
                    )

                    pptx_path = str(Path(pptx_path).relative_to(BASE_DIR))

            except json.JSONDecodeError as e:
                print(f"❌ Erreur JSON section: {e}")
                return JSONResponse({"error": "Erreur parsing JSON"}, status_code=500)

        return {
            "response": f"✅ Section '{section}' générée avec succès ({len(slides)} slides)",
            "slides": slides,
            "pptx_path": pptx_path,
            "section": section,
        }

    except Exception as e:
        print(f"❌ Erreur génération section: {e}")
        import traceback

        traceback.print_exc()
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# === SLIDE EDITOR (HTML-First) ===


@app.get("/slide-editor", response_class=HTMLResponse)
async def slide_editor_page(request: Request):
    return templates.TemplateResponse(
        "slide-editor.html",
        {
            "request": request,
            "active": "slide-editor",
            "consultant_name": CONSULTANT_NAME,
        },
    )


def _extract_images_from_pdf(content_bytes):
    """Extract images from PDF as base64 data URIs"""
    images = []
    try:
        import base64
        import io

        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(content_bytes))
        for page_num, page in enumerate(reader.pages):
            for img_key in page.images:
                img_data = img_key.data
                ext = img_key.name.rsplit(".", 1)[-1].lower()
                mime = {
                    "jpg": "image/jpeg",
                    "jpeg": "image/jpeg",
                    "png": "image/png",
                    "gif": "image/gif",
                }.get(ext, "image/png")
                b64 = base64.b64encode(img_data).decode()
                images.append(
                    {
                        "data_uri": f"data:{mime};base64,{b64}",
                        "name": img_key.name,
                        "page": page_num + 1,
                    }
                )
                if len(images) >= 20:
                    return images
    except Exception as e:
        print(f"Erreur extraction images PDF: {e}")
    return images


def _extract_images_from_pptx(content_bytes):
    """Extract images from PPTX as base64 data URIs"""
    images = []
    try:
        import base64
        import io

        from pptx import Presentation as PptxPres

        prs = PptxPres(io.BytesIO(content_bytes))
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    img = shape.image
                    ext = img.content_type.split("/")[-1]
                    mime = img.content_type
                    b64 = base64.b64encode(img.blob).decode()
                    images.append(
                        {
                            "data_uri": f"data:{mime};base64,{b64}",
                            "name": f"slide_{slide_num + 1}_{shape.name}.{ext}",
                            "slide": slide_num + 1,
                        }
                    )
                    if len(images) >= 20:
                        return images
    except Exception as e:
        print(f"Erreur extraction images PPTX: {e}")
    return images


@app.post("/api/slide-editor/parse-document")
@limiter.limit("20/minute")
async def api_slide_editor_parse_document(request: Request, file: UploadFile = File(...)):
    """Parse un document uploade et retourne son contenu texte + images extraites"""
    try:
        # Valider le fichier uploade (taille, type)
        content = await validate_file_upload(file)
        filename = sanitize_filename(file.filename or "document.txt")
        ext = Path(filename).suffix.lower()
        images = []

        if ext in (".md", ".txt"):
            text = content.decode("utf-8")
        elif ext == ".pdf":
            import io

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            images = _extract_images_from_pdf(content)
        elif ext == ".docx":
            import io

            from docx import Document as DocxDocument

            doc = DocxDocument(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            images = _extract_images_from_pptx(content)
            # Pour PPTX: sauvegarder temporairement le fichier et retourner le
            # chemin
            import tempfile

            temp_dir = Path(tempfile.gettempdir()) / "consulting-tools_uploads"
            temp_dir.mkdir(exist_ok=True)

            temp_path = temp_dir / f"upload_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
            with open(temp_path, "wb") as f:
                f.write(content)

            # Retourner le chemin au lieu du texte
            return {
                "text": str(temp_path),
                "filename": filename,
                "length": len(content),
                "is_pptx": True,
                "images": images,
                "image_count": len(images),
            }
        elif ext == ".html":
            text = content.decode("utf-8")
        else:
            return JSONResponse({"error": f"Format non supporte: {ext}"}, status_code=400)

        return {
            "text": text,
            "filename": filename,
            "length": len(text),
            "images": images,
            "image_count": len(images),
        }
    except Exception as e:
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


def _generate_slide_illustrations(slides, job, topic, gen_type="presentation"):
    """Add Nano Banana image generation prompts to relevant slides (no actual image generation)"""
    try:
        # Slide types that benefit from illustrations
        visual_types = ["content", "highlight", "stat", "diagram", "image", "two_column"]

        # Context adapte au type de generation
        context_map = {
            "formation": "training presentation",
            "proposal": "business proposal",
            "rex": "experience feedback presentation",
            "presentation": "business presentation",
        }
        context_map.get(gen_type, "business presentation")

        for idx, slide in enumerate(slides):
            slide_type = slide.get("type", "")

            # Skip non-visual slides
            if slide_type not in visual_types:
                continue

            # Skip if already has an image or image_prompt
            if slide.get("image") or slide.get("image_prompt"):
                continue

            # Generate prompt for the slide (Nano Banana format)
            slide.get("title", "")
            slide.get("content", "")
            slide.get("bullets", [])
            slide.get("key_points", [])

            # Prompt adapte au type
            if gen_type == "formation":
                prompt = """Create a premium, professional illustration for a training presentation slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: Educational yet professional, Unreal Engine 5 render, engaging and clear.
Colors: Cool blues, warm amber/gold accents, approachable palette.
Mood: Pedagogical, inspiring, professional learning environment.
Format: Wide 16:9, 1792x1024px."""
            elif gen_type == "proposal":
                prompt = """Create a premium, professional illustration for a business proposal slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: High-end corporate, Unreal Engine 5 render, sophisticated and impactful.
Colors: Premium blues, gold/amber accents, executive palette.
Mood: Professional, trustworthy, results-oriented, winning proposal aesthetic.
Format: Wide 16:9, 1792x1024px."""
            else:
                prompt = """Create a premium, professional illustration for a {context} slide.

Topic: {topic}
Slide Title: {title}
Content: {content}
Key Points: {', '.join(bullets[:3] if bullets else key_points[:3])}

Style: Corporate tech aesthetic, Unreal Engine 5 render, clean and modern.
Colors: Cool blues, warm amber/gold accents, professional palette.
Mood: Sophisticated, futuristic but grounded, business presentation quality.
Format: Wide 16:9, 1792x1024px."""

            # Add prompt and dimensions to slide (for Nano Banana generation)
            slide["image_prompt"] = prompt.strip()
            slide["image_dimensions"] = "1792x1024px (16:9)"
            job["steps"].append(
                {"message": f"Prompt Nano Banana ajoute pour slide {idx + 1}", "step": 4}
            )

    except Exception as e:
        print(f"  Erreur ajout prompts illustrations: {e}")
        # Non-blocking - continue sans prompts


def _extract_slides_from_buffer(buffer, job):
    """Parse complete JSON slide objects from a partial LLM response buffer.
    Adds newly found slides to job['slides'] incrementally."""
    # Find the slides array content
    content = buffer
    # Strip markdown fences
    if "```json" in content:
        content = content.split("```json", 1)[1]
        if "```" in content:
            content = content.rsplit("```", 1)[0]
    elif "```" in content:
        parts = content.split("```")
        if len(parts) >= 2:
            content = parts[1]

    # Find the "slides" array or top-level array
    start = -1
    for marker in ['"slides"', "'slides'"]:
        idx = content.find(marker)
        if idx != -1:
            bracket = content.find("[", idx)
            if bracket != -1:
                start = bracket
                break

    if start == -1:
        # Maybe it's a top-level array
        bracket = content.find("[")
        if bracket != -1:
            start = bracket
        else:
            return

    # Extract individual slide objects by tracking brace depth
    existing_count = len(job["slides"])
    obj_start = -1
    depth = 0
    slide_count = 0

    for i in range(start, len(content)):
        c = content[i]
        if c == "{":
            if depth == 0:
                obj_start = i
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0 and obj_start != -1:
                slide_count += 1
                if slide_count > existing_count:
                    # New slide found - try to parse it
                    obj_str = content[obj_start : i + 1]
                    try:
                        slide = json.loads(obj_str)
                        if isinstance(slide, dict) and ("type" in slide or "title" in slide):
                            job["slides"].append(slide)
                    except json.JSONDecodeError:
                        pass
                obj_start = -1


@app.post("/api/slide-editor/start-generate")
@limiter.limit("5/minute")
async def api_slide_editor_start_generate(request: Request):
    """Demarre la generation de slides en arriere-plan et retourne un job_id"""
    data = await request.json()
    topic = data.get("topic", "")
    audience = data.get("audience", "")
    slide_count = min(int(data.get("slide_count", 10)), 40)
    gen_type = data.get("type", "presentation")
    model = data.get("model", "gemini-2.0-flash")
    document_text = data.get("document_text", "")
    feedback = data.get("feedback", "")
    previous_slides = data.get("previous_slides", "")
    provider = "claude" if model.startswith("claude") else "gemini"

    _ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    _uid = id(data) % 10000
    job_id = f"slide_{_ts}_{_uid}"

    jobs[job_id] = {
        "type": "slide-editor",
        "status": "running",
        "steps": [],
        "slides": [],
        "result": None,
        "error": None,
    }

    def run_generation():
        job = jobs[job_id]
        try:
            job["steps"].append(
                {
                    "message": f"Initialisation {model}...",
                    "substep": "Chargement du modele",
                    "step": 0,
                }
            )

            from agents.html_slide_generator import HtmlSlideGeneratorAgent

            agent = HtmlSlideGeneratorAgent(model=model, provider=provider)

            # Load extra context for proposals (references + CVs)
            extra_context = ""
            if gen_type == "proposal":
                try:
                    references_path = Path(BASE_DIR) / "data" / "notebooklm" / "references.json"
                    if references_path.exists():
                        with open(references_path, "r", encoding="utf-8") as f:
                            refs = json.load(f)
                        projects = json.dumps(
                            refs.get("projects", []), indent=2, ensure_ascii=False
                        )[:3000]
                        expertise = json.dumps(refs.get("expertise", []), ensure_ascii=False)[:1000]
                        methodologies = json.dumps(
                            refs.get("methodologies", []), ensure_ascii=False
                        )[:1000]
                        extra_context += (
                            f"\nREFERENCES Consulting Tools :\n{projects}"
                            f"\nEXPERTISE : {expertise}"
                            f"\nMETHODOLOGIES : {methodologies}"
                        )
                except Exception as e:
                    print(f"  Erreur chargement references: {e}")

                try:
                    bio_path = Path(BASE_DIR) / "Biographies - CV All Consulting Tools.pptx"
                    if bio_path.exists():
                        from pptx import Presentation as PptxPres

                        prs = PptxPres(str(bio_path))
                        cvs = []
                        for slide in prs.slides:
                            texts = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for p in shape.text_frame.paragraphs:
                                        if p.text.strip():
                                            texts.append(p.text.strip())
                            full = "\n".join(texts)
                            if len(full) > 50:
                                cvs.append(full)
                        if cvs:
                            extra_context += (
                                "\nCVs EQUIPE Consulting Tools :\n" + "\n---\n".join(cvs[:5])[:3000]
                            )
                except Exception as e:
                    print(f"  Erreur chargement CVs: {e}")

            # Handle feedback/regeneration
            if feedback and previous_slides:
                extra_context += (
                    "\n\nSLIDES ACTUELLES (a modifier selon le feedback) :\n"
                    + previous_slides[:6000]
                    + "\n\nFEEDBACK UTILISATEUR :\n"
                    + feedback
                    + "\n\nModifie les slides en tenant compte du feedback."
                )

            job["steps"].append(
                {
                    "message": "Generation HTML Premium...",
                    "substep": "Analyse du sujet et creation des slides",
                    "step": 1,
                }
            )
            job["head_html"] = ""
            job["html_sections"] = []

            # Stream HTML sections via Gemini 3.1 Pro
            for result in agent.run_streaming(
                topic=topic,
                num_slides=slide_count,
                gen_type=gen_type,
                audience=audience,
                document_text=document_text,
                extra_context=extra_context,
            ):
                if result["head_html"]:
                    job["head_html"] = result["head_html"]

                job["html_sections"].append(result["html_section"])
                slide_data = result["slide_json"]
                slide_data["html_section"] = result["html_section"]
                job["slides"].append(slide_data)

                slide_title = (
                    result.get("slide_json", {}).get("title", "") or f"Slide {result['index'] + 1}"
                )
                job["steps"].append(
                    {
                        "message": f"Slide {result['index'] + 1} generee",
                        "substep": slide_title,
                        "step": 2,
                    }
                )

            if not job["slides"]:
                raise ValueError("Aucune slide generee - reponse LLM invalide")

            job["result"] = {
                "slides": job["slides"],
                "total": len(job["slides"]),
                "head_html": job["head_html"],
            }
            job["status"] = "done"

            # Add image prompts in background (non-blocking)
            _generate_slide_illustrations(job["slides"], job, topic, gen_type)

        except Exception as e:
            print(f"Erreur slide-editor generate: {e}")
            import traceback

            traceback.print_exc()
            job["error"] = safe_error_message(e)
            job["status"] = "error"

    import threading

    threading.Thread(target=run_generation, daemon=True).start()

    return {"job_id": job_id}


@app.get("/api/slide-editor/stream/{job_id}")
async def api_slide_editor_stream(job_id: str):
    """SSE stream pour la progression de generation de slides"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        last_slide_idx = 0

        while True:
            # Send new steps
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("status", step)
                last_step_idx += 1

            # Send slides progressively as they appear
            while last_slide_idx < len(job["slides"]):
                slide = job["slides"][last_slide_idx]
                yield send_sse(
                    "slide",
                    {
                        "index": last_slide_idx,
                        "total": len(job["slides"]),
                        "slide": slide,
                        "head_html": job.get("head_html", "") if last_slide_idx == 0 else "",
                    },
                )
                last_slide_idx += 1
                await asyncio.sleep(0.1)

            if job["status"] == "done":
                # Send any remaining slides
                while last_slide_idx < len(job["slides"]):
                    slide = job["slides"][last_slide_idx]
                    yield send_sse(
                        "slide",
                        {
                            "index": last_slide_idx,
                            "total": len(job["slides"]),
                            "slide": slide,
                            "head_html": job.get("head_html", "") if last_slide_idx == 0 else "",
                        },
                    )
                    last_slide_idx += 1
                yield send_sse(
                    "done",
                    {
                        "total": len(job["slides"]),
                        "head_html": job.get("head_html", ""),
                    },
                )
                return

            if job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.3)

    return StreamingResponse(event_generator(), media_type="text/event-stream")


def _get_slide_json_format():
    return """
IMPORTANT : Tu dois retourner UNIQUEMENT un JSON valide avec cette structure :
{{"slides": [...]}}

Types de slides disponibles et leur format JSON :
- "cover" : {{"type":"cover", "title":"...", "subtitle":"...", "meta":"..."}}
- "section" : {{"type":"section", "title":"...", "number":"01"}}
- "content" : {{"type":"content", "title":"...", "bullets":["...","..."]}}
- "stat" : {{"type":"stat", "stat_value":"73%", "stat_label":"...", "context":"...", "subtitle":"..."}}
- "quote" : {{"type":"quote", "quote_text":"...", "author":"...", "title":"..."}}
- "highlight" : {{"type":"highlight", "title":"...", "key_points":["...","...","..."]}}
- "diagram" : {{"type":"diagram", "title":"...", "diagram_type":"flow|pyramid|timeline|cycle|matrix", "elements":["...","..."]}}
- "two_column" : {{"type":"two_column", "title":"...", "left_title":"...", "left_points":["..."], "right_title":"...", "right_points":["..."]}}
- "table" : {{"type":"table", "title":"...", "headers":["...","..."], "rows":[["...","..."]]}}
- "closing" : {{"type":"closing", "title":"Merci", "subtitle":"..."}}

Optionnel : chaque slide peut avoir un champ "tags": ["TAG1", "TAG2"] pour la categoriser.
"""


def _get_proposal_system_prompt(company_name):
    return """Tu agis en tant que Consultant Senior chez {company_name}. Ta mission est de rediger le contenu et la structure d'une proposition commerciale au format "Slides".

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.
- Style : Tags en en-tete (#TECH #DATA), marqueur vertical "Consulting Tools | 2025-2026", mise en page aeree.

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Utilise des types varies : stat pour les chiffres, diagram pour les processus, highlight pour les points cles, two_column pour les comparaisons
- Le contenu doit etre redige sur un ton expert, convaincant et direct
- Maximum 4 bullet points par slide, phrases courtes et percutantes
- Tu determines toi-meme le nombre de slides adapte au contenu et a la complexite du besoin client
- Ajoute des slides CV equipe et references missions si ces informations sont fournies"""


def _get_proposal_prompt(
    topic,
    audience,
    company_name,
    consultant_name,
    json_format,
    document_text="",
    references="",
    cvs="",
):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT FOURNI (utilise ce contenu comme base) :
{document_text[:6000]}
"""

    ref_section = ""
    if references:
        ref_section = """

REFERENCES DE MISSIONS SIMILAIRES (integre-les dans la proposition) :
{references}
"""

    cv_section = ""
    if cvs:
        cv_section = """

EQUIPE DISPONIBLE (integre les profils pertinents dans la proposition) :
{cvs}
"""

    return """Genere une proposition commerciale pour le sujet suivant : {topic}

PUBLIC CIBLE : {audience or 'Direction et decideurs'}
CONSULTANT : {consultant_name} - {company_name}
{doc_section}{ref_section}{cv_section}
Determine toi-meme le nombre de slides et la structure les plus adaptes au besoin client.
Adapte la profondeur et le nombre de slides a la complexite du sujet.

### ELEMENTS A COUVRIR (ordre et decoupage libres) :
- Couverture (cover) : titre impactant, meta "{company_name} | 2025-2026"
- Comprehension du besoin : montrer que tu as compris la realite du client
- Approche et methodologie : comment tu vas repondre au besoin
- Demarche / Phases : diagram flow ou timeline
- Planning / Feuille de route
- Dispositif et gouvernance : equipe, roles, rituels
- Facteurs cles de succes
- Pourquoi {company_name} : valeur ajoutee, differenciants
- Presentation {company_name} : Cabinet de conseil en transformation digitale (Strategy & Vision, Data & Insights, Product Management, Scalability Platform)
- Si des CVs sont fournis : slides CV equipe (type "cv") avec experiences REFORMULEES pour le besoin client
- Si des references sont fournies : slides references (type "highlight") avec resultats concrets et paralleles avec la mission
- Cloture (closing) : contact {consultant_name}

{json_format}

Genere maintenant la proposition pour : {topic}"""


def _get_presentation_system_prompt(company_name):
    return """Tu es un expert en creation de presentations professionnelles pour {company_name}.
Tu crees des presentations visuelles, impactantes et variees.

REGLES STRICTES :
- VARIETE : Alterne les types de slides. Ne jamais avoir 2 slides "content" consecutives.
- VISUELS : Au moins 40% des slides doivent etre visuelles (stat, diagram, highlight, quote)
- CONCIS : Maximum 3-4 bullet points par slide, phrases courtes
- IMPACT : Chiffres-cles en "stat", principes en "quote", comparaisons en "two_column"

Tu retournes UNIQUEMENT un JSON valide."""


def _get_presentation_prompt(topic, audience, slide_count, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme base) :
{document_text[:6000]}
"""

    return """Genere une presentation de {slide_count} slides sur le sujet suivant :

SUJET : {topic}
PUBLIC CIBLE : {audience or 'Managers et decideurs'}
NOMBRE DE SLIDES : {slide_count}
{doc_section}
Structure attendue : commence par cover, puis alterne sections et contenu, termine par closing.

{json_format}

Reponds UNIQUEMENT avec le JSON."""


def _get_formation_system_prompt(company_name):
    return """Tu es un formateur expert chez {company_name}, specialise en data, IA et transformation digitale.
Tu crees des supports de formation pedagogiques, progressifs et engageants.

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.

### REGLES PEDAGOGIQUES
- Progression : du concept a la pratique (theorie > exemples > exercices)
- 1 idee par slide maximum
- Alterner theorie (content), visualisation (diagram, stat), et recapitulatif (highlight)
- Inclure des slides quiz/exercice avec le type "highlight" (key_points comme options)
- Ajouter un champ "tags" (tableau de strings) pour categoriser chaque slide : ex ["THEORIE"], ["PRATIQUE"], ["QUIZ"], ["DEMO"]
- Maximum 3 bullet points par slide, phrases courtes et pedagogiques

### NOMBRE DE SLIDES
- Tu determines toi-meme le nombre de slides adapte au sujet, a la duree et au contenu du brief
- Chaque module doit avoir : intro, theorie, exercice/quiz, recap

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Au moins 40% de slides visuelles (stat, diagram, highlight, quote)
- Ne jamais avoir 2 slides "content" consecutives"""


def _get_formation_prompt(topic, audience, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme matiere pedagogique) :
{document_text[:6000]}
"""

    return """Genere un support de formation complet sur le sujet suivant :

SUJET : {topic or 'Formation basee sur le document fourni'}
PUBLIC CIBLE : {audience or 'Equipes techniques et managers'}
{doc_section}
Determine le nombre de slides adapte au contenu et a la duree mentionnee dans le brief.
Si aucune duree n est precisee, genere 15-20 slides (1 journee).

### STRUCTURE ATTENDUE :
1. Cover : titre de la formation + sous-titre
2. Section "Objectifs" : ce que les apprenants sauront faire
3. Sections thematiques : alterner content, diagram, stat, highlight
4. Pour chaque section : concept > illustration > exercice/quiz
5. Recapitulatif final (highlight avec les points cles)
6. Closing : prochaines etapes + contact

IMPORTANT : Ajoute un champ "tags" (tableau) a chaque slide pour la categoriser.
Exemples de tags : "THEORIE", "PRATIQUE", "QUIZ", "DEMO", "RECAP", "OBJECTIFS"

{json_format}

Reponds UNIQUEMENT avec le JSON."""


def _get_rex_system_prompt(company_name):
    return """Tu es un consultant senior chez {company_name} qui redige un Retour d Experience (REX) de mission.

### IDENTITE VISUELLE (DA)
- Polices : Titres en "Chakra Petch" (Tech/Impact), Textes courants en "Inter" (Lisibilite).
- Palette Couleurs (Regle 60-30-10) :
  - 60% Dominante : Fond Blanc (#FFFFFF) ou Rose Poudre (#FBF0F4) pour les encadres.
  - 30% Secondaire : Textes en Noir Profond (#1F1F1F) ou Gris Moyen (#474747).
  - 10% Accent : Terracotta (#c0504d) ou Corail (#FF6B58). Maximum 5 elements par slide.

### REGLES DU REX
- Ton factuel et professionnel, oriente resultats
- Chiffres concrets (KPIs, delais, gains, ROI)
- Structure : Contexte > Enjeux > Demarche > Realisations > Resultats > Apprentissages > Recommandations
- Utilise des types visuels varies : stat pour les resultats, diagram pour la demarche, timeline pour le planning
- Maximum 3-4 bullet points par slide
- Ajoute un champ "tags" a chaque slide : "CONTEXTE", "ENJEUX", "DEMARCHE", "RESULTATS", "APPRENTISSAGES"

### REGLES DE GENERATION JSON
- Tu dois retourner UNIQUEMENT du JSON valide, pas de texte explicatif
- Chaque slide a un "type" parmi les types disponibles
- Au moins 40% de slides visuelles (stat, diagram, highlight, timeline)"""


def _get_rex_prompt(topic, client, json_format, document_text=""):
    doc_section = ""
    if document_text:
        doc_section = """

DOCUMENT SOURCE (utilise ce contenu comme base pour le REX) :
{document_text[:6000]}
"""

    return """Genere un Retour d Experience (REX) de mission pour :

MISSION : {topic or 'Mission basee sur le document fourni'}
CLIENT / CONTEXTE : {client or 'Client entreprise'}
{doc_section}
### STRUCTURE OBLIGATOIRE (10 slides) :

**Slide 1 : Couverture (cover)**
- Titre : "REX Mission - [Nom de la mission]"
- Sous-titre avec le client et la periode

**Slide 2 : Contexte Client (content)**
- Qui est le client, son secteur, sa taille
- Contexte de la mission (appel d offre, besoin identifie)
- Tags: ["CONTEXTE"]

**Slide 3 : Enjeux & Objectifs (highlight)**
- 3-4 enjeux business et objectifs mesurables
- Tags: ["ENJEUX"]

**Slide 4 : Notre Demarche (diagram)**
- Type "flow" : phases de la mission (Cadrage > Audit > Design > Build > Run)
- Tags: ["DEMARCHE"]

**Slide 5 : Planning & Jalons (diagram)**
- Type "timeline" : jalons cles de la mission
- Tags: ["DEMARCHE"]

**Slide 6 : Realisations Cles (highlight)**
- 3-4 livrables ou realisations majeures
- Tags: ["RESULTATS"]

**Slide 7 : Resultats Chiffres (stat)**
- KPI principal avec valeur impactante (ex: +35 pourcent, ROI x3, etc.)
- Tags: ["RESULTATS"]

**Slide 8 : Apprentissages (two_column)**
- Colonne gauche : "Ce qui a bien marche" (facteurs de succes)
- Colonne droite : "Points d'amelioration" (lessons learned)
- Tags: ["APPRENTISSAGES"]

**Slide 9 : Recommandations (content)**
- 3-4 recommandations pour la suite ou pour de futures missions similaires
- Tags: ["APPRENTISSAGES"]

**Slide 10 : Cloture (closing)**
- "Merci" + contact

{json_format}

Reponds UNIQUEMENT avec le JSON."""


@app.post("/api/slide-editor/generate")
@limiter.limit("5/minute")
async def api_slide_editor_generate(request: Request):
    """Genere des slides JSON via LLM pour le slide editor (non-streaming fallback)"""
    try:
        data = await request.json()

        # Valider et sanitizer les inputs texte
        topic = sanitize_text_input(data.get("topic", ""), max_length=1000, field_name="topic")
        audience = sanitize_text_input(
            data.get("audience", ""), max_length=500, field_name="audience"
        )
        document_text = sanitize_text_input(
            data.get("document_text", ""), max_length=50000, field_name="document_text"
        )

        slide_count = data.get("slide_count", 10)
        gen_type = data.get("type", "presentation")
        model = data.get("model", "")

        from utils.llm_client import LLMClient

        llm_kwargs = {"max_tokens": 8192}
        if model:
            llm_kwargs["model"] = model
            llm_kwargs["provider"] = "gemini"
        llm = LLMClient(**llm_kwargs)

        consultant_name = os.getenv("CONSULTANT_NAME", "Jean-Sebastien Abessouguie Bayiha")
        company_name = os.getenv("COMPANY_NAME", "Consulting Tools")
        json_format = _get_slide_json_format()

        if gen_type == "proposal":
            system_prompt = _get_proposal_system_prompt(company_name)
            prompt = _get_proposal_prompt(
                topic, audience, company_name, consultant_name, json_format, document_text
            )
        elif gen_type == "formation":
            system_prompt = _get_formation_system_prompt(company_name)
            prompt = _get_formation_prompt(topic, audience, json_format, document_text)
        elif gen_type == "rex":
            system_prompt = _get_rex_system_prompt(company_name)
            prompt = _get_rex_prompt(topic, audience, json_format, document_text)
        else:
            system_prompt = _get_presentation_system_prompt(company_name)
            prompt = _get_presentation_prompt(
                topic, audience, slide_count, json_format, document_text
            )

        response = llm.generate(prompt=prompt, system_prompt=system_prompt, temperature=0.7)

        if "```json" in response:
            json_str = response.split("```json")[1].split("```")[0].strip()
        elif "```" in response:
            json_str = response.split("```")[1].split("```")[0].strip()
        else:
            json_str = response.strip()

        result = json.loads(json_str)
        return JSONResponse(result)

    except json.JSONDecodeError as e:
        print(f"Erreur JSON slide-editor: {e}")
        return JSONResponse({"error": "Erreur de parsing JSON", "slides": []}, status_code=422)
    except Exception as e:
        print(f"Erreur slide-editor generate: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# === DOCUMENT EDITOR ===


@app.get("/document-editor", response_class=HTMLResponse)
async def document_editor_page(request: Request):
    return templates.TemplateResponse(
        "document-editor.html",
        {
            "request": request,
            "active": "document-editor",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.get("/veille", response_class=HTMLResponse)
async def veille_page(request: Request):
    return templates.TemplateResponse(
        "veille.html",
        {
            "request": request,
            "active": "veille",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/document-editor/start-generate")
@limiter.limit("10/minute")
async def api_document_editor_start_generate(request: Request):
    """Demarre la generation de document en arriere-plan"""
    data = await request.json()
    doc_type = data.get("type", "formation")

    # Valider et sanitizer les inputs texte
    topic = sanitize_text_input(data.get("topic", ""), max_length=1000, field_name="topic")
    document_text = sanitize_text_input(
        data.get("document_text", ""), max_length=50000, field_name="document_text"
    )
    audience = sanitize_text_input(data.get("audience", ""), max_length=500, field_name="audience")
    feedback = sanitize_text_input(data.get("feedback", ""), max_length=5000, field_name="feedback")
    previous_content = sanitize_text_input(
        data.get("previous_content", ""), max_length=50000, field_name="previous_content"
    )

    length = data.get("length", "medium")
    model = data.get("model", "")
    use_context = data.get("use_context", False)

    _ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    _uid = id(data) % 10000
    job_id = f"doc_{_ts}_{_uid}"

    jobs[job_id] = {
        "type": "document-editor",
        "status": "running",
        "steps": [],
        "chunks": [],
        "result": None,
        "error": None,
    }

    def run_doc_generation():
        job = jobs[job_id]
        try:
            job["steps"].append({"message": "Initialisation du modele..."})

            from utils.llm_client import LLMClient

            llm_kwargs = {"max_tokens": 8192}
            if model:
                llm_kwargs["model"] = model
                llm_kwargs["provider"] = "gemini"

            result = {"markdown": "", "linkedin_post": ""}

            if feedback and previous_content:
                # Regeneration with feedback
                job["steps"].append({"message": "Application du feedback..."})

                if doc_type == "formation":
                    from agents.formation_generator import FormationGeneratorAgent

                    agent = FormationGeneratorAgent()
                    if model:
                        agent.llm = LLMClient(**llm_kwargs)
                    gen_result = agent.regenerate_with_feedback(previous_content, feedback)
                    result["markdown"] = gen_result.get("markdown", "")
                else:
                    # Article or REX text: generic LLM regeneration
                    llm = LLMClient(**llm_kwargs)
                    system_prompt = "Tu es un expert en redaction chez Consulting Tools. Tu dois modifier un document existant en tenant compte du feedback utilisateur. Conserve la structure markdown et ameliore le contenu."
                    prompt = """Voici le document actuel :
---
{previous_content[:6000]}
---

Feedback utilisateur :
---
{feedback}
---

Modifie le document en tenant compte du feedback. Retourne UNIQUEMENT le document markdown modifie, sans explication."""

                    # Stream response progressively
                    buffer = ""
                    try:
                        for chunk in llm.generate_stream(
                            prompt=prompt, system_prompt=system_prompt, temperature=0.7
                        ):
                            buffer += chunk
                            job["chunks"].append(chunk)
                    except Exception:
                        buffer = llm.generate(
                            prompt=prompt, system_prompt=system_prompt, temperature=0.7
                        )
                        job["chunks"].append(buffer)

                    response = buffer.strip()
                    if response.startswith("```markdown"):
                        response = response[len("```markdown") :]
                    if response.startswith("```"):
                        response = response[3:]
                    if response.endswith("```"):
                        response = response[:-3]
                    result["markdown"] = response.strip()

            elif doc_type == "formation":
                job["steps"].append({"message": "Generation du programme de formation..."})
                from agents.formation_generator import FormationGeneratorAgent

                agent = FormationGeneratorAgent()
                if model:
                    agent.llm = LLMClient(**llm_kwargs)
                input_text = topic
                if document_text:
                    input_text += f"\n\nDocument source :\n{document_text[:6000]}"
                if audience:
                    input_text += f"\n\nPublic cible : {audience}"
                gen_result = agent.generate_programme(input_text)
                result["markdown"] = gen_result.get("markdown", "")

            elif doc_type == "article":
                job["steps"].append({"message": "Generation de l'article..."})
                from agents.article_generator import ArticleGeneratorAgent

                agent = ArticleGeneratorAgent()
                if model:
                    agent.llm = LLMClient(**llm_kwargs)
                idea = topic
                if document_text:
                    idea += f"\n\nDocument source :\n{document_text[:4000]}"
                gen_result = agent.run(idea, target_length=length, use_context=use_context)
                result["markdown"] = gen_result.get("article", "")
                result["linkedin_post"] = gen_result.get("linkedin_post", "")

            elif doc_type == "rex_text":
                job["steps"].append({"message": "Generation du REX..."})
                llm = LLMClient(**llm_kwargs)

                system_prompt = """Tu es un consultant senior chez {
                    os.getenv(
                        'COMPANY_NAME',
                        'Consulting Tools')} qui redige un Retour d Experience (REX) de mission.

TON STYLE :
- Factuel et professionnel, oriente resultats
- Chiffres concrets (KPIs, delais, gains, ROI)
- Structure claire en sections markdown
- Pragmatique et pedagogique

FORMAT MARKDOWN :
- # pour le titre principal
- ## pour les sections
- ### pour les sous-sections
- **Gras** pour les points cles
- Listes a puces pour les enumerations
- > Citations pour les temoignages clients
- Tableaux pour les KPIs et les plannings"""

                doc_section = ""
                if document_text:
                    doc_section = f"\n\nDocument source :\n{document_text[:6000]}"

                prompt = """Redige un Retour d Experience (REX) complet de mission en markdown.

MISSION : {topic or 'Mission basee sur le document fourni'}
CLIENT / CONTEXTE : {audience or 'Client entreprise'}
{doc_section}

### STRUCTURE OBLIGATOIRE :
1. **Titre** (H1) : REX Mission - [Nom]
2. **Contexte** (H2) : client, secteur, contexte de la mission
3. **Enjeux & Objectifs** (H2) : problematique, objectifs mesurables
4. **Demarche** (H2) : phases, methodologie, planning
5. **Realisations** (H2) : livrables, actions menees
6. **Resultats** (H2) : KPIs, metriques, gains (utilise un tableau)
7. **Apprentissages** (H2) : ce qui a marche, axes d amelioration
8. **Recommandations** (H2) : prochaines etapes, conseils
9. **Conclusion** (H2) : synthese

Retourne UNIQUEMENT le document markdown, sans preambule."""

                # Stream response progressively
                buffer = ""
                try:
                    for chunk in llm.generate_stream(
                        prompt=prompt, system_prompt=system_prompt, temperature=0.7
                    ):
                        buffer += chunk
                        job["chunks"].append(chunk)
                except Exception:
                    buffer = llm.generate(
                        prompt=prompt, system_prompt=system_prompt, temperature=0.7
                    )
                    job["chunks"].append(buffer)

                response = buffer.strip()
                if response.startswith("```markdown"):
                    response = response[len("```markdown") :]
                if response.startswith("```"):
                    response = response[3:]
                if response.endswith("```"):
                    response = response[:-3]
                result["markdown"] = response.strip()

            elif doc_type == "linkedin":
                job["steps"].append({"message": "Generation du post LinkedIn..."})
                llm = LLMClient(**llm_kwargs)
                consultant_name = os.getenv("CONSULTANT_NAME", "Jean-Sebastien Abessouguie Bayiha")
                os.getenv("COMPANY_NAME", "Consulting Tools")

                system_prompt = """Tu es {consultant_name}, consultant en strategie data et IA chez {company_name}.
Tu rediges des posts LinkedIn engageants, pragmatiques, orientes resultats, avec une vision critique et constructive.
Style : gen z parisien, oriente resultats, a contre-courant des buzzwords.

STRUCTURE DU POST :
1. Accroche percutante (question provocatrice ou constat surprenant)
2. Corps : 2-3 phrases claires resumant ton point de vue
3. Call to action : question ouverte pour les commentaires
4. 5 hashtags strategiques

REGLES :
- Maximum 1300 caracteres
- Paragraphes courts (1-2 phrases max)
- Pas de emojis excessifs (max 2-3)
- Ton direct et authentique"""

                doc_section = ""
                if document_text:
                    doc_section = f"\n\nContenu source :\n{document_text[:4000]}"

                prompt = """Redige un post LinkedIn sur le sujet suivant :

{topic or 'Sujet base sur le document fourni'}
{doc_section}

Retourne UNIQUEMENT le texte du post, sans explication."""

                # Stream response progressively
                buffer = ""
                try:
                    for chunk in llm.generate_stream(
                        prompt=prompt, system_prompt=system_prompt, temperature=0.7
                    ):
                        buffer += chunk
                        job["chunks"].append(chunk)
                except Exception:
                    buffer = llm.generate(
                        prompt=prompt, system_prompt=system_prompt, temperature=0.7
                    )
                    job["chunks"].append(buffer)

                response = buffer.strip()
                if response.startswith("```"):
                    response = response.split("```", 1)[1]
                    if response.startswith("\n"):
                        response = response[1:]
                    if "```" in response:
                        response = response.rsplit("```", 1)[0]
                result["markdown"] = response.strip()

            elif doc_type == "compte_rendu":
                job["steps"].append({"message": "Analyse du transcript..."})
                from agents.meeting_summarizer import MeetingSummarizerAgent

                agent = MeetingSummarizerAgent()
                if model:
                    agent.llm = LLMClient(**llm_kwargs)

                transcript = topic
                if document_text:
                    transcript += f"\n\n{document_text}"

                # Extract key info
                info_result = agent.extract_key_info(transcript)
                extracted_info = info_result.get("extracted_info", "")

                job["steps"].append({"message": "Generation du compte rendu..."})
                minutes = agent.generate_minutes(transcript, extracted_info)

                job["steps"].append({"message": "Generation de l'email de partage..."})
                email_result = agent.generate_email(extracted_info, minutes)

                result["markdown"] = minutes
                result["email"] = email_result

            elif doc_type == "cv_reference":
                job["steps"].append({"message": "Adaptation du CV/Reference..."})

                from agents.cv_reference_adapter import CVReferenceAdapterAgent

                agent = CVReferenceAdapterAgent()
                if model:
                    agent.llm = LLMClient(**llm_kwargs)

                # Detecter automatiquement si c est un CV ou une reference
                is_cv = any(
                    kw in document_text.lower()
                    for kw in [
                        "experience",
                        "competences",
                        "formation",
                        "cv",
                        "consultant",
                        "diplome",
                        "certification",
                        "parcours professionnel",
                    ]
                )
                type_label = "CV" if is_cv else "Reference"

                job["steps"].append({"message": f"Generation des slides {type_label}..."})

                gen_result = agent.run(
                    document_text=document_text,
                    mission_brief=topic,  # Le "topic" = l appel d offre
                    doc_type=type_label,
                )

                # Retourner slides au lieu de markdown
                job["slides"] = gen_result.get("slides", [])
                result["slides"] = gen_result.get("slides", [])
                result["doc_type"] = type_label

            elif doc_type == "presentation_script":
                job["steps"].append({"message": "Extraction du contenu PPTX..."})

                # document_text contient le chemin vers le fichier PPTX
                # temporaire
                if not document_text or not document_text.endswith(".pptx"):
                    raise ValueError(
                        "Fichier PPTX obligatoire pour generer un script de presentation"
                    )

                from agents.presentation_script_generator import PresentationScriptGenerator

                agent = PresentationScriptGenerator()
                if model:
                    agent.llm = LLMClient(**llm_kwargs)

                job["steps"].append({"message": "Analyse des slides..."})

                # Stream chunks si possible
                gen_result = agent.run(pptx_path=document_text, presentation_context=topic or "")

                result["markdown"] = gen_result.get("markdown", "")
                result["num_slides"] = gen_result.get("num_slides", 0)
                result["estimated_duration"] = gen_result.get("estimated_duration", "")

            job["result"] = result
            job["status"] = "done"
            job["steps"].append({"message": "Generation terminee !"})

        except Exception as e:
            print(f"Erreur document-editor generate: {e}")
            import traceback

            traceback.print_exc()
            job["error"] = safe_error_message(e)
            job["status"] = "error"

    import threading

    threading.Thread(target=run_doc_generation, daemon=True).start()

    return {"job_id": job_id}


@app.get("/api/document-editor/stream/{job_id}")
async def api_document_editor_stream(job_id: str):
    """SSE stream pour la progression de generation de document"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouve"})
            return

        last_step_idx = 0
        last_chunk_idx = 0

        while True:
            # Send new steps
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("status", step)
                last_step_idx += 1

            # Send new chunks progressively
            chunks = job.get("chunks", [])
            while last_chunk_idx < len(chunks):
                yield send_sse("chunk", {"text": chunks[last_chunk_idx]})
                last_chunk_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return

            if job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.3)

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@app.post("/api/document-editor/export-gdocs")
async def api_document_editor_export_gdocs(request: Request):
    """Exporte le markdown vers Google Docs"""
    try:
        data = await request.json()
        markdown = data.get("markdown", "")
        title = data.get("title", "Document Consulting Tools")

        if not markdown:
            return JSONResponse({"error": "Pas de contenu a exporter"}, status_code=400)

        from utils.google_api import GoogleAPIClient

        client = GoogleAPIClient()
        url = client.export_markdown_to_docs(markdown, title)

        return {"url": url}
    except Exception as e:
        print(f"Erreur export Google Docs: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# ===== VEILLE TECHNOLOGIQUE =====


@app.get("/api/veille/articles")
async def api_veille_get_articles(
    limit: int = 50, offset: int = 0, source: str = None, keyword: str = None, days: int = None
):
    """Liste les articles de veille avec filtres"""
    try:
        from utils.article_db import ArticleDatabase

        db = ArticleDatabase()

        articles = db.get_articles(
            limit=limit, offset=offset, source=source, keyword=keyword, days=days
        )

        return {"articles": articles, "count": len(articles)}
    except Exception as e:
        print(f"Erreur lecture articles: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.get("/api/veille/stats")
async def api_veille_stats():
    """Statistiques sur les articles stockes"""
    try:
        from utils.article_db import ArticleDatabase

        db = ArticleDatabase()
        stats = db.get_article_stats()
        return stats
    except Exception as e:
        print(f"Erreur stats: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/veille/generate-digest")
async def api_veille_generate_digest(request: Request):
    """Genere un nouveau digest de veille"""
    try:
        data = await request.json()
        period = data.get("period", "daily")
        days = data.get("days", 1 if period == "daily" else 7)
        keywords = data.get("keywords", [])

        from agents.tech_monitor import TechMonitorAgent

        agent = TechMonitorAgent()

        result = agent.run(keywords=keywords if keywords else None, days=days, period=period)

        return {
            "success": True,
            "digest": result.get("content", result.get("digest", "")),
            "num_articles": result.get("num_articles", 0),
            "digest_id": result.get("digest_id"),
            "md_path": result.get("md_path"),
        }
    except Exception as e:
        print(f"Erreur generation digest: {e}")
        import traceback

        traceback.print_exc()
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.get("/api/veille/digests")
async def api_veille_get_digests(period: str = "daily"):
    """Liste les digests generes"""
    try:
        from utils.article_db import ArticleDatabase

        db = ArticleDatabase()

        # Recuperer le dernier digest de cette periode
        latest = db.get_latest_digest(period=period)

        return {"latest_digest": latest}
    except Exception as e:
        print(f"Erreur lecture digests: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/veille/articles/{article_id}/mark-read")
async def api_veille_mark_read(article_id: int):
    """Marque un article comme lu"""
    try:
        from utils.article_db import ArticleDatabase

        db = ArticleDatabase()
        db.mark_as_read(article_id)
        return {"success": True}
    except Exception as e:
        print(f"Erreur mark read: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


@app.post("/api/veille/articles/{article_id}/toggle-favorite")
async def api_veille_toggle_favorite(article_id: int):
    """Toggle favori sur un article"""
    try:
        from utils.article_db import ArticleDatabase

        db = ArticleDatabase()
        db.toggle_favorite(article_id)
        return {"success": True}
    except Exception as e:
        print(f"Erreur toggle favorite: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# === BUG REPORTS ===

BUG_REPORTS_FILE = Path("data/bug_reports.json")


def _load_bug_reports():
    if BUG_REPORTS_FILE.exists():
        with open(BUG_REPORTS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return []


def _save_bug_reports(reports):
    BUG_REPORTS_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(BUG_REPORTS_FILE, "w", encoding="utf-8") as f:
        json.dump(reports, f, ensure_ascii=False, indent=2)


@app.post("/api/bug-report")
async def api_bug_report(request: Request):
    """Enregistre un signalement de bug"""
    try:
        data = await request.json()
        description = data.get("description", "").strip()
        if not description:
            return JSONResponse({"error": "La description est requise"}, status_code=400)

        severity = data.get("severity", "medium")
        if severity not in ("low", "medium", "high", "critical"):
            severity = "medium"

        page_url = data.get("page_url", "")
        user_agent = data.get("user_agent", "")
        screenshot = data.get("screenshot", "")

        # Generate title with LLM
        title = description[:80]
        try:
            from utils.llm_client import LLMClient

            llm = LLMClient(model="gemini-2.0-flash", provider="gemini", max_tokens=100)
            title = (
                llm.generate_with_context(
                    f"Genere un titre court (max 10 mots) pour ce bug: " f"{description}",
                    system_prompt=(
                        "Tu generes des titres de bug concis en francais. "
                        "Reponds uniquement avec le titre, sans guillemets."
                    ),
                )
                .strip()
                .strip('"')
                .strip("'")
            )
        except Exception:
            pass

        report_id = str(uuid.uuid4())[:8]
        report = {
            "id": report_id,
            "title": title,
            "description": description,
            "severity": severity,
            "page_url": page_url,
            "user_agent": user_agent,
            "screenshot": screenshot[:100] if screenshot else "",
            "has_screenshot": bool(screenshot),
            "created_at": datetime.now().isoformat(),
        }

        reports = _load_bug_reports()
        reports.append(report)
        _save_bug_reports(reports)

        return {"success": True, "id": report_id, "title": title}

    except Exception as e:
        print(f"Erreur bug report: {e}")
        return JSONResponse({"error": safe_error_message(e)}, status_code=500)


# === SKILLS MARKET ===

skills_market_db = ConsultantDatabase()


@app.get("/skills-market", response_class=HTMLResponse)
async def skills_market_page(request: Request):
    return templates.TemplateResponse(
        "skills-market.html",
        {
            "request": request,
            "active": "skills-market",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.get("/api/skills-market/import/status")
async def api_skills_market_import_status():
    """Verifie si les donnees ont deja ete importees"""
    return {
        "imported": skills_market_db.is_imported(),
        "count": skills_market_db.get_consultant_count(),
    }


@app.post("/api/skills-market/import")
@limiter.limit("3/minute")
async def api_skills_market_import(request: Request):
    """Importe les consultants depuis le fichier PPTX"""
    pptx_path = str(BASE_DIR / "Biographies - CV All Consulting Tools.pptx")

    if not os.path.exists(pptx_path):
        return JSONResponse(
            {
                "error": "Fichier PPTX non trouve. "
                "Placez 'Biographies - CV All Consulting Tools.pptx' "
                "a la racine du projet."
            },
            status_code=404,
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "skills-import",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_skills_import,
        args=(job_id, pptx_path),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_skills_import(job_id: str, pptx_path: str):
    """Execute l'import des consultants en background"""
    job = jobs[job_id]

    try:
        skills_market_db.delete_all()
        agent = SkillsMarketAgent()

        job["steps"].append({"step": "reading", "status": "active", "progress": 5})

        def progress_cb(current, total, name):
            pct = int(10 + (current / max(total, 1)) * 80)
            job["steps"].append(
                {
                    "step": f"parsing_{current}",
                    "status": "active",
                    "progress": pct,
                    "message": (f"Analyse: {name} ({current}/{total})"),
                }
            )

        consultants = agent.import_from_pptx(pptx_path, progress_callback=progress_cb)

        job["steps"].append({"step": "saving", "status": "active", "progress": 90})

        for consultant_data in consultants:
            skills_market_db.save_consultant(consultant_data)

        job["steps"].append({"step": "done", "status": "done", "progress": 100})
        job["status"] = "done"
        job["result"] = {
            "imported": len(consultants),
            "message": (f"{len(consultants)} consultants importes"),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)


@app.post("/api/skills-market/upload")
@limiter.limit("10/minute")
async def api_skills_market_upload(
    request: Request,
    file: UploadFile = File(...),
):
    """Upload un CV consultant (PPTX, PDF, ou HTML)"""
    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()
    allowed = {".pptx", ".pdf", ".html", ".htm"}

    if ext not in allowed:
        return JSONResponse(
            {"error": f"Format non supporte: {ext}. " f"Formats acceptes: {', '.join(allowed)}"},
            status_code=400,
        )

    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        return JSONResponse(
            {"error": "Fichier trop volumineux (max 50 Mo)"},
            status_code=400,
        )

    import tempfile

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "skills-upload",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_skills_upload,
        args=(job_id, tmp_path, filename, ext),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_skills_upload(
    job_id: str,
    file_path: str,
    filename: str,
    ext: str,
):
    """Upload et parse un fichier CV en background"""
    job = jobs[job_id]

    try:
        agent = SkillsMarketAgent()
        job["steps"].append(
            {
                "step": "reading",
                "status": "active",
                "progress": 10,
                "message": f"Lecture de {filename}...",
            }
        )

        consultants = []

        if ext == ".pptx":

            def progress_cb(current, total, name):
                pct = int(10 + (current / max(total, 1)) * 80)
                job["steps"].append(
                    {
                        "step": f"parsing_{current}",
                        "status": "active",
                        "progress": pct,
                        "message": (f"Analyse: {name} " f"({current}/{total})"),
                    }
                )

            consultants = agent.import_from_pptx(
                file_path,
                progress_callback=progress_cb,
            )

        elif ext == ".pdf":
            from utils.document_parser import DocumentParser

            text = DocumentParser.parse_file(file_path) or ""
            if not text.strip():
                raise ValueError("Impossible d'extraire du texte du PDF")

            job["steps"].append(
                {
                    "step": "parsing",
                    "status": "active",
                    "progress": 30,
                    "message": "Analyse du contenu PDF...",
                }
            )

            def progress_cb(current, total, name):
                pct = int(30 + (current / max(total, 1)) * 50)
                job["steps"].append(
                    {
                        "step": f"parsing_{current}",
                        "status": "active",
                        "progress": pct,
                        "message": name,
                    }
                )

            consultants = agent.import_from_text(
                text,
                filename,
                progress_callback=progress_cb,
            )

        elif ext in (".html", ".htm"):
            from bs4 import BeautifulSoup

            raw_html = open(file_path, "r", encoding="utf-8").read()
            soup = BeautifulSoup(raw_html, "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)

            if not text.strip():
                raise ValueError("Impossible d'extraire du texte " "du HTML")

            job["steps"].append(
                {
                    "step": "parsing",
                    "status": "active",
                    "progress": 30,
                    "message": "Analyse du contenu HTML...",
                }
            )

            def progress_cb(current, total, name):
                pct = int(30 + (current / max(total, 1)) * 50)
                job["steps"].append(
                    {
                        "step": f"parsing_{current}",
                        "status": "active",
                        "progress": pct,
                        "message": name,
                    }
                )

            consultants = agent.import_from_text(
                text,
                filename,
                progress_callback=progress_cb,
            )

        job["steps"].append(
            {
                "step": "saving",
                "status": "active",
                "progress": 90,
                "message": "Sauvegarde en base...",
            }
        )

        saved = 0
        for consultant_data in consultants:
            skills_market_db.save_consultant(consultant_data)
            saved += 1

        job["steps"].append({"step": "done", "status": "done", "progress": 100})
        job["status"] = "done"
        job["result"] = {
            "imported": saved,
            "message": (f"{saved} consultant(s) ajoute(s) " f"depuis {filename}"),
        }

    except Exception as e:
        job["status"] = "error"
        job["error"] = safe_error_message(e)
    finally:
        try:
            os.unlink(file_path)
        except OSError:
            pass


@app.get("/api/skills-market/stream/{job_id}")
async def api_skills_market_stream(job_id: str):
    """Stream SSE pour l'import"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield (f"data: {json.dumps(step)}\n\n")
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield (f"data: {json.dumps(final)}\n\n")
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


@app.get("/api/skills-market/consultants")
async def api_skills_market_list(
    technical: Optional[str] = None,
    sector: Optional[str] = None,
):
    """Liste les consultants avec filtres optionnels"""
    tech_list = [t.strip() for t in technical.split(",") if t.strip()] if technical else None
    sector_list = [s.strip() for s in sector.split(",") if s.strip()] if sector else None

    if tech_list or sector_list:
        consultants = skills_market_db.search_by_skills(technical=tech_list, sector=sector_list)
    else:
        consultants = skills_market_db.get_all_consultants()

    return {"consultants": consultants}


@app.get("/api/skills-market/consultants/{consultant_id}")
async def api_skills_market_detail(consultant_id: int):
    """Recupere le profil complet d'un consultant"""
    consultant = skills_market_db.get_consultant(consultant_id)
    if not consultant:
        return JSONResponse(
            {"error": "Consultant non trouve"},
            status_code=404,
        )
    return {"consultant": consultant}


@app.delete("/api/skills-market/consultants/{consultant_id}")
async def api_skills_market_delete(consultant_id: int):
    """Supprime un consultant"""
    deleted = skills_market_db.delete_consultant(consultant_id)
    if not deleted:
        return JSONResponse(
            {"error": "Consultant non trouve"},
            status_code=404,
        )
    return {"message": "Consultant supprime"}


@app.get("/api/skills-market/skills")
async def api_skills_market_skills():
    """Recupere toutes les competences uniques"""
    return skills_market_db.get_all_skills()


@app.post("/api/skills-market/consultants" "/{consultant_id}/missions")
@limiter.limit("10/minute")
async def api_skills_market_add_mission(
    request: Request,
    consultant_id: int,
):
    """Ajoute une mission et met a jour les competences"""
    body = await request.json()

    client_name = body.get("client_name", "").strip()
    if not client_name:
        return JSONResponse(
            {"error": "Le nom du client est requis"},
            status_code=400,
        )

    mission_data = {
        "client_name": client_name,
        "context_and_challenges": body.get("context_and_challenges", "").strip(),
        "deliverables": body.get("deliverables", "").strip(),
        "tasks": body.get("tasks", "").strip(),
        "start_date": body.get("start_date", ""),
        "end_date": body.get("end_date", ""),
        "skills_technical": body.get("skills_technical", []),
        "skills_sector": body.get("skills_sector", []),
    }

    mission_id = skills_market_db.add_mission(consultant_id, mission_data)
    return {"mission_id": mission_id, "message": "Experience ajoutee"}


@app.put("/api/skills-market/consultants/{consultant_id}/info")
@limiter.limit("20/minute")
async def api_skills_market_update_consultant_info(
    request: Request,
    consultant_id: int,
):
    """Met a jour les informations de base d'un consultant (nom, prenom, entreprise,
    linkedin, titre, bio, theme couleurs/polices)."""
    body = await request.json()
    updated = skills_market_db.update_consultant_info(consultant_id, body)
    if not updated:
        return JSONResponse(
            {"error": "Consultant introuvable ou aucune modification"},
            status_code=404,
        )
    return {"ok": True, "message": "Profil mis a jour"}


@app.post("/api/skills-market/consultants" "/{consultant_id}/certifications")
@limiter.limit("10/minute")
async def api_skills_market_add_certification(
    request: Request,
    consultant_id: int,
):
    """Ajoute une certification et met a jour les skills"""
    body = await request.json()

    cert_name = body.get("name", "").strip()
    if not cert_name:
        return JSONResponse(
            {"error": "Le nom de la certification est requis"},
            status_code=400,
        )

    cert_data = {
        "name": cert_name,
        "organization": body.get("organization", "").strip(),
        "date_obtained": body.get("date_obtained", "").strip(),
        "description": body.get("description", "").strip(),
        "skills_technical": body.get("skills_technical", []),
        "skills_sector": body.get("skills_sector", []),
    }

    cert_id = skills_market_db.add_certification(consultant_id, cert_data)
    return {"certification_id": cert_id, "message": "Certification ajoutee"}


@app.put("/api/skills-market/consultants" "/{consultant_id}/interests")
@limiter.limit("10/minute")
async def api_skills_market_update_interests(
    request: Request,
    consultant_id: int,
):
    """Met a jour les centres d'interet"""
    body = await request.json()
    interests = body.get("interests", [])

    if not isinstance(interests, list):
        return JSONResponse(
            {"error": "interests doit etre une liste"},
            status_code=400,
        )

    skills_market_db.update_interests(consultant_id, interests)
    return {"message": "Centres d'interet mis a jour"}


@app.put("/api/skills-market/consultants" "/{consultant_id}/disinterests")
@limiter.limit("10/minute")
async def api_skills_market_update_disinterests(
    request: Request,
    consultant_id: int,
):
    """Met a jour les centres de desinteret"""
    body = await request.json()
    disinterests = body.get("disinterests", [])

    if not isinstance(disinterests, list):
        return JSONResponse(
            {"error": "disinterests doit etre une liste"},
            status_code=400,
        )

    skills_market_db.update_disinterests(consultant_id, disinterests)
    return {"message": "Centres de desinteret mis a jour"}


@app.post("/api/skills-market/search")
@limiter.limit("10/minute")
async def api_skills_market_search(request: Request):
    """Recherche en langage naturel"""
    body = await request.json()
    query = body.get("query", "").strip()

    if not query:
        return JSONResponse(
            {"error": "La requete est vide"},
            status_code=400,
        )

    consultants = skills_market_db.get_all_consultants()
    if not consultants:
        return {"results": [], "consultants": []}

    agent = SkillsMarketAgent()
    results = agent.natural_language_search(query, consultants)

    enriched = []
    for r in results:
        consultant = next(
            (c for c in consultants if c["id"] == r["id"]),
            None,
        )
        if consultant:
            enriched.append(
                {
                    **consultant,
                    "score": r.get("score", 0),
                    "explanation": r.get("explanation", ""),
                }
            )

    return {"results": enriched}


@app.get("/api/skills-market/consultants" "/{consultant_id}/analysis")
async def api_skills_market_analysis(
    consultant_id: int,
):
    """Analyse forces/faiblesses via LLM"""
    consultant = skills_market_db.get_consultant(consultant_id)
    if not consultant:
        return JSONResponse(
            {"error": "Consultant non trouve"},
            status_code=404,
        )

    agent = SkillsMarketAgent()
    analysis = agent.analyze_strengths(consultant)

    skills_market_db.update_consultant_analysis(
        consultant_id,
        strengths=analysis.get("strengths", []),
        improvement_areas=analysis.get("improvement_areas", []),
        management_suggestions=analysis.get("management_suggestions", ""),
    )

    return {"analysis": analysis}


# --- Photo de profil consultant ---

PHOTOS_DIR = BASE_DIR / "static" / "photos"
PHOTOS_DIR.mkdir(parents=True, exist_ok=True)

_ALLOWED_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif"}
_MAX_PHOTO_SIZE = 5 * 1024 * 1024  # 5 MB


@app.post("/api/skills-market/consultants/{consultant_id}/photo")
@limiter.limit("5/minute")
async def api_skills_market_upload_photo(
    request: Request,
    consultant_id: int,
    photo: UploadFile = File(...),
):
    """Upload et associe une photo de profil a un consultant"""
    consultant = skills_market_db.get_consultant(consultant_id)
    if not consultant:
        return JSONResponse({"error": "Consultant non trouve"}, status_code=404)

    ext = Path(photo.filename).suffix.lower()
    if ext not in _ALLOWED_IMAGE_EXTS:
        return JSONResponse(
            {"error": f"Format non supporte. Formats acceptes: {', '.join(_ALLOWED_IMAGE_EXTS)}"},
            status_code=400,
        )

    content = await photo.read()
    if len(content) > _MAX_PHOTO_SIZE:
        return JSONResponse({"error": "Fichier trop volumineux (max 5 Mo)"}, status_code=400)

    # Sauvegarder avec un nom stable par consultant
    safe_name = f"consultant_{consultant_id}{ext}"
    photo_path = PHOTOS_DIR / safe_name
    with open(photo_path, "wb") as f:
        f.write(content)

    photo_url = f"/static/photos/{safe_name}"
    updated = skills_market_db.update_photo_url(consultant_id, photo_url)
    if not updated:
        return JSONResponse({"error": "Mise a jour echouee"}, status_code=500)

    return {"photo_url": photo_url, "message": "Photo mise a jour"}


# --- Generation de CV depuis le profil consultant ---


@app.post("/api/skills-market/consultants/{consultant_id}/cv")
@limiter.limit("5/minute")
async def api_skills_market_generate_cv(
    request: Request,
    consultant_id: int,
):
    """Genere un CV HTML/PDF adapte au besoin client depuis le profil consultant"""
    body = await request.json()
    client_need = body.get("client_need", "").strip()
    output_format = body.get("format", "html").lower()  # html ou pdf

    consultant = skills_market_db.get_consultant(consultant_id)
    if not consultant:
        return JSONResponse({"error": "Consultant non trouve"}, status_code=404)

    agent = SkillsMarketAgent()
    try:
        cv_html = agent.generate_cv(consultant, client_need=client_need)
    except Exception as e:
        return JSONResponse(
            {"error": f"Erreur generation CV: {safe_error_message(e)}"}, status_code=500
        )

    if output_format == "pdf":
        try:
            from datetime import datetime as _dt

            from utils.pdf_converter import pdf_converter as _pdf_conv

            timestamp = _dt.now().strftime("%Y%m%d_%H%M%S")
            safe_name_cv = consultant.get("name", "consultant").replace(" ", "_")
            cv_dir = BASE_DIR / "static" / "cvs"
            cv_dir.mkdir(parents=True, exist_ok=True)
            pdf_path = cv_dir / f"cv_{safe_name_cv}_{timestamp}.pdf"
            result = _pdf_conv.html_to_pdf(cv_html, str(pdf_path))
            if not result:
                return JSONResponse(
                    {"error": "Conversion PDF indisponible (WeasyPrint non installe)"},
                    status_code=500,
                )
            return {
                "format": "pdf",
                "download_url": f"/static/cvs/{pdf_path.name}",
            }
        except Exception as e:
            return JSONResponse(
                {"error": f"Erreur export PDF: {safe_error_message(e)}"}, status_code=500
            )

    return {"format": "html", "cv_html": cv_html}


@app.post("/api/skills-market/consultants/{consultant_id}/cover-letter")
@limiter.limit("5/minute")
async def api_skills_market_cover_letter(
    request: Request,
    consultant_id: int,
):
    """Genere une lettre de motivation PDF/HTML a partir du profil consultant et d'une offre"""
    body = await request.json()
    job_offer = sanitize_text_input(body.get("job_offer", ""), max_length=5000)
    output_format = body.get("format", "pdf").lower()

    if not job_offer:
        return JSONResponse({"error": "Champ 'job_offer' requis"}, status_code=400)

    consultant = skills_market_db.get_consultant(consultant_id)
    if not consultant:
        return JSONResponse({"error": "Consultant non trouve"}, status_code=404)

    agent = SkillsMarketAgent()
    try:
        letter_html = agent.generate_cover_letter(consultant, job_offer)
    except Exception as e:
        return JSONResponse(
            {"error": f"Erreur generation lettre: {safe_error_message(e)}"}, status_code=500
        )

    if output_format == "pdf":
        try:
            from datetime import datetime as _dt

            from utils.pdf_converter import pdf_converter as _pdf_conv

            timestamp = _dt.now().strftime("%Y%m%d_%H%M%S")
            safe_name = consultant.get("name", "consultant").replace(" ", "_")
            cv_dir = BASE_DIR / "static" / "cvs"
            cv_dir.mkdir(parents=True, exist_ok=True)
            pdf_path = cv_dir / f"lettre_{safe_name}_{timestamp}.pdf"
            result = _pdf_conv.html_to_pdf(letter_html, str(pdf_path))
            if not result:
                return JSONResponse(
                    {"error": "Conversion PDF indisponible (WeasyPrint non installe)"},
                    status_code=500,
                )
            return {
                "format": "pdf",
                "download_url": f"/static/cvs/{pdf_path.name}",
            }
        except Exception as e:
            return JSONResponse(
                {"error": f"Erreur export PDF: {safe_error_message(e)}"}, status_code=500
            )

    return {"format": "html", "letter_html": letter_html}


# === E-LEARNING ADAPTATIF ===

elearning_db = ElearningDatabase()


@app.get("/elearning", response_class=HTMLResponse)
async def elearning_page(request: Request):
    return templates.TemplateResponse(
        "elearning.html",
        {
            "request": request,
            "active": "elearning",
            "consultant_name": CONSULTANT_NAME,
        },
    )


# --- Sessions ---


@app.post("/api/elearning/session/init")
async def api_elearning_session_init(request: Request):
    """Initialise ou recupere une session etudiant"""
    body = await request.json()
    identifier = body.get("session_identifier")
    session = elearning_db.init_session(identifier)
    return session


# --- Cours: Generation ---


def _run_course_generator(
    job_id,
    topic,
    audience,
    difficulty,
    duration,
    document_content=None,
    mode="free",
    interview_type="",
    consultant_id=0,
):
    """Background job pour generer un cours"""
    try:
        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        # Charger le profil consultant si fourni
        consultant_profile = None
        if consultant_id and int(consultant_id) > 0:
            try:
                from utils.consultant_db import ConsultantDatabase
                cdb = ConsultantDatabase()
                consultant_profile = cdb.get_consultant(int(consultant_id))
            except Exception as e:
                print(f"Erreur chargement consultant {consultant_id}: {e}")

        if document_content:
            result = agent.generate_course_from_document(
                document_content=document_content,
                target_audience=audience,
                difficulty=difficulty,
                duration_hours=duration,
                progress_callback=progress,
                mode=mode,
            )
        else:
            result = agent.generate_course(
                topic=topic,
                target_audience=audience,
                difficulty=difficulty,
                duration_hours=duration,
                progress_callback=progress,
                mode=mode,
                interview_type=interview_type,
                consultant_profile=consultant_profile,
            )

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            course_id = elearning_db.save_course(result)
            result["id"] = course_id
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "course_id": course_id,
                "title": result.get("title", ""),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du cours")


@app.post("/api/elearning/course/generate")
@limiter.limit("3/minute")
async def api_elearning_generate_course(
    request: Request,
    topic: str = Form(""),
    target_audience: str = Form("Professionnels"),
    difficulty: str = Form("intermediate"),
    duration_hours: int = Form(3),
    document_content: str = Form(""),
    mode: str = Form("free"),
    interview_type: str = Form(""),
    consultant_id: int = Form(0),
    auto_duration: str = Form("false"),
):
    """Lance la generation d'un cours (depuis sujet ou document)"""
    topic = sanitize_text_input(topic, max_length=1000)
    target_audience = sanitize_text_input(target_audience)
    document_content = sanitize_text_input(document_content, max_length=50000)
    interview_type = sanitize_text_input(interview_type)

    if not topic and not document_content:
        return JSONResponse(
            {"error": "Sujet ou document requis"},
            status_code=400,
        )

    if difficulty not in ("beginner", "intermediate", "advanced"):
        return JSONResponse({"error": "Difficulte invalide"}, status_code=400)

    valid_modes = ("free", "interview", "certification", "training")
    if mode not in valid_modes:
        mode = "free"

    valid_interview_types = ("rh", "technique", "cas", "fit", "")
    if interview_type not in valid_interview_types:
        interview_type = ""

    # Durée auto : passer 0 à l'agent pour qu'il la détermine lui-même
    if auto_duration.lower() in ("true", "1", "on"):
        duration_hours = 0

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    import threading

    t = threading.Thread(
        target=_run_course_generator,
        args=(
            job_id,
            topic,
            target_audience,
            difficulty,
            duration_hours,
            document_content or None,
            mode,
            interview_type,
            consultant_id,
        ),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@app.post("/api/elearning/course/upload-document")
@limiter.limit("10/minute")
async def api_elearning_upload_document(
    request: Request,
    file: UploadFile = File(...),
):
    """Parse un document uploade pour e-learning"""
    try:
        content = await validate_file_upload(file)
        filename = sanitize_filename(file.filename or "document.txt")
        ext = Path(filename).suffix.lower()

        if ext in (".md", ".txt", ".markdown"):
            text = content.decode("utf-8")
        elif ext == ".pdf":
            import io

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext in (".html", ".htm"):
            text = content.decode("utf-8")
        elif ext == ".pptx":
            import io

            from pptx import Presentation as PptxPres

            prs = PptxPres(io.BytesIO(content))
            slide_texts = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                texts.append(p.text.strip())
                if texts:
                    slide_texts.append("\n".join(texts))
            text = "\n\n---\n\n".join(slide_texts)
        else:
            return JSONResponse(
                {"error": f"Format non supporte: {ext}"},
                status_code=400,
            )

        return {
            "text": text,
            "filename": filename,
            "length": len(text),
        }
    except Exception as e:
        return JSONResponse(
            {"error": safe_error_message(e)},
            status_code=500,
        )


@app.get("/api/elearning/course/stream/{job_id}")
async def api_elearning_course_stream(job_id: str):
    """SSE stream pour la generation de cours"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


# --- Cours: Regeneration ---


def _run_course_regenerator(job_id, course_id, feedback):
    """Background job pour regenerer un cours"""
    try:
        course = elearning_db.get_course(course_id)
        if not course:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Cours non trouve"
            return

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        result = agent.regenerate_with_feedback(course, feedback, progress_callback=progress)

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            elearning_db.delete_course(course_id)
            new_id = elearning_db.save_course(result)
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "course_id": new_id,
                "title": result.get("title", ""),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Regeneration du cours")


@app.post("/api/elearning/course/{course_id}/regenerate")
@limiter.limit("3/minute")
async def api_elearning_regenerate_course(
    request: Request,
    course_id: int,
    feedback: str = Form(...),
):
    """Regenere un cours avec feedback"""
    feedback = sanitize_text_input(feedback)
    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    import threading

    t = threading.Thread(
        target=_run_course_regenerator,
        args=(job_id, course_id, feedback),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


# --- Cours: CRUD ---


@app.get("/api/elearning/courses")
async def api_elearning_list_courses():
    """Liste tous les cours"""
    courses = elearning_db.get_all_courses()
    return {"courses": courses}


@app.get("/api/elearning/course/{course_id}")
async def api_elearning_get_course(course_id: int):
    """Recupere un cours complet"""
    course = elearning_db.get_course(course_id)
    if not course:
        return JSONResponse({"error": "Cours non trouve"}, status_code=404)
    return {"course": course}


@app.delete("/api/elearning/course/{course_id}")
async def api_elearning_delete_course(course_id: int):
    """Supprime un cours"""
    deleted = elearning_db.delete_course(course_id)
    if not deleted:
        return JSONResponse({"error": "Cours non trouve"}, status_code=404)
    return {"ok": True}


@app.post("/api/elearning/lesson/{lesson_id}/solutions")
@limiter.limit("10/minute")
async def api_elearning_lesson_solutions(request: Request, lesson_id: int):
    """Génère les solutions pour les exercices d'une leçon existante."""
    body = await request.json()
    consultant_id = int(body.get("consultant_id", 0))

    # Trouver la leçon dans la DB via le cours
    course_id = int(body.get("course_id", 0))
    topic = sanitize_text_input(body.get("topic", ""), max_length=500)
    lesson_title = sanitize_text_input(body.get("lesson_title", ""), max_length=500)
    exercises = body.get("exercises", [])

    if not exercises:
        return JSONResponse({"error": "Aucun exercice fourni"}, status_code=400)

    consultant_profile = None
    if consultant_id > 0:
        try:
            consultant_db = ConsultantDatabase()
            consultant_profile = consultant_db.get_consultant(consultant_id)
        except Exception:
            pass

    try:
        agent = ElearningAgent()
        updated_exercises = agent.generate_exercise_solutions(
            exercises=exercises,
            topic=topic,
            lesson_title=lesson_title,
            consultant_profile=consultant_profile,
        )
        # Persister en DB si lesson_id fourni
        if lesson_id > 0:
            elearning_db.update_lesson_exercises(lesson_id, updated_exercises)
        return {"exercises": updated_exercises}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# --- Entretien: Chat & Simulation ---


@app.post("/api/elearning/interview/chat")
async def api_elearning_interview_chat(
    topic: str = Form(...),
    role: str = Form(...),
    messages: str = Form(...),
    interviewer_name: str = Form(""),
    interviewer_linkedin: str = Form(""),
):
    """Envoie un message dans la simulation d'entretien"""
    try:
        msg_list = json.loads(messages)
        agent = ElearningAgent()
        response = agent.interview_chat(
            topic=topic,
            role=role,
            messages=msg_list,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
        )
        return {"message": response}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/elearning/interview/analyze")
async def api_elearning_interview_analyze(
    topic: str = Form(...),
    role: str = Form(...),
    messages: str = Form(...),
    interviewer_name: str = Form(""),
    interviewer_linkedin: str = Form(""),
):
    """Analyse la performance de l'entretien"""
    try:
        msg_list = json.loads(messages)
        agent = ElearningAgent()
        analysis = agent.analyze_interview_performance(
            topic=topic,
            role=role,
            messages=msg_list,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
        )
        return {"analysis": analysis}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# --- Quiz: Generation ---


def _run_quiz_generator(job_id, course_id, lesson_id, difficulty, mode="free"):
    """Background job pour generer un quiz"""
    try:
        if lesson_id:
            course = elearning_db.get_course(course_id)
            if not course:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Cours non trouve"
                return

            lesson = None
            for mod in course.get("modules", []):
                for les in mod.get("lessons", []):
                    if les["id"] == lesson_id:
                        lesson = les
                        break

            if not lesson:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Lecon non trouvee"
                return

            lesson_title = lesson.get("title", "")
            lesson_content = lesson.get("content_markdown", "")
        else:
            course = elearning_db.get_course(course_id)
            if not course:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Cours non trouve"
                return

            lesson_title = course.get("title", "")
            all_content = []
            for mod in course.get("modules", []):
                for les in mod.get("lessons", []):
                    all_content.append(les.get("content_markdown", ""))
            lesson_content = "\n\n".join(all_content)

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        result = agent.generate_quiz(
            lesson_title=lesson_title,
            lesson_content=lesson_content,
            difficulty=difficulty,
            progress_callback=progress,
            mode=mode,
        )

        if "error" in result:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = result["error"]
        else:
            result["course_id"] = course_id
            result["lesson_id"] = lesson_id
            result["difficulty_level"] = difficulty
            quiz_id = elearning_db.save_quiz(result)
            jobs[job_id]["status"] = "done"
            jobs[job_id]["result"] = {
                "quiz_id": quiz_id,
                "questions_count": len(result.get("questions", [])),
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du quiz")


@app.post("/api/elearning/quiz/generate")
@limiter.limit("5/minute")
async def api_elearning_generate_quiz(
    request: Request,
    course_id: int = Form(...),
    lesson_id: Optional[int] = Form(None),
    difficulty: str = Form("medium"),
    mode: str = Form("free"),
):
    """Lance la generation d'un quiz"""
    if difficulty not in ("easy", "medium", "hard"):
        return JSONResponse({"error": "Difficulte invalide"}, status_code=400)

    valid_modes = ("free", "interview", "certification", "training")
    if mode not in valid_modes:
        mode = "free"

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    import threading

    t = threading.Thread(
        target=_run_quiz_generator,
        args=(job_id, course_id, lesson_id, difficulty, mode),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@app.get("/api/elearning/quiz/stream/{job_id}")
async def api_elearning_quiz_stream(job_id: str):
    """SSE stream pour la generation de quiz"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


# --- Quiz: Tentatives ---


@app.post("/api/elearning/quiz/start")
async def api_elearning_quiz_start(
    request: Request,
    quiz_id: int = Form(...),
    session_identifier: str = Form(...),
):
    """Demarre une tentative de quiz"""
    session = elearning_db.get_session(session_identifier)
    if not session:
        session = elearning_db.init_session(session_identifier)

    quiz = elearning_db.get_quiz(quiz_id)
    if not quiz:
        return JSONResponse({"error": "Quiz non trouve"}, status_code=404)

    attempt_id = elearning_db.create_attempt(session["id"], quiz_id)

    # Get first question at medium difficulty
    questions = quiz.get("questions", [])
    first_question = questions[0] if questions else None

    if first_question:
        # Remove correct_answer from response
        q = {**first_question}
        q.pop("correct_answer", None)
        q.pop("explanation", None)
    else:
        q = None

    return {
        "attempt_id": attempt_id,
        "total_questions": len(questions),
        "first_question": q,
    }


@app.post("/api/elearning/quiz/answer")
async def api_elearning_quiz_answer(
    request: Request,
    attempt_id: int = Form(...),
    question_id: int = Form(...),
    answer: str = Form(...),
    time_spent: int = Form(0),
):
    """Soumet une reponse et obtient la suivante (adaptatif)"""
    results = elearning_db.get_attempt_results(attempt_id)
    if not results:
        return JSONResponse({"error": "Tentative non trouvee"}, status_code=404)

    quiz = elearning_db.get_quiz(results["quiz_id"])
    if not quiz:
        return JSONResponse({"error": "Quiz non trouve"}, status_code=404)

    # Find the question
    question = None
    for q in quiz.get("questions", []):
        if q["id"] == question_id:
            question = q
            break

    if not question:
        return JSONResponse({"error": "Question non trouvee"}, status_code=404)

    # Evaluate answer
    agent = ElearningAgent()
    evaluation = agent.evaluate_answer(question, answer)

    # Record answer
    elearning_db.record_answer(
        attempt_id,
        question_id,
        answer,
        evaluation["is_correct"],
        time_spent,
    )

    # Adaptive difficulty
    recent = elearning_db.get_recent_answers(attempt_id, 3)
    new_difficulty = agent.adapt_difficulty(recent)
    current_difficulty = results.get("current_difficulty", "medium")

    if new_difficulty:
        elearning_db.update_attempt_difficulty(attempt_id, new_difficulty)
        current_difficulty = new_difficulty

    # Find next unanswered question
    answered_ids = {a["question_id"] for a in results.get("answers", [])}
    answered_ids.add(question_id)

    next_question = None
    for q in quiz.get("questions", []):
        if q["id"] not in answered_ids:
            next_question = {**q}
            next_question.pop("correct_answer", None)
            next_question.pop("explanation", None)
            break

    return {
        "is_correct": evaluation["is_correct"],
        "explanation": evaluation.get("explanation", ""),
        "feedback": evaluation.get("feedback", ""),
        "current_difficulty": current_difficulty,
        "difficulty_changed": new_difficulty is not None,
        "next_question": next_question,
        "questions_remaining": len(quiz.get("questions", [])) - len(answered_ids),
    }


@app.get("/api/elearning/quiz/results/{attempt_id}")
async def api_elearning_quiz_results(attempt_id: int):
    """Recupere les resultats d'une tentative"""
    # Complete the attempt first
    elearning_db.complete_attempt(attempt_id)

    results = elearning_db.get_attempt_results(attempt_id)
    if not results:
        return JSONResponse({"error": "Tentative non trouvee"}, status_code=404)
    return {"results": results}


@app.get("/api/elearning/quizzes/{course_id}")
async def api_elearning_list_quizzes(course_id: int):
    """Liste les quiz d'un cours"""
    quizzes = elearning_db.get_quizzes_for_course(course_id)
    return {"quizzes": quizzes}


# --- Parcours d'apprentissage ---


def _run_learning_path_generator(job_id, session_id, course_id, goals):
    """Background job pour generer un parcours"""
    try:
        course = elearning_db.get_course(course_id)
        if not course:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Cours non trouve"
            return

        session = elearning_db.get_session(session_id)
        if not session:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = "Session non trouvee"
            return

        agent = ElearningAgent()

        def progress(step, detail):
            jobs[job_id]["steps"].append({"step": step, "detail": detail})

        # Get quiz results for gap analysis
        quiz_results = elearning_db.get_session_quiz_results(session["id"], course_id)

        progress("analyzing", "Analyse des performances...")

        gaps = agent.analyze_knowledge_gaps(quiz_results, course.get("modules", []))

        progress("generating", "Generation du parcours...")

        path_data = agent.generate_learning_path(
            gaps=gaps,
            goals=goals,
            course_modules=course.get("modules", []),
            progress_callback=progress,
        )

        # Save learning path
        path_id = elearning_db.save_learning_path(
            {
                "session_id": session["id"],
                "course_id": course_id,
                "stated_goals": goals,
                "knowledge_gaps": path_data.get("knowledge_gaps", []),
                "recommendations": path_data.get("recommendations", []),
            }
        )

        jobs[job_id]["status"] = "done"
        jobs[job_id]["result"] = {
            "path_id": path_id,
            "steps_count": len(path_data.get("path_steps", [])),
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = safe_error_message(e, "Generation du parcours")


@app.post("/api/elearning/learning-path/create")
@limiter.limit("3/minute")
async def api_elearning_create_learning_path(
    request: Request,
    session_identifier: str = Form(...),
    course_id: int = Form(...),
    goals: str = Form(""),
):
    """Cree un parcours d'apprentissage personnalise"""
    try:
        goals_list = json.loads(goals) if goals else []
    except json.JSONDecodeError:
        goals_list = [g.strip() for g in goals.split(",") if g.strip()]

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {"status": "running", "steps": []}

    import threading

    t = threading.Thread(
        target=_run_learning_path_generator,
        args=(job_id, session_identifier, course_id, goals_list),
        daemon=True,
    )
    t.start()

    return {"job_id": job_id}


@app.get("/api/elearning/learning-path/stream/{job_id}")
async def api_elearning_learning_path_stream(job_id: str):
    """SSE stream pour la generation de parcours"""
    if job_id not in jobs:
        return JSONResponse({"error": "Job non trouve"}, status_code=404)

    async def event_generator():
        last_idx = 0
        while True:
            job = jobs.get(job_id, {})
            steps = job.get("steps", [])

            if len(steps) > last_idx:
                for step in steps[last_idx:]:
                    yield f"data: {json.dumps(step)}\n\n"
                last_idx = len(steps)

            if job.get("status") in ("done", "error"):
                final = {
                    "status": job["status"],
                    "result": job.get("result"),
                    "error": job.get("error"),
                }
                yield f"data: {json.dumps(final)}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
    )


@app.get("/api/elearning/learning-path" "/{session_identifier}/{course_id}")
async def api_elearning_get_learning_path(session_identifier: str, course_id: int):
    """Recupere un parcours d'apprentissage"""
    session = elearning_db.get_session(session_identifier)
    if not session:
        return JSONResponse({"error": "Session non trouvee"}, status_code=404)

    path = elearning_db.get_learning_path(session["id"], course_id)
    if not path:
        return JSONResponse({"error": "Parcours non trouve"}, status_code=404)
    return {"path": path}


@app.post("/api/elearning/learning-path/update-progress")
async def api_elearning_update_progress(
    request: Request,
    learning_path_id: int = Form(...),
    module_id: int = Form(...),
    completion_pct: float = Form(...),
    mastery_level: str = Form("none"),
):
    """Met a jour la progression d'un module"""
    valid_levels = (
        "none",
        "beginner",
        "intermediate",
        "proficient",
        "expert",
    )
    if mastery_level not in valid_levels:
        return JSONResponse(
            {"error": "Niveau de maitrise invalide"},
            status_code=400,
        )

    elearning_db.update_module_progress(
        learning_path_id,
        module_id,
        completion_pct,
        mastery_level,
    )
    return {"ok": True}


# === INTERVIEW CHAT ===


@app.post("/api/elearning/interview/chat")
@limiter.limit("30/minute")
async def api_elearning_interview_chat(request: Request):
    """
    Tour de conversation dans un entretien simulé interactif.

    Body JSON :
      {
        "messages": [{"role": "user"|"assistant", "content": "..."}],
        "topic": "Python",          # optionnel
        "role": "Data Engineer"     # optionnel
      }

    Retourne la réponse de l'interviewer.
    """
    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"error": "Corps JSON invalide"}, status_code=400)

    messages = body.get("messages")
    if not messages or not isinstance(messages, list):
        return JSONResponse({"error": "Champ 'messages' requis (liste)"}, status_code=400)

    # Valider les messages
    for msg in messages:
        if not isinstance(msg, dict) or msg.get("role") not in ("user", "assistant"):
            return JSONResponse(
                {"error": "Chaque message doit avoir 'role' (user|assistant) et 'content'"},
                status_code=400,
            )

    topic = str(body.get("topic", ""))[:200]
    role = str(body.get("role", ""))[:200]
    interviewer_name = str(body.get("interviewer_name", ""))[:200]
    interviewer_linkedin = str(body.get("interviewer_linkedin", ""))[:500]
    interview_type = str(body.get("interview_type", ""))[:20]

    try:
        agent = ElearningAgent()
        reply = agent.interview_chat(
            messages,
            topic=topic,
            role=role,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
            interview_type=interview_type,
        )
    except Exception as e:
        return JSONResponse({"error": f"Erreur entretien : {str(e)}"}, status_code=500)

    return {"reply": reply}


@app.post("/api/elearning/interview/analyze")
@limiter.limit("10/minute")
async def api_elearning_interview_analyze(request: Request):
    """
    Analyse les performances d'un entretien simulé complet.

    Body JSON :
      {
        "messages": [{"role": "user"|"assistant", "content": "..."}],
        "topic": "Python",
        "role": "Data Engineer",
        "interview_type": "rh|technique|cas|fit"
      }

    Retourne {overall_score, grade, strengths, improvements, detailed_feedback, ...}
    """
    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"error": "Corps JSON invalide"}, status_code=400)

    messages = body.get("messages")
    if not messages or not isinstance(messages, list):
        return JSONResponse({"error": "Champ 'messages' requis (liste)"}, status_code=400)

    topic = str(body.get("topic", ""))[:200]
    role = str(body.get("role", ""))[:200]
    interviewer_name = str(body.get("interviewer_name", ""))[:200]
    interviewer_linkedin = str(body.get("interviewer_linkedin", ""))[:500]
    interview_type = str(body.get("interview_type", ""))[:20]

    try:
        agent = ElearningAgent()
        analysis = agent.analyze_interview_performance(
            messages,
            topic=topic,
            role=role,
            interviewer_name=interviewer_name,
            interviewer_linkedin=interviewer_linkedin,
            interview_type=interview_type,
        )
    except Exception as e:
        return JSONResponse({"error": f"Erreur analyse : {str(e)}"}, status_code=500)

    return {"analysis": analysis}


# === PREMIUM RESOURCES ===


@app.post("/api/elearning/course/{course_id}/premium-resources")
@limiter.limit("5/minute")
async def api_elearning_premium_resources(
    request: Request,
    course_id: int,
):
    """
    Génère des ressources premium pour un cours :
    fiches mémo, résumés condensés, cheatsheet globale.

    Body JSON optionnel :
      { "resource_types": ["memo", "summary", "cheatsheet"] }
    """
    course = elearning_db.get_course(course_id)
    if not course:
        return JSONResponse({"error": "Cours non trouvé"}, status_code=404)

    try:
        body = await request.json()
    except Exception:
        body = {}

    resource_types = body.get("resource_types") if isinstance(body, dict) else None
    valid_types = {"memo", "summary", "cheatsheet"}
    if resource_types is not None:
        resource_types = [t for t in resource_types if t in valid_types] or None

    try:
        agent = ElearningAgent()
        resources = agent.generate_premium_resources(course, resource_types)
    except Exception as e:
        return JSONResponse(
            {"error": f"Erreur génération ressources premium : {str(e)}"},
            status_code=500,
        )

    return {"resources": resources}


# === MEETING CAPTURE AI ===


@app.get("/meeting-capture", response_class=HTMLResponse)
async def meeting_capture_page(request: Request):
    """Page Meeting Capture AI"""
    user = get_current_user(request)
    if not user:
        return RedirectResponse(url="/login", status_code=302)
    return templates.TemplateResponse(
        "meeting-capture.html",
        {
            "request": request,
            "active": "meeting-capture",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/meeting-capture/upload")
@limiter.limit("10/minute")
async def api_meeting_capture_upload(
    request: Request,
    video: UploadFile = File(...),
    gemini_model: str = Form(default="gemini-2.0-flash"),
):
    """Upload d'une vidéo de réunion et lancement de l'analyse Gemini"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)

    # Valider le modèle
    if gemini_model not in MC_GEMINI_MODELS:
        allowed = ", ".join(MC_GEMINI_MODELS.keys())
        return JSONResponse(
            {"detail": "Modèle non supporté. Modèles disponibles : " + allowed},
            status_code=400,
        )

    # Valider l'extension
    allowed_exts = {".webm", ".mp4", ".mov", ".avi", ".mkv"}
    filename = video.filename or "upload"
    suffix = Path(filename).suffix.lower()
    if suffix not in allowed_exts:
        return JSONResponse(
            {
                "detail": "Extension non supportée. Extensions acceptées : "
                + ", ".join(sorted(allowed_exts))
            },
            status_code=400,
        )

    content = await video.read()
    if not content:
        return JSONResponse({"detail": "Le fichier vidéo est vide."}, status_code=400)

    # Sauvegarder dans temp/
    temp_dir = BASE_DIR / "temp"
    temp_dir.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = sanitize_filename("meeting_" + timestamp + suffix)
    temp_path = temp_dir / safe_name

    with open(str(temp_path), "wb") as f:
        f.write(content)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "meeting-capture",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
        "video_path": str(temp_path),
        "gemini_model": gemini_model,
    }

    thread = threading.Thread(target=_run_meeting_capture, args=(job_id,), daemon=True)
    thread.start()

    return {"job_id": job_id}


def _run_meeting_capture(job_id: str):
    """Exécute l'analyse Meeting Capture en background"""
    job = jobs[job_id]
    video_path = job["video_path"]
    gemini_model = job["gemini_model"]

    try:
        api_key = os.getenv("GEMINI_API_KEY", "") or os.getenv("GOOGLE_API_KEY", "")
        agent = MeetingCaptureAgent(api_key=api_key, model=gemini_model)

        # Étape 1 : Upload
        job["steps"].append(
            {
                "step": "upload",
                "status": "active",
                "progress": 5,
                "message": "Upload de la vidéo vers Gemini…",
            }
        )
        uploaded_file = agent.upload_video(video_path)
        job["steps"].append(
            {"step": "upload", "status": "done", "progress": 20, "message": "Vidéo uploadée"}
        )

        # Étape 2 : Attente traitement
        job["steps"].append(
            {
                "step": "processing",
                "status": "active",
                "progress": 25,
                "message": "Traitement de la vidéo par Gemini…",
            }
        )
        processed_file = agent.wait_for_processing(uploaded_file)
        job["steps"].append(
            {"step": "processing", "status": "done", "progress": 55, "message": "Vidéo traitée"}
        )

        # Étape 3 : Analyse IA
        job["steps"].append(
            {
                "step": "analyze",
                "status": "active",
                "progress": 60,
                "message": "Analyse IA en cours…",
            }
        )
        result = agent.analyze(processed_file)
        job["steps"].append(
            {"step": "analyze", "status": "done", "progress": 95, "message": "Analyse terminée"}
        )

        # Finalisation
        job["steps"].append(
            {"step": "done", "status": "done", "progress": 100, "message": "Résultats prêts"}
        )
        job["result"] = result
        job["status"] = "done"

    except (MCAPIKeyError, VideoProcessingError, TimeoutError) as exc:
        job["status"] = "error"
        job["error"] = safe_error_message(exc)
    except Exception as exc:
        job["status"] = "error"
        job["error"] = safe_error_message(exc)
    finally:
        # Supprimer le fichier temporaire
        try:
            Path(video_path).unlink(missing_ok=True)
        except Exception:
            pass


@app.get("/api/meeting-capture/stream/{job_id}")
async def api_meeting_capture_stream(job_id: str):
    """SSE stream pour la progression de l'analyse Meeting Capture"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouvé"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.post("/api/meeting-capture/gmail-draft")
@limiter.limit("20/minute")
async def api_meeting_capture_gmail_draft(request: Request):
    """Crée un brouillon Gmail à partir du résultat Meeting Capture"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)

    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"detail": "Corps JSON invalide"}, status_code=400)

    to = sanitize_text_input(body.get("to", ""))
    subject = sanitize_text_input(body.get("subject", "Compte rendu de réunion"))
    email_body = body.get("body", "")
    credentials_path = body.get("credentials_path", "")

    if not credentials_path:
        home_dir = Path.home()
        credentials_path = str(home_dir / ".meetingcapture" / "credentials.json")

    token_path = str(Path(credentials_path).parent / "token.json")

    try:
        client = MeetingGmailClient()
        service = client.authenticate(credentials_path, token_path)
        draft = client.create_draft(service, to, subject, email_body)
        return JSONResponse(draft)
    except DraftCreationError as exc:
        return JSONResponse({"detail": str(exc)}, status_code=500)
    except Exception as exc:
        return JSONResponse({"detail": safe_error_message(exc)}, status_code=500)


# ============================================================
# === TENDERSCOUT AI ===
# ============================================================


@app.get("/tenderscout", response_class=HTMLResponse)
async def tenderscout_page(request: Request):
    """Page TenderScout AI"""
    user = get_current_user(request)
    if not user:
        return RedirectResponse(url="/login")
    return templates.TemplateResponse(
        "tenderscout.html",
        {
            "request": request,
            "active": "tenderscout",
            "consultant_name": CONSULTANT_NAME,
        },
    )


@app.post("/api/tenderscout/scan")
@limiter.limit("10/minute")
async def api_tenderscout_scan(request: Request):
    """Lance un scan d'appels d'offres en arrière-plan"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)

    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"detail": "Corps JSON invalide"}, status_code=400)

    keywords = body.get("keywords", [])
    sources = body.get("sources", [])

    if not keywords or not isinstance(keywords, list):
        return JSONResponse(
            {"detail": "Le champ 'keywords' est requis et doit être une liste non vide."},
            status_code=400,
        )
    keywords = [k.strip() for k in keywords if isinstance(k, str) and k.strip()]
    if not keywords:
        return JSONResponse(
            {"detail": "Aucun mot-clé valide fourni."},
            status_code=400,
        )

    if not sources or not isinstance(sources, list):
        return JSONResponse(
            {"detail": "Le champ 'sources' est requis."},
            status_code=400,
        )

    from agents.tender_scout_agent import VALID_SOURCES

    invalid = [s for s in sources if s not in VALID_SOURCES]
    if invalid:
        return JSONResponse(
            {
                "detail": f"Sources invalides : {invalid}. Valeurs acceptées : {sorted(VALID_SOURCES)}."
            },
            status_code=400,
        )

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "type": "tenderscout",
        "status": "running",
        "steps": [],
        "result": None,
        "error": None,
    }

    thread = threading.Thread(
        target=_run_tender_scout,
        args=(job_id, keywords, sources),
        daemon=True,
    )
    thread.start()

    return {"job_id": job_id}


def _run_tender_scout(job_id: str, keywords: list, sources: list):
    """Exécute le scan TenderScout en arrière-plan"""
    job = jobs[job_id]
    db = TenderDatabase()
    agent = TenderScoutAgent()
    all_tenders = []

    try:
        # Étape 1 : scraping BOAMP
        if "boamp" in sources:
            job["steps"].append(
                {
                    "step": "scrape_boamp",
                    "status": "active",
                    "progress": 10,
                    "message": "Scraping BOAMP…",
                }
            )
            boamp_tenders = agent.scrape_boamp(keywords)
            all_tenders.extend(boamp_tenders)
            job["steps"].append(
                {
                    "step": "scrape_boamp",
                    "status": "done",
                    "progress": 30,
                    "message": f"{len(boamp_tenders)} AOs trouvés sur BOAMP",
                }
            )

        # Étape 2 : scraping Francemarches
        if "francemarches" in sources:
            job["steps"].append(
                {
                    "step": "scrape_francemarches",
                    "status": "active",
                    "progress": 35,
                    "message": "Scraping Francemarches.com…",
                }
            )
            fm_tenders = agent.scrape_francemarches(keywords)
            all_tenders.extend(fm_tenders)
            job["steps"].append(
                {
                    "step": "scrape_francemarches",
                    "status": "done",
                    "progress": 50,
                    "message": f"{len(fm_tenders)} AOs trouvés sur Francemarches",
                }
            )

        # Étape 3 : sauvegarde idempotente
        job["steps"].append(
            {"step": "save", "status": "active", "progress": 52, "message": "Sauvegarde en base…"}
        )
        new_count = db.save_tenders(all_tenders)
        job["steps"].append(
            {
                "step": "save",
                "status": "done",
                "progress": 55,
                "message": f"{new_count} nouveau(x) AO(s) enregistré(s)",
            }
        )

        # Étape 4 : analyse Gemini des nouveaux AOs
        refs = [t["reference"] for t in all_tenders]
        unanalyzed_refs = set(db.get_unanalyzed_references(refs))
        to_analyze = [t for t in all_tenders if t["reference"] in unanalyzed_refs]

        analyzed_count = 0
        for i, tender in enumerate(to_analyze):
            progress = 55 + int(40 * (i + 1) / max(len(to_analyze), 1))
            job["steps"].append(
                {
                    "step": "analyze",
                    "status": "active",
                    "progress": progress,
                    "tender": tender["titre"][:60],
                    "message": f"Analyse {i + 1}/{len(to_analyze)}",
                }
            )
            try:
                analysis = agent.analyze_tender(tender)
                db.update_analysis(tender["reference"], analysis)
                analyzed_count += 1
                job["steps"][-1]["status"] = "done"
            except AnalysisError:
                job["steps"][-1]["status"] = "done"

        job["steps"].append(
            {"step": "done", "status": "done", "progress": 100, "message": "Analyse terminée"}
        )
        job["status"] = "done"
        job["result"] = {
            "total": len(all_tenders),
            "new": new_count,
            "analyzed": analyzed_count,
        }

    except ScrapingError as exc:
        job["status"] = "error"
        job["error"] = safe_error_message(exc)
    except Exception as exc:
        job["status"] = "error"
        job["error"] = safe_error_message(exc)


@app.get("/api/tenderscout/stream/{job_id}")
async def api_tenderscout_stream(job_id: str):
    """SSE stream pour la progression du scan TenderScout"""

    async def event_generator():
        job = jobs.get(job_id)
        if not job:
            yield send_sse("error_msg", {"message": "Job non trouvé"})
            return

        last_step_idx = 0
        while True:
            while last_step_idx < len(job["steps"]):
                step = job["steps"][last_step_idx]
                yield send_sse("step", step)
                last_step_idx += 1

            if job["status"] == "done":
                yield send_sse("result", job["result"])
                return
            elif job["status"] == "error":
                yield send_sse("error_msg", {"message": job["error"]})
                return

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"},
    )


@app.get("/api/tenderscout/tenders")
async def api_tenderscout_list(
    request: Request,
    source: Optional[str] = None,
    decision: Optional[str] = None,
):
    """Liste les appels d'offres avec filtres optionnels"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)
    db = TenderDatabase()
    return JSONResponse(db.get_tenders(source=source, decision=decision))


@app.get("/api/tenderscout/export")
async def api_tenderscout_export(request: Request):
    """Exporte les appels d'offres en Excel"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)

    temp_dir = BASE_DIR / "temp"
    temp_dir.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    export_path = str(temp_dir / f"tenders_{timestamp}.xlsx")

    db = TenderDatabase()
    db.export_to_excel(export_path)

    return FileResponse(
        export_path,
        filename=f"tenders_{timestamp}.xlsx",
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


@app.post("/api/tenderscout/notify")
@limiter.limit("10/minute")
async def api_tenderscout_notify(request: Request):
    """Envoie un récap Gmail des AOs GO"""
    user = get_current_user(request)
    if not user:
        return JSONResponse({"detail": "Non authentifié"}, status_code=401)

    try:
        body = await request.json()
    except Exception:
        return JSONResponse({"detail": "Corps JSON invalide"}, status_code=400)

    to = sanitize_text_input(body.get("to", ""))
    credentials_path = body.get("credentials_path", "")

    if not to:
        return JSONResponse({"detail": "Destinataire requis."}, status_code=400)
    if not credentials_path:
        return JSONResponse({"detail": "Chemin credentials.json requis."}, status_code=400)

    # Construire le récap
    db = TenderDatabase()
    go_tenders = db.get_tenders(decision="GO", limit=50)

    if not go_tenders:
        recap_body = "Aucun appel d'offres qualifié GO lors du dernier scan TenderScout."
    else:
        lines = [f"TenderScout AI — Récap des AOs qualifiés GO ({len(go_tenders)} AO(s))\n"]
        for t in go_tenders:
            lines.append(f"• {t['titre']}")
            lines.append(f"  Acheteur : {t.get('acheteur') or '—'}")
            lines.append(f"  Date limite : {t.get('date_limite') or '—'}")
            lines.append(f"  Score : {t.get('score') or '—'}/100")
            lines.append(f"  URL : {t.get('url', '')}")
            lines.append("")
        recap_body = "\n".join(lines)

    subject = f"TenderScout AI — {len(go_tenders)} AO(s) qualifiés GO"
    token_path = str(Path(credentials_path).parent / "tenderscout_token.json")

    try:
        client = MeetingGmailClient()
        service = client.authenticate(credentials_path, token_path)
        draft = client.create_draft(service, to, subject, recap_body)
        return JSONResponse(draft)
    except DraftCreationError as exc:
        return JSONResponse({"detail": str(exc)}, status_code=500)
    except Exception as exc:
        return JSONResponse({"detail": safe_error_message(exc)}, status_code=500)


# === MAIN ===

if __name__ == "__main__":
    import uvicorn

    # Check if SSL certificates exist
    ssl_cert = BASE_DIR / "ssl" / "cert.pem"
    ssl_key = BASE_DIR / "ssl" / "key.pem"
    use_ssl = ssl_cert.exists() and ssl_key.exists()

    protocol = "https" if use_ssl else "http"
    port = 8443 if use_ssl else 8000

    print(f"\n{'=' * 50}")
    print("  Consulting Tools Agents - Interface Web")
    print(f"  {protocol}://localhost:{port}")
    if use_ssl:
        print("  🔒 HTTPS activé (certificat auto-signé)")
    else:
        print("  ⚠️  HTTP uniquement (pas de SSL)")

    # Security warning for default password
    auth_password = os.getenv("AUTH_PASSWORD", "")
    if (
        not auth_password or auth_password == "CHANGE_ME_ON_FIRST_INSTALL"  # pragma: allowlist secret
    ):
        print("\n  🚨 SECURITE : Mot de passe par defaut detecte !")
        print("  → Veuillez configurer AUTH_PASSWORD dans votre fichier .env")
        print("  → Utilisez un mot de passe fort et unique")

    print(f"{'=' * 50}\n")

    if use_ssl:
        uvicorn.run(
            app, host="0.0.0.0", port=port, ssl_keyfile=str(ssl_key), ssl_certfile=str(ssl_cert)
        )
    else:
        uvicorn.run(app, host="0.0.0.0", port=port)
